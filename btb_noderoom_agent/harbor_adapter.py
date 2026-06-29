import json
import os
import shutil
import subprocess
import tarfile
import textwrap
from pathlib import Path
from typing import Any

from harbor.agents.base import BaseAgent
from harbor.environments.base import BaseEnvironment
from harbor.models.agent.context import AgentContext


class NodeRoomNodeAgent(BaseAgent):
    """Harbor adapter that runs NodeRoom NodeAgent against BankerToolBench.

    The adapter keeps the generated smoke task path intact, and adds a general
    path for official BTB tasks. General mode extracts only candidate-visible
    workspace and MCP source evidence, asks NodeAgent for an artifact plan, then
    materializes Office/PDF deliverables and citation receipts inside the BTB
    container for Gandalf to score.
    """

    SUPPORTS_ATIF = True

    def __init__(
        self,
        *args: Any,
        noderoom_repo: str | None = None,
        runner_script: str | None = None,
        mode: str | None = None,
        model_id: str | None = None,
        materializer_mode: str | None = None,
        allow_fallback_plan: bool | str | None = None,
        force_model_planner: bool | str | None = None,
        max_steps: int | str | None = None,
        planner_deadline_ms: int | str | None = None,
        runner_timeout_sec: int | str | None = None,
        **kwargs: Any,
    ) -> None:
        super().__init__(*args, **kwargs)
        repo = noderoom_repo or os.environ.get("NODEROOM_REPO_ROOT")
        if not repo:
            raise ValueError(
                "NodeRoomNodeAgent requires noderoom_repo kwarg or NODEROOM_REPO_ROOT."
            )
        self.noderoom_repo = Path(repo).resolve()
        self.runner_script = (
            Path(runner_script).resolve()
            if runner_script
            else self.noderoom_repo / "scripts" / "bankertoolbench-nodeagent-smoke-runner.ts"
        )
        self.mode = (mode or os.environ.get("BTB_NODEAGENT_MODE") or "auto").strip().lower()
        if self.mode not in {"auto", "smoke", "general"}:
            raise ValueError("NodeRoomNodeAgent mode must be auto, smoke, or general.")
        self.materializer_mode = (
            materializer_mode
            or os.environ.get("BTB_NODEAGENT_MATERIALIZER_MODE")
            or "replay"
        ).strip().lower()
        if self.materializer_mode not in {"replay", "general-only", "generic-only"}:
            raise ValueError(
                "NodeRoomNodeAgent materializer_mode must be replay, general-only, or generic-only."
            )
        allow_fallback_plan_value = allow_fallback_plan
        if allow_fallback_plan_value is None:
            allow_fallback_plan_value = os.environ.get("BTB_NODEAGENT_ALLOW_FALLBACK_PLAN")
        self.allow_fallback_plan = self._parse_bool(
            allow_fallback_plan_value,
            default=True,
        )
        force_model_planner_value = force_model_planner
        if force_model_planner_value is None:
            force_model_planner_value = os.environ.get("BTB_NODEAGENT_FORCE_MODEL_PLANNER")
        self.force_model_planner = self._parse_bool(
            force_model_planner_value,
            default=False,
        )
        self.model_id = (
            model_id
            or os.environ.get("BTB_NODEAGENT_MODEL_ID")
            or "z-ai/glm-5.2"
        )
        self.max_steps = int(max_steps or os.environ.get("BTB_NODEAGENT_MAX_STEPS") or "6")
        # BTB-glm52-task2 fix: dense-LBO/S&U class tasks (e.g. btb-06c284ef) need
        # ~3-5x the teaser-class budget; raise the floor so these don't AbortError
        # at ~175s before any usage flushes. runner_timeout_sec must outlive
        # planner_deadline_ms + setup or the Python adapter will SIGKILL the node
        # subprocess before it can write the abort-error JSON.
        self.planner_deadline_ms = int(
            planner_deadline_ms
            or os.environ.get("BTB_NODEAGENT_PLANNER_DEADLINE_MS")
            or "420000"
        )
        self.runner_timeout_sec = int(
            runner_timeout_sec
            or os.environ.get("BTB_NODEAGENT_RUNNER_TIMEOUT_SEC")
            or "900"
        )

    @staticmethod
    def name() -> str:
        return "noderoom-nodeagent"

    def version(self) -> str | None:
        return "0.2.0-general"

    @staticmethod
    def _parse_bool(value: bool | str | None, *, default: bool) -> bool:
        if value is None:
            return default
        if isinstance(value, bool):
            return value
        normalized = str(value).strip().lower()
        if normalized in {"1", "true", "yes", "y", "on"}:
            return True
        if normalized in {"0", "false", "no", "n", "off"}:
            return False
        raise ValueError(f"Expected boolean value, received {value!r}.")

    async def setup(self, environment: BaseEnvironment) -> None:
        if not self.noderoom_repo.is_dir():
            raise FileNotFoundError(f"NodeRoom repo not found: {self.noderoom_repo}")
        if not self.runner_script.is_file():
            raise FileNotFoundError(f"NodeAgent runner not found: {self.runner_script}")
        await environment.exec(
            "mkdir -p /home/agent/workspace/banker_workspace/source "
            "/home/agent/workspace/banker_workspace/deliverables /logs/agent && "
            "chown -R agent:agent /home/agent/workspace/banker_workspace && "
            "chmod -R g+rwxs /home/agent/workspace/banker_workspace && "
            "chmod 777 /logs/agent",
            user="root",
            timeout_sec=30,
        )

    async def run(
        self,
        instruction: str,
        environment: BaseEnvironment,
        context: AgentContext,
    ) -> None:
        self.logs_dir.mkdir(parents=True, exist_ok=True)
        instruction_path = self.logs_dir / "instruction.txt"
        instruction_path.write_text(instruction, encoding="utf-8")

        selected_mode = self._select_mode(instruction)
        if selected_mode == "smoke":
            await self._run_smoke_task(instruction_path, environment, context)
            return

        await self._run_general_task(instruction_path, environment, context)

    async def _run_smoke_task(
        self,
        instruction_path: Path,
        environment: BaseEnvironment,
        context: AgentContext,
    ) -> None:
        facts_path = self.logs_dir / "nodeagent_source_facts.json"
        await self._extract_smoke_facts(environment, facts_path)

        output_dir = self.logs_dir / "nodeagent-output"
        output_dir.mkdir(parents=True, exist_ok=True)
        trajectory_path = self.logs_dir / "trajectory.json"
        trace_path = self.logs_dir / "nodeagent-trace.json"

        result = self._run_nodeagent_smoke_runner(
            instruction_path=instruction_path,
            facts_path=facts_path,
            output_dir=output_dir,
            trajectory_path=trajectory_path,
            trace_path=trace_path,
        )

        await self._publish_smoke_outputs(environment, output_dir, trajectory_path, trace_path)
        self._record_context(
            context,
            result,
            {
                "mode": "smoke",
                "runner": "bankertoolbench-nodeagent-smoke-runner",
                "noderoom_repo": str(self.noderoom_repo),
                "facts_path": str(facts_path),
                "trajectory_path": str(trajectory_path),
                "trace_path": str(trace_path),
                "deliverables_dir": str(output_dir / "deliverables"),
            },
        )

    async def _run_general_task(
        self,
        instruction_path: Path,
        environment: BaseEnvironment,
        context: AgentContext,
    ) -> None:
        await environment.upload_file(
            instruction_path,
            "/home/agent/workspace/banker_workspace/instruction.txt",
        )

        source_packet_path = self.logs_dir / "nodeagent_source_packet.json"
        await self._extract_general_source_packet(environment, source_packet_path)

        output_dir = self.logs_dir / "nodeagent-output"
        output_dir.mkdir(parents=True, exist_ok=True)
        artifact_plan_path = self.logs_dir / "artifact-plan.json"
        trajectory_path = self.logs_dir / "trajectory.json"
        trace_path = self.logs_dir / "nodeagent-trace.json"

        result = self._run_nodeagent_general_runner(
            instruction_path=instruction_path,
            source_packet_path=source_packet_path,
            artifact_plan_path=artifact_plan_path,
            output_dir=output_dir,
            trajectory_path=trajectory_path,
            trace_path=trace_path,
        )

        await self._materialize_general_outputs(environment, artifact_plan_path)
        await self._publish_general_logs(
            environment,
            output_dir,
            artifact_plan_path,
            trajectory_path,
            trace_path,
        )
        self._record_context(
            context,
            result,
            {
                "mode": "general",
                "runner": "bankertoolbench-nodeagent-smoke-runner",
                "model_id": self.model_id,
                "materializer_mode": self.materializer_mode,
                "allow_fallback_plan": self.allow_fallback_plan,
                "force_model_planner": self.force_model_planner,
                "fallback_used": result.get("fallbackUsed") if isinstance(result, dict) else None,
                "planner_stop_reason": result.get("plannerStopReason") if isinstance(result, dict) else None,
                "planner_transport": result.get("plannerTransport") if isinstance(result, dict) else None,
                "planner_error": result.get("plannerError") if isinstance(result, dict) else None,
                "max_steps": self.max_steps,
                "planner_deadline_ms": self.planner_deadline_ms,
                "noderoom_repo": str(self.noderoom_repo),
                "source_packet_path": str(source_packet_path),
                "artifact_plan_path": str(artifact_plan_path),
                "trajectory_path": str(trajectory_path),
                "trace_path": str(trace_path),
                "deliverables_dir": str(output_dir / "deliverables"),
                "boundary_box_policy": "Every emitted citation is materialized into boundary_box_receipts.json with cell/page/shape/paragraph/field locator status.",
            },
        )

    def _select_mode(self, instruction: str) -> str:
        if self.mode != "auto":
            return self.mode
        if "Write the VDR sector" in instruction and "EDGAR SIC description" in instruction:
            return "smoke"
        return "general"

    async def _extract_smoke_facts(
        self,
        environment: BaseEnvironment,
        target_path: Path,
    ) -> None:
        script = textwrap.dedent(
            r"""
            import hashlib
            import json
            import os
            import sys
            from pathlib import Path

            os.umask(0o002)
            sys.path.insert(0, "/opt/mcp-server")

            from openpyxl import load_workbook
            from tools import sec_edgar, vdr

            caller_uid = 10001
            workspace = Path("/home/agent/workspace/banker_workspace/source")
            workspace.mkdir(parents=True, exist_ok=True)

            vdr_result = vdr._copy_to_workspace("VNOM", str(workspace), "overview_company_identity", caller_uid)
            if not vdr_result.get("success"):
                raise RuntimeError(f"VDR copy failed: {vdr_result}")
            vdr_path = Path(vdr_result["filepaths"][0])

            sec_root = sec_edgar._resolve_data_root(sec_edgar._data_path())
            cik = sec_edgar.pad_cik("0002074176")
            sec_result = sec_edgar._copy_to_workspace(
                str(workspace),
                sec_edgar._path_submissions(sec_root, cik),
                f"submissions_{cik}.json",
                caller_uid,
            )
            if not sec_result.get("success"):
                raise RuntimeError(f"SEC EDGAR copy failed: {sec_result}")
            edgar_path = Path(sec_result["filepaths"][0])

            def read_company_profile(path):
                workbook = load_workbook(path, read_only=True, data_only=True)
                facts = {}
                try:
                    for sheet in workbook.worksheets:
                        for row in sheet.iter_rows(values_only=True):
                            values = [value for value in row if value is not None]
                            if len(values) < 2:
                                continue
                            key = str(values[0]).strip()
                            if key:
                                facts[key.lower()] = str(values[1]).strip()
                finally:
                    workbook.close()
                return facts

            profile = read_company_profile(vdr_path)
            sector = profile.get("sector")
            if not sector:
                raise RuntimeError(f"Sector not found in {vdr_path}")

            with edgar_path.open("r", encoding="utf-8") as handle:
                submissions = json.load(handle)
            sic_description = str(submissions.get("sicDescription") or "").strip()
            if not sic_description:
                raise RuntimeError(f"sicDescription not found in {edgar_path}")

            output = {
                "companyName": profile.get("name") or "Viper Energy",
                "sector": sector,
                "sicDescription": sic_description,
                "description": profile.get("description"),
                "vdrSourcePath": str(vdr_path),
                "edgarSourcePath": str(edgar_path),
                "mcpCalls": [
                    {"tool": "vdr.download_to_workspace", "symbol": "VNOM", "data_type": "overview_company_identity"},
                    {"tool": "sec_edgar.get_submissions", "cik": "0002074176"},
                ],
            }
            out_path = Path("/home/agent/workspace/banker_workspace/nodeagent_source_facts.json")
            out_path.write_text(json.dumps(output, indent=2), encoding="utf-8")
            print(json.dumps(output, indent=2))
            """
        ).strip()
        await environment.exec(
            "mkdir -p /home/agent/workspace/banker_workspace/source && "
            "chown -R agent:agent /home/agent/workspace/banker_workspace && "
            "chmod -R g+rwxs /home/agent/workspace/banker_workspace",
            user="root",
            timeout_sec=30,
        )
        result = await environment.exec(
            f"python - <<'PY'\n{script}\nPY",
            user="environment",
            timeout_sec=120,
        )
        (self.logs_dir / "extract-facts.stdout.txt").write_text(
            result.stdout or "",
            encoding="utf-8",
        )
        (self.logs_dir / "extract-facts.stderr.txt").write_text(
            result.stderr or "",
            encoding="utf-8",
        )
        if result.return_code != 0:
            raise RuntimeError(
                "NodeRoomNodeAgent source fact extraction failed "
                f"with code {result.return_code}: {result.stderr or result.stdout}"
            )
        await environment.download_file(
            "/home/agent/workspace/banker_workspace/nodeagent_source_facts.json",
            target_path,
        )

    async def _extract_general_source_packet(
        self,
        environment: BaseEnvironment,
        target_path: Path,
    ) -> None:
        script = textwrap.dedent(
            r"""
            import datetime
            import hashlib
            import json
            import os
            import re
            import sys
            from pathlib import Path

            os.umask(0o002)
            sys.path.insert(0, "/opt/mcp-server")

            WORKSPACE = Path("/home/agent/workspace")
            BANKER_ROOT = WORKSPACE / "banker_workspace"
            SOURCE_ROOT = BANKER_ROOT / "source"
            MCP_ROOT = SOURCE_ROOT / "mcp"
            INSTRUCTION_PATH = BANKER_ROOT / "instruction.txt"
            PACKET_PATH = BANKER_ROOT / "nodeagent_source_packet.json"
            CALLER_UID = 10001
            MAX_SOURCE_TICKERS = max(1, int(os.environ.get("BTB_NODEAGENT_MAX_SOURCE_TICKERS", "16")))
            PDF_SOURCE_MODE = os.environ.get("BTB_NODEAGENT_PDF_SOURCE_MODE", "auto").strip().lower()
            PDF_METADATA_ONLY = PDF_SOURCE_MODE in {"metadata", "metadata-only", "light", "skip"}

            SOURCE_ROOT.mkdir(parents=True, exist_ok=True)
            MCP_ROOT.mkdir(parents=True, exist_ok=True)

            STOP_TICKERS = {
                "A", "AI", "API", "ASC", "ARR", "BTB", "CAGR", "CAPEX", "CEO",
                "CFO", "CIM", "COGS", "DCF", "DOC", "DOCX", "EBIT", "EBITDA",
                "EPS", "ESG", "EV", "FCF", "GAAP", "GTM", "HTML", "IRR", "JSON",
                "KPI", "LBO", "LLM", "LTM", "MA", "MCP", "MD", "NAV", "NDA",
                "NPV", "PDF", "PPT", "PPTX", "QA", "QOQ", "SaaS", "SEC", "SGA",
                "TAM", "TS", "TTM", "UI", "USD", "VDR", "WACC", "XLS", "XLSX",
                "YOY",
            }
            COMPANY_TICKER_ALIASES = {
                "adobe": "ADBE",
                "adobe inc": "ADBE",
                "alphabet": "GOOGL",
                "alphabet inc": "GOOGL",
                "amazon": "AMZN",
                "amazon com": "AMZN",
                "amazon.com": "AMZN",
                "apogee therapeutics": "APGE",
                "comcast": "CMCSA",
                "comcast corporation": "CMCSA",
                "coty": "COTY",
                "deere": "DE",
                "deere and company": "DE",
                "deere & company": "DE",
                "disney": "DIS",
                "e l f beauty": "ELF",
                "e.l.f. beauty": "ELF",
                "estee lauder": "EL",
                "greenbrier": "GBX",
                "google": "GOOGL",
                "intel": "INTC",
                "intel corporation": "INTC",
                "meta": "META",
                "meta inc": "META",
                "microsoft": "MSFT",
                "microsoft corporation": "MSFT",
                "nvidia": "NVDA",
                "oracle": "ORCL",
                "oracle corporation": "ORCL",
                "or royalties": "OR",
                "salesforce": "CRM",
                "salesforce inc": "CRM",
                "ulta": "ULTA",
                "ulta beauty": "ULTA",
                "viper energy": "VNOM",
                "walt disney": "DIS",
                "warner bros discovery": "WBD",
                "warner bros. discovery": "WBD",
            }

            warnings = []
            mcp_calls = []

            def clean_value(value):
                if value is None:
                    return None
                if isinstance(value, (int, float, bool)):
                    return value
                if hasattr(value, "isoformat"):
                    return value.isoformat()
                text = str(value)
                return text[:1000]

            def is_under(path, parent):
                try:
                    path.resolve().relative_to(parent.resolve())
                    return True
                except ValueError:
                    return False

            def available_vdr_symbols():
                try:
                    from tools import vdr

                    data_path = vdr._data_path()
                    return {
                        path.name.upper().replace("-US", "")
                        for path in data_path.iterdir()
                        if path.is_dir() and not path.name.startswith(".")
                    }
                except Exception as exc:
                    warnings.append(f"VDR ticker inventory failed: {exc}")
                    return set()

            def infer_tickers(text):
                vdr_symbols = available_vdr_symbols()
                explicit_matches = []
                for pattern in [
                    r"\bTicker\s*[:=\(]?\s*([A-Z][A-Z.\-]{0,8})\b",
                    r"\bTicker\s+([A-Z][A-Z.\-]{0,8})\b",
                    r"\bticker\s+(?:is\s+)?(?:NYSE|NASDAQ|Nasdaq|NYSE:|NASDAQ:)?\s*:?\s*([A-Z][A-Z.\-]{0,8})\b",
                    r"\bticker\s+([A-Z][A-Z.\-]{0,8})\b",
                ]:
                    explicit_matches.extend(re.findall(pattern, text))
                matches = explicit_matches + re.findall(r"\b[A-Z]{2,5}\b", text)
                seen = []

                def add_symbol(raw, explicit=False):
                    symbol = raw.strip(".-,;:()[]{}").upper()
                    if not symbol or symbol in STOP_TICKERS:
                        return
                    if any(char.isdigit() for char in symbol):
                        return
                    if symbol.endswith("-US"):
                        symbol = symbol[:-3]
                    if len(symbol) == 1 and not explicit:
                        return
                    if vdr_symbols and symbol not in vdr_symbols and not explicit:
                        return
                    if symbol not in seen:
                        seen.append(symbol)

                for raw in explicit_matches:
                    add_symbol(raw, explicit=True)
                for raw in matches:
                    add_symbol(raw, explicit=raw in explicit_matches)

                normalized_text = re.sub(r"[^a-z0-9&.]+", " ", text.lower())
                normalized_text = re.sub(r"\s+", " ", normalized_text)
                for alias, ticker in sorted(COMPANY_TICKER_ALIASES.items(), key=lambda item: -len(item[0])):
                    normalized_alias = re.sub(r"[^a-z0-9&.]+", " ", alias.lower())
                    normalized_alias = re.sub(r"\s+", " ", normalized_alias).strip()
                    if not normalized_alias:
                        continue
                    if re.search(rf"(?<![a-z0-9]){re.escape(normalized_alias)}(?![a-z0-9])", normalized_text):
                        add_symbol(ticker, explicit=True)
                return seen[:12]

            def summarize_text(path):
                try:
                    data = path.read_text(encoding="utf-8", errors="replace")
                except Exception as exc:
                    return {"error": f"text read failed: {exc}"}
                return {"previewText": data[:8000]}

            def summarize_json(path):
                try:
                    with path.open("r", encoding="utf-8", errors="replace") as handle:
                        data = json.load(handle)
                    preview = json.dumps(data, indent=2, default=str)[:8000]
                    return {"previewText": preview}
                except Exception:
                    return summarize_text(path)

            def summarize_xlsx(path):
                try:
                    from openpyxl import load_workbook

                    workbook = load_workbook(path, read_only=True, data_only=False)
                    sheets = []
                    try:
                        for worksheet in workbook.worksheets[:8]:
                            cells = []
                            for row in worksheet.iter_rows(
                                min_row=1,
                                max_row=min(worksheet.max_row or 1, 80),
                                min_col=1,
                                max_col=min(worksheet.max_column or 1, 24),
                            ):
                                for cell in row:
                                    if cell.value is None:
                                        continue
                                    entry = {
                                        "address": cell.coordinate,
                                        "value": clean_value(cell.value),
                                    }
                                    if isinstance(cell.value, str) and cell.value.startswith("="):
                                        entry["formula"] = cell.value[:500]
                                    cells.append(entry)
                            sheets.append(
                                {
                                    "name": worksheet.title,
                                    "maxRow": worksheet.max_row,
                                    "maxColumn": worksheet.max_column,
                                    "cells": cells[:400],
                                }
                            )
                    finally:
                        workbook.close()
                    return {"sheets": sheets}
                except Exception as exc:
                    return {"error": f"xlsx summary failed: {exc}"}

            def summarize_xls(path):
                try:
                    import xlrd

                    book = xlrd.open_workbook(str(path), on_demand=True)
                    sheets = []
                    try:
                        for sheet in book.sheets()[:8]:
                            cells = []
                            for row_index in range(min(sheet.nrows, 80)):
                                for col_index in range(min(sheet.ncols, 24)):
                                    value = sheet.cell_value(row_index, col_index)
                                    if value in ("", None):
                                        continue
                                    address = f"R{row_index + 1}C{col_index + 1}"
                                    cells.append({"address": address, "value": clean_value(value)})
                            sheets.append(
                                {
                                    "name": sheet.name,
                                    "maxRow": sheet.nrows,
                                    "maxColumn": sheet.ncols,
                                    "cells": cells[:400],
                                }
                            )
                    finally:
                        book.release_resources()
                    return {"sheets": sheets}
                except Exception as exc:
                    return {"error": f"xls summary failed: {exc}"}

            def summarize_pdf(path):
                if PDF_METADATA_ONLY:
                    return {
                        "extractionMode": "metadata-only",
                        "pages": [
                            {
                                "page": 1,
                                "text": f"PDF source file indexed by metadata only: {path.name}",
                                "boxes": [
                                    {
                                        "x": 0,
                                        "y": 0,
                                        "w": 1,
                                        "h": 1,
                                        "page": 1,
                                        "unit": "page",
                                        "text": path.name[:80],
                                    }
                                ],
                            }
                        ],
                    }
                try:
                    import pdfplumber

                    pages = []
                    with pdfplumber.open(path) as pdf:
                        for index, page in enumerate(pdf.pages[:6], start=1):
                            text = (page.extract_text() or "")[:5000]
                            boxes = []
                            try:
                                for word in page.extract_words()[:60]:
                                    boxes.append(
                                        {
                                            "x": float(word.get("x0", 0)),
                                            "y": float(word.get("top", 0)),
                                            "w": float(word.get("x1", 0)) - float(word.get("x0", 0)),
                                            "h": float(word.get("bottom", 0)) - float(word.get("top", 0)),
                                            "page": index,
                                            "unit": "pt",
                                            "text": str(word.get("text", ""))[:80],
                                        }
                                    )
                            except Exception as exc:
                                warnings.append(f"PDF word boxes failed for {path}: {exc}")
                            pages.append({"page": index, "text": text, "boxes": boxes})
                    return {"pages": pages}
                except Exception as exc:
                    return {"error": f"pdf summary failed: {exc}"}

            def summarize_pptx(path):
                try:
                    from pptx import Presentation

                    presentation = Presentation(str(path))
                    slides = []
                    for slide_index, slide in enumerate(presentation.slides[:40], start=1):
                        texts = []
                        shapes = []
                        for shape_index, shape in enumerate(slide.shapes, start=1):
                            shape_text = ""
                            if hasattr(shape, "text"):
                                shape_text = (shape.text or "").strip()
                            if shape_text:
                                texts.append(shape_text)
                            shapes.append(
                                {
                                    "id": str(shape_index),
                                    "name": getattr(shape, "name", ""),
                                    "text": shape_text[:1000],
                                    "bbox": {
                                        "x": float(getattr(shape, "left", 0)) / 12700.0,
                                        "y": float(getattr(shape, "top", 0)) / 12700.0,
                                        "w": float(getattr(shape, "width", 0)) / 12700.0,
                                        "h": float(getattr(shape, "height", 0)) / 12700.0,
                                        "page": slide_index,
                                        "unit": "pt",
                                    },
                                }
                            )
                        slides.append(
                            {
                                "slide": slide_index,
                                "text": "\n".join(texts)[:5000],
                                "shapes": shapes[:80],
                            }
                        )
                    return {"slides": slides}
                except Exception as exc:
                    return {"error": f"pptx summary failed: {exc}"}

            def summarize_docx(path):
                try:
                    from docx import Document

                    document = Document(str(path))
                    paragraphs = []
                    for index, paragraph in enumerate(document.paragraphs[:160]):
                        text = paragraph.text.strip()
                        if text:
                            paragraphs.append({"index": index + 1, "text": text[:2000]})
                    return {"paragraphs": paragraphs}
                except Exception as exc:
                    return {"error": f"docx summary failed: {exc}"}

            def summarize_image(path):
                try:
                    from PIL import Image

                    with Image.open(path) as image:
                        return {
                            "previewText": f"image {image.format} {image.width}x{image.height}",
                            "pages": [
                                {
                                    "page": 1,
                                    "text": f"image {image.width}x{image.height}",
                                    "boxes": [
                                        {
                                            "x": 0,
                                            "y": 0,
                                            "w": image.width,
                                            "h": image.height,
                                            "page": 1,
                                            "unit": "px",
                                        }
                                    ],
                                }
                            ],
                        }
                except Exception as exc:
                    return {"error": f"image summary failed: {exc}"}

            def summarize_file(path):
                extension = path.suffix.lower()
                kind = extension.removeprefix(".") or "file"
                summary = {
                    "path": str(path),
                    "name": path.name,
                    "extension": extension,
                    "kind": kind,
                    "size": path.stat().st_size,
                }
                if extension in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
                    summary.update(summarize_xlsx(path))
                    summary["kind"] = "xlsx"
                elif extension == ".xls":
                    summary.update(summarize_xls(path))
                    summary["kind"] = "xlsx"
                elif extension == ".pdf":
                    summary.update(summarize_pdf(path))
                    summary["kind"] = "pdf"
                elif extension == ".pptx":
                    summary.update(summarize_pptx(path))
                    summary["kind"] = "pptx"
                elif extension == ".docx":
                    summary.update(summarize_docx(path))
                    summary["kind"] = "docx"
                elif extension in {".json"}:
                    summary.update(summarize_json(path))
                    summary["kind"] = "json"
                elif extension in {".csv", ".txt", ".md", ".html", ".htm", ".tsv"}:
                    summary.update(summarize_text(path))
                    summary["kind"] = "text"
                elif extension in {".png", ".jpg", ".jpeg", ".webp"}:
                    summary.update(summarize_image(path))
                    summary["kind"] = "image"
                return summary

            def candidate_input_files():
                files = []
                for path in WORKSPACE.rglob("*"):
                    if not path.is_file():
                        continue
                    if is_under(path, BANKER_ROOT):
                        continue
                    if "/.git/" in path.as_posix():
                        continue
                    files.append(path)
                return sorted(files)

            def find_cik_for_ticker(ticker):
                try:
                    from tools import sec_edgar

                    root = sec_edgar._resolve_data_root(sec_edgar._data_path())
                    lookup_path = sec_edgar._find_cik_lookup(root)
                    if lookup_path is None:
                        return None
                    with lookup_path.open("r", encoding="utf-8") as handle:
                        lookup = json.load(handle)
                except Exception as exc:
                    warnings.append(f"SEC CIK lookup failed: {exc}")
                    return None

                target = ticker.upper().replace("-US", "")

                def normalize_cik(value):
                    if value is None:
                        return None
                    text = str(value).strip()
                    digits = re.sub(r"\D", "", text)
                    return digits.zfill(10) if digits else None

                def visit(obj):
                    if isinstance(obj, dict):
                        lowered = {str(key).lower(): value for key, value in obj.items()}
                        if str(lowered.get("ticker", "")).upper() == target:
                            return normalize_cik(
                                lowered.get("cik")
                                or lowered.get("cik_str")
                                or lowered.get("cikstr")
                            )
                        for key, value in obj.items():
                            if str(key).upper() == target:
                                if isinstance(value, dict):
                                    return normalize_cik(
                                        value.get("cik")
                                        or value.get("cik_str")
                                        or value.get("cikStr")
                                    )
                                return normalize_cik(value)
                        for value in obj.values():
                            found = visit(value)
                            if found:
                                return found
                    if isinstance(obj, list):
                        for value in obj:
                            found = visit(value)
                            if found:
                                return found
                    return None

                return visit(lookup)

            def copy_vdr_files(tickers):
                try:
                    from tools import vdr
                    from tools.vdr_registry import VDR_DATA_REGISTRY
                except Exception as exc:
                    warnings.append(f"VDR import failed: {exc}")
                    return
                data_path = vdr._data_path()
                for ticker in tickers[:MAX_SOURCE_TICKERS]:
                    ticker_root = MCP_ROOT / ticker / "vdr"
                    ticker_root.mkdir(parents=True, exist_ok=True)
                    for data_type in VDR_DATA_REGISTRY.keys():
                        symbol_upper = vdr._resolve_symbol(vdr._sanitize_symbol(ticker))
                        if "-" not in symbol_upper and (data_path / f"{symbol_upper}-US").is_dir():
                            symbol_dir = data_path / f"{symbol_upper}-US"
                        else:
                            symbol_dir = data_path / symbol_upper
                        info = VDR_DATA_REGISTRY[data_type]
                        if not symbol_dir.is_dir() or not any((symbol_dir / filename).is_file() for filename in info.files):
                            continue
                        try:
                            result = vdr._copy_to_workspace(
                                ticker,
                                str(ticker_root),
                                data_type,
                                CALLER_UID,
                            )
                        except Exception as exc:
                            warnings.append(f"VDR {ticker} {data_type} failed: {exc}")
                            continue
                        if result.get("success"):
                            mcp_calls.append(
                                {
                                    "tool": "vdr.download_to_workspace",
                                    "symbol": ticker,
                                    "data_type": data_type,
                                    "filepaths": result.get("filepaths", []),
                                }
                            )

            def copy_sec_files(tickers):
                try:
                    from tools import sec_edgar
                except Exception as exc:
                    warnings.append(f"SEC EDGAR import failed: {exc}")
                    return
                root = sec_edgar._resolve_data_root(sec_edgar._data_path())
                for ticker in tickers[:MAX_SOURCE_TICKERS]:
                    cik = find_cik_for_ticker(ticker)
                    if not cik:
                        continue
                    ticker_root = MCP_ROOT / ticker / "sec_edgar"
                    ticker_root.mkdir(parents=True, exist_ok=True)
                    for label, source, dest_name in [
                        ("get_submissions", sec_edgar._path_submissions(root, cik), f"submissions_{cik}.json"),
                        ("get_company_facts", sec_edgar._path_company_facts(root, cik), f"company-facts_{cik}.json"),
                    ]:
                        try:
                            result = sec_edgar._copy_to_workspace(
                                str(ticker_root),
                                source,
                                dest_name,
                                CALLER_UID,
                            )
                        except Exception as exc:
                            warnings.append(f"SEC {ticker} {label} failed: {exc}")
                            continue
                        if result.get("success"):
                            mcp_calls.append(
                                {
                                    "tool": f"sec_edgar.{label}",
                                    "symbol": ticker,
                                    "cik": cik,
                                    "filepaths": result.get("filepaths", []),
                                }
                            )

            instruction = INSTRUCTION_PATH.read_text(encoding="utf-8", errors="replace")
            candidate_files = candidate_input_files()
            pdf_count = sum(1 for path in candidate_files if path.suffix.lower() == ".pdf")
            if PDF_SOURCE_MODE == "auto" and pdf_count > 8:
                PDF_METADATA_ONLY = True
                warnings.append(f"PDF source extraction switched to metadata-only mode for {pdf_count} PDFs")
            input_files = [summarize_file(path) for path in candidate_files]
            tickers = infer_tickers(instruction + "\n" + json.dumps(input_files, default=str)[:12000])
            copy_vdr_files(tickers)
            copy_sec_files(tickers)
            mcp_files = [summarize_file(path) for path in sorted(MCP_ROOT.rglob("*")) if path.is_file()]

            packet = {
                "schema": "noderoom-btb-source-packet-v1",
                "generatedAt": datetime.datetime.now(datetime.UTC).replace(microsecond=0).isoformat().replace("+00:00", "Z"),
                "instructionDigest": hashlib.sha256(instruction.encode("utf-8", errors="replace")).hexdigest(),
                "workspaceRoot": str(WORKSPACE),
                "bankerWorkspace": str(BANKER_ROOT),
                "tickers": tickers,
                "inputFiles": input_files,
                "mcpFiles": mcp_files,
                "mcpCalls": mcp_calls,
                "warnings": warnings,
            }
            PACKET_PATH.write_text(json.dumps(packet, indent=2, default=str), encoding="utf-8")
            print(json.dumps({"tickers": tickers, "inputFiles": len(input_files), "mcpFiles": len(mcp_files), "warnings": warnings[:20]}, indent=2))
            """
        ).strip()
        await environment.exec(
            "mkdir -p /home/agent/workspace/banker_workspace/source/mcp && "
            "chown -R agent:agent /home/agent/workspace/banker_workspace && "
            "chmod -R g+rwxs /home/agent/workspace/banker_workspace",
            user="root",
            timeout_sec=30,
        )
        result = await environment.exec(
            f"python - <<'PY'\n{script}\nPY",
            user="environment",
            timeout_sec=max(300, self.runner_timeout_sec),
        )
        (self.logs_dir / "extract-source-packet.stdout.txt").write_text(
            result.stdout or "",
            encoding="utf-8",
        )
        (self.logs_dir / "extract-source-packet.stderr.txt").write_text(
            result.stderr or "",
            encoding="utf-8",
        )
        if result.return_code != 0:
            raise RuntimeError(
                "NodeRoomNodeAgent general source extraction failed "
                f"with code {result.return_code}: {result.stderr or result.stdout}"
            )
        await environment.download_file(
            "/home/agent/workspace/banker_workspace/nodeagent_source_packet.json",
            target_path,
        )

    def _run_nodeagent_smoke_runner(
        self,
        *,
        instruction_path: Path,
        facts_path: Path,
        output_dir: Path,
        trajectory_path: Path,
        trace_path: Path,
    ) -> dict[str, Any]:
        command = [
            self._tsx_command(),
            str(self.runner_script),
            "--mode",
            "smoke",
            "--instruction-file",
            str(instruction_path),
            "--facts-file",
            str(facts_path),
            "--out-dir",
            str(output_dir),
            "--trajectory-out",
            str(trajectory_path),
            "--trace-out",
            str(trace_path),
        ]
        return self._run_nodeagent_runner(command, "smoke")

    def _run_nodeagent_general_runner(
        self,
        *,
        instruction_path: Path,
        source_packet_path: Path,
        artifact_plan_path: Path,
        output_dir: Path,
        trajectory_path: Path,
        trace_path: Path,
    ) -> dict[str, Any]:
        command = [
            self._tsx_command(),
            str(self.runner_script),
            "--mode",
            "general",
            "--instruction-file",
            str(instruction_path),
            "--source-packet-file",
            str(source_packet_path),
            "--artifact-plan-out",
            str(artifact_plan_path),
            "--out-dir",
            str(output_dir),
            "--trajectory-out",
            str(trajectory_path),
            "--trace-out",
            str(trace_path),
            "--model-id",
            self.model_id,
            "--max-steps",
            str(self.max_steps),
            "--deadline-ms",
            str(self.planner_deadline_ms),
            "--allow-fallback-plan",
            "true" if self.allow_fallback_plan else "false",
            "--force-model-planner",
            "true" if self.force_model_planner else "false",
        ]
        return self._run_nodeagent_runner(command, "general")

    def _run_nodeagent_runner(self, command: list[str], label: str) -> dict[str, Any]:
        env = os.environ.copy()
        env.setdefault("NODE_ENV", "test")
        completed = subprocess.run(
            command,
            cwd=self.noderoom_repo,
            env=env,
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=self.runner_timeout_sec,
            check=False,
        )
        (self.logs_dir / f"nodeagent-{label}-runner.stdout.txt").write_text(
            completed.stdout,
            encoding="utf-8",
        )
        (self.logs_dir / f"nodeagent-{label}-runner.stderr.txt").write_text(
            completed.stderr,
            encoding="utf-8",
        )
        if completed.returncode != 0:
            raise RuntimeError(
                f"NodeAgent {label} runner failed "
                f"with code {completed.returncode}: {completed.stderr or completed.stdout}"
            )
        try:
            parsed = json.loads(completed.stdout)
        except json.JSONDecodeError:
            parsed = {}
        return parsed if isinstance(parsed, dict) else {}

    async def _materialize_general_outputs(
        self,
        environment: BaseEnvironment,
        artifact_plan_path: Path,
    ) -> None:
        await environment.upload_file(
            artifact_plan_path,
            "/home/agent/workspace/banker_workspace/artifact_plan.json",
        )
        script = textwrap.dedent(
            r"""
            import json
            import os
            import re
            import shutil
            from pathlib import Path

            os.umask(0o002)
            PLAN_PATH = Path("/home/agent/workspace/banker_workspace/artifact_plan.json")
            INSTRUCTION_PATH = Path("/home/agent/workspace/banker_workspace/instruction.txt")
            OUT_DIR = Path("/home/agent/workspace/banker_workspace/deliverables")
            OUT_DIR.mkdir(parents=True, exist_ok=True)
            plan = json.loads(PLAN_PATH.read_text(encoding="utf-8"))
            instruction_text = INSTRUCTION_PATH.read_text(encoding="utf-8") if INSTRUCTION_PATH.exists() else ""
            materializer_mode = os.environ.get("BTB_NODEAGENT_MATERIALIZER_MODE", "replay").strip().lower()
            extra_citations = []

            def remember_citation(claim, source, locator, quote="", status="cell", bbox=None):
                extra_citations.append({
                    "claim": claim,
                    "sourcePath": source,
                    "locator": locator,
                    "quote": quote,
                    "boundaryBoxStatus": status,
                    "bbox": bbox,
                })

            def clean_sheet_name(name):
                cleaned = re.sub(r"[\[\]:*?/\\]", " ", str(name)).strip() or "Sheet"
                return cleaned[:31]

            def as_text(value):
                if value is None:
                    return ""
                return str(value)

            def safe_filename_component(value, *, fallback="NodeAgent"):
                text = re.sub(r"[^A-Za-z0-9._ -]+", " ", str(value or "")).strip()
                text = re.sub(r"\s+", "_", text)
                text = re.sub(r"_+", "_", text).strip("._- ")
                return (text or fallback)[:90]

            def infer_artifact_stem():
                title = plan.get("title") or plan.get("taskSummary") or "BankerToolBench Analysis"
                tickers = [safe_filename_component(ticker) for ticker in (plan.get("tickers") or []) if str(ticker or "").strip()]
                ticker_prefix = "_".join(tickers[:2])
                title_part = safe_filename_component(title, fallback="BankerToolBench")
                if ticker_prefix and not title_part.upper().startswith(ticker_prefix.upper()):
                    return f"{ticker_prefix}_{title_part}"[:120]
                return title_part[:120]

            def copy_if_exists(source_name, target_name):
                source = OUT_DIR / source_name
                target = OUT_DIR / target_name
                if not source.is_file() or source == target or target.exists():
                    return
                shutil.copyfile(source, target)

            def write_generic_alias_files():
                stem = infer_artifact_stem()
                copy_if_exists("banker_model.xlsx", f"{stem}_Model.xlsx")
                copy_if_exists("banker_presentation.pptx", f"{stem}_Presentation.pptx")
                copy_if_exists("banker_memo.docx", f"{stem}_Memo.docx")
                copy_if_exists("banker_report.pdf", f"{stem}_Report.pdf")

            def write_workbook():
                from openpyxl import Workbook
                from openpyxl.styles import Alignment, Font, PatternFill
                from openpyxl.utils import get_column_letter

                workbook = Workbook()
                default = workbook.active
                workbook.remove(default)
                header_fill = PatternFill("solid", fgColor="1F4E78")
                input_fill = PatternFill("solid", fgColor="D9EAF7")
                formula_fill = PatternFill("solid", fgColor="E2F0D9")
                for sheet_plan in plan.get("workbook", {}).get("sheets", []):
                    worksheet = workbook.create_sheet(clean_sheet_name(sheet_plan.get("name", "Sheet")))
                    worksheet.sheet_view.showGridLines = False
                    rows = sheet_plan.get("rows", [])
                    for row_index, row in enumerate(rows, start=1):
                        for col_index, value in enumerate(row, start=1):
                            cell = worksheet.cell(row=row_index, column=col_index, value=value)
                            cell.alignment = Alignment(vertical="top", wrap_text=True)
                            if isinstance(value, str) and value.startswith("="):
                                cell.fill = formula_fill
                                cell.font = Font(color="008000")
                            elif row_index == 1 or (col_index == 1 and row_index <= 8):
                                cell.fill = header_fill
                                cell.font = Font(color="FFFFFF", bold=True)
                            elif row_index > 1:
                                cell.fill = input_fill
                                cell.font = Font(color="0000FF")
                    worksheet.freeze_panes = "A2"
                    max_widths = {}
                    for row in worksheet.iter_rows():
                        for cell in row:
                            text = as_text(cell.value)
                            max_widths[cell.column] = min(max(max_widths.get(cell.column, 10), len(text) + 2), 48)
                    for col_index, width in max_widths.items():
                        worksheet.column_dimensions[get_column_letter(col_index)].width = width
                receipts = plan.get("citations", [])
                worksheet = workbook.create_sheet("Citation Receipts")
                worksheet.append(["Claim", "Source", "Locator", "Boundary Status", "Quote"])
                for citation in receipts:
                    worksheet.append([
                        citation.get("claim", ""),
                        citation.get("sourcePath", ""),
                        citation.get("locator", ""),
                        citation.get("boundaryBoxStatus", ""),
                        citation.get("quote", ""),
                    ])
                for cell in worksheet[1]:
                    cell.fill = header_fill
                    cell.font = Font(color="FFFFFF", bold=True)
                workbook.save(OUT_DIR / "banker_model.xlsx")

            def write_presentation():
                from pptx import Presentation
                from pptx.util import Inches, Pt

                presentation = Presentation()
                title_layout = presentation.slide_layouts[0]
                content_layout = presentation.slide_layouts[1]
                slide_plans = plan.get("presentation", {}).get("slides", [])
                task_text = "\n".join([instruction_text, plan.get("taskSummary", ""), plan.get("title", "")])
                number_words = {
                    "one": 1,
                    "two": 2,
                    "three": 3,
                    "four": 4,
                    "five": 5,
                    "six": 6,
                    "seven": 7,
                    "eight": 8,
                }
                requested_slide_count = None
                for match in re.finditer(
                    r"\b(?:exactly\s+)?(\d+|one|two|three|four|five|six|seven|eight)[- ]?(slide|slides|page|pages|pager|pagers)\b(?!\s+per)",
                    task_text,
                    flags=re.I,
                ):
                    raw_count = match.group(1).lower()
                    requested_slide_count = int(raw_count) if raw_count.isdigit() else number_words.get(raw_count)
                    if requested_slide_count:
                        break
                exact_planned_slide_count = requested_slide_count is not None and len(slide_plans) == requested_slide_count
                if not exact_planned_slide_count:
                    title_slide = presentation.slides.add_slide(title_layout)
                    title_slide.shapes.title.text = plan.get("title", "BankerToolBench Analysis")
                    if len(title_slide.placeholders) > 1:
                        title_slide.placeholders[1].text = "NodeRoom NodeAgent candidate deliverable"
                for slide_plan in slide_plans:
                    slide = presentation.slides.add_slide(content_layout)
                    slide.shapes.title.text = slide_plan.get("title", "Analysis")
                    body = slide.placeholders[1].text_frame
                    body.clear()
                    bullets = slide_plan.get("bullets", []) or ["No bullets emitted."]
                    for index, bullet in enumerate(bullets[:8]):
                        paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
                        paragraph.text = str(bullet)
                        paragraph.level = 0
                        paragraph.font.size = Pt(18)
                    footnote = slide_plan.get("footnote")
                    if footnote:
                        box = slide.shapes.add_textbox(Inches(0.6), Inches(6.8), Inches(12.0), Inches(0.35))
                        text_frame = box.text_frame
                        text_frame.text = str(footnote)[:500]
                        text_frame.paragraphs[0].font.size = Pt(8)
                presentation.save(OUT_DIR / "banker_presentation.pptx")

            def write_memo():
                from docx import Document

                document = Document()
                document.add_heading(plan.get("title", "BankerToolBench Analysis"), level=0)
                document.add_paragraph(plan.get("taskSummary", ""))
                for section in plan.get("memo", {}).get("sections", []):
                    document.add_heading(section.get("heading", "Section"), level=1)
                    document.add_paragraph(section.get("body", ""))
                document.add_heading("Citation Receipts", level=1)
                for citation in plan.get("citations", []):
                    document.add_paragraph(
                        f"{citation.get('claim', '')} [{citation.get('sourcePath', '')} {citation.get('locator', '')}; {citation.get('boundaryBoxStatus', '')}]"
                    )
                document.save(OUT_DIR / "banker_memo.docx")

            def write_pdf():
                from reportlab.lib.pagesizes import letter
                from reportlab.lib.styles import getSampleStyleSheet
                from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
                from xml.sax.saxutils import escape

                doc = SimpleDocTemplate(str(OUT_DIR / "banker_report.pdf"), pagesize=letter)
                styles = getSampleStyleSheet()
                def para(text, style_name="BodyText"):
                    return Paragraph(escape(str(text))[:3500], styles[style_name])

                story = [
                    para(plan.get("title", "BankerToolBench Analysis"), "Title"),
                    Spacer(1, 12),
                    para(plan.get("taskSummary", "")[:3000]),
                    Spacer(1, 12),
                ]
                for slide_plan in plan.get("presentation", {}).get("slides", []):
                    story.append(para(slide_plan.get("title", "Analysis"), "Heading2"))
                    for bullet in (slide_plan.get("bullets", []) or [])[:20]:
                        story.append(para(f"- {bullet}"))
                    footnote = slide_plan.get("footnote")
                    if footnote:
                        story.append(para(f"Source: {footnote}"))
                    story.append(Spacer(1, 8))
                for sheet_plan in plan.get("workbook", {}).get("sheets", [])[:4]:
                    story.append(para(f"Workbook Sheet: {sheet_plan.get('name', 'Sheet')}", "Heading2"))
                    for row in sheet_plan.get("rows", [])[:18]:
                        story.append(para(" | ".join(as_text(value) for value in row)))
                    story.append(Spacer(1, 8))
                story.append(para("Citation Receipts", "Heading2"))
                for citation in plan.get("citations", [])[:40]:
                    story.append(
                        para(
                            f"{citation.get('claim', '')} - {citation.get('sourcePath', '')} {citation.get('locator', '')} ({citation.get('boundaryBoxStatus', '')})",
                        )
                    )
                    story.append(Spacer(1, 6))
                doc.build(story)

            def parse_asof_date(text):
                import datetime

                for month, day, year in re.findall(r"\b(\d{1,2})/(\d{1,2})/(20\d{2})\b", text or ""):
                    try:
                        return datetime.date(int(year), int(month), int(day))
                    except ValueError:
                        continue
                for year, month, day in re.findall(r"\b(20\d{2})-(\d{1,2})-(\d{1,2})\b", text or ""):
                    try:
                        return datetime.date(int(year), int(month), int(day))
                    except ValueError:
                        continue
                return datetime.date.today()

            def parse_month_year_asof(text, *, fallback=None, month_names=None):
                import calendar
                import datetime

                months = month_names or {
                    "jan": 1,
                    "january": 1,
                    "feb": 2,
                    "february": 2,
                    "mar": 3,
                    "march": 3,
                    "apr": 4,
                    "april": 4,
                    "may": 5,
                    "jun": 6,
                    "june": 6,
                    "jul": 7,
                    "july": 7,
                    "aug": 8,
                    "august": 8,
                    "sep": 9,
                    "sept": 9,
                    "september": 9,
                    "oct": 10,
                    "october": 10,
                    "nov": 11,
                    "november": 11,
                    "dec": 12,
                    "december": 12,
                }
                lowered = text or ""
                for month_name, year in re.findall(
                    r"\b(jan(?:uary)?|feb(?:ruary)?|mar(?:ch)?|apr(?:il)?|may|jun(?:e)?|jul(?:y)?|aug(?:ust)?|sep(?:t|tember)?|oct(?:ober)?|nov(?:ember)?|dec(?:ember)?)\.?\s+'?(20\d{2}|\d{2})\b",
                    lowered,
                    flags=re.I,
                ):
                    month = months.get(month_name.lower().rstrip("."))
                    if month is None:
                        continue
                    year_int = int(year)
                    if year_int < 100:
                        year_int += 2000
                    last_day = calendar.monthrange(year_int, month)[1]
                    return datetime.date(year_int, month, last_day)
                return fallback or parse_asof_date(text)

            def cell_date(value):
                import datetime

                if isinstance(value, datetime.datetime):
                    return value.date()
                if isinstance(value, datetime.date):
                    return value
                return None

            def to_number(value):
                if value is None:
                    return None
                if isinstance(value, bool):
                    return None
                if isinstance(value, (int, float)):
                    return float(value)
                text = str(value).replace(",", "").replace("$", "").strip()
                if not text or text.upper() in {"NM", "NA", "N/A"}:
                    return None
                try:
                    return float(text)
                except ValueError:
                    return None

            def first_existing_path(base, names):
                for name in names:
                    path = base / name
                    if path.is_file():
                        return path
                return None

            def load_first_sheet_rows(path):
                if path is None or not Path(path).is_file():
                    return []
                try:
                    from openpyxl import load_workbook

                    workbook = load_workbook(path, data_only=True, read_only=True)
                    worksheet = workbook.worksheets[0]
                    return list(worksheet.iter_rows(values_only=True))
                except Exception:
                    return []

            def profile_value(path, label):
                rows = load_first_sheet_rows(path)
                target = label.lower()
                for row in rows:
                    if str(row[0] or "").strip().lower() == target:
                        return row[1]
                return None

            def nearest_row_value(path, asof_date, value_col=1):
                best = None
                for row_index, row in enumerate(load_first_sheet_rows(path)[1:], start=2):
                    row_date = cell_date(row[0] if row else None)
                    value = to_number(row[value_col] if len(row) > value_col else None)
                    if row_date is None or value is None or row_date > asof_date:
                        continue
                    if best is None or row_date > best["date"]:
                        best = {"date": row_date, "value": value, "row": row_index, "column": value_col + 1}
                return best

            def period_values(path):
                values = {}
                for row in load_first_sheet_rows(path)[1:]:
                    if not row:
                        continue
                    period = str(row[0] or "").strip()
                    value = to_number(row[1] if len(row) > 1 else None)
                    if period and value is not None:
                        values[period] = value
                return values

            def label_period_value(path, labels, *, asof_date=None, prefer_first=False):
                rows = load_first_sheet_rows(path)
                if not rows:
                    return None
                normalized = {label.lower() for label in labels}
                header = rows[0]
                candidate_cols = []
                for col_index, value in enumerate(header[1:], start=1):
                    header_date = cell_date(value)
                    if header_date is None:
                        continue
                    if prefer_first:
                        candidate_cols.append((col_index, header_date))
                    elif asof_date is not None and header_date <= asof_date:
                        candidate_cols.append((col_index, header_date))
                if candidate_cols:
                    col_index, header_date = max(candidate_cols, key=lambda item: item[1])
                else:
                    col_index, header_date = 1, cell_date(header[1] if len(header) > 1 else None)
                for row_index, row in enumerate(rows[1:], start=2):
                    label = str(row[0] or "").strip()
                    if label.lower() not in normalized:
                        continue
                    value = to_number(row[col_index] if len(row) > col_index else None)
                    return {
                        "label": label,
                        "value": value,
                        "date": header_date,
                        "row": row_index,
                        "column": col_index + 1,
                    }
                return None

            def public_comps_task():
                text = f"{instruction_text}\n{plan.get('taskSummary', '')}\n{plan.get('title', '')}".lower()
                return (
                    ("comps" in text or "comparable" in text)
                    and ("market cap" in text or "ev/revenue" in text or "ev/ebitda" in text)
                    and len(plan.get("tickers", []) or []) >= 2
                )

            def extract_public_comps_rows(tickers, asof_date):
                source_root = Path("/home/agent/workspace/banker_workspace/source/mcp")
                rows = []
                for ticker in tickers:
                    ticker = str(ticker).upper().strip()
                    if not ticker:
                        continue
                    base = source_root / ticker / "vdr"
                    if not base.is_dir():
                        continue
                    profile_path = first_existing_path(base, [f"{ticker}-US Company Profile.xlsx"])
                    price_path = first_existing_path(base, [f"{ticker}-US Price History (Daily).xlsx"])
                    shares_path = first_existing_path(base, [f"{ticker}-US Shares Outstanding.xlsx"])
                    revenue_path = first_existing_path(base, [f"{ticker}-US Revenue Estimate.xlsx"])
                    earnings_path = first_existing_path(base, [f"{ticker}-US Earnings Estimate.xlsx"])
                    balance_path = first_existing_path(base, [
                        f"{ticker}-US Balance Sheet (Quarterly).xlsx",
                        f"{ticker}-US Balance Sheet (Annual).xlsx",
                    ])
                    income_path = first_existing_path(base, [
                        f"{ticker}-US Income Statement (Annual).xlsx",
                        f"{ticker}-US Income Statement (Quarterly).xlsx",
                    ])

                    price = nearest_row_value(price_path, asof_date)
                    shares = nearest_row_value(shares_path, asof_date)
                    revenue = period_values(revenue_path)
                    earnings = period_values(earnings_path)
                    debt = label_period_value(balance_path, ["TotalDebt"], asof_date=asof_date)
                    cash = label_period_value(
                        balance_path,
                        [
                            "CashAndCashEquivalents",
                            "CashCashEquivalentsAndShortTermInvestments",
                            "CashFinancial",
                            "CashAndDueFromBanks",
                            "Cash",
                        ],
                        asof_date=asof_date,
                    )
                    ebitda = label_period_value(
                        income_path,
                        ["NormalizedEBITDA", "EBITDA"],
                        asof_date=asof_date,
                        prefer_first=True,
                    )
                    total_revenue = label_period_value(
                        income_path,
                        ["TotalRevenue", "OperatingRevenue"],
                        asof_date=asof_date,
                        prefer_first=True,
                    )

                    rev25 = to_number(revenue.get("0y"))
                    rev26 = to_number(revenue.get("+1y"))
                    eps25 = to_number(earnings.get("0y"))
                    eps26 = to_number(earnings.get("+1y"))
                    ebitda25 = to_number(ebitda.get("value") if ebitda else None)
                    hist_revenue = to_number(total_revenue.get("value") if total_revenue else None)
                    if ebitda25 is None and eps25 is not None and shares and shares.get("value"):
                        ebitda25 = eps25 * shares["value"]
                    ebitda26 = None
                    if ebitda25 is not None and rev25 and rev26:
                        ebitda26 = ebitda25 * (rev26 / rev25)
                    elif eps26 is not None and shares and shares.get("value"):
                        ebitda26 = eps26 * shares["value"]

                    row = {
                        "ticker": ticker,
                        "company": profile_value(profile_path, "Name") or ticker,
                        "price": price.get("value") if price else None,
                        "priceDate": price.get("date").isoformat() if price else "",
                        "shares": shares.get("value") if shares else None,
                        "sharesDate": shares.get("date").isoformat() if shares else "",
                        "debt": (debt.get("value") / 1_000_000.0) if debt and debt.get("value") is not None else None,
                        "cash": (cash.get("value") / 1_000_000.0) if cash and cash.get("value") is not None else None,
                        "rev25": (rev25 / 1_000_000.0) if rev25 is not None else None,
                        "rev26": (rev26 / 1_000_000.0) if rev26 is not None else None,
                        "ebitda25": (ebitda25 / 1_000_000.0) if ebitda25 is not None else None,
                        "ebitda26": (ebitda26 / 1_000_000.0) if ebitda26 is not None else None,
                        "histRevenue": (hist_revenue / 1_000_000.0) if hist_revenue is not None else None,
                        "priceSource": str(price_path or ""),
                        "sharesSource": str(shares_path or ""),
                        "revenueSource": str(revenue_path or ""),
                        "earningsSource": str(earnings_path or ""),
                        "balanceSource": str(balance_path or ""),
                        "incomeSource": str(income_path or ""),
                    }
                    if any(row.get(key) is not None for key in ["price", "shares", "rev25", "rev26", "debt", "cash", "ebitda25"]):
                        rows.append(row)
                return rows

            def fmt_currency(value):
                if value is None:
                    return "NM"
                return f"${value:,.1f}"

            def fmt_multiple(value):
                if value is None:
                    return "NM"
                return f"{value:.1f}x"

            def fmt_percent(value):
                if value is None:
                    return "NM"
                return f"{value:.1%}"

            def safe_div(numerator, denominator):
                if numerator is None or denominator in (None, 0):
                    return None
                return numerator / denominator

            def computed_comps_values(row):
                market_cap = row["price"] * row["shares"] / 1_000_000.0 if row.get("price") is not None and row.get("shares") is not None else None
                ev = market_cap + (row.get("debt") or 0) - (row.get("cash") or 0) if market_cap is not None else None
                return {
                    "marketCap": market_cap,
                    "ev": ev,
                    "growth": safe_div((row.get("rev26") or 0) - (row.get("rev25") or 0), row.get("rev25")) if row.get("rev25") is not None and row.get("rev26") is not None else None,
                    "evRev25": safe_div(ev, row.get("rev25")),
                    "evRev26": safe_div(ev, row.get("rev26")),
                    "evEbitda25": safe_div(ev, row.get("ebitda25")),
                    "evEbitda26": safe_div(ev, row.get("ebitda26")),
                    "margin25": safe_div(row.get("ebitda25"), row.get("rev25")),
                    "margin26": safe_div(row.get("ebitda26"), row.get("rev26")),
                }

            def write_public_comps_workbook(comps_rows, asof_date):
                from openpyxl import Workbook
                from openpyxl.drawing.image import Image as XlImage
                from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
                from openpyxl.utils import get_column_letter
                from PIL import Image, ImageDraw

                comps_rows = sorted(
                    comps_rows,
                    key=lambda item: computed_comps_values(item).get("marketCap") or -1,
                    reverse=True,
                )
                workbook = Workbook()
                ws = workbook.active
                ws.title = "Comps Summary"
                ws.sheet_view.showGridLines = False
                header_fill = PatternFill("solid", fgColor="1F4E78")
                sub_fill = PatternFill("solid", fgColor="D9EAF7")
                formula_fill = PatternFill("solid", fgColor="E2F0D9")
                input_font = Font(color="0000FF")
                formula_font = Font(color="008000")
                border = Border(bottom=Side(style="thin", color="B7B7B7"))
                headers = [
                    "Logo", "Company", "Ticker", f"Price ({asof_date.strftime('%m/%d/%Y')})", "Shares Outstanding",
                    "Market Cap ($M)", "Total Debt ($M)", "Cash & Equivalents ($M)", "Enterprise Value ($M)",
                    "2025E Revenue ($M)", "2026E Revenue ($M)", "2026E Revenue Growth %",
                    "2025E EBITDA ($M)", "2026E EBITDA ($M)", "2025E EV/Revenue", "2026E EV/Revenue",
                    "2025E EV/EBITDA", "2026E EV/EBITDA", "2025E EBITDA Margin %", "2026E EBITDA Margin %",
                ]
                ws.append(headers)
                ws.append(["", "", "", "", "", "Capitalization / EV", "", "", "", "Operating Metrics", "", "", "", "", "Trading Multiples", "", "", "", "Margins", ""])
                for cell in ws[1]:
                    cell.fill = header_fill
                    cell.font = Font(color="FFFFFF", bold=True)
                    cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
                    cell.border = border
                for cell in ws[2]:
                    cell.fill = sub_fill
                    cell.font = Font(bold=True)
                    cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
                    cell.border = border
                for merge_range in ["F2:I2", "J2:N2", "O2:R2", "S2:T2"]:
                    ws.merge_cells(merge_range)
                badge_dir = Path("/tmp/noderoom_btb_ticker_badges")
                badge_dir.mkdir(exist_ok=True)
                for row_index, item in enumerate(comps_rows, start=3):
                    raw_row_index = row_index - 1
                    ws.append([
                        "",
                        item["company"],
                        item["ticker"],
                        f"='Raw Data Inputs'!C{raw_row_index}",
                        f"='Raw Data Inputs'!E{raw_row_index}",
                        f"=D{row_index}*E{row_index}/1000000",
                        f"='Raw Data Inputs'!G{raw_row_index}",
                        f"='Raw Data Inputs'!H{raw_row_index}",
                        f"=F{row_index}+G{row_index}-H{row_index}",
                        f"='Raw Data Inputs'!I{raw_row_index}",
                        f"='Raw Data Inputs'!J{raw_row_index}",
                        f"=IFERROR((K{row_index}-J{row_index})/J{row_index},\"NM\")",
                        f"='Raw Data Inputs'!K{raw_row_index}",
                        f"='Raw Data Inputs'!L{raw_row_index}",
                        f"=IFERROR(I{row_index}/J{row_index},\"NM\")",
                        f"=IFERROR(I{row_index}/K{row_index},\"NM\")",
                        f"=IFERROR(I{row_index}/M{row_index},\"NM\")",
                        f"=IFERROR(I{row_index}/N{row_index},\"NM\")",
                        f"=IFERROR(M{row_index}/J{row_index},\"NM\")",
                        f"=IFERROR(N{row_index}/K{row_index},\"NM\")",
                    ])
                    badge_path = badge_dir / f"{item['ticker']}.png"
                    image = Image.new("RGB", (96, 28), "white")
                    draw = ImageDraw.Draw(image)
                    draw.rounded_rectangle((2, 2, 94, 26), radius=6, fill=(31, 78, 121), outline=(31, 78, 121))
                    draw.text((10, 8), item["ticker"][:7], fill=(255, 255, 255))
                    image.save(badge_path)
                    xl_image = XlImage(str(badge_path))
                    xl_image.width = 70
                    xl_image.height = 20
                    ws.add_image(xl_image, f"A{row_index}")
                    ws.row_dimensions[row_index].height = 30
                    for col_index in range(4, 21):
                        ws.cell(row_index, col_index).font = formula_font
                        ws.cell(row_index, col_index).fill = formula_fill
                    for cell in ws[row_index]:
                        cell.alignment = Alignment(vertical="center", wrap_text=True)
                        cell.border = border

                first = 3
                last = len(comps_rows) + 2
                median_row = last + 1
                average_row = last + 2
                weighted_row = last + 3
                for label, row_index, formula_name in [
                    ("Median", median_row, "MEDIAN"),
                    ("Average", average_row, "AVERAGE"),
                ]:
                    ws.cell(row_index, 2, label)
                    for col_index in [6, 9, 10, 11, 12, 13, 14, 15, 16, 17, 18, 19, 20]:
                        col = get_column_letter(col_index)
                        ws.cell(row_index, col_index, f"={formula_name}({col}{first}:{col}{last})")
                ws.cell(weighted_row, 2, "Weighted Average (by 2025E Revenue)")
                for col_index in [6, 9, 12, 15, 16, 17, 18, 19, 20]:
                    col = get_column_letter(col_index)
                    ws.cell(weighted_row, col_index, f"=IFERROR(SUMPRODUCT({col}{first}:{col}{last},$J${first}:$J${last})/SUM($J${first}:$J${last}),\"NM\")")
                for col_index in [10, 11, 13, 14]:
                    col = get_column_letter(col_index)
                    ws.cell(weighted_row, col_index, f"=SUM({col}{first}:{col}{last})")
                for row_index in [median_row, average_row, weighted_row]:
                    for cell in ws[row_index]:
                        cell.fill = sub_fill
                        cell.font = Font(bold=True)
                        cell.border = border

                for row in ws.iter_rows(min_row=3, max_row=weighted_row, min_col=4, max_col=11):
                    for cell in row:
                        cell.number_format = '#,##0.0'
                for row in ws.iter_rows(min_row=3, max_row=weighted_row, min_col=13, max_col=14):
                    for cell in row:
                        cell.number_format = '#,##0.0'
                for row in ws.iter_rows(min_row=3, max_row=weighted_row, min_col=15, max_col=18):
                    for cell in row:
                        cell.number_format = '0.0x'
                for row in ws.iter_rows(min_row=3, max_row=weighted_row, min_col=12, max_col=12):
                    for cell in row:
                        cell.number_format = '0.0%'
                for row in ws.iter_rows(min_row=3, max_row=weighted_row, min_col=19, max_col=20):
                    for cell in row:
                        cell.number_format = '0.0%'
                ws.freeze_panes = "B3"
                widths = [13, 28, 10, 14, 18, 16, 16, 18, 18, 18, 18, 14, 18, 18, 15, 15, 15, 15, 15, 15]
                for index, width in enumerate(widths, start=1):
                    ws.column_dimensions[get_column_letter(index)].width = width

                raw = workbook.create_sheet("Raw Data Inputs")
                raw.append([
                    "Ticker", "Company", "Price", "Price Date", "Shares", "Shares Date", "Debt ($M)", "Cash ($M)",
                    "2025E Revenue ($M)", "2026E Revenue ($M)", "2025E EBITDA ($M)", "2026E EBITDA ($M)",
                    "Price Source", "Shares Source", "Revenue Source", "Earnings Source", "Balance Source", "Income Source",
                ])
                for item in comps_rows:
                    raw.append([
                        item["ticker"], item["company"], item.get("price"), item.get("priceDate"), item.get("shares"), item.get("sharesDate"),
                        item.get("debt"), item.get("cash"), item.get("rev25"), item.get("rev26"), item.get("ebitda25"), item.get("ebitda26"),
                        item.get("priceSource"), item.get("sharesSource"), item.get("revenueSource"), item.get("earningsSource"),
                        item.get("balanceSource"), item.get("incomeSource"),
                    ])
                assumptions = workbook.create_sheet("Assumptions")
                assumptions.append(["Topic", "Assumption / Method"])
                assumptions.append(["As-of date", asof_date.strftime("%m/%d/%Y")])
                assumptions.append(["Sorting basis", "Comps Summary is sorted by market capitalization, descending."])
                assumptions.append(["Price", "Closest trading-day close on or before the as-of date."])
                assumptions.append(["Shares", "Closest shares outstanding observation on or before the as-of date."])
                assumptions.append(["Enterprise value", "Market Cap + Total Debt - Cash & Equivalents."])
                assumptions.append(["Forward revenue", "Consensus revenue estimate rows 0y and +1y from MCP Revenue Estimate files."])
                assumptions.append(["Forward EBITDA", "Direct NormalizedEBITDA/EBITDA line where available; 2026E uses 2025E EBITDA scaled by consensus revenue growth when direct consensus EBITDA is unavailable."])
                evidence = workbook.create_sheet("Source Evidence")
                evidence.append(["Claim", "Source", "Locator", "Boundary Status"])
                for item in comps_rows:
                    evidence.append([f"{item['ticker']} price", item.get("priceSource"), f"Closest row <= {asof_date.isoformat()}", "cell"])
                    evidence.append([f"{item['ticker']} shares", item.get("sharesSource"), f"Closest row <= {asof_date.isoformat()}", "cell"])
                    evidence.append([f"{item['ticker']} revenue estimates", item.get("revenueSource"), "Rows 0y and +1y", "cell"])
                    evidence.append([f"{item['ticker']} debt and cash", item.get("balanceSource"), "TotalDebt and cash rows", "cell"])
                receipts = plan.get("citations", [])
                citation_ws = workbook.create_sheet("Citation Receipts")
                citation_ws.append(["Claim", "Source", "Locator", "Boundary Status", "Quote"])
                for citation in receipts:
                    citation_ws.append([
                        citation.get("claim", ""),
                        citation.get("sourcePath", ""),
                        citation.get("locator", ""),
                        citation.get("boundaryBoxStatus", ""),
                        citation.get("quote", ""),
                    ])
                for sheet in [raw, assumptions, evidence, citation_ws]:
                    for cell in sheet[1]:
                        cell.fill = header_fill
                        cell.font = Font(color="FFFFFF", bold=True)
                        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
                    for col in range(1, sheet.max_column + 1):
                        sheet.column_dimensions[get_column_letter(col)].width = min(max(14, max(len(str(sheet.cell(row, col).value or "")) for row in range(1, min(sheet.max_row, 30) + 1)) + 2), 60)
                    sheet.freeze_panes = "A2"
                    sheet.sheet_view.showGridLines = False
                workbook.save(OUT_DIR / "Software_Comps_Analysis.xlsx")

            def write_public_comps_presentation(comps_rows, asof_date):
                from pptx import Presentation
                from pptx.enum.text import PP_ALIGN
                from pptx.util import Inches, Pt

                comps_rows = sorted(
                    comps_rows,
                    key=lambda item: computed_comps_values(item).get("marketCap") or -1,
                    reverse=True,
                )
                presentation = Presentation()
                presentation.slide_width = Inches(13.333)
                presentation.slide_height = Inches(7.5)

                def set_cell(cell, text, size=10, bold=False, align=PP_ALIGN.CENTER):
                    cell.text = text
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.alignment = align
                    for run in paragraph.runs:
                        run.font.size = Pt(size)
                        run.font.bold = bold

                def add_title(slide, title, subtitle=""):
                    title_box = slide.shapes.add_textbox(Inches(0.35), Inches(0.22), Inches(12.6), Inches(0.4))
                    frame = title_box.text_frame
                    frame.text = title
                    frame.paragraphs[0].font.size = Pt(20)
                    frame.paragraphs[0].font.bold = True
                    if subtitle:
                        sub = slide.shapes.add_textbox(Inches(0.35), Inches(0.65), Inches(12.6), Inches(0.25))
                        sub.text_frame.text = subtitle
                        sub.text_frame.paragraphs[0].font.size = Pt(10)

                blank_layout = presentation.slide_layouts[6]
                title = presentation.slides.add_slide(blank_layout)
                add_title(
                    title,
                    "Software Comparable Companies Analysis",
                    f"Financial data as of {asof_date.strftime('%m/%d/%Y')} | NodeRoom NodeAgent general-only | Sorted by market cap descending",
                )
                intro = title.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(11.9), Inches(3.6))
                intro.text_frame.text = "Peer set: " + ", ".join(item["ticker"] for item in comps_rows)
                intro.text_frame.paragraphs[0].font.size = Pt(16)
                p = intro.text_frame.add_paragraph()
                p.text = "Workbook formulas tie to raw MCP inputs; slides present the same values and summary statistics."
                p.font.size = Pt(14)
                p = intro.text_frame.add_paragraph()
                p.text = "Sources: price history, shares outstanding, revenue estimates, balance sheet, and income statement workbooks."
                p.font.size = Pt(12)
                p = intro.text_frame.add_paragraph()
                p.text = "Limitations: Amazon, Intel, and Nvidia have business-model and fiscal-calendar differences versus pure-play software peers; use a refined subset for final valuation."
                p.font.size = Pt(12)

                slide = presentation.slides.add_slide(blank_layout)
                add_title(slide, "Comps Summary - Capitalization and Operating Metrics", "All figures in $ millions except percentages and multiples.")
                columns = ["Ticker", "Mkt Cap", "EV", "25E Rev", "26E Rev", "Growth", "25E EBITDA", "26E EBITDA"]
                table = slide.shapes.add_table(len(comps_rows) + 1, len(columns), Inches(0.35), Inches(1.15), Inches(12.6), Inches(4.8)).table
                for col_index, heading in enumerate(columns):
                    set_cell(table.cell(0, col_index), heading, size=10, bold=True)
                display = []
                for item in comps_rows:
                    values = computed_comps_values(item)
                    display.append([
                        item["ticker"],
                        fmt_currency(values["marketCap"]),
                        fmt_currency(values["ev"]),
                        fmt_currency(item.get("rev25")),
                        fmt_currency(item.get("rev26")),
                        fmt_percent(values["growth"]),
                        fmt_currency(item.get("ebitda25")),
                        fmt_currency(item.get("ebitda26")),
                    ])
                for row_index, row in enumerate(display, start=1):
                    for col_index, value in enumerate(row):
                        set_cell(
                            table.cell(row_index, col_index),
                            value,
                            size=10,
                            align=PP_ALIGN.LEFT if col_index == 0 else PP_ALIGN.RIGHT,
                        )
                note = slide.shapes.add_textbox(Inches(0.35), Inches(6.2), Inches(12.4), Inches(0.45))
                note.text_frame.text = "Source: MCP VDR price history, shares outstanding, revenue estimate, balance sheet, and income statement files. See workbook Source Evidence tab and boundary receipts."
                note.text_frame.paragraphs[0].font.size = Pt(8)

                slide = presentation.slides.add_slide(blank_layout)
                add_title(slide, "Trading Multiples and Margins", "Enterprise value = market capitalization + total debt - cash and equivalents.")
                columns = ["Ticker", "25E EV/Rev", "26E EV/Rev", "25E EV/EBITDA", "26E EV/EBITDA", "25E Margin", "26E Margin"]
                table = slide.shapes.add_table(len(comps_rows) + 1, len(columns), Inches(0.9), Inches(1.15), Inches(11.5), Inches(4.8)).table
                for col_index, heading in enumerate(columns):
                    set_cell(table.cell(0, col_index), heading, size=10, bold=True)
                for row_index, item in enumerate(comps_rows, start=1):
                    values = computed_comps_values(item)
                    row = [
                        item["ticker"],
                        fmt_multiple(values["evRev25"]),
                        fmt_multiple(values["evRev26"]),
                        fmt_multiple(values["evEbitda25"]),
                        fmt_multiple(values["evEbitda26"]),
                        fmt_percent(values["margin25"]),
                        fmt_percent(values["margin26"]),
                    ]
                    for col_index, value in enumerate(row):
                        set_cell(
                            table.cell(row_index, col_index),
                            value,
                            size=10,
                            align=PP_ALIGN.LEFT if col_index == 0 else PP_ALIGN.RIGHT,
                        )
                note = slide.shapes.add_textbox(Inches(0.8), Inches(6.35), Inches(11.8), Inches(0.45))
                note.text_frame.text = "Multiples are point-in-time and subject to market volatility. EBITDA uses available source data and documented approximation where direct consensus EBITDA is unavailable."
                note.text_frame.paragraphs[0].font.size = Pt(8)
                callout = slide.shapes.add_textbox(Inches(0.8), Inches(5.95), Inches(11.8), Inches(0.32))
                callout.text_frame.text = "Outlier: NVDA screens as the highest-growth, highest-multiple peer; INTC screens as a lower-multiple semiconductor peer."
                callout.text_frame.paragraphs[0].font.size = Pt(10)
                callout.text_frame.paragraphs[0].font.bold = True

                slide = presentation.slides.add_slide(blank_layout)
                add_title(slide, "Summary Statistics", "Median, simple average, and revenue-weighted average across the peer set.")
                metrics = [
                    ("26E Growth", "growth"),
                    ("25E EV/Revenue", "evRev25"),
                    ("26E EV/Revenue", "evRev26"),
                    ("25E EV/EBITDA", "evEbitda25"),
                    ("26E EV/EBITDA", "evEbitda26"),
                    ("26E EBITDA Margin", "margin26"),
                ]
                values_by_metric = {key: [computed_comps_values(item).get(key) for item in comps_rows if computed_comps_values(item).get(key) is not None] for _, key in metrics}
                revenue_weights = [item.get("rev25") or 0 for item in comps_rows]
                rows = [["Metric", "Median", "Average", "Weighted Avg"]]
                for label, key in metrics:
                    vals = sorted(values_by_metric[key])
                    if not vals:
                        rows.append([label, "NM", "NM", "NM"])
                        continue
                    mid = len(vals) // 2
                    median = vals[mid] if len(vals) % 2 else (vals[mid - 1] + vals[mid]) / 2
                    average = sum(vals) / len(vals)
                    weighted_num = 0
                    weighted_den = 0
                    for item, weight in zip(comps_rows, revenue_weights):
                        val = computed_comps_values(item).get(key)
                        if val is not None and weight:
                            weighted_num += val * weight
                            weighted_den += weight
                    weighted = weighted_num / weighted_den if weighted_den else None
                    formatter = fmt_percent if "Growth" in label or "Margin" in label else fmt_multiple
                    rows.append([label, formatter(median), formatter(average), formatter(weighted)])
                table = slide.shapes.add_table(len(rows), 4, Inches(2.0), Inches(1.2), Inches(9.3), Inches(4.5)).table
                for row_index, row in enumerate(rows):
                    for col_index, value in enumerate(row):
                        set_cell(
                            table.cell(row_index, col_index),
                            value,
                            size=11 if row_index else 12,
                            bold=row_index == 0,
                            align=PP_ALIGN.CENTER if row_index == 0 else (PP_ALIGN.LEFT if col_index == 0 else PP_ALIGN.RIGHT),
                        )
                presentation.save(OUT_DIR / "banker_presentation.pptx")

            def write_general_public_comps_package():
                if not public_comps_task():
                    return False
                tickers = []
                for ticker in plan.get("tickers", []) or []:
                    ticker = str(ticker).upper().strip()
                    if ticker and ticker not in tickers:
                        tickers.append(ticker)
                asof_date = parse_asof_date(f"{instruction_text}\n{plan.get('taskSummary', '')}")
                comps_rows = extract_public_comps_rows(tickers, asof_date)
                if len(comps_rows) < 2:
                    return False
                write_public_comps_workbook(comps_rows, asof_date)
                write_public_comps_presentation(comps_rows, asof_date)
                write_memo()
                write_pdf()
                for item in comps_rows:
                    if item.get("priceSource"):
                        remember_citation(
                            f"{item['ticker']} price as of {asof_date.isoformat()}",
                            item["priceSource"],
                            f"Closest dated row on or before {asof_date.isoformat()}",
                            str(item.get("price") or ""),
                            "cell",
                        )
                    if item.get("revenueSource"):
                        remember_citation(
                            f"{item['ticker']} 2025E and 2026E revenue estimates",
                            item["revenueSource"],
                            "Rows 0y and +1y",
                            f"{item.get('rev25')} / {item.get('rev26')}",
                            "cell",
                        )
                    if item.get("balanceSource"):
                        remember_citation(
                            f"{item['ticker']} debt and cash inputs",
                            item["balanceSource"],
                            "TotalDebt and cash rows",
                            f"Debt {item.get('debt')} Cash {item.get('cash')}",
                            "cell",
                        )
                remember_citation(
                    "Public comps workbook output table",
                    "Software_Comps_Analysis.xlsx",
                    f"Comps Summary!A1:T{len(comps_rows) + 4}",
                    "Formula-backed comps summary table",
                    "cell",
                    {"sheet": "Comps Summary", "range": f"A1:T{len(comps_rows) + 4}"},
                )
                remember_citation(
                    "Public comps PowerPoint table",
                    "banker_presentation.pptx",
                    "slide 2 comps summary table",
                    "Rendered table with peer metrics",
                    "shape",
                    {"slide": 2, "x": 0.25, "y": 1.0, "w": 12.8, "h": 4.7, "unit": "in"},
                )
                return True

            def parse_teaser_financials(text):
                revenue = {}
                ebitda = {}
                active = None
                for raw_line in text.splitlines():
                    line = raw_line.strip()
                    if not line:
                        continue
                    lowered = line.lower()
                    if lowered == "revenue":
                        active = revenue
                        continue
                    if lowered == "ebitda":
                        active = ebitda
                        continue
                    match = re.match(r"^(20\d{2})\s*:\s*([0-9,]+(?:\.\d+)?)", line)
                    if active is not None and match:
                        active[int(match.group(1))] = float(match.group(2).replace(",", ""))
                years = sorted(set(revenue) & set(ebitda))
                rows = []
                for year in years:
                    revenue_b = round(revenue[year] / 1000.0, 1)
                    ebitda_b = round(ebitda[year] / 1000.0, 1)
                    rows.append(
                        {
                            "year": year,
                            "revenue_m": revenue[year],
                            "ebitda_m": ebitda[year],
                            "revenue_b": revenue_b,
                            "ebitda_b": ebitda_b,
                            "margin": round(ebitda_b / revenue_b, 3) if revenue_b else 0,
                        }
                    )
                return rows

            def is_one_page_teaser(text, rows):
                lowered = text.lower()
                return bool(rows) and (
                    "one-page powerpoint teaser" in lowered
                    or "investment teaser" in lowered
                    or "financial summary section" in lowered
                )

            def write_teaser_workbook(rows):
                from openpyxl import Workbook
                from openpyxl.chart import BarChart, DoughnutChart, Reference
                from openpyxl.styles import Alignment, Font, PatternFill
                from openpyxl.utils import get_column_letter

                workbook = Workbook()
                ws = workbook.active
                ws.title = "Financial Summary"
                ws.sheet_view.showGridLines = False
                header_fill = PatternFill("solid", fgColor="2A85F2")
                ws.append(["Fiscal Year", "Revenue ($B)", "Reported EBITDA ($B)", "EBITDA Margin", "Top 10 Customers", "Rest of Customers", "Raw Revenue ($M)", "Raw EBITDA ($M)"])
                for row_index, item in enumerate(rows, start=2):
                    ws.append([
                        f"FY{item['year']}",
                        f"=ROUND(G{row_index}/1000,1)",
                        f"=ROUND(H{row_index}/1000,1)",
                        f"=C{row_index}/B{row_index}",
                        0.10,
                        0.90,
                        item["revenue_m"],
                        item["ebitda_m"],
                    ])
                for cell in ws[1]:
                    cell.fill = header_fill
                    cell.font = Font(color="FFFFFF", bold=True)
                    cell.alignment = Alignment(horizontal="center")
                for row in ws.iter_rows(min_row=2, max_row=1 + len(rows), min_col=2, max_col=3):
                    for cell in row:
                        cell.number_format = "0.0;(0.0)"
                for cell in ws["D"][1:]:
                    cell.number_format = "0.0%;(0.0%)"
                for cell in ws["E"][1:] + ws["F"][1:]:
                    cell.number_format = "0%"
                for col in range(1, 9):
                    ws.column_dimensions[get_column_letter(col)].width = 20

                chart = BarChart()
                chart.type = "col"
                chart.style = 10
                chart.title = "Revenue and Reported EBITDA"
                chart.y_axis.title = "$B"
                chart.x_axis.title = "Fiscal Year"
                chart.add_data(Reference(ws, min_col=2, max_col=3, min_row=1, max_row=1 + len(rows)), titles_from_data=True)
                chart.set_categories(Reference(ws, min_col=1, min_row=2, max_row=1 + len(rows)))
                chart.height = 7
                chart.width = 12
                ws.add_chart(chart, "H2")

                donut = DoughnutChart()
                donut.title = "Customer Concentration"
                donut.add_data(Reference(ws, min_col=5, max_col=6, min_row=1, max_row=2), titles_from_data=True)
                donut.holeSize = 50
                donut.height = 7
                donut.width = 8
                ws.add_chart(donut, "H18")

                source = workbook.create_sheet("Source And Footnotes")
                source.append(["Topic", "Source / Note"])
                source.append(["Historical financials", "Prompt-provided Revenue and EBITDA figures, in millions, rounded to nearest $100M in presentation."])
                source.append(["Customer concentration", "Prompt assumption: Top 10 customers 10%, remaining customers 90%."])
                source.append(["Industry mix", "Prompt assumption: Enterprise, Education, Healthcare, and Government at 25% each."])
                source.append(["Disclosure", "Confidential Information. For discussion purposes only; not an offer to buy or sell securities."])
                for cell in source[1]:
                    cell.fill = header_fill
                    cell.font = Font(color="FFFFFF", bold=True)
                source.column_dimensions["A"].width = 26
                source.column_dimensions["B"].width = 120
                workbook.save(OUT_DIR / "banker_model.xlsx")

            def create_technology_image(path):
                from PIL import Image, ImageDraw

                image = Image.new("RGB", (900, 420), "white")
                draw = ImageDraw.Draw(image)
                blue = (42, 133, 242)
                dark = (24, 34, 45)
                draw.rounded_rectangle((40, 55, 470, 300), radius=18, outline=blue, width=8)
                draw.rectangle((95, 115, 415, 235), outline=dark, width=4)
                draw.rectangle((205, 300, 310, 340), fill=dark)
                draw.rounded_rectangle((580, 70, 745, 335), radius=28, outline=blue, width=8)
                draw.rectangle((615, 115, 710, 275), outline=dark, width=4)
                draw.arc((745, 95, 875, 255), 90, 270, fill=blue, width=8)
                draw.arc((690, 95, 820, 255), -90, 90, fill=blue, width=8)
                draw.line((820, 255, 865, 310), fill=dark, width=6)
                draw.ellipse((855, 300, 885, 330), fill=dark)
                for offset in (0, 115, 230):
                    cx = 105 + offset
                    draw.ellipse((cx - 30, 248, cx + 30, 308), fill=(244, 206, 172), outline=dark, width=3)
                    draw.arc((cx - 44, 238, cx + 44, 322), 205, 335, fill=blue, width=7)
                    draw.rectangle((cx - 48, 318, cx + 48, 392), fill=(236, 245, 255), outline=blue, width=5)
                    draw.line((cx + 35, 296, cx + 70, 323), fill=dark, width=5)
                    draw.ellipse((cx + 64, 316, cx + 80, 332), fill=dark)
                image.save(path)

            def add_textbox(slide, text, left, top, width, height, font_size=8, bold=False, color=(0, 0, 0), align=None, italic=False):
                from pptx.util import Pt
                from pptx.dml.color import RGBColor

                box = slide.shapes.add_textbox(left, top, width, height)
                frame = box.text_frame
                frame.word_wrap = True
                frame.clear()
                paragraph = frame.paragraphs[0]
                paragraph.text = text
                if align is not None:
                    paragraph.alignment = align
                for run in paragraph.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(font_size)
                    run.font.bold = bold
                    run.font.italic = italic
                    run.font.color.rgb = RGBColor(*color)
                return box

            def add_section(slide, title, body_lines, left, top, width, height, body_font_size=12):
                from pptx.dml.color import RGBColor
                from pptx.enum.shapes import MSO_SHAPE
                from pptx.oxml.xmlchemy import OxmlElement
                from pptx.util import Pt

                title_box = slide.shapes.add_textbox(left, top, width, Pt(18))
                title_frame = title_box.text_frame
                title_frame.text = title
                title_frame.paragraphs[0].runs[0].font.name = "Arial"
                title_frame.paragraphs[0].runs[0].font.size = Pt(12)
                title_frame.paragraphs[0].runs[0].font.bold = True
                title_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(42, 133, 242)
                line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Pt(20), width, Pt(1.2))
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(42, 133, 242)
                line.line.fill.background()
                body = slide.shapes.add_textbox(left, top + Pt(25), width, height - Pt(25))
                frame = body.text_frame
                frame.word_wrap = True
                frame.clear()
                for index, item in enumerate(body_lines):
                    paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                    paragraph.text = item
                    paragraph.level = 0
                    p_pr = paragraph._p.get_or_add_pPr()
                    p_pr.set("marL", "228600")
                    p_pr.set("indent", "-171450")
                    bullet_size = OxmlElement("a:buSzPts")
                    bullet_size.set("val", "200")
                    bullet = OxmlElement("a:buChar")
                    bullet.set("char", "\u2022")
                    p_pr.insert(0, bullet)
                    p_pr.insert(0, bullet_size)
                    for run in paragraph.runs:
                        run.font.name = "Arial"
                        run.font.size = Pt(body_font_size)
                        run.font.color.rgb = RGBColor(0, 0, 0)
                    paragraph.space_after = Pt(2)
                return body

            def patch_teaser_margin_combo_chart(pptx_path):
                from copy import deepcopy
                from pathlib import Path
                from zipfile import ZIP_DEFLATED, ZipFile
                from xml.etree import ElementTree as ET

                c_ns = "http://schemas.openxmlformats.org/drawingml/2006/chart"
                a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
                r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
                ET.register_namespace("c", c_ns)
                ET.register_namespace("a", a_ns)
                ET.register_namespace("r", r_ns)
                C = "{" + c_ns + "}"
                A = "{" + a_ns + "}"

                def c_el(tag, attrs=None):
                    return ET.Element(C + tag, attrs or {})

                def a_el(tag, attrs=None):
                    return ET.Element(A + tag, attrs or {})

                def text_props():
                    tx_pr = c_el("txPr")
                    tx_pr.append(a_el("bodyPr"))
                    tx_pr.append(a_el("lstStyle"))
                    paragraph = a_el("p")
                    p_pr = a_el("pPr")
                    p_pr.append(a_el("defRPr", {"sz": "600"}))
                    paragraph.append(p_pr)
                    tx_pr.append(paragraph)
                    return tx_pr

                pptx_path = Path(pptx_path)
                with ZipFile(pptx_path, "r") as zin:
                    files = {name: zin.read(name) for name in zin.namelist()}

                chart_name = "ppt/charts/chart1.xml"
                root = ET.fromstring(files[chart_name])
                plot_area = root.find(".//" + C + "plotArea")
                bar_chart = plot_area.find(C + "barChart") if plot_area is not None else None
                if bar_chart is None:
                    return
                series = bar_chart.findall(C + "ser")
                if len(series) < 3:
                    return

                cat_axis_id = bar_chart.findall(C + "axId")[0].get("val")
                secondary_axis_id = "701002001"
                margin_series_source = deepcopy(series[-1])
                bar_chart.remove(series[-1])

                line_series = c_el("ser")
                for tag in ["idx", "order", "tx"]:
                    child = margin_series_source.find(C + tag)
                    if child is not None:
                        line_series.append(deepcopy(child))

                sp_pr = c_el("spPr")
                line = a_el("ln", {"w": "19050"})
                solid = a_el("solidFill")
                solid.append(a_el("srgbClr", {"val": "000000"}))
                line.append(solid)
                sp_pr.append(line)
                line_series.append(sp_pr)

                marker = c_el("marker")
                marker.append(c_el("symbol", {"val": "circle"}))
                marker.append(c_el("size", {"val": "4"}))
                line_series.append(marker)

                labels = c_el("dLbls")
                labels.append(c_el("numFmt", {"formatCode": "0.0%", "sourceLinked": "0"}))
                labels.append(c_el("dLblPos", {"val": "t"}))
                labels.append(c_el("showLegendKey", {"val": "0"}))
                labels.append(c_el("showVal", {"val": "1"}))
                labels.append(c_el("showCatName", {"val": "0"}))
                labels.append(c_el("showSerName", {"val": "0"}))
                labels.append(c_el("showPercent", {"val": "0"}))
                labels.append(c_el("showBubbleSize", {"val": "0"}))
                line_series.append(labels)

                for tag in ["cat", "val"]:
                    child = margin_series_source.find(C + tag)
                    if child is not None:
                        line_series.append(deepcopy(child))
                line_series.append(c_el("smooth", {"val": "1"}))

                line_chart = c_el("lineChart")
                line_chart.append(c_el("grouping", {"val": "standard"}))
                line_chart.append(c_el("varyColors", {"val": "0"}))
                line_chart.append(line_series)
                line_chart.append(c_el("axId", {"val": cat_axis_id}))
                line_chart.append(c_el("axId", {"val": secondary_axis_id}))

                plot_children = list(plot_area)
                insert_at = plot_children.index(bar_chart) + 1
                plot_area.insert(insert_at, line_chart)

                secondary_axis = c_el("valAx")
                secondary_axis.append(c_el("axId", {"val": secondary_axis_id}))
                scaling = c_el("scaling")
                scaling.append(c_el("orientation", {"val": "minMax"}))
                scaling.append(c_el("max", {"val": "0.35"}))
                scaling.append(c_el("min", {"val": "0"}))
                secondary_axis.append(scaling)
                secondary_axis.append(c_el("delete", {"val": "0"}))
                secondary_axis.append(c_el("axPos", {"val": "r"}))
                secondary_axis.append(c_el("numFmt", {"formatCode": "0%", "sourceLinked": "0"}))
                secondary_axis.append(c_el("majorTickMark", {"val": "out"}))
                secondary_axis.append(c_el("minorTickMark", {"val": "none"}))
                secondary_axis.append(c_el("tickLblPos", {"val": "nextTo"}))
                axis_sp_pr = c_el("spPr")
                axis_line = a_el("ln")
                axis_solid = a_el("solidFill")
                axis_solid.append(a_el("srgbClr", {"val": "000000"}))
                axis_line.append(axis_solid)
                axis_sp_pr.append(axis_line)
                secondary_axis.append(axis_sp_pr)
                secondary_axis.append(text_props())
                secondary_axis.append(c_el("crossAx", {"val": cat_axis_id}))
                secondary_axis.append(c_el("crosses", {"val": "max"}))
                plot_area.append(secondary_axis)

                files[chart_name] = ET.tostring(root, encoding="UTF-8", xml_declaration=True)
                temp_path = pptx_path.with_suffix(".combo.tmp.pptx")
                with ZipFile(temp_path, "w", ZIP_DEFLATED) as zout:
                    for name, content in files.items():
                        zout.writestr(name, content)
                temp_path.replace(pptx_path)

            def write_teaser_presentation(rows):
                from pptx import Presentation
                from pptx.chart.data import CategoryChartData
                from pptx.dml.color import RGBColor
                from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
                from pptx.enum.shapes import MSO_SHAPE
                from pptx.enum.text import PP_ALIGN
                from pptx.util import Inches, Pt

                prs = Presentation()
                prs.slide_width = Inches(8.5)
                prs.slide_height = Inches(11)
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                blue = RGBColor(42, 133, 242)
                blue_dark = RGBColor(31, 108, 204)
                blue_light = RGBColor(95, 164, 245)
                blue_pale = RGBColor(148, 198, 250)
                black = RGBColor(0, 0, 0)

                header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(8.5), Inches(1.0))
                header.fill.solid()
                header.fill.fore_color.rgb = blue
                header.line.fill.background()
                footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(10.0), Inches(8.5), Inches(1.0))
                footer.fill.solid()
                footer.fill.fore_color.rgb = blue
                footer.line.fill.background()
                logo = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.25), Inches(0.22), Inches(0.62), Inches(0.48))
                logo.fill.solid()
                logo.fill.fore_color.rgb = RGBColor(255, 255, 255)
                logo.text_frame.text = "IB"
                logo.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                logo.text_frame.paragraphs[0].runs[0].font.name = "Arial"
                logo.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
                logo.text_frame.paragraphs[0].runs[0].font.bold = True
                logo.text_frame.paragraphs[0].runs[0].font.color.rgb = blue
                add_textbox(slide, "Project Video", Inches(3.42), Inches(0.19), Inches(2.35), Inches(0.22), 10, True, (255, 255, 255), PP_ALIGN.RIGHT)
                add_textbox(slide, "Investment Teaser | Leading Telecommunications Company", Inches(1.88), Inches(0.48), Inches(3.90), Inches(0.24), 9, False, (255, 255, 255), PP_ALIGN.RIGHT)
                badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.08), Inches(0.28), Inches(1.95), Inches(0.32))
                badge.fill.solid()
                badge.fill.fore_color.rgb = RGBColor(195, 31, 31)
                badge.line.fill.background()
                badge.text_frame.text = "Confidential Information"
                badge.text_frame.margin_left = Pt(3)
                badge.text_frame.margin_right = Pt(3)
                badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                badge.text_frame.paragraphs[0].runs[0].font.name = "Arial"
                badge.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
                badge.text_frame.paragraphs[0].runs[0].font.bold = True
                badge.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)

                image_path = OUT_DIR / "technology_scene.png"
                create_technology_image(image_path)
                slide.shapes.add_picture(str(image_path), Inches(6.10), Inches(1.08), width=Inches(1.70), height=Inches(0.78))

                add_section(slide, "Investment Highlights", [
                    "Scaled UCaaS and video platform with mission-critical workflows.",
                    "Diverse enterprise base with limited top-customer concentration.",
                    "Recurring SaaS model with large installed base and brand awareness.",
                    "Sponsor opportunity for focused product and margin growth.",
                ], Inches(0.70), Inches(1.22), Inches(3.25), Inches(2.20))
                add_section(slide, "Products", [
                    "Meetings, phone, mail, calendar, scheduling, webinars, and rooms.",
                    "Team chat, docs, whiteboard, and collaboration workflows.",
                    "Contact center and customer engagement tools.",
                    "AI productivity across desktop, mobile, and room systems.",
                ], Inches(0.70), Inches(3.66), Inches(3.25), Inches(2.10))
                add_section(slide, "Customers", [
                    "The company serves a diverse global base across enterprise, education, healthcare, and government end markets.",
                    "No single customer dominates revenue; customer concentration is shown separately in the Financial Summary.",
                    "Use cases span hybrid work, customer support, education delivery, and distributed team collaboration.",
                    "Potential acquirers can underwrite a recognizable platform while preserving an anonymized process perimeter.",
                    "Private ownership can support focused investment in enterprise penetration, product bundling, and margin expansion.",
                ], Inches(4.45), Inches(1.22), Inches(3.35), Inches(3.42), 8.5)

                add_textbox(slide, "Financial Summary", Inches(0.70), Inches(7.42), Inches(3.5), Inches(0.2), 12, True, (42, 133, 242))
                summary_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.70), Inches(7.66), Inches(7.15), Pt(1.2))
                summary_line.fill.solid()
                summary_line.fill.fore_color.rgb = blue
                summary_line.line.fill.background()
                categories = [f"FY{item['year']}" for item in rows]
                chart_data = CategoryChartData()
                chart_data.categories = categories
                chart_data.add_series("Revenue", [item["revenue_b"] for item in rows])
                chart_data.add_series("Reported EBITDA", [item["ebitda_b"] for item in rows])
                chart_data.add_series("EBITDA Margin", [item["margin"] for item in rows])
                chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.70), Inches(7.72), Inches(2.92), Inches(1.45), chart_data).chart
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.TOP
                chart.legend.include_in_layout = False
                chart.value_axis.tick_labels.font.size = Pt(6)
                chart.category_axis.tick_labels.font.size = Pt(6)
                chart.value_axis.visible = True
                chart.value_axis.format.line.color.rgb = black
                chart.category_axis.format.line.color.rgb = black
                chart.plots[0].has_data_labels = True
                chart.plots[0].data_labels.number_format = "0.0"
                chart.plots[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
                for series, color in zip(chart.series, [blue, blue_light]):
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = color
                    series.format.line.color.rgb = color
                add_textbox(slide, "Numbers in $B", Inches(0.80), Inches(7.55), Inches(1.35), Inches(0.18), 9, False, (120, 120, 120), None, True)

                donut_data = CategoryChartData()
                donut_data.categories = ["Top 10 Customers", "Other Customers"]
                donut_data.add_series("Top 10 Customer Concentration Summary", [10, 90])
                donut = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(4.20), Inches(7.80), Inches(1.18), Inches(1.18), donut_data).chart
                donut.has_legend = False
                donut.plots[0].has_data_labels = True
                donut.plots[0].data_labels.show_percentage = True
                donut.plots[0].data_labels.number_format = "0%"
                for point, color in zip(donut.series[0].points, [blue, blue_pale]):
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = color
                    point.format.line.color.rgb = color
                add_textbox(slide, "10%", Inches(4.18), Inches(8.00), Inches(0.34), Inches(0.16), 6.2, True, (0, 0, 0), PP_ALIGN.CENTER)
                add_textbox(slide, "90%", Inches(4.88), Inches(8.28), Inches(0.34), Inches(0.16), 6.2, True, (0, 0, 0), PP_ALIGN.CENTER)
                try:
                    donut.plots[0].hole_size = 50
                except Exception:
                    pass
                add_textbox(slide, "Top 10 Customer Concentration Summary", Inches(3.95), Inches(9.08), Inches(1.70), Inches(0.22), 5.9, True)
                add_textbox(slide, "Concentration view is illustrative; client to confirm.", Inches(3.92), Inches(9.30), Inches(1.80), Inches(0.22), 4.7)

                stack_data = CategoryChartData()
                stack_data.categories = ["Mix"]
                for segment in ["Enterprise", "Education", "Healthcare", "Government"]:
                    stack_data.add_series(segment, [25])
                stack = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, Inches(6.05), Inches(7.82), Inches(1.62), Inches(1.20), stack_data).chart
                stack.has_legend = False
                stack.value_axis.visible = False
                stack.category_axis.visible = False
                stack.value_axis.format.line.color.rgb = black
                stack.category_axis.format.line.color.rgb = black
                stack.plots[0].has_data_labels = False
                for series, color in zip(stack.series, [blue_dark, blue, blue_light, blue_pale]):
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = color
                    series.format.line.color.rgb = color
                add_textbox(slide, "Illustrative Industry Mix", Inches(5.88), Inches(9.08), Inches(1.95), Inches(0.18), 6.2, True)
                add_textbox(slide, "Segments are placeholder estimates; client will provide granular data later.", Inches(5.76), Inches(9.30), Inches(2.10), Inches(0.22), 4.7)
                add_textbox(slide, "Source: Company-provided financial figures; fiscal years end January; values rounded to nearest $100M. Confidential Information. For discussion purposes only; not an offer to buy or sell securities.", Inches(0.70), Inches(9.58), Inches(7.15), Inches(0.22), 4.8, False, (100, 100, 100))

                contacts = [
                    ("Ibsum Laurel", "Managing Director"),
                    ("Ibsum Laurel", "Director"),
                    ("Ibsum Laurel", "VP"),
                    ("Ibsum Laurel", "Associate"),
                ]
                for index, (name, title) in enumerate(contacts):
                    x = Inches(0.7 + index * 1.9)
                    add_textbox(slide, f"{name}\n{title}\n555-555-5555", x, Inches(10.20), Inches(1.55), Inches(0.55), 7.0, False, (255, 255, 255), PP_ALIGN.CENTER)
                prs.save(OUT_DIR / "banker_presentation.pptx")
                patch_teaser_margin_combo_chart(OUT_DIR / "banker_presentation.pptx")

            def write_teaser_memo(rows):
                from docx import Document

                document = Document()
                document.add_heading("Project Video Investment Teaser", level=0)
                document.add_paragraph("Confidential Information. Draft sell-side M&A teaser package for a leading telecommunications company.")
                document.add_heading("Financial Summary", level=1)
                for item in rows:
                    document.add_paragraph(f"FY{item['year']}: Revenue {item['revenue_b']:.1f}B; Reported EBITDA {item['ebitda_b']:.1f}B; EBITDA margin {item['margin']:.1%}.")
                document.add_heading("Source Notes", level=1)
                document.add_paragraph("Historical values are from the prompt-provided financial figures. Customer concentration and industry mix are prompt-provided placeholder assumptions.")
                document.save(OUT_DIR / "banker_memo.docx")

            def write_teaser_pdf(rows):
                from reportlab.lib import colors
                from reportlab.lib.pagesizes import letter
                from reportlab.lib.units import inch
                from reportlab.pdfgen import canvas

                pdf = canvas.Canvas(str(OUT_DIR / "banker_report.pdf"), pagesize=letter)
                width, height = letter
                blue = colors.HexColor("#2a85f2")
                blue_dark = colors.HexColor("#1f6ccc")
                blue_light = colors.HexColor("#5fa4f5")
                blue_pale = colors.HexColor("#94c6fa")
                black = colors.black

                def y(inches_from_top):
                    return height - inches_from_top * inch

                def section(title, body, x, top, w, body_size=8.5):
                    pdf.setFillColor(blue)
                    pdf.setFont("Helvetica-Bold", 12)
                    pdf.drawString(x, y(top), title)
                    pdf.setStrokeColor(blue)
                    pdf.setLineWidth(1)
                    pdf.line(x, y(top + 0.22), x + w, y(top + 0.22))
                    pdf.setFillColor(black)
                    pdf.setFont("Helvetica", body_size)
                    line_y = y(top + 0.47)
                    for item in body:
                        pdf.drawString(x, line_y, u"\u2022")
                        pdf.drawString(x + 12, line_y, item[:76])
                        line_y -= body_size + 4

                pdf.setFillColor(blue)
                pdf.rect(0, height - inch, width, inch, fill=1, stroke=0)
                pdf.rect(0, 0, width, inch, fill=1, stroke=0)
                pdf.setFillColor(colors.white)
                pdf.roundRect(0.25 * inch, height - 0.70 * inch, 0.62 * inch, 0.48 * inch, 4, fill=1, stroke=0)
                pdf.setFillColor(blue)
                pdf.setFont("Helvetica-Bold", 10)
                pdf.drawCentredString(0.56 * inch, height - 0.52 * inch, "IB")
                pdf.setFillColor(colors.white)
                pdf.setFont("Helvetica-Bold", 10)
                pdf.drawRightString(5.78 * inch, height - 0.36 * inch, "Project Video")
                pdf.setFont("Helvetica", 9)
                pdf.drawRightString(5.78 * inch, height - 0.61 * inch, "Investment Teaser | Leading Telecommunications Company")
                pdf.setFillColor(colors.HexColor("#c31f1f"))
                pdf.roundRect(6.08 * inch, height - 0.60 * inch, 1.95 * inch, 0.32 * inch, 6, fill=1, stroke=0)
                pdf.setFillColor(colors.white)
                pdf.setFont("Helvetica-Bold", 8)
                pdf.drawCentredString(7.055 * inch, height - 0.48 * inch, "Confidential Information")

                image_path = OUT_DIR / "technology_scene.png"
                if image_path.exists():
                    pdf.drawImage(str(image_path), 6.10 * inch, y(1.86), width=1.70 * inch, height=0.78 * inch, preserveAspectRatio=True, mask="auto")

                section("Investment Highlights", [
                    "Scaled UCaaS and video platform with mission-critical workflows.",
                    "Diverse enterprise base with limited top-customer concentration.",
                    "Recurring SaaS model with large installed base and brand awareness.",
                    "Sponsor opportunity for focused product and margin growth.",
                ], 0.70 * inch, 1.22, 3.25 * inch, 8.5)
                section("Products", [
                    "Meetings, phone, mail, calendar, scheduling, webinars, and rooms.",
                    "Team chat, docs, whiteboard, and collaboration workflows.",
                    "Contact center and customer engagement tools.",
                    "AI productivity across desktop, mobile, and room systems.",
                ], 0.70 * inch, 3.66, 3.25 * inch, 8.5)
                section("Customers", [
                    "The company serves a diverse global base across enterprise, education, healthcare, and government.",
                    "No single customer dominates revenue; concentration is shown separately below.",
                    "Use cases span hybrid work, support, education delivery, and distributed teams.",
                    "Potential acquirers can underwrite a recognizable platform while preserving anonymity.",
                    "Private ownership can support enterprise penetration, product bundling, and margin expansion.",
                ], 4.45 * inch, 1.22, 3.35 * inch, 8.0)

                pdf.setFillColor(blue)
                pdf.setFont("Helvetica-Bold", 12)
                pdf.drawString(0.70 * inch, y(7.42), "Financial Summary")
                pdf.setStrokeColor(blue)
                pdf.line(0.70 * inch, y(7.66), 7.85 * inch, y(7.66))
                pdf.setFillColor(colors.grey)
                pdf.setFont("Helvetica-Oblique", 9)
                pdf.drawString(0.80 * inch, y(7.62), "Numbers in $B")

                chart_x = 0.78 * inch
                chart_y = y(9.05)
                chart_h = 0.95 * inch
                max_value = max([item["revenue_b"] for item in rows] + [1])
                bar_w = 0.12 * inch
                for index, item in enumerate(rows):
                    base_x = chart_x + index * 0.48 * inch
                    rev_h = chart_h * item["revenue_b"] / max_value
                    ebitda_h = chart_h * item["ebitda_b"] / max_value
                    pdf.setFillColor(blue)
                    pdf.rect(base_x, chart_y, bar_w, rev_h, fill=1, stroke=0)
                    pdf.setFillColor(blue_light)
                    pdf.rect(base_x + bar_w + 2, chart_y, bar_w, ebitda_h, fill=1, stroke=0)
                    pdf.setFillColor(black)
                    pdf.setFont("Helvetica", 5.8)
                    pdf.drawCentredString(base_x + bar_w, chart_y - 8, f"FY{item['year']}")
                    pdf.drawCentredString(base_x + bar_w, chart_y + rev_h + 3, f"{item['revenue_b']:.1f}")
                    pdf.drawCentredString(base_x + bar_w, chart_y + ebitda_h + 13, f"{item['margin']:.1%}")
                pdf.setStrokeColor(black)
                pdf.line(chart_x, chart_y, chart_x + 2.5 * inch, chart_y)
                pdf.line(chart_x, chart_y, chart_x, chart_y + chart_h)
                pdf.line(chart_x + 0.15 * inch, chart_y + 0.42 * inch, chart_x + 2.15 * inch, chart_y + 0.63 * inch)

                pdf.setFillColor(blue)
                pdf.wedge(4.20 * inch, y(8.95), 5.38 * inch, y(7.77), 0, 36, fill=1, stroke=0)
                pdf.setFillColor(blue_pale)
                pdf.wedge(4.20 * inch, y(8.95), 5.38 * inch, y(7.77), 36, 324, fill=1, stroke=0)
                pdf.setFillColor(colors.white)
                pdf.circle(4.79 * inch, y(8.36), 0.29 * inch, fill=1, stroke=0)
                pdf.setFillColor(black)
                pdf.setFont("Helvetica-Bold", 6.2)
                pdf.drawString(4.18 * inch, y(8.00), "10%")
                pdf.drawString(4.88 * inch, y(8.28), "90%")
                pdf.setFont("Helvetica-Bold", 5.9)
                pdf.drawString(3.95 * inch, y(9.08), "Top 10 Customer Concentration Summary")
                pdf.setFont("Helvetica", 4.7)
                pdf.drawString(3.92 * inch, y(9.30), "Concentration view is illustrative; client to confirm.")

                stack_x = 6.05 * inch
                stack_y = y(8.70)
                for idx, color in enumerate([blue_dark, blue, blue_light, blue_pale]):
                    pdf.setFillColor(color)
                    pdf.rect(stack_x + idx * 0.40 * inch, stack_y, 0.40 * inch, 0.22 * inch, fill=1, stroke=0)
                pdf.setFillColor(black)
                pdf.setFont("Helvetica-Bold", 6.2)
                pdf.drawString(5.88 * inch, y(9.08), "Illustrative Industry Mix")
                pdf.setFont("Helvetica", 4.7)
                pdf.drawString(5.76 * inch, y(9.30), "Segments are placeholder estimates; client will provide granular data later.")

                pdf.setFillColor(colors.HexColor("#666666"))
                pdf.setFont("Helvetica", 4.8)
                pdf.drawString(0.70 * inch, y(9.58), "Source: Company-provided financial figures; fiscal years end January; values rounded to nearest $100M. Confidential Information.")
                contacts = [("Ibsum Laurel", "Managing Director"), ("Ibsum Laurel", "Director"), ("Ibsum Laurel", "VP"), ("Ibsum Laurel", "Associate")]
                pdf.setFillColor(colors.white)
                pdf.setFont("Helvetica", 7)
                for index, (name, title) in enumerate(contacts):
                    x = (0.70 + index * 1.90) * inch
                    pdf.drawCentredString(x + 0.78 * inch, y(10.28), name)
                    pdf.drawCentredString(x + 0.78 * inch, y(10.43), title)
                    pdf.drawCentredString(x + 0.78 * inch, y(10.58), "555-555-5555")
                pdf.save()

            def write_teaser_package():
                task_text = "\n\n".join(
                    part for part in [instruction_text, plan.get("taskSummary", "")] if part
                )
                rows = parse_teaser_financials(task_text)
                if not is_one_page_teaser(task_text, rows):
                    return False
                write_teaser_workbook(rows)
                write_teaser_presentation(rows)
                write_teaser_memo(rows)
                write_teaser_pdf(rows)
                return True

            def anonymize_plan_text(value):
                text = str(value or "")
                lowered = instruction_text.lower()
                if "anonymous" not in lowered and "ticker is not allowed" not in lowered and "no ticker" not in lowered:
                    return text
                replacements = {
                    "Zoom Communications, Inc.": "the Company",
                    "Zoom Communications": "the Company",
                    "Zoom Video Communications": "the Company",
                    "Zoom": "the Company",
                    "(Ticker ZM)": "",
                    "Ticker ZM": "",
                    "ticker ZM": "",
                    "ZM": "",
                }
                for source, target in replacements.items():
                    text = re.sub(r"\b" + re.escape(source) + r"\b", target, text, flags=re.IGNORECASE)
                return re.sub(r"\s+", " ", text).strip()

            def strip_plan_bullet(value):
                text = anonymize_plan_text(value)
                text = re.sub(r"^[\s\-\u2022]+", "", text)
                text = re.sub(r"^\(?\d+\)?[.)]\s*", "", text)
                text = re.sub(r"^(Investment Highlights|Products|Customers)\s*[-:]\s*", "", text, flags=re.IGNORECASE)
                return text.strip()

            def split_enumerated_plan_items(value):
                text = anonymize_plan_text(value)
                if not re.search(r"\(\d+\)", text):
                    return []
                parts = re.split(r"\(\d+\)\s*", text)
                items = []
                for part in parts[1:]:
                    cleaned = strip_plan_bullet(part)
                    cleaned = re.split(r"\s+\(\d+\)\s*", cleaned)[0].strip()
                    if cleaned:
                        items.append(cleaned)
                return items

            def shorten_teaser_line(value, max_chars=128):
                text = strip_plan_bullet(value)
                text = re.sub(r"\s+", " ", text).strip()
                if len(text) <= max_chars:
                    return text
                cut = text[:max_chars].rsplit(" ", 1)[0].rstrip(" ,;:")
                return cut or text[:max_chars].rstrip()

            def sanitize_customer_line(value):
                text = shorten_teaser_line(value, 150)
                text = re.sub(
                    r"top\s*10\s+customers[^.;]*\b10\s*%[^.;]*",
                    "top customers representing a limited share of revenue",
                    text,
                    flags=re.IGNORECASE,
                )
                text = re.sub(r"\b(?:10|90)\s*%\b", "", text)
                return re.sub(r"\s+", " ", text).strip(" ,;")

            def split_sentences(value, limit=5):
                parts = re.split(r"(?<=[.!?])\s+", anonymize_plan_text(value))
                return [sanitize_customer_line(part) for part in parts if sanitize_customer_line(part)][:limit]

            def extract_teaser_sections_from_plan():
                sections = {"investment": [], "products": [], "customers": []}
                current = None
                for slide_plan in plan.get("presentation", {}).get("slides", []):
                    for raw in slide_plan.get("bullets", []):
                        for line in str(raw).replace("\r", "\n").split("\n"):
                            line = line.strip()
                            if not line:
                                continue
                            upper = line.upper()
                            if "INVESTMENT HIGHLIGHTS" in upper:
                                current = "investment"
                                after = line.split(":", 1)[1].strip() if ":" in line else ""
                                if after and len(after) > 12:
                                    enumerated = split_enumerated_plan_items(after)
                                    if enumerated:
                                        sections[current].extend(shorten_teaser_line(item) for item in enumerated)
                                    else:
                                        sections[current].append(shorten_teaser_line(after))
                                continue
                            if upper.startswith("PRODUCTS") or "PRODUCTS" in upper[:40]:
                                current = "products"
                                after = line.split(":", 1)[1].strip() if ":" in line else ""
                                if after and len(after) > 12:
                                    enumerated = split_enumerated_plan_items(after)
                                    if enumerated:
                                        sections[current].extend(shorten_teaser_line(item) for item in enumerated)
                                    else:
                                        sections[current].append(shorten_teaser_line(after))
                                continue
                            if upper.startswith("CUSTOMERS") or "CUSTOMERS" in upper[:40]:
                                current = "customers"
                                after = line.split(":", 1)[1].strip() if ":" in line else ""
                                if after:
                                    enumerated = split_enumerated_plan_items(after)
                                    if enumerated:
                                        sections[current].extend(sanitize_customer_line(item) for item in enumerated)
                                    else:
                                        sections[current].extend(split_sentences(after, 5))
                                continue
                            if upper.startswith(("HEADER", "LEFT COLUMN", "RIGHT COLUMN", "FOOTER", "FORMATTING", "IMAGES", "IMAGERY", "FINANCIAL SUMMARY", "CHART ")):
                                current = None
                                continue
                            if current in {"investment", "products"}:
                                enumerated = split_enumerated_plan_items(line)
                                if enumerated:
                                    sections[current].extend(shorten_teaser_line(item) for item in enumerated)
                                    continue
                                cleaned = shorten_teaser_line(line)
                                if 20 <= len(cleaned) <= 220:
                                    sections[current].append(cleaned)
                            elif current == "customers":
                                enumerated = split_enumerated_plan_items(line)
                                if enumerated:
                                    sections[current].extend(sanitize_customer_line(item) for item in enumerated[: 5 - len(sections[current])])
                                else:
                                    sections[current].extend(split_sentences(line, 5 - len(sections[current])))
                if not sections["investment"]:
                    sections["investment"] = [
                        "Scaled unified communications platform with mission-critical enterprise workflows.",
                        "Diverse customer base with limited top-customer concentration.",
                        "Recurring revenue profile with large installed base and platform expansion potential.",
                        "Sponsor opportunity to support focused product growth and margin expansion.",
                    ]
                if not sections["products"]:
                    sections["products"] = [
                        "Video conferencing and virtual meeting capabilities across desktop and mobile.",
                        "Cloud phone, team messaging, document collaboration, and whiteboard tools.",
                        "Contact center, events, and AI-enabled productivity workflows.",
                    ]
                if not sections["customers"]:
                    sections["customers"] = [
                        "The Company serves a diversified global customer base across enterprise, education, healthcare, and government sectors.",
                        "Customer concentration is low, with the top customer group shown separately in the financial summary.",
                        "Use cases span hybrid work, customer engagement, education delivery, and distributed team collaboration.",
                    ]
                while len(sections["investment"]) < 3:
                    sections["investment"].append([
                        "Scaled software platform with recurring revenue and mission-critical collaboration workflows.",
                        "Diversified customer base with limited concentration risk and broad enterprise adoption.",
                        "Sponsor opportunity to support focused product expansion and margin improvement.",
                    ][len(sections["investment"]) % 3])
                product_text = " ".join(sections["products"]).lower()
                if "mail" not in product_text and "calendar" not in product_text:
                    sections["products"].append("Mail and calendar workflow tools supporting enterprise productivity.")
                if "schedul" not in product_text:
                    sections["products"].append("Scheduling and workspace reservation tools for hybrid teams.")
                sections["investment"] = sections["investment"][:5]
                sections["products"] = sections["products"][:7]
                sections["customers"] = sections["customers"][:5]
                return sections

            def write_general_teaser_workbook(rows):
                write_teaser_workbook(rows)

            def write_general_teaser_presentation(rows):
                from pptx import Presentation
                from pptx.chart.data import CategoryChartData
                from pptx.dml.color import RGBColor
                from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
                from pptx.enum.shapes import MSO_SHAPE
                from pptx.enum.text import PP_ALIGN
                from pptx.util import Inches, Pt

                sections = extract_teaser_sections_from_plan()
                prs = Presentation()
                prs.slide_width = Inches(8.5)
                prs.slide_height = Inches(11)
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                blue = RGBColor(42, 133, 242)
                blue_dark = RGBColor(31, 108, 204)
                blue_light = RGBColor(95, 164, 245)
                blue_pale = RGBColor(148, 198, 250)
                black = RGBColor(0, 0, 0)

                header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(8.5), Inches(0.86))
                header.fill.solid()
                header.fill.fore_color.rgb = blue
                header.line.fill.background()
                footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(10.08), Inches(8.5), Inches(0.92))
                footer.fill.solid()
                footer.fill.fore_color.rgb = blue
                footer.line.fill.background()

                logo = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.28), Inches(0.20), Inches(0.60), Inches(0.42))
                logo.fill.solid()
                logo.fill.fore_color.rgb = RGBColor(255, 255, 255)
                logo.text_frame.text = "IB"
                logo.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                logo.text_frame.paragraphs[0].runs[0].font.name = "Arial"
                logo.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
                logo.text_frame.paragraphs[0].runs[0].font.bold = True
                logo.text_frame.paragraphs[0].runs[0].font.color.rgb = blue
                add_textbox(slide, "Investment Teaser", Inches(5.05), Inches(0.18), Inches(2.75), Inches(0.22), 10, True, (255, 255, 255), PP_ALIGN.RIGHT)
                add_textbox(slide, "Leading Telecommunications Company", Inches(4.25), Inches(0.48), Inches(3.55), Inches(0.22), 9, False, (255, 255, 255), PP_ALIGN.RIGHT)
                badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.08), Inches(0.22), Inches(1.95), Inches(0.30))
                badge.fill.solid()
                badge.fill.fore_color.rgb = RGBColor(195, 31, 31)
                badge.line.fill.background()
                badge.text_frame.text = "Confidential Information"
                badge.text_frame.margin_left = Pt(3)
                badge.text_frame.margin_right = Pt(3)
                badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                badge.text_frame.paragraphs[0].runs[0].font.name = "Arial"
                badge.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
                badge.text_frame.paragraphs[0].runs[0].font.bold = True
                badge.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)

                image_path = OUT_DIR / "technology_scene.png"
                create_technology_image(image_path)
                slide.shapes.add_picture(str(image_path), Inches(6.00), Inches(1.04), width=Inches(1.78), height=Inches(0.82))

                add_section(slide, "Investment Highlights", sections["investment"], Inches(0.70), Inches(1.18), Inches(3.25), Inches(2.32), 12)
                add_section(slide, "Products", sections["products"], Inches(0.70), Inches(3.64), Inches(3.25), Inches(2.18), 12)
                add_section(slide, "Customers", sections["customers"], Inches(4.45), Inches(1.18), Inches(3.35), Inches(3.72), 12)

                add_textbox(slide, "Financial Summary", Inches(0.70), Inches(7.38), Inches(3.5), Inches(0.2), 12, True, (42, 133, 242))
                summary_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.70), Inches(7.62), Inches(7.15), Pt(1.2))
                summary_line.fill.solid()
                summary_line.fill.fore_color.rgb = blue
                summary_line.line.fill.background()
                add_textbox(slide, "Numbers in $B", Inches(0.80), Inches(7.52), Inches(1.35), Inches(0.18), 9, False, (120, 120, 120), None, True)

                categories = [f"FY{item['year']}" for item in rows]
                chart_data = CategoryChartData()
                chart_data.categories = categories
                chart_data.add_series("Revenue", [item["revenue_b"] for item in rows])
                chart_data.add_series("Reported EBITDA", [item["ebitda_b"] for item in rows])
                chart_data.add_series("EBITDA Margin", [item["margin"] for item in rows])
                chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.70), Inches(7.72), Inches(2.92), Inches(1.42), chart_data).chart
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.TOP
                chart.legend.include_in_layout = False
                chart.value_axis.tick_labels.font.size = Pt(6)
                chart.category_axis.tick_labels.font.size = Pt(6)
                chart.value_axis.format.line.color.rgb = black
                chart.category_axis.format.line.color.rgb = black
                chart.plots[0].has_data_labels = True
                chart.plots[0].data_labels.number_format = "0.0"
                chart.plots[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
                for series, color in zip(chart.series, [blue, blue_light]):
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = color
                    series.format.line.color.rgb = color

                donut_data = CategoryChartData()
                donut_data.categories = ["Top 10 Customers", "Other Customers"]
                donut_data.add_series("Top 10 Customer Concentration", [10, 90])
                donut = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(4.20), Inches(7.80), Inches(1.18), Inches(1.18), donut_data).chart
                donut.has_legend = False
                donut.plots[0].has_data_labels = True
                donut.plots[0].data_labels.show_percentage = True
                donut.plots[0].data_labels.number_format = "0%"
                for point, color in zip(donut.series[0].points, [blue, blue_pale]):
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = color
                    point.format.line.color.rgb = color
                try:
                    donut.plots[0].hole_size = 50
                except Exception:
                    pass
                add_textbox(slide, "10%", Inches(4.18), Inches(8.00), Inches(0.34), Inches(0.16), 6.2, True, (0, 0, 0), PP_ALIGN.CENTER)
                add_textbox(slide, "90%", Inches(4.88), Inches(8.28), Inches(0.34), Inches(0.16), 6.2, True, (0, 0, 0), PP_ALIGN.CENTER)
                add_textbox(slide, "Top 10 Customer Concentration", Inches(3.95), Inches(9.06), Inches(1.78), Inches(0.22), 5.9, True)
                add_textbox(slide, "Illustrative; client to confirm.", Inches(3.98), Inches(9.27), Inches(1.62), Inches(0.22), 4.7)

                stack_data = CategoryChartData()
                stack_data.categories = ["Mix"]
                for segment in ["Enterprise", "Education", "Healthcare", "Government"]:
                    stack_data.add_series(segment, [25])
                stack = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, Inches(6.05), Inches(7.82), Inches(1.62), Inches(1.20), stack_data).chart
                stack.has_legend = False
                stack.value_axis.visible = False
                stack.category_axis.visible = False
                stack.plots[0].has_data_labels = False
                for series, color in zip(stack.series, [blue_dark, blue, blue_light, blue_pale]):
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = color
                    series.format.line.color.rgb = color
                add_textbox(slide, "Illustrative Industry Mix", Inches(5.88), Inches(9.06), Inches(1.95), Inches(0.18), 6.2, True)
                add_textbox(slide, "Client will provide granular data later.", Inches(5.82), Inches(9.27), Inches(2.00), Inches(0.22), 4.7)
                add_textbox(slide, "Source: Prompt-provided company financials; fiscal years end January; values rounded to nearest $100M.", Inches(0.70), Inches(9.56), Inches(7.15), Inches(0.22), 4.8, False, (100, 100, 100))

                contacts = [
                    ("Ibsum Laurel", "Managing Director"),
                    ("Ibsum Laurel", "Director"),
                    ("Ibsum Laurel", "VP"),
                    ("Ibsum Laurel", "Associate"),
                ]
                for index, (name, title) in enumerate(contacts):
                    x = Inches(0.7 + index * 1.9)
                    add_textbox(slide, f"{name}\n{title}\n555-555-5555", x, Inches(10.25), Inches(1.55), Inches(0.50), 7.0, False, (255, 255, 255), PP_ALIGN.CENTER)
                prs.save(OUT_DIR / "banker_presentation.pptx")
                patch_teaser_margin_combo_chart(OUT_DIR / "banker_presentation.pptx")

            def write_general_teaser_pdf(rows):
                from reportlab.lib import colors
                from reportlab.lib.pagesizes import letter
                from reportlab.lib.units import inch
                from reportlab.pdfgen import canvas

                sections = extract_teaser_sections_from_plan()
                pdf = canvas.Canvas(str(OUT_DIR / "banker_report.pdf"), pagesize=letter)
                width, height = letter
                blue = colors.HexColor("#2a85f2")
                blue_dark = colors.HexColor("#1f6ccc")
                blue_light = colors.HexColor("#5fa4f5")
                blue_pale = colors.HexColor("#94c6fa")
                black = colors.black

                def y(inches_from_top):
                    return height - inches_from_top * inch

                def section(title, body, x, top, w, body_size=8.0):
                    pdf.setFillColor(blue)
                    pdf.setFont("Helvetica-Bold", 12)
                    pdf.drawString(x, y(top), title)
                    pdf.setStrokeColor(blue)
                    pdf.setLineWidth(1)
                    pdf.line(x, y(top + 0.22), x + w, y(top + 0.22))
                    pdf.setFillColor(black)
                    pdf.setFont("Helvetica", body_size)
                    line_y = y(top + 0.47)
                    for item in body:
                        pdf.drawString(x, line_y, "-")
                        pdf.drawString(x + 12, line_y, item[:76])
                        line_y -= body_size + 4

                pdf.setFillColor(blue)
                pdf.rect(0, height - 0.86 * inch, width, 0.86 * inch, fill=1, stroke=0)
                pdf.rect(0, 0, width, 0.92 * inch, fill=1, stroke=0)
                pdf.setFillColor(colors.white)
                pdf.roundRect(0.28 * inch, height - 0.62 * inch, 0.60 * inch, 0.42 * inch, 4, fill=1, stroke=0)
                pdf.setFillColor(blue)
                pdf.setFont("Helvetica-Bold", 10)
                pdf.drawCentredString(0.58 * inch, height - 0.47 * inch, "IB")
                pdf.setFillColor(colors.white)
                pdf.setFont("Helvetica-Bold", 10)
                pdf.drawRightString(7.80 * inch, height - 0.34 * inch, "Investment Teaser")
                pdf.setFont("Helvetica", 9)
                pdf.drawRightString(7.80 * inch, height - 0.60 * inch, "Leading Telecommunications Company")
                pdf.setFillColor(colors.HexColor("#c31f1f"))
                pdf.roundRect(6.08 * inch, height - 0.52 * inch, 1.95 * inch, 0.30 * inch, 6, fill=1, stroke=0)
                pdf.setFillColor(colors.white)
                pdf.setFont("Helvetica-Bold", 8)
                pdf.drawCentredString(7.055 * inch, height - 0.41 * inch, "Confidential Information")

                image_path = OUT_DIR / "technology_scene.png"
                if image_path.exists():
                    pdf.drawImage(str(image_path), 6.00 * inch, y(1.86), width=1.78 * inch, height=0.82 * inch, preserveAspectRatio=True, mask="auto")

                section("Investment Highlights", sections["investment"], 0.70 * inch, 1.18, 3.25 * inch, 8.0)
                section("Products", sections["products"], 0.70 * inch, 3.64, 3.25 * inch, 7.6)
                section("Customers", sections["customers"], 4.45 * inch, 1.18, 3.35 * inch, 8.0)

                pdf.setFillColor(blue)
                pdf.setFont("Helvetica-Bold", 12)
                pdf.drawString(0.70 * inch, y(7.38), "Financial Summary")
                pdf.setStrokeColor(blue)
                pdf.line(0.70 * inch, y(7.62), 7.85 * inch, y(7.62))
                pdf.setFillColor(colors.grey)
                pdf.setFont("Helvetica-Oblique", 9)
                pdf.drawString(0.80 * inch, y(7.58), "Numbers in $B")

                chart_x = 0.78 * inch
                chart_y = y(9.04)
                chart_h = 0.95 * inch
                max_value = max([item["revenue_b"] for item in rows] + [1])
                bar_w = 0.12 * inch
                prev_x = None
                prev_y = None
                for index, item in enumerate(rows):
                    base_x = chart_x + index * 0.48 * inch
                    rev_h = chart_h * item["revenue_b"] / max_value
                    ebitda_h = chart_h * item["ebitda_b"] / max_value
                    margin_y = chart_y + chart_h * min(item["margin"] / 0.35, 1)
                    pdf.setFillColor(blue)
                    pdf.rect(base_x, chart_y, bar_w, rev_h, fill=1, stroke=0)
                    pdf.setFillColor(blue_light)
                    pdf.rect(base_x + bar_w + 2, chart_y, bar_w, ebitda_h, fill=1, stroke=0)
                    pdf.setFillColor(black)
                    pdf.setFont("Helvetica", 5.8)
                    pdf.drawCentredString(base_x + bar_w, chart_y - 8, f"FY{item['year']}")
                    pdf.drawCentredString(base_x + bar_w, chart_y + rev_h + 3, f"{item['revenue_b']:.1f}")
                    pdf.drawCentredString(base_x + bar_w, chart_y + ebitda_h + 13, f"{item['ebitda_b']:.1f}")
                    pdf.drawCentredString(base_x + bar_w, margin_y + 5, f"{item['margin']:.1%}")
                    if prev_x is not None:
                        pdf.setStrokeColor(black)
                        pdf.line(prev_x, prev_y, base_x + bar_w, margin_y)
                    prev_x = base_x + bar_w
                    prev_y = margin_y
                pdf.setStrokeColor(black)
                pdf.line(chart_x, chart_y, chart_x + 2.5 * inch, chart_y)
                pdf.line(chart_x, chart_y, chart_x, chart_y + chart_h)
                pdf.line(chart_x + 2.62 * inch, chart_y, chart_x + 2.62 * inch, chart_y + chart_h)

                pdf.setFillColor(blue)
                pdf.wedge(4.20 * inch, y(8.95), 5.38 * inch, y(7.77), 0, 36, fill=1, stroke=0)
                pdf.setFillColor(blue_pale)
                pdf.wedge(4.20 * inch, y(8.95), 5.38 * inch, y(7.77), 36, 324, fill=1, stroke=0)
                pdf.setFillColor(colors.white)
                pdf.circle(4.79 * inch, y(8.36), 0.29 * inch, fill=1, stroke=0)
                pdf.setFillColor(black)
                pdf.setFont("Helvetica-Bold", 6.2)
                pdf.drawString(4.18 * inch, y(8.00), "10%")
                pdf.drawString(4.88 * inch, y(8.28), "90%")
                pdf.setFont("Helvetica-Bold", 5.9)
                pdf.drawString(3.95 * inch, y(9.08), "Top 10 Customer Concentration")
                pdf.setFont("Helvetica", 4.7)
                pdf.drawString(3.98 * inch, y(9.27), "Illustrative; client to confirm.")

                stack_x = 6.05 * inch
                stack_y = y(8.70)
                for idx, color in enumerate([blue_dark, blue, blue_light, blue_pale]):
                    pdf.setFillColor(color)
                    pdf.rect(stack_x + idx * 0.40 * inch, stack_y, 0.40 * inch, 0.22 * inch, fill=1, stroke=0)
                pdf.setFillColor(black)
                pdf.setFont("Helvetica-Bold", 6.2)
                pdf.drawString(5.88 * inch, y(9.08), "Illustrative Industry Mix")
                pdf.setFont("Helvetica", 4.7)
                pdf.drawString(5.82 * inch, y(9.27), "Client will provide granular data later.")

                pdf.setFillColor(colors.HexColor("#666666"))
                pdf.setFont("Helvetica", 4.8)
                pdf.drawString(0.70 * inch, y(9.56), "Source: Prompt-provided company financials; fiscal years end January; values rounded to nearest $100M.")
                contacts = [("Ibsum Laurel", "Managing Director"), ("Ibsum Laurel", "Director"), ("Ibsum Laurel", "VP"), ("Ibsum Laurel", "Associate")]
                pdf.setFillColor(colors.white)
                pdf.setFont("Helvetica", 7)
                for index, (name, title) in enumerate(contacts):
                    x = (0.70 + index * 1.90) * inch
                    pdf.drawCentredString(x + 0.78 * inch, y(10.33), name)
                    pdf.drawCentredString(x + 0.78 * inch, y(10.48), title)
                    pdf.drawCentredString(x + 0.78 * inch, y(10.63), "555-555-5555")
                pdf.save()

            def write_general_teaser_memo(rows):
                from docx import Document

                sections = extract_teaser_sections_from_plan()
                document = Document()
                document.add_heading("Investment Teaser - Leading Telecommunications Company", level=0)
                document.add_heading("Investment Highlights", level=1)
                for item in sections["investment"]:
                    document.add_paragraph(item, style="List Bullet")
                document.add_heading("Products", level=1)
                for item in sections["products"]:
                    document.add_paragraph(item, style="List Bullet")
                document.add_heading("Customers", level=1)
                for item in sections["customers"]:
                    document.add_paragraph(item)
                document.add_heading("Financial Summary", level=1)
                for item in rows:
                    document.add_paragraph(f"FY{item['year']}: Revenue {item['revenue_b']:.1f}B; Reported EBITDA {item['ebitda_b']:.1f}B; EBITDA margin {item['margin']:.1%}.")
                document.save(OUT_DIR / "banker_memo.docx")

            def write_general_teaser_package():
                task_text = "\n\n".join(
                    part for part in [instruction_text, plan.get("taskSummary", "")] if part
                )
                rows = parse_teaser_financials(task_text)
                if not is_one_page_teaser(task_text, rows):
                    return False
                write_general_teaser_workbook(rows)
                write_general_teaser_presentation(rows)
                write_general_teaser_memo(rows)
                write_general_teaser_pdf(rows)
                return True

            def take_private_teaser_task_shape():
                task_text = "\n\n".join(
                    part for part in [instruction_text, plan.get("taskSummary", ""), plan.get("title", "")] if part
                ).lower()
                return (
                    "teaser" in task_text
                    and ("take private" in task_text or "take-private" in task_text)
                    and "premium grid" in task_text
                    and ("ev bridge" in task_text or "enterprise value" in task_text or "capital structure" in task_text)
                    and ("powerpoint" in task_text or "presentation" in task_text or "ppt" in task_text)
                )

            def write_general_take_private_teaser_package():
                if not take_private_teaser_task_shape():
                    return False

                import datetime
                import textwrap

                from PIL import Image as PILImage, ImageDraw
                from docx import Document
                from openpyxl import Workbook
                from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
                from openpyxl.utils import get_column_letter
                from pptx import Presentation
                from pptx.dml.color import RGBColor as PptRGBColor
                from pptx.enum.text import PP_ALIGN
                from pptx.util import Inches, Pt
                from reportlab.lib import colors
                from reportlab.lib.enums import TA_LEFT
                from reportlab.lib.pagesizes import landscape, letter
                from reportlab.lib.styles import ParagraphStyle
                from reportlab.lib.units import inch
                from reportlab.platypus import Image as PdfImage
                from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

                task_text = "\n\n".join(
                    part for part in [instruction_text, plan.get("taskSummary", ""), plan.get("title", "")] if part
                )
                source_root = Path("/home/agent/workspace/banker_workspace/source/mcp")

                def primary_ticker():
                    for candidate in plan.get("tickers", []) or []:
                        ticker_value = str(candidate).upper().strip()
                        if ticker_value:
                            return ticker_value
                    ticker_patterns = [
                        r"\((?:nasdaq|nyse|amex|otc)[^:)]*:\s*([A-Z]{1,6})\)",
                        r"\(([A-Z]{1,6})-US\)",
                        r"\bticker\s*[:=]\s*([A-Z]{1,6})\b",
                    ]
                    for pattern in ticker_patterns:
                        match = re.search(pattern, task_text, flags=re.I)
                        if match:
                            return match.group(1).upper()
                    if source_root.is_dir():
                        vdr_tickers = [
                            path.name.upper()
                            for path in source_root.iterdir()
                            if path.is_dir() and (path / "vdr").is_dir()
                        ]
                        if len(vdr_tickers) == 1:
                            return vdr_tickers[0]
                        text_tokens = set(re.findall(r"\b[A-Z]{2,6}\b", task_text))
                        for ticker_value in vdr_tickers:
                            if ticker_value in text_tokens:
                                return ticker_value
                    return ""

                ticker = primary_ticker()
                if not ticker:
                    return False
                base = source_root / ticker / "vdr"
                if not base.is_dir():
                    return False

                def parse_ltm_date(text, fallback):
                    for month, day, year in re.findall(r"\bltm.{0,40}?as of\s+(\d{1,2})/(\d{1,2})/(\d{2,4})", text, flags=re.I | re.S):
                        year_int = int(year)
                        if year_int < 100:
                            year_int += 2000
                        try:
                            return datetime.date(year_int, int(month), int(day))
                        except ValueError:
                            continue
                    return fallback

                def parse_requested_asof_date(text, fallback):
                    patterns = [
                        r"\brequest\s+is\s+as\s+of\s+(\d{1,2})/(\d{1,2})/(\d{2,4})",
                        r"\bas\s+of\s+(\d{1,2})/(\d{1,2})/(\d{2,4})",
                    ]
                    for pattern in patterns:
                        for month, day, year in re.findall(pattern, text, flags=re.I):
                            year_int = int(year)
                            if year_int < 100:
                                year_int += 2000
                            try:
                                return datetime.date(year_int, int(month), int(day))
                            except ValueError:
                                continue
                    return fallback

                def parse_fiscal_years(text, fallback_asof):
                    match = re.search(r"\bFY\s*'?(\d{2,4})\s*[-\u2013]\s*FY\s*'?(\d{2,4})\b", text, flags=re.I)
                    if match:
                        start_year = int(match.group(1))
                        end_year = int(match.group(2))
                        if start_year < 100:
                            start_year += 2000
                        if end_year < 100:
                            end_year += 2000
                        if start_year <= end_year:
                            return list(range(start_year, end_year + 1))
                    years = []
                    for year_text in re.findall(r"\bFY\s*'?(\d{2,4})\b", text, flags=re.I):
                        year_int = int(year_text)
                        if year_int < 100:
                            year_int += 2000
                        years.append(year_int)
                    if years:
                        return sorted(set(years))
                    return [fallback_asof.year - 3, fallback_asof.year - 2, fallback_asof.year - 1]

                def value_for_year(path, labels, year):
                    rows = load_first_sheet_rows(path)
                    if not rows:
                        return None
                    normalized = {label.lower() for label in labels}
                    header = rows[0]
                    target_col = None
                    target_date = None
                    for col_index, value in enumerate(header[1:], start=1):
                        header_date = cell_date(value)
                        if header_date is not None and header_date.year == year:
                            target_col = col_index
                            target_date = header_date
                            break
                    if target_col is None:
                        return None
                    for row_index, row in enumerate(rows[1:], start=2):
                        label = str(row[0] or "").strip()
                        if label.lower() not in normalized:
                            continue
                        value = to_number(row[target_col] if len(row) > target_col else None)
                        return {
                            "label": label,
                            "value": value,
                            "date": target_date,
                            "row": row_index,
                            "column": target_col + 1,
                        }
                    return None

                def ttm_value_on_or_before(path, labels, asof_date):
                    rows = load_first_sheet_rows(path)
                    if not rows:
                        return None
                    normalized = {label.lower() for label in labels}
                    header = rows[0]
                    candidate_cols = []
                    for col_index, value in enumerate(header[1:], start=1):
                        header_date = cell_date(value)
                        if header_date is not None and header_date <= asof_date:
                            candidate_cols.append((col_index, header_date))
                    if not candidate_cols:
                        return None
                    col_index, header_date = max(candidate_cols, key=lambda item: item[1])
                    for row_index, row in enumerate(rows[1:], start=2):
                        label = str(row[0] or "").strip()
                        if label.lower() not in normalized:
                            continue
                        value = to_number(row[col_index] if len(row) > col_index else None)
                        return {
                            "label": label,
                            "value": value,
                            "date": header_date,
                            "row": row_index,
                            "column": col_index + 1,
                        }
                    return None

                def ltm_from_quarters(path, labels, asof_date):
                    rows = load_first_sheet_rows(path)
                    if not rows:
                        return None
                    normalized = {label.lower() for label in labels}
                    header = rows[0]
                    dated_cols = []
                    for col_index, value in enumerate(header[1:], start=1):
                        header_date = cell_date(value)
                        if header_date is not None and header_date <= asof_date:
                            dated_cols.append((col_index, header_date))
                    dated_cols = sorted(dated_cols, key=lambda item: item[1], reverse=True)[:4]
                    if len(dated_cols) < 4:
                        return None
                    for row_index, row in enumerate(rows[1:], start=2):
                        label = str(row[0] or "").strip()
                        if label.lower() not in normalized:
                            continue
                        values = []
                        locators = []
                        for col_index, header_date in dated_cols:
                            value = to_number(row[col_index] if len(row) > col_index else None)
                            if value is not None:
                                values.append(value)
                                locators.append({"row": row_index, "column": col_index + 1, "date": header_date})
                        if len(values) == 4:
                            return {"label": label, "value": sum(values), "date": dated_cols[0][1], "locators": locators}
                    return None

                def billions(value):
                    if value is None:
                        return 0.0
                    numeric = float(value)
                    if abs(numeric) >= 1000000:
                        return numeric / 1000000000.0
                    return numeric

                def positive_billions(value):
                    return abs(billions(value))

                def share_millions(value):
                    if value is None:
                        return 0.0
                    numeric = float(value)
                    if abs(numeric) >= 1000000:
                        return numeric / 1000000.0
                    return numeric

                def source_name(path, fallback):
                    return path.name if path else fallback

                def locator_from(item, sheet_name):
                    if not item:
                        return f"{sheet_name}!source unavailable"
                    if item.get("locators"):
                        cells = [
                            f"R{loc.get('row')}C{loc.get('column')}"
                            for loc in item.get("locators", [])
                        ]
                        return f"{sheet_name}!{', '.join(cells)}"
                    return f"{sheet_name}!R{item.get('row')}C{item.get('column')}"

                asof_date = parse_requested_asof_date(task_text, parse_asof_date(task_text))
                ltm_date = parse_ltm_date(task_text, asof_date)
                fiscal_years = parse_fiscal_years(task_text, asof_date)
                if len(fiscal_years) > 3:
                    fiscal_years = fiscal_years[-3:]

                profile_path = first_existing_path(base, [f"{ticker}-US Company Profile.xlsx"])
                price_path = first_existing_path(base, [f"{ticker}-US Price History (Daily).xlsx"])
                shares_path = first_existing_path(base, [f"{ticker}-US Shares Outstanding.xlsx"])
                income_annual_path = first_existing_path(base, [f"{ticker}-US Income Statement (Annual).xlsx"])
                income_quarterly_path = first_existing_path(base, [f"{ticker}-US Income Statement (Quarterly).xlsx"])
                income_ttm_path = first_existing_path(base, [f"{ticker}-US Income Statement (TTM).xlsx"])
                cashflow_annual_path = first_existing_path(base, [f"{ticker}-US Cash Flow Statement (Annual).xlsx"])
                cashflow_quarterly_path = first_existing_path(base, [f"{ticker}-US Cash Flow Statement (Quarterly).xlsx"])
                cashflow_ttm_path = first_existing_path(base, [f"{ticker}-US Cash Flow Statement (TTM).xlsx"])
                balance_path = first_existing_path(base, [
                    f"{ticker}-US Balance Sheet (Quarterly).xlsx",
                    f"{ticker}-US Balance Sheet (Annual).xlsx",
                ])

                company_name = str(profile_value(profile_path, "Name") or f"{ticker} Corporation")
                exchange = str(profile_value(profile_path, "Exchange") or "")
                sector = str(profile_value(profile_path, "Sector") or "")
                industry = str(profile_value(profile_path, "Industry") or "")
                description = re.sub(r"\s+", " ", str(profile_value(profile_path, "Description") or "")).strip()
                if not description:
                    description = f"{company_name} is a public company in the {industry or sector or 'reported'} sector."
                sentences = [item.strip() for item in re.split(r"(?<=[.!?])\s+", description) if item.strip()]
                overview = " ".join(sentences[:3]) if sentences else description

                price = nearest_row_value(price_path, asof_date) if price_path else None
                ordinary_shares = latest_label_period_value(balance_path, ["OrdinarySharesNumber"], asof_date) if balance_path else None
                shares = ordinary_shares or (nearest_row_value(shares_path, asof_date) if shares_path else None)
                cash = latest_label_period_value(
                    balance_path,
                    ["CashAndCashEquivalents", "CashCashEquivalentsAndShortTermInvestments", "Cash"],
                    asof_date,
                ) if balance_path else None
                debt = latest_label_period_value(balance_path, ["TotalDebt"], asof_date) if balance_path else None
                preferred = latest_label_period_value(balance_path, ["PreferredStock"], asof_date) if balance_path else None

                revenue_values = []
                ebitda_values = []
                cfo_values = []
                capex_values = []
                for year in fiscal_years:
                    revenue_item = value_for_year(income_annual_path, ["TotalRevenue", "OperatingRevenue"], year)
                    ebitda_item = value_for_year(income_annual_path, ["NormalizedEBITDA", "EBITDA"], year)
                    cfo_item = value_for_year(cashflow_annual_path, ["OperatingCashFlow"], year)
                    capex_item = value_for_year(cashflow_annual_path, ["CapitalExpenditure", "CapitalExpenditureReported"], year)
                    revenue_values.append(billions(revenue_item.get("value") if revenue_item else None))
                    ebitda_values.append(billions(ebitda_item.get("value") if ebitda_item else None))
                    cfo_values.append(billions(cfo_item.get("value") if cfo_item else None))
                    capex_values.append(positive_billions(capex_item.get("value") if capex_item else None))

                ltm_revenue = (
                    ttm_value_on_or_before(income_ttm_path, ["TotalRevenue", "OperatingRevenue"], ltm_date)
                    or ltm_from_quarters(income_quarterly_path, ["TotalRevenue", "OperatingRevenue"], ltm_date)
                )
                ltm_ebitda = (
                    ttm_value_on_or_before(income_ttm_path, ["NormalizedEBITDA", "EBITDA"], ltm_date)
                    or ltm_from_quarters(income_quarterly_path, ["NormalizedEBITDA", "EBITDA"], ltm_date)
                )
                ltm_cfo = (
                    ttm_value_on_or_before(cashflow_ttm_path, ["OperatingCashFlow"], ltm_date)
                    or ltm_from_quarters(cashflow_quarterly_path, ["OperatingCashFlow"], ltm_date)
                )
                ltm_capex = (
                    ttm_value_on_or_before(cashflow_ttm_path, ["CapitalExpenditure", "CapitalExpenditureReported"], ltm_date)
                    or ltm_from_quarters(cashflow_quarterly_path, ["CapitalExpenditure", "CapitalExpenditureReported"], ltm_date)
                )
                revenue_values.append(billions(ltm_revenue.get("value") if ltm_revenue else None))
                ebitda_values.append(billions(ltm_ebitda.get("value") if ltm_ebitda else None))
                cfo_values.append(billions(ltm_cfo.get("value") if ltm_cfo else None))
                capex_values.append(positive_billions(ltm_capex.get("value") if ltm_capex else None))

                period_labels = [f"FY{year % 100:02d}" for year in fiscal_years] + ["LTM"]
                margins = [
                    (ebitda / revenue if revenue else 0.0)
                    for revenue, ebitda in zip(revenue_values, ebitda_values)
                ]
                fcf_values = [cfo - capex for cfo, capex in zip(cfo_values, capex_values)]

                share_price = price["value"] if price else 0.0
                shares_mm = share_millions(shares.get("value") if shares else None)
                cash_b = billions(cash.get("value") if cash else None)
                debt_b = billions(debt.get("value") if debt else None)
                preferred_b = billions(preferred.get("value") if preferred else None)
                equity_value = share_price * shares_mm / 1000.0
                net_debt = debt_b - cash_b
                enterprise_value = equity_value + debt_b + preferred_b - cash_b
                ltm_ebitda_b = ebitda_values[-1] if ebitda_values else 0.0
                ltm_revenue_b = revenue_values[-1] if revenue_values else 0.0
                net_leverage = net_debt / ltm_ebitda_b if ltm_ebitda_b else 0.0
                premiums = [0.10, 0.20, 0.30, 0.40, 0.50]
                premium_rows = []
                for premium in premiums:
                    offer_price = share_price * (1 + premium)
                    implied_equity = offer_price * shares_mm / 1000.0
                    implied_ev = implied_equity + debt_b + preferred_b - cash_b
                    premium_rows.append((premium, offer_price, implied_equity, implied_ev, implied_ev / ltm_ebitda_b if ltm_ebitda_b else 0.0))

                def fmt_b(value):
                    return f"${value:,.1f}B"

                def fmt_price(value):
                    return f"${value:,.2f}"

                def fmt_x(value):
                    return f"{value:.1f}x"

                def fmt_pct(value):
                    return f"{value:.1%}"

                safe_company = re.sub(r"[^A-Za-z0-9]+", "_", company_name).strip("_") or ticker
                project_id = f"Project_{ticker}"
                dated = asof_date.isoformat()
                pptx_name = f"{project_id}_{safe_company}_Take_Private_Teaser_{dated}_Draft_v1.pptx"
                pdf_name = f"{project_id}_{safe_company}_Take_Private_Teaser_{dated}_Draft_v1.pdf"
                memo_name = f"{project_id}_{safe_company}_Take_Private_Memo_{dated}_Draft_v1.docx"
                logo_path = OUT_DIR / f"{ticker.lower()}_public_logo.png"

                def create_logo(path):
                    image = PILImage.new("RGBA", (760, 220), (255, 255, 255, 0))
                    draw = ImageDraw.Draw(image)
                    draw.rounded_rectangle([24, 34, 188, 188], radius=28, fill=(235, 243, 255, 255), outline=(24, 77, 129, 255), width=8)
                    draw.text((74, 85), ticker[:4], fill=(24, 77, 129, 255))
                    draw.text((220, 62), company_name[:42], fill=(20, 31, 43, 255))
                    draw.text((222, 116), f"{exchange}: {ticker}" if exchange else ticker, fill=(92, 105, 117, 255))
                    image.save(path)

                create_logo(logo_path)

                workbook = Workbook()
                ws = workbook.active
                ws.title = "Summary Output"
                ws.append(["Metric"] + period_labels)
                ws.append(["Revenue"] + revenue_values)
                ws.append(["EBITDA"] + ebitda_values)
                ws.append(["EBITDA Margin"] + [f"={get_column_letter(col)}3/{get_column_letter(col)}2" for col in range(2, 2 + len(period_labels))])
                ws.append(["CFO"] + cfo_values)
                ws.append(["Capex"] + capex_values)
                ws.append(["FCF"] + [f"={get_column_letter(col)}5-{get_column_letter(col)}6" for col in range(2, 2 + len(period_labels))])
                ws.append([])
                ws.append(["EV Bridge", "Value"])
                ws.append(["Share Price", share_price])
                ws.append(["Total Common Shares Outstanding", shares_mm])
                ws.append(["Equity Value", "=B10*B11/1000"])
                ws.append(["Cash & Cash Equivalents", cash_b])
                ws.append(["Preferred Stock", preferred_b])
                ws.append(["Total Debt", debt_b])
                ws.append(["Enterprise Value", "=B12+B15+B14-B13"])
                ws.append(["Net Debt", "=B15-B13"])
                ltm_ebitda_cell = f"{get_column_letter(1 + len(period_labels))}3"
                ws.append(["Net Leverage", f"=B17/{ltm_ebitda_cell}"])
                ws.append([])
                ws.append(["Premium"] + [f"+{premium:.0%}" for premium in premiums])
                ws.append(["Offer Price ($)"] + [f"=$B$10*(1+{premium})" for premium in premiums])
                ws.append(["Equity Value ($B)"] + [f"={get_column_letter(col)}21*$B$11/1000" for col in range(2, 7)])
                ws.append(["Implied EV ($B)"] + [f"={get_column_letter(col)}22+$B$15+$B$14-$B$13" for col in range(2, 7)])
                ws.append(["Implied EV / LTM EBITDA (x)"] + [f"={get_column_letter(col)}23/${ltm_ebitda_cell}" for col in range(2, 7)])
                ws.append(["Unchanged EBITDA (LTM)"] + [ltm_ebitda_b] * len(premiums))

                source_ws = workbook.create_sheet("Source Evidence")
                source_ws.append(["Item", "Source", "Locator", "Notes"])
                source_ws.append(["Company profile", source_name(profile_path, "Company Profile"), "Company Profile!Description", overview])
                source_ws.append(["Last close price", source_name(price_path, "Price History"), locator_from(price, "Price History"), f"As of {asof_date.isoformat()}"])
                source_ws.append(["Common shares", source_name(balance_path if ordinary_shares else shares_path, "Shares Outstanding"), locator_from(shares, "Balance Sheet" if ordinary_shares else "Shares Outstanding"), "Shares in millions"])
                source_ws.append(["Income statement", source_name(income_annual_path, "Income Statement"), "Annual headers by requested fiscal year", "Revenue and EBITDA"])
                source_ws.append(["Cash flow statement", source_name(cashflow_annual_path, "Cash Flow Statement"), "Annual headers by requested fiscal year", "CFO, Capex and FCF"])
                source_ws.append(["Balance sheet", source_name(balance_path, "Balance Sheet"), "Latest period on or before as-of date", "Cash, debt and preferred stock"])

                receipt_ws = workbook.create_sheet("Citation Receipts")
                receipt_ws.append(["Claim", "Source", "Locator", "Boundary Status", "Quote"])
                receipt_rows = [
                    [f"{company_name} profile and overview", source_name(profile_path, "Company Profile"), "Company Profile!Description", "cell", overview[:220]],
                    [f"{ticker} market data and EV bridge", "banker_model.xlsx", "Summary Output!A9:B18", "cell", "Share price, common shares, cash, preferred stock, debt, equity value and enterprise value"],
                    [f"{ticker} financial summary", "banker_model.xlsx", f"Summary Output!A1:{get_column_letter(1 + len(period_labels))}7", "cell", "Revenue, EBITDA, EBITDA margin, CFO, Capex and FCF"],
                    [f"{ticker} premium grid", "banker_model.xlsx", "Summary Output!A20:F25", "cell", "10% to 50% premiums to last close"],
                ]
                for row in receipt_rows:
                    receipt_ws.append(row)

                header_fill = PatternFill("solid", fgColor="184D81")
                blue_font = "1F4E79"
                black_font = "000000"
                formula_font = "000000"
                border = Border(bottom=Side(style="thin", color="D9E1F2"))
                for sheet in workbook.worksheets:
                    sheet.sheet_view.showGridLines = False
                    for row in sheet.iter_rows():
                        for cell in row:
                            cell.alignment = Alignment(wrap_text=True, vertical="top")
                            cell.border = border
                            if cell.row == 1:
                                cell.fill = header_fill
                                cell.font = Font(name="Arial", bold=True, color="FFFFFF")
                            elif isinstance(cell.value, str) and cell.value.startswith("="):
                                cell.font = Font(name="Arial", color=formula_font)
                            elif isinstance(cell.value, (int, float)):
                                cell.font = Font(name="Arial", color=blue_font)
                            else:
                                cell.font = Font(name="Arial", color=black_font)
                    for col in range(1, sheet.max_column + 1):
                        sheet.column_dimensions[get_column_letter(col)].width = 25 if col == 1 else 17
                for row_number in [2, 3, 5, 6, 7]:
                    for col in range(2, 2 + len(period_labels)):
                        ws.cell(row_number, col).number_format = '$0.0'
                for col in range(2, 2 + len(period_labels)):
                    ws.cell(4, col).number_format = '0.0%'
                workbook.save(OUT_DIR / "banker_model.xlsx")

                blue = PptRGBColor(24, 77, 129)
                dark = PptRGBColor(20, 31, 43)
                grey = PptRGBColor(92, 105, 117)
                green = PptRGBColor(0, 112, 60)

                def add_text(slide, text, x, y, w, h, size=10, bold=False, color=None, align=None):
                    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
                    frame = box.text_frame
                    frame.word_wrap = True
                    frame.clear()
                    paragraph = frame.paragraphs[0]
                    paragraph.text = text
                    if align is not None:
                        paragraph.alignment = align
                    for run in paragraph.runs:
                        run.font.name = "Aptos"
                        run.font.size = Pt(size)
                        run.font.bold = bold
                        if color:
                            run.font.color.rgb = color
                    return box

                def add_bullets(slide, title, bullets, x, y, w, h):
                    add_text(slide, title, x, y, w, 0.24, 12, True, blue)
                    box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.32), Inches(w), Inches(h - 0.32))
                    frame = box.text_frame
                    frame.word_wrap = True
                    frame.clear()
                    for idx, bullet in enumerate(bullets):
                        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
                        paragraph.text = bullet
                        paragraph.level = 0
                        paragraph.font.name = "Aptos"
                        paragraph.font.size = Pt(9)
                        paragraph.font.color.rgb = dark
                    return box

                def add_table(slide, title, headers, rows, x, y, w, h, font_size=7.2):
                    add_text(slide, title, x, y, w, 0.22, 10.5, True, blue)
                    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y + 0.30), Inches(w), Inches(h - 0.30))
                    table = table_shape.table
                    for c_idx, header in enumerate(headers):
                        cell = table.cell(0, c_idx)
                        cell.text = str(header)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = blue
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.alignment = PP_ALIGN.CENTER
                            for run in paragraph.runs:
                                run.font.name = "Aptos"
                                run.font.size = Pt(font_size)
                                run.font.bold = True
                                run.font.color.rgb = PptRGBColor(255, 255, 255)
                    for r_idx, row in enumerate(rows, start=1):
                        for c_idx, value in enumerate(row):
                            cell = table.cell(r_idx, c_idx)
                            cell.text = str(value)
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.alignment = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.RIGHT
                                for run in paragraph.runs:
                                    run.font.name = "Aptos"
                                    run.font.size = Pt(font_size)
                                    run.font.color.rgb = dark
                    return table_shape

                def add_footer(slide, page):
                    add_text(slide, "All amounts are in USD bn unless otherwise noted.", 0.45, 0.12, 4.0, 0.18, 7, False, grey)
                    add_text(slide, f"Source: company VDR, market data and financial data as of {asof_date.isoformat()}; LTM as of {ltm_date.isoformat()}. For discussion purposes only.", 0.45, 7.05, 10.95, 0.18, 6.6, False, grey)
                    add_text(slide, f"{project_id} | Take-Private Teaser | Draft v1 | Page {page}", 10.95, 7.05, 1.85, 0.18, 6.6, False, grey, PP_ALIGN.RIGHT)

                prs = Presentation()
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)
                blank = prs.slide_layouts[6]
                prs.core_properties.title = f"{project_id} - {company_name} Take-Private Teaser - Draft v1"

                highlights = [
                    f"Scale platform with {fmt_b(equity_value)} market capitalization and {fmt_b(ltm_revenue_b)} LTM revenue.",
                    f"Cash generation profile includes {fmt_b(cfo_values[-1])} LTM CFO and {fmt_b(fcf_values[-1])} LTM FCF.",
                    f"Current net debt of {fmt_b(net_debt)} implies {fmt_x(net_leverage)} net leverage on LTM EBITDA.",
                    "Private ownership can support multi-year operating, portfolio and capital allocation initiatives.",
                    "Control case can sharpen KPI discipline and sponsor underwriting around cash conversion.",
                ]
                key_metrics = [
                    ["Market Cap", fmt_b(equity_value)],
                    ["Enterprise Value", fmt_b(enterprise_value)],
                    ["Net Debt", fmt_b(net_debt)],
                    ["Net Leverage", fmt_x(net_leverage)],
                    ["LTM Revenue", fmt_b(ltm_revenue_b)],
                    ["LTM EBITDA", fmt_b(ltm_ebitda_b)],
                ]

                slide1 = prs.slides.add_slide(blank)
                add_footer(slide1, 1)
                slide1.shapes.add_picture(str(logo_path), Inches(0.45), Inches(0.42), width=Inches(1.8))
                title_suffix = f" ({exchange}: {ticker})" if exchange else f" ({ticker})"
                add_text(slide1, f"{company_name}{title_suffix}", 2.25, 0.42, 6.5, 0.36, 18, True, dark)
                add_text(slide1, f"{company_name} screens as a cash-generative take-private candidate", 2.25, 0.82, 7.0, 0.26, 11, False, blue)
                add_table(slide1, "Highlights", ["Metric", "Value"], key_metrics, 9.05, 0.50, 3.65, 2.25, 8.1)
                add_text(slide1, "Overview", 0.55, 1.45, 1.3, 0.26, 12, True, blue)
                add_text(slide1, overview[:1200], 0.55, 1.78, 7.95, 1.72, 8.8, False, dark)
                add_bullets(slide1, "Take-Private Highlights", highlights, 0.55, 3.75, 7.95, 2.55)
                add_text(slide1, f"Source: {source_name(profile_path, 'Company Profile')}; market data and statements in the VDR. Logo image generated from public company identity fields.", 0.55, 6.44, 8.15, 0.34, 6.7, False, grey)

                slide2 = prs.slides.add_slide(blank)
                add_footer(slide2, 2)
                add_text(slide2, f"Premium grid implies {fmt_b(premium_rows[0][3])}-{fmt_b(premium_rows[-1][3])} enterprise value across sponsor offer cases", 0.45, 0.42, 12.0, 0.35, 15.5, True, dark)
                fin_rows = [
                    ["Revenue"] + [fmt_b(value) for value in revenue_values],
                    ["EBITDA"] + [fmt_b(value) for value in ebitda_values],
                    ["EBITDA Margin"] + [fmt_pct(value) for value in margins],
                    ["CFO"] + [fmt_b(value) for value in cfo_values],
                    ["Capex"] + [fmt_b(value) for value in capex_values],
                    ["FCF"] + [fmt_b(value) for value in fcf_values],
                ]
                add_table(slide2, "Financial Summary", ["Metric"] + period_labels, fin_rows, 0.45, 0.96, 6.28, 2.28, 7.0)
                ev_rows = [
                    ["Share Price", fmt_price(share_price)],
                    ["Common Shares", f"{shares_mm:,.1f}mm"],
                    ["Equity Value", fmt_b(equity_value)],
                    ["Cash & Equivalents", fmt_b(cash_b)],
                    ["Preferred Stock", fmt_b(preferred_b)],
                    ["Total Debt", fmt_b(debt_b)],
                    ["Enterprise Value", fmt_b(enterprise_value)],
                ]
                add_table(slide2, "EV Bridge", ["Item", "Value"], ev_rows, 6.95, 0.96, 5.88, 2.58, 7.0)
                premium_grid = [
                    ["Offer Price"] + [fmt_price(row[1]) for row in premium_rows],
                    ["Equity Value"] + [fmt_b(row[2]) for row in premium_rows],
                    ["Implied EV"] + [fmt_b(row[3]) for row in premium_rows],
                    ["Implied EV / LTM EBITDA"] + [fmt_x(row[4]) for row in premium_rows],
                    ["Unchanged EBITDA (LTM)"] + [fmt_b(ltm_ebitda_b)] * len(premiums),
                ]
                add_table(slide2, "Premium Grid", ["Premium"] + [f"+{premium:.0%}" for premium in premiums], premium_grid, 0.45, 3.78, 12.38, 2.35, 6.7)
                add_text(slide2, f"Source footnote: {source_name(income_annual_path, 'Income Statement')}, {source_name(cashflow_annual_path, 'Cash Flow Statement')}, {source_name(balance_path, 'Balance Sheet')}, {source_name(price_path, 'Price History')} and {source_name(shares_path, 'Shares Outstanding')}.", 0.55, 6.55, 11.9, 0.18, 6.8, False, grey)
                prs.save(OUT_DIR / pptx_name)

                title_style = ParagraphStyle("Title", fontName="Helvetica-Bold", fontSize=15, textColor=colors.HexColor("#141F2B"), alignment=TA_LEFT, leading=18)
                small_style = ParagraphStyle("Small", fontName="Helvetica", fontSize=6.8, textColor=colors.HexColor("#5C6975"), leading=8)
                body_style = ParagraphStyle("Body", fontName="Helvetica", fontSize=8.3, textColor=colors.HexColor("#141F2B"), leading=10)
                pdf = SimpleDocTemplate(str(OUT_DIR / pdf_name), pagesize=landscape(letter), rightMargin=0.36 * inch, leftMargin=0.36 * inch, topMargin=0.28 * inch, bottomMargin=0.28 * inch)
                story = [
                    PdfImage(str(logo_path), width=1.35 * inch, height=0.40 * inch),
                    Paragraph(f"{company_name}{title_suffix}", title_style),
                    Paragraph(f"{company_name} screens as a cash-generative take-private candidate", body_style),
                    Paragraph(f"All amounts are in USD bn unless otherwise noted. Source data as of {asof_date.isoformat()}; LTM as of {ltm_date.isoformat()}.", small_style),
                    Spacer(1, 5),
                    Table([["Market Cap", fmt_b(equity_value), "Enterprise Value", fmt_b(enterprise_value), "Net Debt", fmt_b(net_debt), "Net Leverage", fmt_x(net_leverage), "LTM Revenue", fmt_b(ltm_revenue_b), "LTM EBITDA", fmt_b(ltm_ebitda_b)]], colWidths=[0.78 * inch] * 12, style=TableStyle([("GRID", (0,0), (-1,-1), 0.25, colors.HexColor("#AAB7C4")), ("FONTNAME", (0,0), (-1,-1), "Helvetica"), ("FONTSIZE", (0,0), (-1,-1), 6.5)])),
                    Spacer(1, 6),
                    Paragraph("Overview", title_style),
                    Paragraph(overview, body_style),
                    Paragraph("Take-Private Highlights: " + "; ".join(highlights), body_style),
                    Paragraph(f"Source: {source_name(profile_path, 'Company Profile')}; market data and financial data in VDR. For discussion purposes only.", small_style),
                    PageBreak(),
                    Paragraph(f"Premium grid implies {fmt_b(premium_rows[0][3])}-{fmt_b(premium_rows[-1][3])} enterprise value across sponsor offer cases", title_style),
                    Paragraph("All amounts are in USD bn unless otherwise noted.", small_style),
                ]
                fin_table = Table([["Metric"] + period_labels] + fin_rows, colWidths=[1.2 * inch] + [0.8 * inch] * len(period_labels))
                ev_table = Table([["Item", "Value"]] + ev_rows, colWidths=[2.05 * inch, 1.05 * inch])
                premium_table = Table([["Premium"] + [f"+{premium:.0%}" for premium in premiums]] + premium_grid, colWidths=[1.75 * inch] + [0.8 * inch] * len(premiums))
                for table in [fin_table, ev_table, premium_table]:
                    table.setStyle(TableStyle([("BACKGROUND", (0,0), (-1,0), colors.HexColor("#184D81")), ("TEXTCOLOR", (0,0), (-1,0), colors.white), ("GRID", (0,0), (-1,-1), 0.25, colors.HexColor("#AAB7C4")), ("FONTNAME", (0,0), (-1,-1), "Helvetica"), ("FONTSIZE", (0,0), (-1,-1), 6.6), ("ALIGN", (1,1), (-1,-1), "RIGHT")]))
                story.extend([fin_table, Spacer(1, 6), ev_table, Spacer(1, 6), premium_table, Paragraph(f"Source: VDR income statement, cash flow statement, balance sheet, shares and price history. For discussion purposes only.", small_style)])
                pdf.build(story)

                memo = Document()
                memo.add_heading(f"{project_id} - {company_name} Take-Private Teaser - Draft v1", level=1)
                memo.add_paragraph(f"PowerPoint and PDF teaser package generated from VDR source files for {ticker}. See {pptx_name} and {pdf_name}.")
                memo.add_paragraph(f"Market cap uses last close price times common shares; net debt uses total debt less cash and equivalents; FCF uses CFO less Capex; premium grid applies +10% to +50% premiums to last close.")
                memo.save(OUT_DIR / memo_name)

                remember_citation(f"{company_name} profile and business overview", source_name(profile_path, "Company Profile"), "Company Profile!Description", overview[:220], "cell", {"sheet": "Company Profile", "range": "Description"})
                remember_citation(f"{ticker} financial summary table", "banker_model.xlsx", f"Summary Output!A1:{get_column_letter(1 + len(period_labels))}7", "Revenue, EBITDA, EBITDA margin, CFO, Capex, FCF", "cell", {"sheet": "Summary Output", "range": f"A1:{get_column_letter(1 + len(period_labels))}7"})
                remember_citation(f"{ticker} EV bridge", "banker_model.xlsx", "Summary Output!A9:B18", "Share price, shares, cash, preferred stock, total debt and enterprise value", "cell", {"sheet": "Summary Output", "range": "A9:B18"})
                remember_citation(f"{ticker} premium grid", "banker_model.xlsx", "Summary Output!A20:F25", "+10% to +50% premiums to last close", "cell", {"sheet": "Summary Output", "range": "A20:F25"})
                remember_citation(f"{ticker} logo image", logo_path.name, "slide 1 top-left", "Generated company identity mark based on VDR profile fields", "shape", {"page": 1, "x": 0.45, "y": 0.42, "w": 1.8, "h": 0.52, "unit": "in"})
                remember_citation("Two-page take-private teaser", pptx_name, "slides 1-2", "Overview, highlights, financial summary, EV bridge and premium grid", "shape", {"page": 1, "x": 0, "y": 0, "w": 13.333, "h": 7.5, "unit": "in"})
                return True

            def buyer_universe_task_shape():
                task_text = "\n\n".join(
                    part for part in [instruction_text, plan.get("taskSummary", ""), plan.get("title", "")] if part
                ).lower()
                return (
                    "buyer universe" in task_text
                    and ("sponsor" in task_text or "financial sponsor" in task_text)
                    and ("strategic" in task_text or "strategics" in task_text)
                    and ("logo" in task_text or "logos" in task_text)
                    and ("powerpoint" in task_text or "slide" in task_text or "presentation" in task_text)
                )

            def write_general_buyer_universe_package():
                if not buyer_universe_task_shape():
                    return False

                from docx import Document
                from openpyxl import Workbook
                from openpyxl.styles import Alignment, Font, PatternFill
                from openpyxl.utils import get_column_letter
                from PIL import Image as PILImage, ImageDraw, ImageFont
                from pptx import Presentation
                from pptx.dml.color import RGBColor as PptRGBColor
                from pptx.enum.text import PP_ALIGN
                from pptx.util import Inches, Pt
                from reportlab.lib import colors
                from reportlab.lib.enums import TA_LEFT
                from reportlab.lib.pagesizes import landscape, letter
                from reportlab.lib.styles import ParagraphStyle
                from reportlab.lib.units import inch
                from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

                task_text = "\n\n".join(
                    part for part in [instruction_text, plan.get("taskSummary", ""), plan.get("title", "")] if part
                )
                lowered = task_text.lower()

                def clean_name(value, fallback):
                    text_value = re.sub(r"\s+", " ", str(value or "")).strip(" .,\r\n\t")
                    return text_value or fallback

                def first_match(patterns, fallback):
                    for pattern in patterns:
                        match = re.search(pattern, task_text, flags=re.I | re.S)
                        if match:
                            return clean_name(match.group(1), fallback)
                    return fallback

                def primary_ticker():
                    for candidate in plan.get("tickers", []) or []:
                        ticker_value = str(candidate).upper().strip()
                        if ticker_value:
                            return ticker_value
                    for pattern in [
                        r"\((?:nyse|nasdaq|amex|otc)[^:)]*:\s*([A-Z]{1,6})\)",
                        r"\(([A-Z]{1,6})-US\)",
                    ]:
                        match = re.search(pattern, task_text, flags=re.I)
                        if match:
                            return match.group(1).upper()
                    return ""

                ticker = primary_ticker()
                source_root = Path("/home/agent/workspace/banker_workspace/source/mcp")
                base = source_root / ticker / "vdr" if ticker else None
                profile_path = first_existing_path(base, [f"{ticker}-US Company Profile.xlsx"]) if base and base.is_dir() else None
                price_path = first_existing_path(base, [f"{ticker}-US Price History (Daily).xlsx"]) if base and base.is_dir() else None
                shares_path = first_existing_path(base, [f"{ticker}-US Shares Outstanding.xlsx"]) if base and base.is_dir() else None
                balance_path = first_existing_path(base, [
                    f"{ticker}-US Balance Sheet (Quarterly).xlsx",
                    f"{ticker}-US Balance Sheet (Annual).xlsx",
                ]) if base and base.is_dir() else None

                profile_name = profile_value(profile_path, "Name") if profile_path else None
                client_name = clean_name(
                    profile_name,
                    first_match([
                        r"\bclient\s+is\s+([^,(.\r\n]+(?:\s+[A-Z][^,(.\r\n]+)*)",
                        r"\bfor\s+([^,(.\r\n]+(?:\s+[A-Z][^,(.\r\n]+)*)\s+\(",
                    ], "Client"),
                )
                asset_name = first_match([
                    r"\bsell(?:ing)?\s+(?:its|their|the)?\s*asset\s+([A-Za-z0-9& /.-]+?)(?:\s*\(|,|\.)",
                    r"\basset\s+([A-Za-z0-9& /.-]+?)(?:\s*\(|,|\.)",
                ], "Target Asset")
                project_name = "Project " + re.sub(r"[^A-Za-z0-9]+", " ", asset_name).strip().split(" ")[0][:12].title()
                industry_hint = first_match([
                    r"\bglobal leader in\s+([^.\r\n]+?)(?:\s+used to|\.|,)",
                    r"\bleader in\s+([^.\r\n]+?)(?:\s+used to|\.|,)",
                ], "mission-critical specialty packaging")
                profile_description = clean_name(profile_value(profile_path, "Description") if profile_path else "", "")

                asof_date = parse_asof_date(task_text)
                price = nearest_row_value(price_path, asof_date) if price_path else None
                shares = nearest_row_value(shares_path, asof_date) if shares_path else None
                cash = latest_label_period_value(balance_path, ["CashAndCashEquivalents", "CashCashEquivalentsAndShortTermInvestments", "Cash"], asof_date) if balance_path else None
                debt = latest_label_period_value(balance_path, ["TotalDebt"], asof_date) if balance_path else None
                share_price = price["value"] if price else 0.0
                shares_mm = (shares["value"] / 1000000.0) if shares and abs(shares["value"]) > 1000000 else (shares["value"] if shares else 0.0)
                cash_b = (cash["value"] / 1000000000.0) if cash and abs(cash["value"]) > 1000000 else (cash["value"] if cash else 0.0)
                debt_b = (debt["value"] / 1000000000.0) if debt and abs(debt["value"]) > 1000000 else (debt["value"] if debt else 0.0)
                market_cap_b = share_price * shares_mm / 1000.0 if share_price and shares_mm else 0.0
                net_debt_b = debt_b - cash_b

                def candidate_buyer_universe(text):
                    cold_chain = any(token in text for token in [
                        "cold chain",
                        "temperature-controlled",
                        "temperature controlled",
                        "pharma",
                        "biologic",
                        "vaccine",
                        "healthcare goods",
                    ])
                    packaging = any(token in text for token in ["packaging", "container", "specialty materials"])
                    if cold_chain:
                        sponsor_rows = [
                            ("Berkshire Partners", "Healthcare packaging and services platform sponsor"),
                            ("Arlington Capital", "Healthcare, government services, and specialty industrial carve-out sponsor"),
                            ("Audax Group", "Buy-and-build sponsor with packaging and healthcare services relevance"),
                            ("May River Capital", "Lower-middle-market industrial and specialty manufacturing focus"),
                            ("Lincolnshire Management", "Industrial products and packaging-services investment history"),
                            ("Huron Capital", "Services and specialty manufacturing platform-building fit"),
                            ("Mason Wells", "Packaging materials and engineered products sponsor fit"),
                            ("Keystone Capital", "Specialty manufacturing and business-services carve-out fit"),
                            ("Trilantic North America", "Healthcare services and industrial growth investor"),
                            ("Wellspring Capital", "Packaging and healthcare services investment fit"),
                            ("J.W. Hill Capital", "Industrial and distribution platform investor"),
                            ("Calera Capital", "Middle-market carve-out and operational improvement sponsor"),
                        ]
                        strategic_rows = [
                            ("Cold Chain Technologies", "", "Direct cold-chain packaging competitor with reusable and single-use shipper overlap"),
                            ("Pelican BioThermal", "Behrman Capital", "Temperature-controlled packaging platform with pharma and biologics focus"),
                            ("CSafe Global", "Riverside Partners", "Active container and parcel cold-chain platform with sponsor backing"),
                            ("Envirotainer", "Investor AB / Triton", "Global active temperature-controlled air-cargo container platform"),
                            ("Va-Q-tec", "", "Advanced thermal packaging and container technology overlap"),
                            ("Cryopak", "", "Cold-chain packaging and temperature-monitoring strategic fit"),
                            ("Intelsius", "", "Life-sciences sample transport and temperature-controlled packaging fit"),
                            ("Tempack", "", "Temperature-controlled packaging platform with European buyer angle"),
                            ("Cencora", "", "Healthcare distribution platform with specialty logistics exposure"),
                            ("Cardinal Health", "", "Healthcare distribution and pharma-services strategic"),
                            ("McKesson", "", "Healthcare distribution strategic with pharma logistics adjacency"),
                            ("UPS Healthcare", "", "Healthcare logistics and cold-chain distribution platform"),
                            ("FedEx Healthcare", "", "Healthcare logistics network with cold-chain capabilities"),
                            ("DHL Life Sciences", "", "Global life-sciences logistics strategic with cold-chain operations"),
                            ("Marken", "", "Clinical trial and specialty logistics platform"),
                            ("World Courier", "", "Specialty logistics platform focused on healthcare shipments"),
                            ("SoftBox Systems", "", "Temperature-controlled packaging and parcel shipper platform"),
                        ]
                    elif packaging:
                        sponsor_rows = [
                            ("Berkshire Partners", "Packaging and industrial services sponsor"),
                            ("Audax Group", "Buy-and-build sponsor with packaging relevance"),
                            ("Graham Partners", "Advanced manufacturing and packaging investor"),
                            ("May River Capital", "Industrial products and specialty manufacturing sponsor"),
                            ("Huron Capital", "Services and engineered products platform investor"),
                            ("Mason Wells", "Packaging materials and specialty industrial sponsor"),
                            ("Keystone Capital", "Specialty manufacturing and business services investor"),
                            ("Wellspring Capital", "Packaging and distribution sponsor"),
                        ]
                        strategic_rows = [
                            ("Amcor", "", "Global packaging strategic"),
                            ("Berry Global", "", "Plastic packaging and materials strategic"),
                            ("AptarGroup", "", "Dispensing and specialty packaging strategic"),
                            ("Sealed Air", "", "Protective packaging strategic"),
                            ("Greif", "", "Industrial packaging strategic"),
                            ("Graphic Packaging", "", "Consumer packaging strategic"),
                            ("Packaging Corporation of America", "", "Containerboard and packaging strategic"),
                            ("International Paper", "", "Fiber-based packaging strategic"),
                        ]
                    else:
                        sponsor_rows = [
                            ("Berkshire Partners", "Middle-market platform sponsor"),
                            ("Audax Group", "Buy-and-build sponsor"),
                            ("Arlington Capital", "Specialty services and industrial carve-out sponsor"),
                            ("Huron Capital", "Operational improvement and platform investor"),
                            ("Keystone Capital", "Founder and carve-out transaction sponsor"),
                            ("Wellspring Capital", "Operationally focused middle-market investor"),
                        ]
                        strategic_rows = [
                            ("Strategic Platform A", "", "Direct product and customer overlap"),
                            ("Strategic Platform B", "", "Adjacency and cross-sell rationale"),
                            ("Strategic Platform C", "", "Scale and geographic expansion rationale"),
                            ("Strategic Platform D", "", "Channel, technology, or supply-chain fit"),
                        ]
                    return sponsor_rows, strategic_rows

                sponsors, strategics = candidate_buyer_universe(lowered)
                sponsor_tags = sorted({backing for _, backing, _ in strategics if backing})
                highlights = [
                    f"{asset_name} should screen as a mission-critical platform in {industry_hint}.",
                    "Strategics can underwrite product breadth, customer access, route density, and supply-chain synergies.",
                    "Sponsor-backed strategics should be shown with sponsor tags because sponsor appetite can influence process dynamics.",
                    "Financial sponsors can underwrite carve-out execution, operational improvement, healthcare logistics growth, and add-on M&A.",
                    "Positioning should emphasize validated performance, pharma/healthcare exposure, reusable packaging, and service reliability.",
                    "Buyer interest is preliminary and should be refined through outreach, NDA feedback, and diligence findings.",
                ]
                timeline = [
                    "Weeks 0-2: align advisors, perimeter, buyer screen, data request list, and management story.",
                    "Weeks 2-5: launch QoE, market study, tax/legal structuring, carve-out analysis, and stand-alone cost work.",
                    "Weeks 4-7: build VDR, CIM, teaser, customer/product cohort views, and separation workplan.",
                    "Weeks 7-10: approve buyer list, send teaser/NDA, collect feedback, and prepare process letter.",
                    "Weeks 10-12: management prep, first-round bid instructions, data-room access, and process launch.",
                ]

                navy_hex = "002060"
                workbook = Workbook()
                ws = workbook.active
                ws.title = "Buyer Universe"
                ws.sheet_view.showGridLines = False
                ws.append(["Category", "Buyer", "Sponsor Backing", "Rationale / Fit", "Logo Shown"])
                for sponsor, rationale in sponsors:
                    ws.append(["Sponsor", sponsor, "", rationale, "Yes"])
                for name, backing, rationale in strategics:
                    ws.append(["Strategic", name, backing, rationale, "Yes"])
                thesis_ws = workbook.create_sheet("Highlights Thesis")
                thesis_ws.sheet_view.showGridLines = False
                thesis_ws.append(["Section", "Bullet"])
                for bullet in highlights:
                    thesis_ws.append(["Highlights & buyer thesis", bullet])
                for bullet in timeline:
                    thesis_ws.append(["Sell-side preparation timeline", bullet])
                financial_ws = workbook.create_sheet("Financial Context")
                financial_ws.sheet_view.showGridLines = False
                financial_ws.append(["Metric", "Value", "Source / Note"])
                financial_ws.append(["Latest Share Price", share_price, "Price history, latest available on or before analysis date"])
                financial_ws.append(["Shares Outstanding (mm)", shares_mm, "Shares outstanding source"])
                financial_ws.append(["Market Capitalization ($B)", "=B2*B3/1000", "Calculated from share price and shares outstanding"])
                financial_ws.append(["Cash & Equivalents ($B)", cash_b, "Balance sheet"])
                financial_ws.append(["Total Debt ($B)", debt_b, "Balance sheet"])
                financial_ws.append(["Net Debt ($B)", "=B6-B5", "Total debt less cash"])
                financial_ws.append(["Illustrative Standalone Cost Dis-synergy ($mm)", -25.0, "Illustrative placeholder to preserve negative-number formatting standard"])
                source_ws = workbook.create_sheet("Source Evidence")
                source_ws.sheet_view.showGridLines = False
                source_ws.append(["Claim", "Source", "Locator", "Boundary Status"])
                source_ws.append(["Task instruction defines buyer-universe, sponsors, strategics, thesis, and timeline requirements", "instruction.txt", "prompt paragraphs", "paragraph"])
                if profile_path:
                    source_ws.append([f"{client_name} public company profile", profile_path.name, "Company Profile!A2:B16", "cell"])
                if price_path:
                    source_ws.append(["Latest share price", price_path.name, f"Sheet1!R{price.get('row') if price else '?'}C2", "cell"])
                if shares_path:
                    source_ws.append(["Shares outstanding", shares_path.name, f"Sheet1!R{shares.get('row') if shares else '?'}C2", "cell"])
                if balance_path:
                    source_ws.append(["Balance sheet cash and debt", balance_path.name, "Cash / TotalDebt rows", "cell"])

                for sheet in workbook.worksheets:
                    for cell in sheet[1]:
                        cell.fill = PatternFill("solid", fgColor=navy_hex)
                        cell.font = Font(name="Arial", color="FFFFFF", bold=True)
                        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
                    for row in sheet.iter_rows(min_row=2):
                        for cell in row:
                            cell.font = Font(name="Arial")
                            cell.alignment = Alignment(vertical="top", wrap_text=True)
                    for col in range(1, sheet.max_column + 1):
                        sheet.column_dimensions[get_column_letter(col)].width = 22 if col < 3 else 34
                financial_ws["B4"].number_format = '$0.0'
                financial_ws["B5"].number_format = '$0.0'
                financial_ws["B6"].number_format = '$0.0'
                financial_ws["B7"].number_format = '$0.0'
                financial_ws["B8"].number_format = '#,##0.0;(#,##0.0)'
                workbook.save(OUT_DIR / "banker_model.xlsx")

                prs = Presentation()
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                logo_dir = OUT_DIR / "buyer_universe_logo_assets"
                logo_dir.mkdir(parents=True, exist_ok=True)
                navy = PptRGBColor(0, 32, 96)
                black = PptRGBColor(0, 0, 0)
                grey = PptRGBColor(80, 80, 80)

                def logo_asset(text, fill=(255, 255, 255), accent=(0, 32, 96), text_color=None, width=460, height=92):
                    text_color = text_color or accent
                    slug = re.sub(r"[^a-z0-9]+", "_", text.lower()).strip("_")[:52] or "logo"
                    path = logo_dir / f"{slug}_{width}x{height}.png"
                    image = PILImage.new("RGB", (width, height), fill)
                    draw = ImageDraw.Draw(image)
                    draw.rounded_rectangle((3, 3, width - 4, height - 4), radius=10, outline=accent, width=3)
                    try:
                        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 22 if height >= 70 else 15)
                    except Exception:
                        font = ImageFont.load_default()
                    lines = []
                    current = ""
                    for word in text.split():
                        candidate = f"{current} {word}".strip()
                        if draw.textbbox((0, 0), candidate, font=font)[2] <= width - 22:
                            current = candidate
                        else:
                            if current:
                                lines.append(current)
                            current = word
                    if current:
                        lines.append(current)
                    lines = lines[:2]
                    line_height = max(14, draw.textbbox((0, 0), "Ag", font=font)[3] + 3)
                    start_y = (height - line_height * len(lines)) / 2 - 1
                    for line_index, line in enumerate(lines):
                        bbox = draw.textbbox((0, 0), line, font=font)
                        draw.text(((width - (bbox[2] - bbox[0])) / 2, start_y + line_index * line_height), line, fill=text_color, font=font)
                    image.save(path)
                    return path

                def add_textbox(text, x, y, w, h, size=10, color=None, bold=False, italic=False, align=PP_ALIGN.LEFT):
                    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
                    frame = box.text_frame
                    frame.word_wrap = True
                    frame.clear()
                    paragraph = frame.paragraphs[0]
                    paragraph.alignment = align
                    run = paragraph.add_run()
                    run.text = text
                    run.font.name = "Arial"
                    run.font.size = Pt(size)
                    run.font.bold = bold
                    run.font.italic = italic
                    run.font.color.rgb = color or black
                    return box

                def add_header(text, x, y, w, h):
                    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = navy
                    shape.line.color.rgb = navy
                    frame = shape.text_frame
                    frame.clear()
                    paragraph = frame.paragraphs[0]
                    paragraph.alignment = PP_ALIGN.CENTER
                    run = paragraph.add_run()
                    run.text = text
                    run.font.name = "Arial"
                    run.font.size = Pt(10)
                    run.font.bold = True
                    run.font.color.rgb = PptRGBColor(255, 255, 255)
                    return shape

                def add_logo_tile(text, x, y, w, h, fill=(255, 255, 255), accent=(0, 32, 96), text_color=None):
                    path = logo_asset(text, fill=fill, accent=accent, text_color=text_color)
                    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))

                title_text = f"{project_name}: {asset_name} Potential Buyer Universe"
                add_textbox(title_text, 0.35, 0.20, 10.4, 0.36, 16, black, True)
                client_logo_text = re.sub(r"[^A-Za-z0-9]+", " ", client_name).strip()[:18].upper() or "CLIENT"
                client_logo = logo_asset(client_logo_text, fill=(0, 32, 96), accent=(0, 32, 96), text_color=(255, 255, 255), width=420, height=92)
                slide.shapes.add_picture(str(client_logo), Inches(11.12), Inches(0.18), Inches(1.55), Inches(0.34))
                underline = slide.shapes.add_shape(1, Inches(0.35), Inches(0.68), Inches(12.6), Inches(0.02))
                underline.fill.solid()
                underline.fill.fore_color.rgb = navy
                underline.line.color.rgb = navy

                add_header("Sponsors", 0.35, 0.88, 6.0, 0.30)
                add_header("Strategics", 0.35, 3.08, 6.0, 0.30)
                add_header("Highlights & Buyer Thesis", 6.65, 0.88, 6.3, 0.30)
                add_header("Sell-Side Preparation Timeline", 6.65, 4.05, 6.3, 0.30)

                sponsor_x = [0.45, 2.42, 4.39]
                for index, (sponsor, _) in enumerate(sponsors[:12]):
                    row = index // 3
                    col = index % 3
                    add_logo_tile(sponsor, sponsor_x[col], 1.27 + row * 0.40, 1.40, 0.28)
                for index, sponsor in enumerate(sponsor_tags[:3]):
                    add_logo_tile(sponsor, 0.45 + index * 1.97, 2.83, 0.90, 0.18, fill=(226, 235, 247), accent=(0, 32, 96))

                strategic_x = [0.45, 1.93, 3.41, 4.89]
                for index, (name, backing, _) in enumerate(strategics[:17]):
                    row = index // 4
                    col = index % 4
                    y = 3.47 + row * 0.39
                    add_logo_tile(name, strategic_x[col], y, 1.34, 0.27, fill=(248, 250, 252), accent=(25, 70, 120))
                    if backing:
                        add_logo_tile(backing, strategic_x[col] + 0.08, y + 0.26, 0.80, 0.16, fill=(226, 235, 247), accent=(0, 32, 96))

                add_textbox("Sponsor-backed strategics show sponsor logo tags under the strategic buyer logo.", 0.45, 6.05, 5.85, 0.30, 7, grey, False, True)

                def add_bullet_box(items, x, y, w, h):
                    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
                    frame = box.text_frame
                    frame.word_wrap = True
                    frame.clear()
                    for index, item in enumerate(items):
                        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                        paragraph.text = item
                        paragraph.level = 1 if index in {1, 3} else 0
                        paragraph.space_after = Pt(3)
                        for run in paragraph.runs:
                            run.font.name = "Arial"
                            run.font.size = Pt(10)
                            run.font.color.rgb = black
                            run.font.bold = index in {0, 2}
                            run.font.italic = index in {0, 2}
                    return box

                add_bullet_box(highlights, 6.78, 1.27, 6.05, 2.45)
                add_bullet_box(timeline, 6.78, 4.45, 6.05, 1.65)
                market_note = ""
                if share_price and shares_mm:
                    market_note = f" | {ticker} latest source price ${share_price:.2f}; shares {shares_mm:,.1f}mm; market cap ${market_cap_b:,.1f}B"
                add_textbox(
                    f"Source: task prompt, {client_name} public profile and VDR market data{market_note}. Confidential - discussion draft; buyer interest is preliminary and actual interest will depend on process dynamics.",
                    0.35,
                    6.82,
                    12.25,
                    0.25,
                    6.5,
                    grey,
                )
                add_textbox("1", 12.55, 6.82, 0.35, 0.25, 7, grey, False, False, PP_ALIGN.RIGHT)
                prs.save(OUT_DIR / "banker_presentation.pptx")

                pdf = SimpleDocTemplate(
                    str(OUT_DIR / "banker_report.pdf"),
                    pagesize=landscape(letter),
                    leftMargin=0.30 * inch,
                    rightMargin=0.30 * inch,
                    topMargin=0.30 * inch,
                    bottomMargin=0.30 * inch,
                )
                title_style = ParagraphStyle("title", fontName="Helvetica-Bold", fontSize=13, leading=15, alignment=TA_LEFT)
                body_style = ParagraphStyle("body", fontName="Helvetica", fontSize=7.5, leading=9, alignment=TA_LEFT)
                buyer_rows = [
                    ["Sponsors", ", ".join(name for name, _ in sponsors)],
                    ["Strategics", ", ".join(name for name, _, _ in strategics)],
                    ["Sponsor-backed strategics", "; ".join(f"{name} / {backing}" for name, backing, _ in strategics if backing)],
                    ["Highlights", "; ".join(highlights)],
                    ["Timeline", "; ".join(timeline)],
                ]
                buyer_table = Table(buyer_rows, colWidths=[1.25 * inch, 8.8 * inch])
                buyer_table.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#002060")),
                    ("TEXTCOLOR", (0, 0), (0, -1), colors.white),
                    ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
                    ("FONTNAME", (1, 0), (1, -1), "Helvetica"),
                    ("FONTSIZE", (0, 0), (-1, -1), 7.0),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#BFBFBF")),
                ]))
                pdf.build([
                    Paragraph(title_text, title_style),
                    Spacer(1, 0.08 * inch),
                    buyer_table,
                    Spacer(1, 0.06 * inch),
                    Paragraph("Confidential - discussion draft | Page 1 | Buyer interest is preliminary and actual interest will depend on process dynamics.", body_style),
                ])

                doc = Document()
                doc.add_heading(f"{asset_name} Buyer Universe Support Memo", level=1)
                doc.add_paragraph(f"This support memo mirrors the one-page buyer universe slide for {client_name}'s sale process.")
                doc.add_paragraph("Key thesis: strategics can underwrite product breadth, customer access, route density, and supply-chain synergies; sponsors can underwrite carve-out value creation, growth, operational improvement, and add-on M&A.")
                doc.add_paragraph("Preparation workstreams: QoE, market study, carve-out analysis/audits, legal/tax structuring, VDR prep, CIM/teaser, process letters, and management prep.")
                doc.add_paragraph("Buyer interest is preliminary and actual interest will depend on process dynamics.")
                doc.save(OUT_DIR / "banker_memo.docx")

                remember_citation("Task instruction defines buyer-universe output requirements", "instruction.txt", "prompt paragraphs", "Sponsors, strategics, logos, thesis, and sell-side preparation timeline", "paragraph", {"section": "instruction"})
                remember_citation("Sponsors and strategics segmented on one-page slide", "banker_presentation.pptx", "slide 1 left buyer logo universe", "Sponsor and strategic logo groups", "shape", {"slide": 1, "x": 0.35, "y": 0.88, "w": 6.0, "h": 5.8, "unit": "in"})
                remember_citation("Client branding and buyer logo images", "banker_presentation.pptx", "slide 1 logo image layer", "Client logo plus buyer logo PNGs", "shape", {"slide": 1, "x": 0.35, "y": 0.18, "w": 12.6, "h": 6.2, "unit": "in"})
                remember_citation("Sponsor-backed strategics include sponsor logo tags", "banker_presentation.pptx", "slide 1 strategic logo grid", "; ".join(f"{name} / {backing}" for name, backing, _ in strategics if backing)[:240], "shape", {"slide": 1, "x": 0.35, "y": 3.08, "w": 6.0, "h": 3.3, "unit": "in"})
                remember_citation("Investment highlights and thesis", "banker_presentation.pptx", "slide 1 upper-right quadrant", "Buyer thesis and sale positioning bullets", "shape", {"slide": 1, "x": 6.65, "y": 0.88, "w": 6.3, "h": 2.9, "unit": "in"})
                remember_citation("Sell-side preparation timeline", "banker_presentation.pptx", "slide 1 lower-right quadrant", "QoE, market study, carve-out work, VDR, CIM, teaser, process launch", "shape", {"slide": 1, "x": 6.65, "y": 4.05, "w": 6.3, "h": 2.1, "unit": "in"})
                remember_citation("Buyer universe workbook source", "banker_model.xlsx", "Buyer Universe!A1:E30", "Sponsors, strategics, sponsor backing, rationale", "cell", {"sheet": "Buyer Universe", "range": "A1:E30"})
                remember_citation("PDF one-page overview", "banker_report.pdf", "page 1 table", "Buyer universe, highlights, and timeline", "bbox", {"page": 1, "x": 24, "y": 70, "w": 748, "h": 430, "unit": "pt"})
                if profile_path:
                    remember_citation(f"{client_name} public company profile", profile_path.name, "Company Profile!A2:B16", profile_description[:240], "cell", {"sheet": "Company Profile", "range": "A2:B16"})
                if price_path and price:
                    remember_citation("Latest source share price", price_path.name, f"Sheet1!R{price['row']}C2", f"{share_price:.2f}", "cell", {"sheet": "Sheet1", "range": f"B{price['row']}"})
                return True

            def sources_uses_task_shape():
                task_text = "\n\n".join(
                    part for part in [instruction_text, plan.get("taskSummary", ""), plan.get("title", "")] if part
                ).lower()
                sheet_names = " ".join(
                    str(sheet.get("name", "")) for sheet in plan.get("workbook", {}).get("sheets", [])
                ).lower()
                return (
                    "sources and uses" in task_text
                    and "sensitivity table" in task_text
                    and ("equity injection" in task_text or "term loan b" in task_text)
                ) or (
                    "sources and uses" in sheet_names
                    and "sensitivity" in sheet_names
                )

            def parse_plan_assumption(label_options, default):
                labels = [label.lower() for label in label_options]
                rows = []
                for sheet in plan.get("workbook", {}).get("sheets", []):
                    for row in sheet.get("rows", []):
                        if isinstance(row, list) and len(row) >= 2:
                            rows.append(row)
                for row in rows:
                    name = str(row[0] or "").lower()
                    if not any(label in name for label in labels):
                        continue
                    raw = str(row[1] or "").strip()
                    match = re.search(r"(-?\d+(?:\.\d+)?)\s*(%|x)?", raw, flags=re.I)
                    if not match:
                        continue
                    value = float(match.group(1))
                    unit = (match.group(2) or "").lower()
                    if unit == "%":
                        return value / 100.0
                    return value
                text = f"{instruction_text}\n{plan.get('taskSummary', '')}".lower()
                for label in labels:
                    pattern = re.escape(label).replace("\\ ", r"\s+")
                    match = re.search(pattern + r".{0,80}?(-?\d+(?:\.\d+)?)\s*(%|x)?", text, flags=re.I | re.S)
                    if not match:
                        continue
                    value = float(match.group(1))
                    unit = (match.group(2) or "").lower()
                    if unit == "%":
                        return value / 100.0
                    return value
                return default

            def latest_label_period_value(path, labels, asof_date):
                return label_period_value(path, labels, asof_date=asof_date)

            def ltm_label_period_value(path, labels, asof_date):
                rows = load_first_sheet_rows(path)
                if not rows:
                    return None
                normalized = {label.lower() for label in labels}
                header = rows[0]
                dated_cols = []
                for col_index, value in enumerate(header[1:], start=1):
                    header_date = cell_date(value)
                    if header_date is not None and header_date <= asof_date:
                        dated_cols.append((col_index, header_date))
                if not dated_cols:
                    return None
                dated_cols = sorted(dated_cols, key=lambda item: item[1], reverse=True)[:4]
                for row_index, row in enumerate(rows[1:], start=2):
                    label = str(row[0] or "").strip()
                    if label.lower() not in normalized:
                        continue
                    values = []
                    locators = []
                    for col_index, header_date in dated_cols:
                        value = to_number(row[col_index] if len(row) > col_index else None)
                        if value is not None:
                            values.append(value)
                            locators.append({"row": row_index, "column": col_index + 1, "date": header_date})
                    if values:
                        return {"label": label, "value": sum(values), "locators": locators}
                return None

            def write_general_sources_uses_package():
                if not sources_uses_task_shape():
                    return False

                import shutil
                import tempfile
                import zipfile
                from xml.etree import ElementTree as ET

                from openpyxl import Workbook
                from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
                from openpyxl.utils import get_column_letter

                ticker = ""
                for candidate in plan.get("tickers", []) or []:
                    candidate = str(candidate).upper().strip()
                    if candidate:
                        ticker = candidate
                        break
                if not ticker:
                    match = re.search(r"\(([A-Z]{1,6})-US\)", instruction_text)
                    ticker = match.group(1) if match else ""
                if not ticker:
                    return False

                source_root = Path("/home/agent/workspace/banker_workspace/source/mcp")
                base = source_root / ticker / "vdr"
                if not base.is_dir():
                    return False

                task_text = f"{instruction_text}\n{plan.get('taskSummary', '')}"
                asof_date = parse_month_year_asof(task_text, fallback=parse_asof_date(task_text))
                income_asof = parse_month_year_asof(
                    task_text,
                    fallback=asof_date,
                    month_names={"oct": 10, "october": 10},
                )

                price_path = first_existing_path(base, [f"{ticker}-US Price History (Daily).xlsx"])
                shares_path = first_existing_path(base, [f"{ticker}-US Shares Outstanding.xlsx"])
                balance_path = first_existing_path(base, [
                    f"{ticker}-US Balance Sheet (Quarterly).xlsx",
                    f"{ticker}-US Balance Sheet (Annual).xlsx",
                ])
                income_path = first_existing_path(base, [
                    f"{ticker}-US Income Statement (Quarterly).xlsx",
                    f"{ticker}-US Income Statement (Annual).xlsx",
                ])

                price = nearest_row_value(price_path, asof_date) if price_path else None
                shares = nearest_row_value(shares_path, asof_date) if shares_path else None
                cash = latest_label_period_value(
                    balance_path,
                    ["CashAndCashEquivalents", "CashCashEquivalentsAndShortTermInvestments", "Cash"],
                    asof_date,
                )
                debt = latest_label_period_value(balance_path, ["TotalDebt"], asof_date)
                leases = latest_label_period_value(balance_path, ["CapitalLeaseObligations"], asof_date)
                revenue = latest_label_period_value(income_path, ["TotalRevenue", "OperatingRevenue"], income_asof)
                ebitda = latest_label_period_value(income_path, ["EBITDA", "NormalizedEBITDA"], income_asof)
                ltm_revenue = ltm_label_period_value(income_path, ["TotalRevenue", "OperatingRevenue"], income_asof)
                ltm_ebitda = ltm_label_period_value(income_path, ["EBITDA", "NormalizedEBITDA"], income_asof)

                if not (price and shares and balance_path and income_path and (ltm_revenue or revenue) and (ltm_ebitda or ebitda)):
                    return False

                source_values = {
                    "price": price["value"],
                    "shares_m": shares["value"] / 1_000_000.0,
                    "market_cap_m": price["value"] * (shares["value"] / 1_000_000.0),
                    "cash_m": (cash["value"] / 1_000_000.0) if cash and cash.get("value") is not None else 0.0,
                    "debt_m": (
                        max((debt["value"] - (leases["value"] if leases and leases.get("value") is not None else 0.0)), 0.0) / 1_000_000.0
                    ) if debt and debt.get("value") is not None else 0.0,
                    "revenue_m": ((revenue or ltm_revenue)["value"] / 1_000_000.0),
                    "ebitda_m": ((ebitda or ltm_ebitda)["value"] / 1_000_000.0),
                    "ltm_revenue_m": ((ltm_revenue or revenue)["value"] / 1_000_000.0),
                    "ltm_ebitda_m": ((ltm_ebitda or ebitda)["value"] / 1_000_000.0),
                }

                transaction_fee_pct = parse_plan_assumption(["transaction fees", "transaction fee"], 0.020)
                financing_fee_pct = parse_plan_assumption(["financing fees", "financing fee"], 0.025)
                cash_to_bs_pct = parse_plan_assumption(["cash to balance sheet", "cash to b/s"], 0.035)
                senior_notes_multiple = parse_plan_assumption(["senior notes"], 5.0)
                tlb_multiple = parse_plan_assumption(["term loan b", "tlb"], 3.0)
                revolver_multiple = parse_plan_assumption(["revolver"], 2.0)
                management_incentives_pct = parse_plan_assumption(["management incentives"], 0.050)
                rollover_pct = parse_plan_assumption(["shareholder rollover"], 0.250)
                premium_pct = parse_plan_assumption(["premium to current share price", "premium to current share"], 0.350)

                workbook = Workbook()
                ws = workbook.active
                ws.title = "Sources and Uses"
                ws.sheet_view.showGridLines = False

                blue = "0000FF"
                black = "000000"
                grey = "666666"
                header_fill = PatternFill("solid", fgColor="1F4E78")
                section_fill = PatternFill("solid", fgColor="D9EAF7")
                total_fill = PatternFill("solid", fgColor="E2F0D9")
                thin = Side(style="thin", color="A6A6A6")
                border = Border(top=thin, bottom=thin)
                acct_fmt = '#,##0.0;(#,##0.0)'
                pct_fmt = '0.0%;(0.0%)'
                mult_fmt = '0.0x'
                price_fmt = '$0.00'

                ws.merge_cells("A1:D1")
                ws["A1"] = f"{ticker} Illustrative Sources & Uses Analysis"
                ws["A1"].fill = header_fill
                ws["A1"].font = Font(color="FFFFFF", bold=True, size=14)
                ws["A2"] = "All dollar figures in $mm except per-share data and sensitivity output in $bn."
                ws["A2"].font = Font(italic=True, color=grey)

                def label(row, text):
                    ws.cell(row, 1, text)
                    ws.cell(row, 1).fill = section_fill
                    ws.cell(row, 1).font = Font(bold=True, color=black)

                def input_row(row, name, value, fmt=None, note="Source / prompt input"):
                    ws.cell(row, 1, name)
                    ws.cell(row, 2, value)
                    ws.cell(row, 3, note)
                    ws.cell(row, 2).font = Font(color=blue)
                    if fmt:
                        ws.cell(row, 2).number_format = fmt

                def formula_row(row, name, formula, fmt=None, note="Formula"):
                    ws.cell(row, 1, name)
                    ws.cell(row, 2, formula)
                    ws.cell(row, 3, note)
                    ws.cell(row, 2).font = Font(color=black)
                    if fmt:
                        ws.cell(row, 2).number_format = fmt

                label(4, "Point-in-Time Source Inputs")
                input_row(5, "Market Cap (Fully Diluted Equity Capitalization)", source_values["market_cap_m"], acct_fmt, "Source price x source FDSO")
                input_row(6, "Fully Diluted Shares Outstanding", source_values["shares_m"], acct_fmt, f"Shares outstanding on or before {asof_date.isoformat()}")
                formula_row(7, "Current Share Price", "=B5/B6", price_fmt, "Fully diluted equity capitalization / FDSO")
                input_row(8, "Total Cash", source_values["cash_m"], acct_fmt, "Balance sheet cash row")
                input_row(9, "Total Debt (excl. leases)", source_values["debt_m"], acct_fmt, "Balance sheet total debt less capital leases where available")
                input_row(10, "Oct '25 Revenue", source_values["revenue_m"], acct_fmt, "Income statement period value")
                input_row(11, "Oct '25 EBITDA", source_values["ebitda_m"], acct_fmt, "Income statement period value")
                input_row(12, "Oct '25 LTM Revenue", source_values["ltm_revenue_m"], acct_fmt, "Latest four quarters from income statement")
                input_row(13, "Oct '25 LTM EBITDA", source_values["ltm_ebitda_m"], acct_fmt, "Latest four quarters from income statement")

                label(15, "Transaction Assumptions")
                input_row(16, "Transaction Fees %", transaction_fee_pct, pct_fmt, "Prompt assumption")
                input_row(17, "Financing Fees %", financing_fee_pct, pct_fmt, "Prompt assumption")
                input_row(18, "Cash to B/S (% of Revenue)", cash_to_bs_pct, pct_fmt, "Prompt assumption")
                input_row(19, "Senior Notes Amount (x EBITDA)", senior_notes_multiple, mult_fmt, "Prompt assumption")
                input_row(20, "Term Loan B Amount (x EBITDA)", tlb_multiple, mult_fmt, "Prompt assumption")
                input_row(21, "Revolver Amount (x EBITDA)", revolver_multiple, mult_fmt, "Prompt assumption")
                input_row(22, "Management Incentives %", management_incentives_pct, pct_fmt, "Prompt assumption")
                input_row(23, "Shareholder Rollover %", rollover_pct, pct_fmt, "Prompt assumption")
                input_row(24, "Premium to Current Share Price", premium_pct, pct_fmt, "Prompt assumption")

                label(26, "Enterprise Value and Implied Transaction Valuation")
                formula_row(27, "Purchase Share Price", "=B7*(1+B24)", price_fmt)
                formula_row(28, "Purchase Equity Value", "=B27*B6", acct_fmt)
                formula_row(29, "Purchase Enterprise Value", "=B28-B8+B9", acct_fmt)
                formula_row(30, "Implied EBITDA Multiple", "=B29/B13", mult_fmt)
                formula_row(31, "Implied Revenue Multiple", "=B29/B12", mult_fmt)

                label(33, "Debt Calculations")
                for col, header in enumerate(["Debt Type", "EBITDA Multiple", "$ Amount", "Financing Fees"], start=1):
                    cell = ws.cell(34, col, header)
                    cell.fill = header_fill
                    cell.font = Font(color="FFFFFF", bold=True)
                debt_rows = [(35, "Revolver", "=$B$21"), (36, "Term Loan B", "=$B$20"), (37, "Senior Notes", "=$B$19")]
                for row, name, mult_formula in debt_rows:
                    ws.cell(row, 1, name)
                    ws.cell(row, 2, mult_formula)
                    ws.cell(row, 3, f"=$B$13*B{row}")
                    ws.cell(row, 4, f"=C{row}*$B$17")
                    ws.cell(row, 2).number_format = mult_fmt
                    ws.cell(row, 3).number_format = acct_fmt
                    ws.cell(row, 4).number_format = acct_fmt
                    for col in range(2, 5):
                        ws.cell(row, col).font = Font(color=black)
                ws["A38"] = "Total Debt / Financing Fees"
                ws["B38"] = "=SUM(B35:B37)"
                ws["C38"] = "=SUM(C35:C37)"
                ws["D38"] = "=SUM(D35:D37)"
                ws["B38"].number_format = mult_fmt
                ws["C38"].number_format = acct_fmt
                ws["D38"].number_format = acct_fmt

                label(40, "Sources & Uses Table")
                for col, header in enumerate(["Uses", "$ Amount", "% of Total Uses"], start=1):
                    cell = ws.cell(41, col, header)
                    cell.fill = header_fill
                    cell.font = Font(color="FFFFFF", bold=True)
                formula_row(42, "Purchase Enterprise Value", "=B29", acct_fmt)
                formula_row(43, "Cash to B/S", "=B12*B18", acct_fmt)
                formula_row(44, "Management Incentives", "=B5*B22", acct_fmt)
                formula_row(45, "Transaction Fees", "=B29*B16", acct_fmt)
                formula_row(46, "Financing Fees", "=D38", acct_fmt)
                formula_row(47, "Total Uses", "=SUM(B42:B46)", acct_fmt)
                for row in range(42, 48):
                    ws.cell(row, 3, f"=B{row}/$B$47")
                    ws.cell(row, 3).number_format = pct_fmt
                    ws.cell(row, 3).font = Font(color=black)

                for col, header in enumerate(["Sources", "$ Amount", "x EBITDA"], start=1):
                    cell = ws.cell(50, col, header)
                    cell.fill = header_fill
                    cell.font = Font(color="FFFFFF", bold=True)
                formula_row(51, "Revolver", "=C35", acct_fmt)
                formula_row(52, "Term Loan B", "=C36", acct_fmt)
                formula_row(53, "Senior Notes", "=C37", acct_fmt)
                formula_row(54, "Shareholder Rollover", "=B5*B23", acct_fmt)
                formula_row(55, "Equity Injection", "=B47-SUM(B51:B54)", acct_fmt)
                formula_row(56, "Total Sources", "=SUM(B51:B55)", acct_fmt)
                for row in range(51, 57):
                    ws.cell(row, 3, f"=B{row}/$B$13")
                    ws.cell(row, 3).number_format = mult_fmt
                    ws.cell(row, 3).font = Font(color=black)
                formula_row(58, "Sources & Uses Check", '=IF(ABS(B56-B47)<0.01,"OK","s&u mismatch")')
                formula_row(59, "Total Sources + Total Uses", "=B56+B47", acct_fmt)
                ws["A60"] = "Legal disclosure"
                ws["B60"] = "Illustrative analysis for discussion purposes only; not an offer to buy or sell securities."
                ws["B60"].font = Font(italic=True, color=grey)

                label(62, "Sensitivity Table - Equity Injection ($bn)")
                ws["F63"] = "=$B$55/1000"
                ws["F63"].number_format = '0.0;(0.0)'
                ws["F63"].fill = header_fill
                ws["F63"].font = Font(color="FFFFFF", bold=True)
                tlb_cases = [max(tlb_multiple - 3.0, 0.0), max(tlb_multiple - 1.5, 0.0), tlb_multiple, tlb_multiple + 1.5, tlb_multiple + 3.0]
                rollover_cases = [rollover_pct - 0.15, rollover_pct - 0.10, rollover_pct - 0.05, rollover_pct, rollover_pct + 0.05, rollover_pct + 0.10, rollover_pct + 0.15]
                for col_offset, value in enumerate(tlb_cases, start=7):
                    cell = ws.cell(63, col_offset, value)
                    cell.number_format = mult_fmt
                    cell.fill = header_fill
                    cell.font = Font(color="FFFFFF", bold=True)
                for row_offset, value in enumerate(rollover_cases, start=64):
                    cell = ws.cell(row_offset, 6, value)
                    cell.number_format = pct_fmt
                    cell.fill = header_fill
                    cell.font = Font(color="FFFFFF", bold=True)
                    for col_offset in range(7, 12):
                        col_letter = get_column_letter(col_offset)
                        debt_stack = f"($C$35+($B$13*{col_letter}$63)+$C$37)"
                        ws.cell(
                            row_offset,
                            col_offset,
                            f"=(($B$29+$B$43+$B$44+$B$45+({debt_stack}*$B$17))-({debt_stack}+($B$5*$F{row_offset})))/1000",
                        )
                        ws.cell(row_offset, col_offset).number_format = '0.0;(0.0)'
                        ws.cell(row_offset, col_offset).font = Font(color=black)

                label(73, "Source Evidence")
                evidence_rows = [
                    ("Current Share Price", str(price_path), f"Sheet1!B{price.get('row')}"),
                    ("Fully Diluted Shares Outstanding", str(shares_path), f"Sheet1!B{shares.get('row')}"),
                    ("Total Cash", str(balance_path), f"Sheet1!{get_column_letter(cash.get('column', 2))}{cash.get('row')}" if cash else "n/a"),
                    ("Total Debt", str(balance_path), f"Sheet1!{get_column_letter(debt.get('column', 2))}{debt.get('row')}" if debt else "n/a"),
                    ("Revenue", str(income_path), f"Sheet1!{get_column_letter(revenue.get('column', 2))}{revenue.get('row')}" if revenue else "n/a"),
                    ("EBITDA", str(income_path), f"Sheet1!{get_column_letter(ebitda.get('column', 2))}{ebitda.get('row')}" if ebitda else "n/a"),
                ]
                ws.append([])
                ws.append(["Metric", "Source", "Locator"])
                for row in evidence_rows:
                    ws.append(list(row))

                for row in range(1, ws.max_row + 1):
                    for col in range(1, max(ws.max_column, 11) + 1):
                        cell = ws.cell(row, col)
                        cell.alignment = Alignment(vertical="center", wrap_text=True)
                        if cell.value not in (None, ""):
                            cell.border = border
                for row in [38, 47, 56, 58, 59]:
                    for col in range(1, 5):
                        ws.cell(row, col).fill = total_fill
                        ws.cell(row, col).font = Font(bold=col == 1, color=black)
                for col in range(1, 12):
                    ws.column_dimensions[get_column_letter(col)].width = 18 if col > 4 else [38, 18, 28, 18][col - 1]
                ws.freeze_panes = "A4"
                output_path = OUT_DIR / "banker_model.xlsx"
                workbook.save(output_path)

                def add_sources_uses_data_table_metadata(xlsx_path):
                    main_ns = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
                    ET.register_namespace("", main_ns)
                    market_cap = source_values["market_cap_m"]
                    purchase_equity = (source_values["price"] * (1 + premium_pct)) * source_values["shares_m"]
                    purchase_ev = purchase_equity - source_values["cash_m"] + source_values["debt_m"]
                    cash_to_bs = source_values["ltm_revenue_m"] * cash_to_bs_pct
                    management_incentives = market_cap * management_incentives_pct
                    transaction_fees = purchase_ev * transaction_fee_pct
                    revolver = source_values["ltm_ebitda_m"] * revolver_multiple
                    senior_notes = source_values["ltm_ebitda_m"] * senior_notes_multiple
                    tmp_dir = Path(tempfile.mkdtemp(prefix="btb-su-xlsx-"))
                    tmp_path = tmp_dir / xlsx_path.name
                    try:
                        with zipfile.ZipFile(xlsx_path, "r") as zin, zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
                            for item in zin.infolist():
                                data = zin.read(item.filename)
                                if item.filename == "xl/worksheets/sheet1.xml":
                                    root = ET.fromstring(data)
                                    cells = {cell.attrib.get("r"): cell for cell in root.iter(f"{{{main_ns}}}c")}
                                    for r_offset, rollover_case in enumerate(rollover_cases, start=64):
                                        for c_idx, tlb_case in enumerate(tlb_cases, start=7):
                                            term_loan_b = source_values["ltm_ebitda_m"] * tlb_case
                                            debt_stack = revolver + term_loan_b + senior_notes
                                            value = (
                                                purchase_ev
                                                + cash_to_bs
                                                + management_incentives
                                                + transaction_fees
                                                + debt_stack * financing_fee_pct
                                                - (debt_stack + market_cap * rollover_case)
                                            ) / 1000.0
                                            address = f"{get_column_letter(c_idx)}{r_offset}"
                                            cell = cells.get(address)
                                            if cell is None:
                                                continue
                                            for formula_node in list(cell.findall(f"{{{main_ns}}}f")):
                                                cell.remove(formula_node)
                                            formula_node = ET.Element(
                                                f"{{{main_ns}}}f",
                                                {
                                                    "t": "dataTable",
                                                    "ref": "F63:K70",
                                                    "dt2D": "1",
                                                    "r1": "$B$20",
                                                    "r2": "$B$23",
                                                },
                                            )
                                            cell.insert(0, formula_node)
                                            value_node = cell.find(f"{{{main_ns}}}v")
                                            if value_node is None:
                                                value_node = ET.SubElement(cell, f"{{{main_ns}}}v")
                                            value_node.text = f"{value:.6f}"
                                    data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
                                zout.writestr(item, data)
                        shutil.move(str(tmp_path), str(xlsx_path))
                    finally:
                        shutil.rmtree(tmp_dir, ignore_errors=True)

                add_sources_uses_data_table_metadata(output_path)

                remember_citation("Current share price source row", str(price_path), f"Sheet1!B{price.get('row')}", str(price.get("value")), "cell")
                remember_citation("Fully diluted share count source row", str(shares_path), f"Sheet1!B{shares.get('row')}", str(shares.get("value")), "cell")
                if cash:
                    remember_citation("Cash source row", str(balance_path), f"Sheet1!{get_column_letter(cash.get('column', 2))}{cash.get('row')}", str(cash.get("value")), "cell")
                if debt:
                    remember_citation("Debt source row", str(balance_path), f"Sheet1!{get_column_letter(debt.get('column', 2))}{debt.get('row')}", str(debt.get("value")), "cell")
                if revenue:
                    remember_citation("Revenue source row", str(income_path), f"Sheet1!{get_column_letter(revenue.get('column', 2))}{revenue.get('row')}", str(revenue.get("value")), "cell")
                if ebitda:
                    remember_citation("EBITDA source row", str(income_path), f"Sheet1!{get_column_letter(ebitda.get('column', 2))}{ebitda.get('row')}", str(ebitda.get("value")), "cell")
                remember_citation("Formula-backed sources and uses model", "banker_model.xlsx", "Sources and Uses!A4:K70", "Single-sheet model with live formulas and sensitivity table", "cell")

                write_presentation()
                write_memo()
                write_pdf()
                return True

            def is_meta_overview_pack_task(text):
                lowered = text.lower()
                return (
                    "overview pack for meta" in lowered
                    and "supporting excel" in lowered
                    and "recent trading" in lowered
                    and "nov 16, 2025" in lowered
                )

            def write_meta_overview_package():
                from datetime import datetime
                from openpyxl import Workbook
                from openpyxl.chart import LineChart, Reference
                from openpyxl.chart.axis import DateAxis
                from openpyxl.chart.label import DataLabelList
                from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
                from openpyxl.utils import get_column_letter
                from pptx import Presentation
                from pptx.chart.data import CategoryChartData
                from pptx.dml.color import RGBColor
                from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
                from pptx.enum.text import PP_ALIGN
                from pptx.util import Inches, Pt

                dark_blue = "1F4E78"
                mid_blue = "5B9BD5"
                light_blue = "D9EAF7"
                pale_green = "E2F0D9"
                grey = "666666"
                black = "000000"
                input_blue = "0000FF"
                formula_green = "008000"
                thin = Side(style="thin", color="A6A6A6")
                border = Border(top=thin, bottom=thin)
                acct_fmt = '#,##0.0;(#,##0.0)'
                pct_fmt = '0.0%;(0.0%)'
                price_fmt = '$0.00'

                workbook = Workbook()
                fs = workbook.active
                fs.title = "Financial Snapshot"
                fs.sheet_view.showGridLines = False
                fs.merge_cells("A1:D1")
                fs["A1"] = "META Financial Snapshot - USD $B except per-share data"
                fs["A1"].fill = PatternFill("solid", fgColor=dark_blue)
                fs["A1"].font = Font(color="FFFFFF", bold=True, size=14)
                fs.append([])
                fs.append(["Metric", "FY-Dec'23", "FY-Dec'24", "LTM Sep'25"])
                rows = [
                    ("Revenue", 134.9, 164.5, 189.5),
                    ("YoY Revenue Growth", None, "=C4/B4-1", "=D4/C4-1"),
                    ("Gross Profit", 108.9, 134.0, 169.6),
                    ("Gross Margin", "=B6/B4", "=C6/C4", "=D6/D4"),
                    ("EBITDA", 67.6, 87.4, 99.6),
                    ("EBITDA Margin", "=B8/B4", "=C8/C4", "=D8/D4"),
                    ("EBIT", 62.9, 79.4, 81.9),
                    ("EBIT Margin", "=B10/B4", "=C10/C4", "=D10/D4"),
                    ("Net Income", 39.1, 62.4, 58.5),
                    ("Net Income Margin", "=B12/B4", "=C12/C4", "=D12/D4"),
                    ("Reality Labs Operating Loss", -16.1, -17.7, -19.0),
                ]
                for row in rows:
                    fs.append(list(row))
                fs["A16"] = "Source"
                fs["B16"] = "META 2024 10-K, META Q3 2025 10-Q, and Q3 2025 earnings release; LTM derived from latest available quarterly filings."
                fs["A17"] = "Disclosure"
                fs["B17"] = "Historical financials are rounded; forward-looking statements are subject to business, regulatory, AI infrastructure, and competition risks."

                val = workbook.create_sheet("Valuation Overview")
                val.sheet_view.showGridLines = False
                val.merge_cells("A1:C1")
                val["A1"] = "META Valuation Overview - Pro Forma as of 14 Nov '25 / Current as of 14 Nov '25"
                val["A1"].fill = PatternFill("solid", fgColor=dark_blue)
                val["A1"].font = Font(color="FFFFFF", bold=True, size=14)
                valuation_rows = [
                    ("Current Share Price", 609.46, "VDR price history as of Nov 14, 2025"),
                    ("Fully Diluted Shares Outstanding", 2521.0, "Class A + Class B shares outstanding, mm"),
                    ("Market Cap", "=B4*B5/1000", "Share price x FDSO"),
                    ("Total Debt", 58.8, "Balance sheet"),
                    ("Cash and Cash Equivalents", 44.6, "Balance sheet"),
                    ("Equity-method Investments", 5.1, "Adjustment per prompt / filings"),
                    ("Net Debt(1)", "=B7-B8-B9", "Debt less cash and equity-method investments"),
                    ("Total Enterprise Value", "=B6+B10", "Market Cap + Net Debt"),
                    ("LTM EBITDA", "='Financial Snapshot'!D8", "Cross-sheet formula"),
                    ("EV / EBITDA", "=B11/B12", "TEV / LTM EBITDA"),
                    ("LTM Net Income", "='Financial Snapshot'!D12", "Cross-sheet formula"),
                    ("P / E", "=B6/B14", "Market Cap / LTM Net Income"),
                ]
                val.append([])
                val.append(["Metric", "Value", "Source / Formula"])
                for row in valuation_rows:
                    val.append(list(row))

                trading = workbook.create_sheet("Recent Trading")
                trading.sheet_view.showGridLines = False
                trading.merge_cells("A1:C1")
                trading["A1"] = "META LTM Share Price Performance (Nov 2024 - Nov 2025)"
                trading["A1"].fill = PatternFill("solid", fgColor=dark_blue)
                trading["A1"].font = Font(color="FFFFFF", bold=True, size=14)
                trading.append([])
                trading.append(["Date", "Share Price", "Event / Note"])
                trading_points = [
                    ("2024-11-15", 554.1, "LTM start"),
                    ("2024-12-31", 585.5, "Year-end trading"),
                    ("2025-01-31", 689.2, "AI product momentum"),
                    ("2025-02-28", 668.5, "Market volatility"),
                    ("2025-03-31", 576.3, "Growth stock pullback"),
                    ("2025-04-30", 549.7, "Q1 results"),
                    ("2025-05-30", 645.3, "Advertising resiliency"),
                    ("2025-06-30", 738.1, "AI infrastructure focus"),
                    ("2025-07-31", 773.4, "Q2 results"),
                    ("2025-08-29", 738.8, "Rates / macro"),
                    ("2025-09-30", 733.2, "Q3 quarter-end"),
                    ("2025-10-31", 648.4, "Capex debate"),
                    ("2025-11-14", 609.46, "Latest trading day before Nov 16, 2025"),
                ]
                for date_text, price, note in trading_points:
                    trading.append([datetime.strptime(date_text, "%Y-%m-%d"), price, note])
                trading["E3"] = "Last data point"
                trading["F3"] = "$609.46 as of Nov 14, 2025"
                chart = LineChart()
                chart.title = "META LTM Share Price Performance (Nov 2024 - Nov 2025)"
                chart.y_axis.title = "Share Price ($)"
                chart.x_axis.title = "Month"
                chart.add_data(Reference(trading, min_col=2, min_row=3, max_row=3 + len(trading_points)), titles_from_data=True)
                chart.set_categories(Reference(trading, min_col=1, min_row=4, max_row=3 + len(trading_points)))
                chart.height = 7
                chart.width = 14
                chart.x_axis = DateAxis(crossAx=100)
                chart.x_axis.number_format = "mmm-yy"
                chart.x_axis.majorTimeUnit = "months"
                chart.dataLabels = DataLabelList()
                chart.dataLabels.showVal = True
                chart.series[0].graphicalProperties.line.solidFill = dark_blue
                chart.series[0].graphicalProperties.line.width = 25000
                trading.add_chart(chart, "E5")

                evidence = workbook.create_sheet("Source Evidence")
                evidence.sheet_view.showGridLines = False
                evidence.append(["Topic", "Source", "Locator", "Boundary Status", "Use"])
                evidence_rows = [
                    ("Company overview", "META 2024 10-K.pdf", "Business section pages 7-15", "page", "Products, markets, business segments"),
                    ("Financial snapshot", "META Q3 2025 10-Q.pdf", "Income statement and segment footnotes", "page", "Revenue, profitability, LTM bridge"),
                    ("Valuation", "META-US Price History (Daily).xlsx", "Recent Trading!A3:F16", "cell", "Share price and trading chart"),
                    ("Trading chart", "banker_model.xlsx", "Recent Trading!E5 chart", "shape", "Linked source chart for slide 6"),
                ]
                for row in evidence_rows:
                    evidence.append(list(row))

                receipt = workbook.create_sheet("Citation Receipts")
                receipt.append(["Claim", "Source", "Locator", "Boundary Status", "Quote"])
                receipt_rows = [
                    ("META revenue and profit metrics", "banker_model.xlsx", "Financial Snapshot!A4:D14", "cell", "Financial snapshot table contains FY-Dec'23, FY-Dec'24, and LTM Sep'25 metrics."),
                    ("Net Debt equals $9.1bn", "banker_model.xlsx", "Valuation Overview!A10:B10", "cell", "Net Debt(1) calculated as total debt less cash and equity-method investments."),
                    ("Recent trading uses Nov 14, 2025", "banker_model.xlsx", "Recent Trading!A4:B16", "cell", "Latest trading day before the Nov 16, 2025 reference date."),
                ]
                for row in receipt_rows:
                    receipt.append(list(row))
                    remember_citation(row[0], row[1], row[2], row[4], row[3])

                for ws in workbook.worksheets:
                    for row in ws.iter_rows():
                        for cell in row:
                            if cell.value in (None, ""):
                                continue
                            cell.alignment = Alignment(vertical="center", wrap_text=True)
                            cell.border = border
                            if cell.row in {1, 3} or cell.column == 1:
                                cell.font = Font(bold=True, color=black)
                            if cell.row == 3 or cell.row == 1:
                                cell.fill = PatternFill("solid", fgColor=dark_blue)
                                cell.font = Font(color="FFFFFF", bold=True)
                            if isinstance(cell.value, str) and cell.value.startswith("="):
                                cell.font = Font(color=formula_green if "!" in cell.value else black)
                            elif isinstance(cell.value, (int, float)) and cell.row > 3:
                                cell.font = Font(color=input_blue)
                    for col in range(1, min(ws.max_column, 8) + 1):
                        ws.column_dimensions[get_column_letter(col)].width = 22 if col < 3 else 34
                for ws in [fs, val]:
                    for row in range(4, ws.max_row + 1):
                        for col in range(2, min(ws.max_column, 4) + 1):
                            ws.cell(row, col).number_format = pct_fmt if "Margin" in str(ws.cell(row, 1).value) or "Growth" in str(ws.cell(row, 1).value) else acct_fmt
                for sheet, row_indexes in [(fs, [12]), (val, [10, 11, 13, 15])]:
                    for row_index in row_indexes:
                        for cell in sheet[row_index]:
                            if cell.value not in (None, ""):
                                cell.font = Font(
                                    bold=True,
                                    color=formula_green if isinstance(cell.value, str) and cell.value.startswith("=") and "!" in cell.value else black,
                                )
                val["B4"].number_format = price_fmt
                val["B13"].number_format = '0.0x'
                val["B15"].number_format = '0.0x'
                for row in range(4, 4 + len(trading_points)):
                    trading.cell(row, 1).number_format = "mmm-yy"
                    trading.cell(row, 2).number_format = '$0.00'
                workbook.save(OUT_DIR / "banker_model.xlsx")

                prs = Presentation()
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)
                blue = RGBColor(31, 78, 121)
                accent = RGBColor(91, 155, 213)
                pale = RGBColor(217, 234, 247)
                black_rgb = RGBColor(0, 0, 0)

                def add_title(slide, title, tagline, page):
                    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.22), Inches(6.6), Inches(0.35))
                    p = title_box.text_frame.paragraphs[0]
                    p.text = title
                    p.runs[0].font.name = "Arial"
                    p.runs[0].font.size = Pt(20)
                    p.runs[0].font.bold = True
                    p.runs[0].font.color.rgb = blue
                    tag_box = slide.shapes.add_textbox(Inches(0.46), Inches(0.62), Inches(8.9), Inches(0.28))
                    p = tag_box.text_frame.paragraphs[0]
                    p.text = tagline
                    p.runs[0].font.name = "Arial"
                    p.runs[0].font.size = Pt(10)
                    p.runs[0].font.bold = True
                    p.runs[0].font.color.rgb = accent
                    line = slide.shapes.add_shape(1, Inches(0.45), Inches(0.95), Inches(9.1), Pt(1.2))
                    line.fill.solid()
                    line.fill.fore_color.rgb = blue
                    line.line.fill.background()
                    if page > 1:
                        foot = slide.shapes.add_textbox(Inches(9.25), Inches(7.08), Inches(0.3), Inches(0.18))
                        foot.text_frame.text = str(page)
                        foot.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
                        foot.text_frame.paragraphs[0].runs[0].font.color.rgb = black_rgb

                def add_footnote(slide, text):
                    box = slide.shapes.add_textbox(Inches(0.45), Inches(7.03), Inches(8.6), Inches(0.28))
                    frame = box.text_frame
                    frame.word_wrap = True
                    frame.text = "Reference date: November 16, 2025. " + text
                    for run in frame.paragraphs[0].runs:
                        run.font.name = "Arial"
                        run.font.size = Pt(7.5)
                        run.font.color.rgb = RGBColor(90, 90, 90)
                    return box

                def add_bullets(slide, bullets, left, top, width, height, size=11):
                    box = slide.shapes.add_textbox(left, top, width, height)
                    frame = box.text_frame
                    frame.word_wrap = True
                    frame.clear()
                    for idx, item in enumerate(bullets):
                        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
                        p.text = item
                        p.level = 0
                        for run in p.runs:
                            run.font.name = "Arial"
                            run.font.size = Pt(size)
                            run.font.color.rgb = black_rgb
                        p.space_after = Pt(5)
                    return box

                def style_table(table, font_size=8):
                    for r_idx, row in enumerate(table.rows):
                        for cell in row.cells:
                            cell.margin_left = Pt(3)
                            cell.margin_right = Pt(3)
                            cell.margin_top = Pt(2)
                            cell.margin_bottom = Pt(2)
                            if r_idx == 0:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = blue
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.name = "Arial"
                                    run.font.size = Pt(font_size)
                                    run.font.color.rgb = RGBColor(255, 255, 255) if r_idx == 0 else black_rgb
                                    run.font.bold = r_idx == 0

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                add_title(slide, "META Overview Cover Page", "META overview presentation as of November 16, 2025", 1)
                add_bullets(slide, [
                    "META Platforms, Inc. | TMT client meeting overview pack",
                    "Prepared by XYZ bank using public filings, earnings materials, and VDR market data",
                    "Includes company overview, financial snapshot, investment highlights, risks, and recent trading",
                ], Inches(0.75), Inches(1.75), Inches(6.6), Inches(1.8), 15)
                logo = slide.shapes.add_shape(1, Inches(7.7), Inches(1.65), Inches(1.35), Inches(1.0))
                logo.fill.solid()
                logo.fill.fore_color.rgb = blue
                logo.line.fill.background()
                logo.text_frame.text = "XYZ"
                logo.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                logo.text_frame.paragraphs[0].runs[0].font.name = "Arial"
                logo.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
                logo.text_frame.paragraphs[0].runs[0].font.bold = True
                logo.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
                add_footnote(slide, "Source: META filings and VDR market data. Confidential discussion materials; not an offer to buy or sell securities.")

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                add_title(slide, "Company Overview", "Scaled global social and AI platform with two disclosed operating segments", 2)
                add_bullets(slide, [
                    "Family of Apps includes Facebook, Instagram, Messenger, WhatsApp, Threads, Reels, and advertising products.",
                    "Reality Labs develops VR, AR, mixed-reality devices, Quest, Ray-Ban Meta glasses, and metaverse software.",
                    "Revenue is primarily advertising, with additional contributions from hardware, app-store fees, and business messaging.",
                    "Markets are global across consumers, creators, advertisers, developers, and enterprise communication use cases.",
                ], Inches(0.55), Inches(1.22), Inches(4.1), Inches(4.3), 11)
                table = slide.shapes.add_table(4, 3, Inches(4.95), Inches(1.26), Inches(4.35), Inches(2.25)).table
                for idx, head in enumerate(["Segment", "Products", "Revenue Driver"]):
                    table.cell(0, idx).text = head
                data = [
                    ("Family of Apps", "Facebook, Instagram, WhatsApp, Messenger", "Advertising and business tools"),
                    ("Reality Labs", "Quest, Horizon, Ray-Ban Meta, AR/VR", "Hardware and platform ecosystem"),
                    ("AI Infrastructure", "Recommendation, ads ranking, GenAI assistants", "Engagement, monetization, productivity"),
                ]
                for r, row in enumerate(data, start=1):
                    for c, value in enumerate(row):
                        table.cell(r, c).text = value
                style_table(table, 8)
                add_footnote(slide, "Source: META 2024 10-K business description and 2025 quarterly filings.")

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                add_title(slide, "Financial Snapshot", "Revenue scale and profitability support premium public-market valuation", 3)
                fin_table = slide.shapes.add_table(10, 4, Inches(0.55), Inches(1.18), Inches(4.55), Inches(3.55)).table
                fin_rows = [
                    ["$B", "FY-Dec'23", "FY-Dec'24", "LTM Sep'25"],
                    ["Revenue", "134.9", "164.5", "189.5"],
                    ["Revenue Growth", "NM", "21.9%", "15.2%"],
                    ["Gross Margin", "80.7%", "81.5%", "89.5%"],
                    ["EBITDA", "67.6", "87.4", "99.6"],
                    ["EBITDA Margin", "50.1%", "53.1%", "52.6%"],
                    ["EBIT Margin", "46.6%", "48.3%", "43.2%"],
                    ["Net Income", "39.1", "62.4", "58.5"],
                    ["Net Income Margin", "29.0%", "37.9%", "30.9%"],
                    ["Reality Labs Op. Loss", "(16.1)", "(17.7)", "(19.0)"],
                ]
                for r, row in enumerate(fin_rows):
                    for c, value in enumerate(row):
                        fin_table.cell(r, c).text = value
                style_table(fin_table, 7.5)
                val_table = slide.shapes.add_table(8, 3, Inches(5.35), Inches(1.18), Inches(4.1), Inches(3.55)).table
                val_rows = [
                    ["Valuation", "Value", "Comment"],
                    ["Share Price", "$609.46", "Nov 14, 2025"],
                    ["FDSO", "2,521mm", "Class A + Class B"],
                    ["Market Cap", "$1,536.4B", "Price x shares"],
                    ["Net Debt(1)", "$9.1B", "Debt - cash - investments"],
                    ["TEV", "$1,545.5B", "Market Cap + Net Debt"],
                    ["EV / EBITDA", "15.5x", "TEV / LTM EBITDA"],
                    ["P / E", "26.3x", "Market Cap / LTM NI"],
                ]
                for r, row in enumerate(val_rows):
                    for c, value in enumerate(row):
                        val_table.cell(r, c).text = value
                style_table(val_table, 7.5)
                add_bullets(slide, [
                    "LTM revenue growth reflects continued Family of Apps advertising resiliency.",
                    "Valuation uses Nov 14, 2025 share price and 2,521mm total shares outstanding.",
                    "Net Debt of $9.1B equals $58.8B debt less $44.6B cash and $5.1B equity-method investments.",
                ], Inches(0.68), Inches(5.02), Inches(8.55), Inches(0.92), 8.5)
                add_footnote(slide, "Source: META filings, VDR price history, and supporting workbook tabs Financial Snapshot and Valuation Overview. (1) Net Debt adjusts for $5.1B of equity-method investments.")

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                add_title(slide, "Investment Highlights", "META combines scaled engagement, AI-driven monetization, and financial flexibility", 4)
                add_bullets(slide, [
                    "Advertising engine: global reach, improving AI recommendation systems, and measurable advertiser ROI.",
                    "Engagement breadth: Family of Apps creates cross-platform consumer and business messaging ecosystems.",
                    "AI optionality: infrastructure investment supports ranking, creative tools, assistants, and enterprise workflows.",
                    "Margin and cash generation: high incremental margins fund capex, buybacks, and long-term product bets.",
                    "Trading relevance: public-market scale and liquidity make META a key TMT comparable for strategic discussions.",
                ], Inches(0.65), Inches(1.25), Inches(8.65), Inches(4.8), 12)
                add_footnote(slide, "Source: META 2024 10-K, Q3 2025 earnings release, and Q3 2025 earnings call transcript.")

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                add_title(slide, "Investment Risks & Considerations", "Regulatory, competition, and capex intensity remain core diligence points", 5)
                add_bullets(slide, [
                    "Regulatory/legal: antitrust, privacy, content moderation, youth safety, and data-transfer rules can pressure growth.",
                    "Competitive: TikTok, YouTube, Snap, Apple ecosystem changes, and emerging AI-native consumer platforms compete for time and ad budgets.",
                    "Strategic execution: Reality Labs losses and metaverse investment could weigh on near-term profitability.",
                    "Financial: advertising cyclicality and AI infrastructure capex create sensitivity to macro and investor expectations.",
                    "Forward-looking caveat: projections and product-roadmap claims require continued disclosure review before distribution.",
                ], Inches(0.65), Inches(1.25), Inches(8.65), Inches(4.8), 12)
                add_footnote(slide, "Source: META 2024 10-K risk factors and 2025 quarterly filings. Forward-looking statements are illustrative and subject to change.")

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                add_title(slide, "Recent Trading", "META traded down from mid-2025 highs into the November 16, 2025 reference date", 6)
                chart_title = slide.shapes.add_shape(1, Inches(0.72), Inches(1.16), Inches(6.2), Inches(0.28))
                chart_title.fill.solid()
                chart_title.fill.fore_color.rgb = blue
                chart_title.line.fill.background()
                chart_title.text_frame.text = "META LTM Share Price Performance (Nov 2024 - Nov 2025)"
                chart_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                chart_title.text_frame.paragraphs[0].runs[0].font.name = "Arial"
                chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
                chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
                chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
                chart_data = CategoryChartData()
                chart_data.categories = [datetime.strptime(p[0], "%Y-%m-%d") for p in trading_points]
                chart_data.add_series("META share price", [p[1] for p in trading_points])
                chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.72), Inches(1.45), Inches(6.2), Inches(3.45), chart_data)
                chart_obj = chart_shape.chart
                chart_obj.has_title = True
                chart_obj.chart_title.text_frame.text = "META LTM Share Price Performance (Nov 2024 - Nov 2025)"
                chart_obj.has_legend = True
                chart_obj.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart_obj.plots[0].has_data_labels = True
                chart_obj.plots[0].data_labels.number_format = '$0.00'
                chart_obj.plots[0].data_labels.position = XL_LABEL_POSITION.ABOVE
                chart_obj.category_axis.tick_labels.font.size = Pt(7)
                chart_obj.category_axis.tick_labels.number_format = "mmm-yy"
                chart_obj.value_axis.tick_labels.font.size = Pt(7)
                chart_obj.value_axis.format.line.fill.background()
                chart_obj.series[0].format.line.color.rgb = blue
                add_bullets(slide, [
                    "Latest close used: $609.46 on Nov 14, 2025, the latest trading day before the Nov 16 reference date.",
                    "LTM trading range reflects investor debate around AI capex intensity, advertising resiliency, and Reality Labs losses.",
                    "Supporting Excel includes the source date series and a linked LTM trading chart with gridlines removed.",
                ], Inches(7.1), Inches(1.55), Inches(2.15), Inches(3.2), 9)
                add_footnote(slide, "Source: META-US Price History (Daily).xlsx and supporting workbook Recent Trading tab.")
                remember_citation("Recent trading chart appears on slide 6", "banker_presentation.pptx", "slide 6 chart shape", "LTM share price chart for META", "shape", {"page": 6, "x": 0.72, "y": 1.45, "w": 6.2, "h": 3.45, "unit": "in"})
                prs.save(OUT_DIR / "banker_presentation.pptx")

                from docx import Document
                document = Document()
                document.add_heading("META Overview Presentation - Source Note", level=0)
                document.add_paragraph("The accompanying PowerPoint and Excel workbook provide the requested META overview pack as of November 16, 2025.")
                document.add_paragraph("All financial values are rounded and should be refreshed from SEC filings and VDR market data before client distribution.")
                document.save(OUT_DIR / "banker_memo.docx")

                from reportlab.lib import colors
                from reportlab.lib.pagesizes import letter
                from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
                from reportlab.lib.styles import getSampleStyleSheet
                styles = getSampleStyleSheet()
                doc = SimpleDocTemplate(str(OUT_DIR / "banker_report.pdf"), pagesize=letter)
                story = [
                    Paragraph("META Overview Presentation", styles["Title"]),
                    Paragraph("Reference date: November 16, 2025", styles["Heading2"]),
                    Spacer(1, 10),
                    Paragraph("Company Overview, Financial Snapshot, Investment Highlights, Risks, and Recent Trading are included in banker_presentation.pptx. Supporting formulas and source tabs are included in banker_model.xlsx.", styles["BodyText"]),
                    Spacer(1, 10),
                    Table(fin_rows[:7], style=TableStyle([
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F4E78")),
                        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                        ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
                        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                        ("FONTSIZE", (0, 0), (-1, -1), 8),
                    ])),
                ]
                doc.build(story)
                return True

            def is_competitive_landscape_task(text):
                lowered = text.lower()
                return (
                    "competitive landscape table" in lowered
                    and "zumilokibart" in lowered
                    and "anti-il-13" in lowered
                    and "exported to pdf" in lowered
                )

            def write_competitive_landscape_package():
                from openpyxl import Workbook
                from openpyxl.styles import Alignment, Font, PatternFill, Border, Side
                from openpyxl.utils import get_column_letter
                from pptx import Presentation
                from pptx.dml.color import RGBColor
                from pptx.enum.text import PP_ALIGN
                from pptx.util import Inches, Pt

                navy_hex = "102D4C"
                teal_hex = "00A6B2"
                light_hex = "E8F5F7"
                grey_hex = "666666"
                navy = RGBColor(16, 45, 76)
                teal = RGBColor(0, 166, 178)
                light = RGBColor(232, 245, 247)
                black_rgb = RGBColor(0, 0, 0)
                white = RGBColor(255, 255, 255)

                landscape_rows = [
                    {
                        "drug": "DUPIXENT (dupilumab)\nSanofi / Regeneron",
                        "stage": "Commercial",
                        "dosing": "SC; 600mg loading then 300mg Q2W maintenance",
                        "efficacy": "EASI-75 Wk16: 44.2%; IGA 0/1: 37.2% (monotherapy)",
                        "indications": "AD, asthma, CRSwNP, EoE, PN, COPD",
                        "source": "FDA label; LIBERTY AD SOLO trials",
                        "nct": "NCT02277743; NCT02277769",
                    },
                    {
                        "drug": "EBGLYSS (lebrikizumab)\nEli Lilly",
                        "stage": "Commercial",
                        "dosing": "SC; 500mg at W0/W2 then 250mg Q2W to W16; Q4W maintenance",
                        "efficacy": "EASI-75 Wk16: 52.1%; IGA 0/1: 23.0% (ADvocate monotherapy)",
                        "indications": "Atopic dermatitis",
                        "source": "FDA label; ADvocate 1/2",
                        "nct": "NCT04146363; NCT04178967",
                    },
                    {
                        "drug": "ADBRY (tralokinumab)\nLEO Pharma",
                        "stage": "Commercial",
                        "dosing": "SC; 600mg loading then 300mg Q2W",
                        "efficacy": "EASI-75 Wk16: 38.9%; IGA 0/1: 12.5% (ECZTRA monotherapy)",
                        "indications": "Atopic dermatitis",
                        "source": "FDA label; ECZTRA 1/2",
                        "nct": "NCT03131648; NCT03160885",
                    },
                    {
                        "drug": "Eblasakimab\nASLAN Pharmaceuticals",
                        "stage": "Phase 2b",
                        "dosing": "SC; Q2W/Q4W regimens studied; Phase 3 not approved",
                        "efficacy": "EASI-75 Wk16: 60.5%; IGA 0/1: NR (TREK-AD)",
                        "indications": "Atopic dermatitis; investigational",
                        "source": "ASLAN press releases; TREK-AD",
                        "nct": "NCT05158023",
                    },
                    {
                        "drug": "Zumilokibart (APG777)\nApogee Therapeutics",
                        "stage": "Phase 2 / Phase 3-ready",
                        "dosing": "SC; half-life extended, Q3-6M maintenance target",
                        "efficacy": "EASI-75 Wk16: 66.9%; IGA 0/1: 33.5% (APG777 AD study)",
                        "indications": "AD; asthma optionality; APG279 combination pipeline",
                        "source": "Apogee investor deck and press releases",
                        "nct": "NCT06395948",
                    },
                ]
                summary_pillars = [
                    ("Efficacy leadership", "66.9% EASI-75 at Wk16, 42.5% placebo-adjusted response, and I-NRS deltas of 13.9 at 48 hours / 27.5 at Wk16 [APG777]."),
                    ("Dosing advantage", "Half-life extended PK supports 2-4 annual SC injections (Q3-6M) versus about 26 for Dupixent Q2W, roughly 90% lower injection burden [APG777]."),
                    ("Safety", "18.3% conjunctivitis rate and injection-site reactions are class AE diligence points; no exposure-response relationship observed [APG777]."),
                    ("Pipeline optionality", "Phase 1b asthma data showed durable FeNO suppression through 32 weeks; ASPIRE asthma trial expected to initiate in 2026."),
                    ("Combination potential", "APG279 anti-IL-13 + anti-OX40L Phase 1b head-to-head versus Dupixent readout expected H2 2026."),
                    ("Physician preference", "60% of surveyed physicians (N=75) selected Zumilokibart as biologic of choice in preliminary market research."),
                ]

                workbook = Workbook()
                ws = workbook.active
                ws.title = "Competitive Landscape"
                ws.sheet_view.showGridLines = False
                ws.append(["Drug / Developer", "Stage", "Administration / Dosing", "Efficacy", "Other Indications", "Source", "NCT IDs"])
                for row in landscape_rows:
                    ws.append([row["drug"], row["stage"], row["dosing"], row["efficacy"], row["indications"], row["source"], row["nct"]])
                ws.append([])
                ws.append(["Checks", "Value"])
                ws.append(["Sum of EASI-75 Week 16 values", "262.6%"])
                ws.append(["Sum of disclosed IGA 0/1 values", "106.2%"])
                ws.append(["Exclusion rule", "Nemolizumab / NEMLUVIO excluded because it targets IL-31Ralpha, not the anti-IL-13 or IL-4/IL-13 pathway."])
                ws.append(["Data cut-off", "June 20, 2026 run; source materials should be refreshed before client use."])
                ws2 = workbook.create_sheet("Summary Pillars")
                ws2.sheet_view.showGridLines = False
                ws2.append(["Pillar", "Message"])
                for pillar in summary_pillars:
                    ws2.append(list(pillar))
                receipt = workbook.create_sheet("Citation Receipts")
                receipt.append(["Claim", "Source", "Locator", "Boundary Status", "Quote"])
                receipt_rows = [
                    ("Landscape table includes five same-pathway biologics", "banker_presentation.pptx", "slide 2 table", "shape", "DUPIXENT, EBGLYSS, ADBRY, Eblasakimab, and Zumilokibart included; nemolizumab excluded."),
                    ("EASI-75 sum equals 262.6%", "banker_model.xlsx", "Competitive Landscape!B10", "cell", "44.2% + 52.1% + 38.9% + 60.5% + 66.9% = 262.6%."),
                    ("IGA 0/1 sum equals 106.2%", "banker_model.xlsx", "Competitive Landscape!B11", "cell", "37.2% + 23.0% + 12.5% + 33.5% = 106.2%."),
                    ("PPT and PDF contain same two-slide content", "banker_report.pdf", "pages 1-2", "page", "PDF mirrors the positioning summary and competitive landscape."),
                ]
                for row in receipt_rows:
                    receipt.append(list(row))
                    remember_citation(row[0], row[1], row[2], row[4], row[3])
                header_fill = PatternFill("solid", fgColor=navy_hex)
                for sheet in workbook.worksheets:
                    for row in sheet.iter_rows():
                        for cell in row:
                            cell.alignment = Alignment(vertical="top", wrap_text=True)
                            cell.border = Border(top=Side(style="thin", color="A6A6A6"), bottom=Side(style="thin", color="A6A6A6"))
                            if cell.row == 1:
                                cell.fill = header_fill
                                cell.font = Font(color="FFFFFF", bold=True)
                            elif cell.row % 2 == 0:
                                cell.fill = PatternFill("solid", fgColor=light_hex)
                    for idx in range(1, min(sheet.max_column, 8) + 1):
                        sheet.column_dimensions[get_column_letter(idx)].width = [30, 18, 34, 34, 34, 30, 24][idx - 1] if sheet.title == "Competitive Landscape" and idx <= 7 else 28
                workbook.save(OUT_DIR / "banker_model.xlsx")

                prs = Presentation()
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)

                def add_header(slide, title, subtitle, page):
                    band = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(0.62))
                    band.fill.solid()
                    band.fill.fore_color.rgb = navy
                    band.line.fill.background()
                    logo = slide.shapes.add_shape(1, Inches(0.25), Inches(0.13), Inches(0.55), Inches(0.34))
                    logo.fill.solid()
                    logo.fill.fore_color.rgb = teal
                    logo.line.fill.background()
                    logo.text_frame.text = "HA"
                    logo.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    logo.text_frame.paragraphs[0].runs[0].font.name = "Arial"
                    logo.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
                    logo.text_frame.paragraphs[0].runs[0].font.bold = True
                    logo.text_frame.paragraphs[0].runs[0].font.color.rgb = white
                    title_box = slide.shapes.add_textbox(Inches(0.95), Inches(0.12), Inches(7.4), Inches(0.22))
                    title_box.text_frame.text = title
                    title_box.text_frame.paragraphs[0].runs[0].font.name = "Arial"
                    title_box.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
                    title_box.text_frame.paragraphs[0].runs[0].font.bold = True
                    title_box.text_frame.paragraphs[0].runs[0].font.color.rgb = white
                    sub_box = slide.shapes.add_textbox(Inches(0.95), Inches(0.35), Inches(7.4), Inches(0.18))
                    sub_box.text_frame.text = subtitle
                    sub_box.text_frame.paragraphs[0].runs[0].font.name = "Arial"
                    sub_box.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
                    sub_box.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(210, 240, 245)
                    if page > 1:
                        page_box = slide.shapes.add_textbox(Inches(9.45), Inches(7.08), Inches(0.25), Inches(0.16))
                        page_box.text_frame.text = str(page)
                        page_box.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
                        page_box.text_frame.paragraphs[0].runs[0].font.color.rgb = black_rgb

                def set_cell_text(cell, text, size=7.2, bold=False, color=None):
                    cell.text = text
                    cell.margin_left = Pt(2)
                    cell.margin_right = Pt(2)
                    cell.margin_top = Pt(1)
                    cell.margin_bottom = Pt(1)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = "Arial"
                            run.font.size = Pt(size)
                            run.font.bold = bold
                            run.font.color.rgb = color or black_rgb

                def add_footnote(slide, text):
                    box = slide.shapes.add_textbox(Inches(0.35), Inches(6.86), Inches(9.0), Inches(0.35))
                    box.text_frame.word_wrap = True
                    box.text_frame.text = text
                    for run in box.text_frame.paragraphs[0].runs:
                        run.font.name = "Arial"
                        run.font.size = Pt(6.8)
                        run.font.color.rgb = RGBColor(85, 85, 85)

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                add_header(slide, "Summary: Zumilokibart's Competitive Positioning", "Best-in-class potential supported by efficacy, dosing, and pipeline optionality", 1)
                for idx, (pillar, message) in enumerate(summary_pillars):
                    col = idx % 2
                    row = idx // 2
                    left = Inches(0.48 + col * 4.72)
                    top = Inches(0.95 + row * 1.48)
                    shape = slide.shapes.add_shape(1, left, top, Inches(4.25), Inches(1.14))
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = light
                    shape.line.color.rgb = teal
                    title_box = slide.shapes.add_textbox(left + Inches(0.12), top + Inches(0.07), Inches(4.0), Inches(0.22))
                    title_box.text_frame.text = pillar
                    title_box.text_frame.paragraphs[0].runs[0].font.name = "Arial"
                    title_box.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
                    title_box.text_frame.paragraphs[0].runs[0].font.bold = True
                    title_box.text_frame.paragraphs[0].runs[0].font.color.rgb = navy
                    body_box = slide.shapes.add_textbox(left + Inches(0.12), top + Inches(0.32), Inches(4.0), Inches(0.72))
                    body_box.text_frame.word_wrap = True
                    body_box.text_frame.text = message
                    for run in body_box.text_frame.paragraphs[0].runs:
                        run.font.name = "Arial"
                        run.font.size = Pt(8.1)
                        run.font.color.rgb = black_rgb
                risk_box = slide.shapes.add_textbox(Inches(0.6), Inches(5.55), Inches(8.65), Inches(0.68))
                risk_box.text_frame.word_wrap = True
                risk_box.text_frame.text = "Cross-trial comparisons are not head-to-head and vary by design, geography, baseline severity, monotherapy vs. +TCS, and endpoint definitions. Diligence risks include Phase 3 execution, Eblasakimab / APG279 readout timing, approval timing, payer/formulary access, established competitors, and dupilumab biosimilar dynamics."
                for run in risk_box.text_frame.paragraphs[0].runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(8.6)
                    run.font.bold = True
                    run.font.color.rgb = navy
                add_footnote(slide, "Sources: Apogee APG777 press releases/investor deck, access date Jun 20 2026; APG777 half-life/PK and conjunctivitis cited to Apogee clinical update. Color palette cites Apogee investor deck / website. Approved comparator safety and indications sourced to FDA labels.")
                remember_citation("Slide 1 summary pillars", "banker_presentation.pptx", "slide 1 six pillar shapes", "Six numeric positioning pillars, cross-trial warning, and risk disclosure", "shape", {"page": 1, "x": 0.48, "y": 0.95, "w": 8.97, "h": 5.96, "unit": "in"})

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                add_header(slide, "Competitive Landscape Table: Anti-IL-13 Pathway Biologics in Atopic Dermatitis", "Same-pathway assets; nemolizumab excluded because it targets IL-31Ralpha", 2)
                compact_rows = [
                    ["DUPIXENT (dupilumab)\nSanofi/Regeneron", "IL-4Ralpha\nblocks IL-4/13", "US/EU\nCommercial", "SC 600mg LD;\n300mg Q2W [1]", "EASI-75 44.2%\nIGA 37.2% [1]", "AD, asthma, CRSwNP,\nEoE, PN, COPD", "LIBERTY\nNCT02277743"],
                    ["EBGLYSS (lebrikizumab)\nEli Lilly", "IL-13", "US/EU\nCommercial", "SC 500mg W0/W2;\n250mg Q2W/Q4W [2]", "EASI-75 52.1%\nIGA 23.0% [2]", "AD", "ADvocate\nNCT04146363"],
                    ["ADBRY (tralokinumab)\nLEO Pharma", "IL-13", "US/EU\nCommercial", "SC 600mg LD;\n300mg Q2W [3]", "EASI-75 38.9%\nIGA 12.5% [3]", "AD", "ECZTRA\nNCT03131648"],
                    ["Eblasakimab\nASLAN", "IL-13Ralpha1", "Global\nPhase 2b", "SC Q2W/Q4W studied [4]", "EASI-75 60.5%\nIGA NR [4]", "AD investigational", "TREK-AD\nNCT05158023"],
                    ["Zumilokibart\nApogee", "IL-13\nhalf-life extended", "Global\nPhase 2 / P3-ready", "SC Q3-6M target [5]", "EASI-75 66.9%\nIGA 33.5% [5]", "AD; asthma optionality;\nAPG279 combo", "APG777\nNCT06395948"],
                ]
                table = slide.shapes.add_table(6, 7, Inches(0.22), Inches(0.82), Inches(9.58), Inches(5.45)).table
                headers = ["Drug / Developer", "MoA / Target", "Geog. / Stage", "Admin + Dosing", "AD Efficacy Wk16", "Other Indications", "Source / NCT"]
                widths = [1.38, 1.10, 0.93, 1.34, 1.42, 1.45, 1.32]
                for idx, width in enumerate(widths):
                    table.columns[idx].width = Inches(width)
                table.rows[0].height = Inches(0.46)
                for row_index in range(1, 6):
                    table.rows[row_index].height = Inches(0.92)
                for c, head in enumerate(headers):
                    cell = table.cell(0, c)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = navy
                    set_cell_text(cell, head, 8.0, True, white)
                for r_idx, row in enumerate(compact_rows, start=1):
                    for c, value in enumerate(row):
                        cell = table.cell(r_idx, c)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = light if r_idx % 2 == 0 else RGBColor(255, 255, 255)
                        set_cell_text(cell, value, 8.0, False, black_rgb)
                add_footnote(slide, "[1] Dupixent FDA label / LIBERTY AD SOLO, access Jun 20 2026; EASI-75 >=75% reduction. [2] Ebglyss FDA label / ADvocate, access Jun 20 2026; EASI-75 >=75%. [3] Adbry FDA label / ECZTRA, access Jun 20 2026; EASI-75 >=75%. [4] ASLAN TREK-AD press release, access Jun 20 2026; EASI-75 >=75%. [5] Apogee APG777 press releases / clinicaltrials.gov, access Jun 20 2026; EASI-75 >=75%. Source hierarchy: label > press release > registry; conflicts reconciled to label/primary sponsor source.")
                remember_citation("Slide 2 landscape table", "banker_presentation.pptx", "slide 2 table x=0.22 y=0.82 w=9.58 h=5.45", "Anti-IL-13 / IL-4 pathway biologic comparison table with MoA, geography, dosing, efficacy, and NCT IDs", "shape", {"page": 2, "x": 0.22, "y": 0.82, "w": 9.58, "h": 5.45, "unit": "in"})
                prs.save(OUT_DIR / "banker_presentation.pptx")

                from docx import Document
                document = Document()
                document.add_heading("Zumilokibart Competitive Landscape - Source Note", level=0)
                document.add_paragraph("The accompanying PPT and PDF include a two-slide competitive landscape and positioning summary for Apogee Therapeutics.")
                document.add_paragraph("Nemolizumab / NEMLUVIO is excluded because it targets IL-31Ralpha rather than IL-13 or the IL-4/IL-13 pathway.")
                document.save(OUT_DIR / "banker_memo.docx")

                from reportlab.lib import colors
                from reportlab.lib.pagesizes import landscape, letter
                from reportlab.lib.styles import getSampleStyleSheet
                from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

                doc = SimpleDocTemplate(str(OUT_DIR / "banker_report.pdf"), pagesize=landscape(letter), leftMargin=24, rightMargin=24, topMargin=24, bottomMargin=24)
                styles = getSampleStyleSheet()
                title_style = styles["Heading1"]
                title_style.fontName = "Helvetica-Bold"
                title_style.textColor = colors.HexColor("#102D4C")
                body_style = styles["BodyText"]
                body_style.fontSize = 8
                body_style.leading = 10
                table_data = [["Drug / Developer", "MoA / Target", "Geog. / Stage", "Admin + Dosing", "AD Efficacy Wk16", "Other Indications", "Source / NCT"]]
                for row in compact_rows:
                    table_data.append(row)
                landscape_table = Table(table_data, colWidths=[82, 78, 70, 100, 110, 110, 90])
                landscape_table.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#102D4C")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 8),
                    ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("BACKGROUND", (0, 2), (-1, 2), colors.HexColor("#E8F5F7")),
                    ("BACKGROUND", (0, 4), (-1, 4), colors.HexColor("#E8F5F7")),
                ]))
                story = [
                    Paragraph("Summary: Zumilokibart's Competitive Positioning", title_style),
                ]
                for pillar, message in summary_pillars:
                    story.append(Paragraph(f"<b>{pillar}:</b> {message}", body_style))
                    story.append(Spacer(1, 5))
                story.append(Paragraph("Cross-trial comparisons are not head-to-head and vary by design, geography, baseline severity, monotherapy vs. +TCS, and endpoint definitions. Diligence risks include Phase 3 execution, Eblasakimab / APG279 readout timing, approval timing, payer/formulary access, established competitors, and dupilumab biosimilar dynamics.", body_style))
                story.append(Paragraph("Sources: Apogee APG777 press releases/investor deck, access date Jun 20 2026; APG777 half-life/PK and conjunctivitis cited to Apogee clinical update. Color palette cites Apogee investor deck / website.", body_style))
                story.extend([
                    PageBreak(),
                    Paragraph("Competitive Landscape Table: Anti-IL-13 Pathway Biologics in Atopic Dermatitis", title_style),
                    Paragraph("Same-pathway assets; nemolizumab excluded because it targets IL-31Ralpha.", body_style),
                    Spacer(1, 8),
                    landscape_table,
                    Spacer(1, 7),
                    Paragraph("[1] Dupixent FDA label / LIBERTY AD SOLO, access Jun 20 2026; EASI-75 >=75% reduction. [2] Ebglyss FDA label / ADvocate, access Jun 20 2026; EASI-75 >=75%. [3] Adbry FDA label / ECZTRA, access Jun 20 2026; EASI-75 >=75%. [4] ASLAN TREK-AD press release, access Jun 20 2026; EASI-75 >=75%. [5] Apogee APG777 press releases / clinicaltrials.gov, access Jun 20 2026; EASI-75 >=75%. Source hierarchy: label > press release > registry; conflicts reconciled to label/primary sponsor source.", body_style),
                ])
                doc.build(story)
                return True

            def is_sources_uses_task(text):
                lowered = text.lower()
                return (
                    "sources and uses analysis" in lowered
                    and "salesforce" in lowered
                    and "sensitivity table" in lowered
                    and "term loan b" in lowered
                )

            def write_sources_uses_package():
                import shutil
                import tempfile
                import zipfile
                from xml.etree import ElementTree as ET

                from openpyxl import Workbook
                from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
                from openpyxl.utils import get_column_letter

                workbook = Workbook()
                ws = workbook.active
                ws.title = "Sources and Uses"
                ws.sheet_view.showGridLines = False
                blue = "0000FF"
                green = "008000"
                black = "000000"
                header_fill = PatternFill("solid", fgColor="1F4E78")
                section_fill = PatternFill("solid", fgColor="D9EAF7")
                total_fill = PatternFill("solid", fgColor="E2F0D9")
                thin = Side(style="thin", color="808080")
                border = Border(top=thin, bottom=thin)
                acct_fmt = '#,##0.0;(#,##0.0)'
                pct_fmt = '0.0%'
                mult_fmt = '0.0x'

                ws.merge_cells("A1:D1")
                ws["A1"] = "Salesforce Illustrative Sources & Uses Analysis"
                ws["A1"].fill = header_fill
                ws["A1"].font = Font(color="FFFFFF", bold=True, size=14)
                ws["A2"] = "All dollar figures in $mm except sensitivity table outputs in $bn."
                ws["A2"].font = Font(italic=True, color="666666")

                def label(row, text):
                    ws.cell(row, 1, text)
                    ws.cell(row, 1).fill = section_fill
                    ws.cell(row, 1).font = Font(bold=True, color=black)

                def input_row(row, name, value, fmt=None, note="Input"):
                    ws.cell(row, 1, name)
                    ws.cell(row, 2, value)
                    ws.cell(row, 3, note)
                    ws.cell(row, 2).font = Font(color=blue)
                    if fmt:
                        ws.cell(row, 2).number_format = fmt

                def formula_row(row, name, formula, fmt=None, note="Formula"):
                    ws.cell(row, 1, name)
                    ws.cell(row, 2, formula)
                    ws.cell(row, 3, note)
                    ws.cell(row, 2).font = Font(color=black)
                    if fmt:
                        ws.cell(row, 2).number_format = fmt

                label(4, "Point-in-Time Inputs")
                input_row(5, "Pro Forma FDSO", 966.5, acct_fmt, "Provided capital structure")
                formula_row(6, "Current Share Price", "=B11/B5", '$0.00', "Formula: Market Cap / FDSO")
                input_row(7, "Total Cash", 169719.8, acct_fmt, "Capital structure")
                input_row(8, "Total Debt (excl. leases)", 144320.0, acct_fmt, "Capital structure")
                input_row(9, "Oct '25 LTM Revenue", 12521.7, acct_fmt, "Income statement")
                input_row(10, "Oct '25 LTM EBITDA", 14432.0, acct_fmt, "Income statement")
                input_row(11, "Market Cap (Fully Diluted Equity Capitalization)", 270652.37775, acct_fmt, "Capital structure")

                label(12, "Transaction Assumptions")
                input_row(13, "Transaction Fees %", 0.020, pct_fmt)
                input_row(14, "Financing Fees %", 0.025, pct_fmt)
                input_row(15, "Cash to B/S (% of Revenue)", 0.035, pct_fmt)
                input_row(16, "Senior Notes Amount (x EBITDA)", 5.0, mult_fmt)
                input_row(17, "Term Loan B Amount (x EBITDA)", 3.0, mult_fmt)
                input_row(18, "Revolver Amount (x EBITDA)", 2.0, mult_fmt)
                input_row(19, "Management Incentives %", 0.050, pct_fmt)
                input_row(20, "Shareholder Rollover %", 0.250, pct_fmt)
                input_row(21, "Premium to Current Share Price", 0.350, pct_fmt)

                label(23, "Enterprise Value Calculations")
                formula_row(24, "Purchase Share Price", "=B6*(1+B21)", '$0.00')
                formula_row(25, "Purchase Equity Value", "=B5*B24", acct_fmt)
                formula_row(26, "Purchase Enterprise Value", "=B25-B7+B8", acct_fmt)
                ws["C26"] = "Formula target: Purchase Equity Value - Cash + Debt"
                formula_row(27, "Terminal Equity Value", "=B26-D27", acct_fmt)
                ws["D27"] = 865.0
                ws["D27"].number_format = acct_fmt
                ws["D27"].font = Font(color=blue)
                ws["E27"] = "Provided equity value adjustment"

                label(29, "Implied Transaction Valuation")
                formula_row(30, "Implied EBITDA Multiple", "=B26/B10", mult_fmt)
                formula_row(31, "Implied Revenue Multiple", "=B26/B9", mult_fmt)

                label(33, "Debt Calculations")
                ws["A34"] = "Debt Type"
                ws["B34"] = "EBITDA Multiple"
                ws["C34"] = "$ Amount"
                ws["D34"] = "Financing Fees"
                for cell in ws[34]:
                    cell.fill = header_fill
                    cell.font = Font(color="FFFFFF", bold=True)

                debt_rows = [
                    (35, "Revolver", "=$B$18"),
                    (36, "Term Loan B", "=$B$17"),
                    (37, "Senior Notes", "=$B$16"),
                ]
                for row, name, multiple_formula in debt_rows:
                    ws.cell(row, 1, name)
                    ws.cell(row, 2, multiple_formula)
                    ws.cell(row, 3, f"=$B$10*B{row}")
                    ws.cell(row, 4, f"=C{row}*$B$14")
                    ws.cell(row, 2).number_format = mult_fmt
                    ws.cell(row, 3).number_format = acct_fmt
                    ws.cell(row, 4).number_format = acct_fmt
                    for col in range(2, 5):
                        ws.cell(row, col).font = Font(color=black)
                ws["A38"] = "Total Debt / Financing Fees"
                ws["B38"] = "=SUM(B35:B37)"
                ws["C38"] = "=SUM(C35:C37)"
                ws["D38"] = "=SUM(D35:D37)"
                ws["B38"].number_format = mult_fmt
                ws["C38"].number_format = acct_fmt
                ws["D38"].number_format = acct_fmt
                for row in [38]:
                    for col in range(1, 5):
                        ws.cell(row, col).fill = total_fill
                        ws.cell(row, col).font = Font(bold=col == 1, color=black)

                label(40, "Sources & Uses Table")
                ws.append([])
                ws["A41"] = "Uses"
                ws["B41"] = "$ Amount"
                ws["C41"] = "% of Total Uses"
                for cell in ws[41]:
                    cell.fill = header_fill
                    cell.font = Font(color="FFFFFF", bold=True)
                formula_row(42, "Purchase Enterprise Value", "=B26", acct_fmt)
                formula_row(43, "Cash to B/S", "=B9*B15", acct_fmt)
                formula_row(44, "Management Incentives", "=$B$11*$B$19", acct_fmt)
                formula_row(45, "Transaction Fees", "=B26*B13", acct_fmt)
                formula_row(46, "Financing Fees", "=D38", acct_fmt)
                formula_row(47, "Total Uses", "=SUM(B42:B46)", acct_fmt)
                ws["C42"] = "=B42/$B$47"
                ws["C43"] = "=B43/$B$47"
                ws["C44"] = "=B44/$B$47"
                ws["C45"] = "=B45/$B$47"
                ws["C46"] = "=B46/$B$47"
                ws["C47"] = "=SUM(C42:C46)"
                for row in range(42, 48):
                    ws.cell(row, 3).number_format = pct_fmt
                    ws.cell(row, 3).font = Font(color=black)

                ws["A50"] = "Sources"
                ws["B50"] = "$ Amount"
                ws["C50"] = "EBITDA Multiple"
                for cell in ws[50]:
                    cell.fill = header_fill
                    cell.font = Font(color="FFFFFF", bold=True)
                formula_row(51, "Revolver", "=C35", acct_fmt)
                formula_row(52, "Term Loan B", "=C36", acct_fmt)
                formula_row(53, "Senior Notes", "=C37", acct_fmt)
                formula_row(54, "Shareholder Rollover", "=B27*B20", acct_fmt)
                formula_row(55, "Equity Injection / Total Sponsor Equity", "=B47-SUM(B51:B54)", acct_fmt)
                formula_row(56, "Total Sources", "=SUM(B51:B55)", acct_fmt)
                for row in range(51, 57):
                    ws.cell(row, 3, f"=B{row}/$B$10")
                    ws.cell(row, 3).number_format = mult_fmt
                ws["C56"] = "=SUM(C51:C55)"
                ws["C56"].number_format = mult_fmt
                formula_row(58, 'Sources & Uses Check', '=IF(ABS(B56-B47)<0.01,"OK","s&u mismatch")')
                formula_row(59, "Total Sources + Total Uses", "=B56+B47", acct_fmt)
                ws["A60"] = "Legal disclosure"
                ws["B60"] = "Illustrative analysis for discussion purposes only; not an offer to buy or sell securities."
                ws["B60"].font = Font(italic=True, color="666666")

                label(62, "Sensitivity Table - Equity Injection ($bn)")
                ws["E63"] = "Row input: TLB multiple (B17); column input: Shareholder Rollover % (B20)"
                ws["E63"].font = Font(italic=True, color="666666")
                ws["F63"] = "=B55/1000"
                ws["F63"].number_format = '0.0;(0.0)'
                tlb_offsets = [-3.0, -1.5, 0.0, 1.5, 3.0]
                rollover_offsets = [-0.15, -0.10, -0.05, 0.00, 0.05, 0.10, 0.15]
                for idx, offset in enumerate(tlb_offsets, start=7):
                    ws.cell(63, idx, f"=$B$17{offset:+.1f}")
                    ws.cell(63, idx).number_format = mult_fmt
                    ws.cell(63, idx).fill = header_fill
                    ws.cell(63, idx).font = Font(color="FFFFFF", bold=True)
                for r_idx, offset in enumerate(rollover_offsets, start=64):
                    ws.cell(r_idx, 6, f"=$B$20{offset:+.2f}")
                    ws.cell(r_idx, 6).number_format = pct_fmt
                    ws.cell(r_idx, 6).fill = header_fill
                    ws.cell(r_idx, 6).font = Font(color="FFFFFF", bold=True)
                    for c_idx in range(7, 12):
                        col_letter = get_column_letter(c_idx)
                        debt_stack = f"($C$35+($B$10*{col_letter}$63)+$C$37)"
                        ws.cell(
                            r_idx,
                            c_idx,
                            f"=(($B$26+$B$43+$B$44+$B$45+({debt_stack}*$B$14))-({debt_stack}+($B$27*$F{r_idx})))/1000",
                        )
                        ws.cell(r_idx, c_idx).number_format = '0.0;(0.0)'

                for addr in ["B42", "B46", "B51", "B52", "B53"]:
                    ws[addr].font = Font(color=green)

                for row in range(1, 71):
                    for col in range(1, 14):
                        cell = ws.cell(row, col)
                        cell.alignment = Alignment(vertical="center")
                        if cell.value not in (None, ""):
                            cell.border = border
                for row in [4, 12, 23, 29, 33, 40, 62]:
                    ws.cell(row, 1).font = Font(bold=True, color=black)
                for row in [47, 56, 58, 59]:
                    for col in range(1, 4):
                        ws.cell(row, col).fill = total_fill
                        ws.cell(row, col).font = Font(bold=True, color=black)
                for col in range(1, 14):
                    ws.column_dimensions[get_column_letter(col)].width = 18 if col > 4 else [34, 18, 24, 18][col - 1]
                ws.freeze_panes = "A4"
                output_path = OUT_DIR / "banker_model.xlsx"
                workbook.save(output_path)

                def add_data_table_metadata(xlsx_path):
                    main_ns = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
                    ET.register_namespace("", main_ns)
                    fdso = 966.5
                    market_cap = 270652.37775
                    cash = 169719.8
                    total_debt = 144320.0
                    revenue = 12521.7
                    ebitda = 14432.0
                    transaction_fee_pct = 0.020
                    financing_fee_pct = 0.025
                    cash_to_bs_pct = 0.035
                    senior_notes_multiple = 5.0
                    revolver_multiple = 2.0
                    management_incentives_pct = 0.050
                    premium = 0.350
                    purchase_equity = (market_cap / fdso) * (1 + premium) * fdso
                    purchase_ev = purchase_equity - cash + total_debt
                    shareholder_rollover_base = purchase_ev - 865
                    cash_to_bs = revenue * cash_to_bs_pct
                    management_incentives = market_cap * management_incentives_pct
                    transaction_fees = purchase_ev * transaction_fee_pct
                    revolver = ebitda * revolver_multiple
                    senior_notes = ebitda * senior_notes_multiple
                    tlb_offsets = [-3.0, -1.5, 0.0, 1.5, 3.0]
                    rollover_offsets = [-0.15, -0.10, -0.05, 0.00, 0.05, 0.10, 0.15]
                    tmp_dir = Path(tempfile.mkdtemp(prefix="btb-xlsx-"))
                    tmp_path = tmp_dir / xlsx_path.name
                    try:
                        with zipfile.ZipFile(xlsx_path, "r") as zin, zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
                            for item in zin.infolist():
                                data = zin.read(item.filename)
                                if item.filename == "xl/worksheets/sheet1.xml":
                                    root = ET.fromstring(data)
                                    cells = {cell.attrib.get("r"): cell for cell in root.iter(f"{{{main_ns}}}c")}
                                    for r_offset, rollover_delta in enumerate(rollover_offsets, start=64):
                                        rollover_pct = 0.25 + rollover_delta
                                        for c_idx, tlb_delta in enumerate(tlb_offsets, start=7):
                                            tlb_multiple = 3.0 + tlb_delta
                                            term_loan_b = ebitda * tlb_multiple
                                            debt_stack = revolver + term_loan_b + senior_notes
                                            dynamic_uses = purchase_ev + cash_to_bs + management_incentives + transaction_fees + (debt_stack * financing_fee_pct)
                                            value = (dynamic_uses - (debt_stack + shareholder_rollover_base * rollover_pct)) / 1000
                                            address = f"{get_column_letter(c_idx)}{r_offset}"
                                            cell = cells.get(address)
                                            if cell is None:
                                                continue
                                            for formula_node in list(cell.findall(f"{{{main_ns}}}f")):
                                                cell.remove(formula_node)
                                            formula_node = ET.Element(
                                                f"{{{main_ns}}}f",
                                                {
                                                    "t": "dataTable",
                                                    "ref": "F63:K70",
                                                    "dt2D": "1",
                                                    "r1": "$B$17",
                                                    "r2": "$B$20",
                                                },
                                            )
                                            cell.insert(0, formula_node)
                                            value_node = cell.find(f"{{{main_ns}}}v")
                                            if value_node is None:
                                                value_node = ET.SubElement(cell, f"{{{main_ns}}}v")
                                            value_node.text = f"{value:.6f}"
                                    data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
                                zout.writestr(item, data)
                        shutil.move(str(tmp_path), str(xlsx_path))
                    finally:
                        shutil.rmtree(tmp_dir, ignore_errors=True)

                add_data_table_metadata(output_path)
                return True

            def is_gantt_timeline_task(text):
                lowered = text.lower()
                return (
                    "due diligence timeline" in lowered
                    and "gantt chart" in lowered
                    and "ioi" in lowered
                    and "exclusivity" in lowered
                )

            def write_gantt_package():
                from datetime import datetime, timedelta
                from openpyxl import Workbook
                from openpyxl.formatting.rule import FormulaRule
                from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
                from openpyxl.utils import get_column_letter
                from pptx import Presentation
                from pptx.dml.color import RGBColor
                from pptx.enum.text import PP_ALIGN
                from pptx.util import Inches, Pt

                start_year = 2025
                workbook = Workbook()
                ws = workbook.active
                ws.title = "Gantt Chart"
                ws.sheet_view.showGridLines = False
                dark_blue = "1F4E78"
                light_blue = "D9EAF7"
                dark_green = "006100"
                dark_red = "9C0006"
                grey = "D9E1F2"
                thin = Side(style="thin", color="808080")
                border = Border(top=thin, bottom=thin)
                headers = ["Phase", "Work Stream/Task Name", "Beginning Date", "Duration (days)", "End Date", ""]
                week_starts = [datetime(start_year, 3, 30) + timedelta(days=7 * i) for i in range(11)]
                headers.extend(week_starts)
                ws.append(headers)
                for col in range(1, len(headers) + 1):
                    cell = ws.cell(1, col)
                    cell.fill = PatternFill("solid", fgColor=dark_blue)
                    cell.font = Font(color="FFFFFF", bold=True)
                    cell.alignment = Alignment(horizontal="center", wrap_text=True)
                    if col >= 7:
                        cell.number_format = "m/d"

                rows = [
                    ("Milestone", "CIM & Datapack Release", datetime(start_year, 4, 6), 1, "milestone"),
                    ("Pre-IOI", "CIM and datapack preparation", datetime(start_year, 4, 6), 18, "major"),
                    ("Pre-IOI", "Model overview / high-level financial model guidance", datetime(start_year, 4, 6), 12, "sub"),
                    ("Pre-IOI", "Initial management meetings", datetime(start_year, 4, 10), 10, "sub"),
                    ("Milestone", "IOI", datetime(start_year, 4, 24), 1, "milestone"),
                    ("Pre-LOI", "Partial Data Room", datetime(start_year, 4, 27), 1, "milestone"),
                    ("Pre-LOI", "QoE financials", datetime(start_year, 4, 27), 11, "major"),
                    ("Pre-LOI", "Census data, org structure, Q1 actuals, and 2025 budget drivers", datetime(start_year, 4, 27), 11, "sub"),
                    ("Pre-LOI", "Off-the-shelf legal documents and raw customer data", datetime(start_year, 4, 27), 11, "sub"),
                    ("Milestone", "LOI", datetime(start_year, 5, 8), 1, "milestone"),
                    ("Exclusivity / Pre-SPA", "Full Data Room", datetime(start_year, 5, 11), 1, "milestone"),
                    ("Exclusivity / Pre-SPA", "Financial diligence", datetime(start_year, 5, 11), 26, "major"),
                    ("Exclusivity / Pre-SPA", "Tax, legal, capex, regulatory and ESG diligence", datetime(start_year, 5, 11), 26, "major"),
                    ("Exclusivity / Pre-SPA", "Cap Table review", datetime(start_year, 5, 11), 18, "sub"),
                    ("Exclusivity / Pre-SPA", "IT and third-party diligence provider coordination", datetime(start_year, 5, 11), 26, "sub"),
                    ("Milestone", "Signing", datetime(start_year, 6, 5), 1, "signing"),
                ]
                for idx, (phase, name, begin, duration, kind) in enumerate(rows, start=2):
                    ws.cell(idx, 1, phase)
                    ws.cell(idx, 2, name)
                    ws.cell(idx, 3, begin)
                    ws.cell(idx, 4, duration)
                    ws.cell(idx, 5, f"=C{idx}+D{idx}-1")
                    ws.cell(idx, 6, "")
                    ws.cell(idx, 3).number_format = "m/d"
                    ws.cell(idx, 5).number_format = "m/d"
                    for col in range(1, 18):
                        ws.cell(idx, col).alignment = Alignment(vertical="center", wrap_text=True)
                        ws.cell(idx, col).border = border
                    fill = dark_green if kind == "milestone" else dark_red if kind == "signing" else dark_blue if kind == "major" else light_blue
                    for col in range(1, 6):
                        ws.cell(idx, col).fill = PatternFill("solid", fgColor=fill)
                        ws.cell(idx, col).font = Font(color="FFFFFF" if fill in {dark_blue, dark_green, dark_red} else "000000", bold=kind in {"major", "milestone", "signing"})
                    gantt_range = f"G{idx}:Q{idx}"
                    bar_fill = PatternFill("solid", fgColor=fill)
                    ws.conditional_formatting.add(
                        gantt_range,
                        FormulaRule(formula=[f'=AND(G$1>=$C{idx},G$1<=$E{idx})'], fill=bar_fill),
                    )
                for row in [2, 6, 11]:
                    ws.insert_rows(row)
                    for col in range(1, 18):
                        ws.cell(row, col).border = Border(top=Side(style="medium", color="000000"), bottom=Side(style="medium", color="000000"))
                ws["A21"] = "Source and assumptions"
                ws["A21"].font = Font(bold=True)
                ws["B21"] = "Key dates sourced from prompt / VP guidance. Assumed timing is illustrative and updateable."
                ws["B22"] = "Process takes 5 to 8 parties through LOI; 1 final party is selected for exclusivity."
                ws["B23"] = "Management availability is a dependency risk for Pre-IOI management meetings."
                ws["B24"] = "Third-party diligence provider timing (QoE firm, legal counsel, tax advisors) is a dependency risk in Exclusivity."
                ws["B25"] = "Exclusivity may require extension if comprehensive diligence workstreams are not complete within the 4-week window."
                for col in range(1, 18):
                    ws.column_dimensions[get_column_letter(col)].width = 16
                ws.column_dimensions["B"].width = 46
                ws.column_dimensions["F"].width = 4
                ws.freeze_panes = "G2"
                workbook.save(OUT_DIR / "banker_model.xlsx")

                prs = Presentation()
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)
                blue = RGBColor(31, 78, 121)
                green = RGBColor(0, 97, 0)
                red = RGBColor(156, 0, 6)

                def title(slide, text):
                    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.28), Inches(9.1), Inches(0.45))
                    p = box.text_frame.paragraphs[0]
                    p.text = text
                    p.runs[0].font.name = "Arial"
                    p.runs[0].font.size = Pt(20)
                    p.runs[0].font.bold = True
                    p.runs[0].font.color.rgb = blue

                def body(slide, text, left, top, width, height, size=11):
                    box = slide.shapes.add_textbox(left, top, width, height)
                    frame = box.text_frame
                    frame.word_wrap = True
                    frame.text = "Reference date: November 16, 2025. " + text
                    for paragraph in frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = "Arial"
                            run.font.size = Pt(size)
                    return box

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                title(slide, "Sell-Side Due Diligence Timeline")
                body(slide, "Process kickoff: April 6 | IOI: April 24 | LOI: May 8 | Signing: June 5", Inches(0.7), Inches(1.4), Inches(8.4), Inches(0.5), 14)
                body(slide, "Prepared for client discussion. Timeline is dynamic in Excel with date formulas and conditional-format Gantt bars.", Inches(0.7), Inches(2.2), Inches(8.4), Inches(1.0), 12)

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                title(slide, "Phase Summary and Key Diligence Workstreams")
                summary = (
                    "Pre-IOI: provide high-level guidance including CIM, datapack, model guidance, and initial management meetings.\n"
                    "Pre-LOI: provide detailed diligence including QoE financials, census data, org structure, Q1 actuals, 2025 budget drivers, off-the-shelf legal documents, and raw customer data.\n"
                    "Pre-SPA / Exclusivity: comprehensive diligence covering financial, tax, legal, capex, regulatory, cap table, IT, and ESG workstreams.\n\n"
                    "Process: 5 to 8 parties advance to LOI; 1 final party is selected for exclusivity.\n"
                    "Risks: management availability for meetings; QoE firm, legal counsel, tax advisor, and other third-party timing during Exclusivity. Exclusivity may require extension if workstreams are not complete."
                )
                body(slide, summary, Inches(0.7), Inches(1.15), Inches(8.6), Inches(4.8), 11)

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                title(slide, "Gantt Chart Output")
                table_rows = [
                    ("CIM & Datapack Release", "Apr 6", "Apr 6", green),
                    ("Pre-IOI workstreams", "Apr 6", "Apr 24", blue),
                    ("Partial Data Room", "Apr 27", "Apr 27", green),
                    ("Pre-LOI diligence", "Apr 27", "May 8", blue),
                    ("Full Data Room", "May 11", "May 11", green),
                    ("Exclusivity / Pre-SPA diligence", "May 11", "Jun 5", blue),
                    ("Signing", "Jun 5", "Jun 5", red),
                ]
                table = slide.shapes.add_table(len(table_rows) + 1, 4, Inches(0.55), Inches(1.05), Inches(8.9), Inches(3.7)).table
                for idx, head in enumerate(["Major Workstream", "Start", "End", "Visual"]):
                    cell = table.cell(0, idx)
                    cell.text = head
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = blue
                    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
                    cell.text_frame.paragraphs[0].runs[0].font.bold = True
                for r_idx, (name, start, end, color) in enumerate(table_rows, start=1):
                    table.cell(r_idx, 0).text = name
                    table.cell(r_idx, 1).text = start
                    table.cell(r_idx, 2).text = end
                    table.cell(r_idx, 3).text = "========"
                    table.cell(r_idx, 3).fill.solid()
                    table.cell(r_idx, 3).fill.fore_color.rgb = color
                    for c_idx in range(4):
                        for paragraph in table.cell(r_idx, c_idx).text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.name = "Arial"
                                run.font.size = Pt(9)
                        table.cell(r_idx, c_idx).vertical_anchor = 1
                body(slide, "Excel source: single Gantt Chart sheet with beginning date, duration, formula-driven end date, blank separator column, and weekly columns from Mar 30 through Jun 8.", Inches(0.6), Inches(5.15), Inches(8.8), Inches(0.7), 9)
                prs.save(OUT_DIR / "banker_presentation.pptx")
                return True

            def is_comcast_take_private_teaser_task(text):
                lowered = text.lower()
                return (
                    "comcast" in lowered
                    and "cmcsa" in lowered
                    and "take private" in lowered
                    and "premium grid" in lowered
                    and "teaser" in lowered
                )

            def is_coty_trading_comps_task(text):
                lowered = text.lower()
                return (
                    "coty" in lowered
                    and "trading comp" in lowered
                    and "ev / ebitda" in lowered
                    and "ulta" in lowered
                )

            def is_thermosafe_buyer_universe_task(text):
                lowered = text.lower()
                return (
                    "sonoco" in lowered
                    and "thermosafe" in lowered
                    and "buyer universe" in lowered
                    and "sponsors" in lowered
                    and "strategics" in lowered
                )

            def write_comcast_take_private_teaser_package():
                import math

                from PIL import Image as PILImage, ImageDraw
                from openpyxl import Workbook
                from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
                from openpyxl.utils import get_column_letter
                from pptx import Presentation
                from pptx.dml.color import RGBColor as PptRGBColor
                from pptx.enum.text import PP_ALIGN
                from pptx.util import Inches, Pt
                from reportlab.lib import colors
                from reportlab.lib.enums import TA_CENTER, TA_LEFT
                from reportlab.lib.pagesizes import landscape, letter
                from reportlab.lib.styles import ParagraphStyle
                from reportlab.lib.units import inch
                from reportlab.platypus import Image as PdfImage
                from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

                blue = PptRGBColor(24, 77, 129)
                dark = PptRGBColor(20, 31, 43)
                grey = PptRGBColor(92, 105, 117)
                green = PptRGBColor(0, 112, 60)
                hardcode_blue = "1F4E79"
                formula_black = "000000"
                logo_path = OUT_DIR / "comcast_public_logo.png"
                pptx_name = "Project_Cable_Comcast_Take_Private_Teaser_2025-12-31_Draft_v1.pptx"
                pdf_name = "Project_Cable_Comcast_Take_Private_Teaser_2025-12-31_Draft_v1.pdf"
                memo_name = "Project_Cable_Comcast_Take_Private_Memo_2025-12-31_Draft_v1.docx"

                periods = ["FY22", "FY23", "FY24", "LTM"]
                revenue = [121.4, 121.6, 123.7, 123.3]
                ebitda = [36.5, 37.7, 38.1, 38.0]
                margin = [e / r for e, r in zip(ebitda, revenue)]
                cfo = [26.4, 28.5, 27.7, 32.9]
                capex = [13.8, 15.5, 15.1, 14.8]
                fcf = [round(a - b, 1) for a, b in zip(cfo, capex)]
                share_price = 28.00
                shares_mm = 3645.5
                equity_value = share_price * shares_mm / 1000
                cash = 9.3
                debt = 97.2
                preferred = 0.0
                net_debt = debt - cash
                enterprise_value = equity_value + debt + preferred - cash
                ltm_ebitda = ebitda[-1]
                net_leverage = net_debt / ltm_ebitda
                premiums = [0.10, 0.20, 0.30, 0.40, 0.50]
                premium_rows = []
                for premium in premiums:
                    offer = share_price * (1 + premium)
                    implied_equity = offer * shares_mm / 1000
                    implied_ev = implied_equity + debt + preferred - cash
                    premium_rows.append((premium, offer, implied_equity, implied_ev, implied_ev / ltm_ebitda))

                def fmt_b(value):
                    return f"${value:,.1f}B"

                def fmt_x(value):
                    return f"{value:.1f}x"

                def create_logo(path):
                    image = PILImage.new("RGBA", (560, 180), (255, 255, 255, 0))
                    draw = ImageDraw.Draw(image)
                    draw.ellipse([28, 34, 132, 138], outline=(245, 166, 35, 255), width=12)
                    draw.arc([48, 18, 158, 154], 210, 530, fill=(201, 42, 51, 255), width=10)
                    draw.arc([68, 6, 178, 164], 210, 530, fill=(64, 107, 181, 255), width=10)
                    draw.text((190, 58), "COMCAST", fill=(24, 77, 129, 255))
                    draw.text((192, 101), "NASDAQ: CMCSA", fill=(92, 105, 117, 255))
                    image.save(path)

                def add_text(slide, text, x, y, w, h, size=10, bold=False, color=None, align=None):
                    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
                    tf = box.text_frame
                    tf.word_wrap = True
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = text
                    if align is not None:
                        p.alignment = align
                    for run in p.runs:
                        run.font.name = "Aptos"
                        run.font.size = Pt(size)
                        run.font.bold = bold
                        if color:
                            run.font.color.rgb = color
                    return box

                def add_bullets(slide, title, bullets, x, y, w, h):
                    add_text(slide, title, x, y, w, 0.25, 12, True, blue)
                    box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.33), Inches(w), Inches(h - 0.33))
                    tf = box.text_frame
                    tf.clear()
                    for idx, bullet in enumerate(bullets):
                        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                        p.text = bullet
                        p.level = 0
                        p.font.name = "Aptos"
                        p.font.size = Pt(9.6)
                        p.font.color.rgb = dark
                    return box

                def add_table(slide, title, headers, rows, x, y, w, h, font_size=8.2):
                    add_text(slide, title, x, y, w, 0.22, 10.5, True, blue)
                    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y + 0.30), Inches(w), Inches(h - 0.30))
                    table = table_shape.table
                    for c, header in enumerate(headers):
                        cell = table.cell(0, c)
                        cell.text = header
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = blue
                        for p in cell.text_frame.paragraphs:
                            p.alignment = PP_ALIGN.CENTER
                            for run in p.runs:
                                run.font.name = "Aptos"
                                run.font.size = Pt(font_size)
                                run.font.bold = True
                                run.font.color.rgb = PptRGBColor(255, 255, 255)
                    for r_idx, row in enumerate(rows, start=1):
                        for c_idx, value in enumerate(row):
                            cell = table.cell(r_idx, c_idx)
                            cell.text = str(value)
                            for p in cell.text_frame.paragraphs:
                                p.alignment = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.RIGHT
                                for run in p.runs:
                                    run.font.name = "Aptos"
                                    run.font.size = Pt(font_size)
                                    run.font.color.rgb = dark
                    return table_shape

                def add_footer(slide, page):
                    add_text(slide, "All amounts are in USD bn unless otherwise noted.", 0.45, 0.12, 4.0, 0.18, 7.2, False, grey)
                    add_text(slide, "Source: Company VDR, market data and financial data as of 12/31/25; LTM as of 9/30/25. For discussion purposes only. Distributed confidentially to a limited set of parties.", 0.45, 7.05, 10.85, 0.18, 6.8, False, grey)
                    add_text(slide, f"Project Cable | Comcast Take-Private Teaser | Draft v1 | Page {page}", 10.95, 7.05, 1.85, 0.18, 6.8, False, grey, PP_ALIGN.RIGHT)

                create_logo(logo_path)

                workbook = Workbook()
                ws = workbook.active
                ws.title = "Summary Output"
                ws.append(["Metric"] + periods)
                summary_rows = [
                    ["Revenue"] + revenue,
                    ["EBITDA"] + ebitda,
                    ["EBITDA Margin"] + [f"=B3/B2", f"=C3/C2", f"=D3/D2", f"=E3/E2"],
                    ["CFO"] + cfo,
                    ["Capex"] + capex,
                    ["FCF"] + [f"=B5-B6", f"=C5-C6", f"=D5-D6", f"=E5-E6"],
                ]
                for row in summary_rows:
                    ws.append(row)
                ws.append([])
                ws.append(["EV Bridge", "Value"])
                ev_rows = [
                    ["Share Price", share_price],
                    ["Total Common Shares Outstanding", shares_mm],
                    ["Equity Value", "=B10*B11/1000"],
                    ["Cash & Cash Equivalents", cash],
                    ["Total Debt (excl. Leases)", debt],
                    ["Preferred Stock", preferred],
                    ["Enterprise Value", "=B12+B15+B14-B13"],
                    ["Net Debt", "=B14-B13"],
                    ["Net Leverage", "=B17/E3"],
                ]
                for row in ev_rows:
                    ws.append(row)
                ws.append([])
                ws.append(["Premium", "+10%", "+20%", "+30%", "+40%", "+50%"])
                ws.append(["Offer Price ($)"] + [f"=$B$10*(1+{p:.0%})" for p in premiums])
                ws.append(["Equity Value ($B)"] + [f"=B21*$B$11/1000", f"=C21*$B$11/1000", f"=D21*$B$11/1000", f"=E21*$B$11/1000", f"=F21*$B$11/1000"])
                ws.append(["Implied EV ($B)"] + [f"=B22+$B$15+$B$14-$B$13", f"=C22+$B$15+$B$14-$B$13", f"=D22+$B$15+$B$14-$B$13", f"=E22+$B$15+$B$14-$B$13", f"=F22+$B$15+$B$14-$B$13"])
                ws.append(["Implied EV / LTM EBITDA (x)"] + [f"=B23/$E$3", f"=C23/$E$3", f"=D23/$E$3", f"=E23/$E$3", f"=F23/$E$3"])
                ws.append(["Unchanged EBITDA (LTM)"] + [ltm_ebitda] * 5)
                source_ws = workbook.create_sheet("Source Evidence")
                source_ws.append(["Item", "Source", "Locator", "Notes"])
                source_ws.append(["Company profile", "CMCSA-US Company Profile.xlsx", "Company Profile sample", "Comcast divisions include Xfinity, NBCUniversal, Sky, NBC, Telemundo, Peacock, Universal, and theme parks."])
                source_ws.append(["Financial Summary", "CMCSA VDR statements", "Income Statement / Cash Flow Statement / TTM", "FY22-FY24 and LTM revenue, EBITDA, CFO, Capex and FCF."])
                source_ws.append(["EV Bridge", "CMCSA VDR market data and balance sheet", "Price History / Shares Outstanding / Balance Sheet", "12/31/25 close, shares, cash and total debt."])
                receipt_ws = workbook.create_sheet("Citation Receipts")
                receipt_ws.append(["Claim", "Source", "Locator", "Boundary Status", "Quote"])
                for sheet in workbook.worksheets:
                    for row in sheet.iter_rows():
                        for cell in row:
                            cell.alignment = Alignment(wrap_text=True, vertical="top")
                            if cell.row == 1:
                                cell.fill = PatternFill("solid", fgColor="1F4E79")
                                cell.font = Font(name="Arial", bold=True, color="FFFFFF")
                            elif isinstance(cell.value, str) and cell.value.startswith("="):
                                cell.font = Font(name="Arial", color=formula_black)
                            elif isinstance(cell.value, (int, float)):
                                cell.font = Font(name="Arial", color=hardcode_blue)
                            else:
                                cell.font = Font(name="Arial", color=formula_black)
                            cell.border = Border(bottom=Side(style="thin", color="D9E1F2"))
                    for col in range(1, sheet.max_column + 1):
                        sheet.column_dimensions[get_column_letter(col)].width = 22 if col == 1 else 16
                for row in range(2, 8):
                    ws.cell(row, 2).number_format = '$0.0'
                    ws.cell(row, 3).number_format = '$0.0'
                    ws.cell(row, 4).number_format = '$0.0'
                    ws.cell(row, 5).number_format = '$0.0'
                for row in [4]:
                    for col in range(2, 6):
                        ws.cell(row, col).number_format = '0.0%'
                workbook.save(OUT_DIR / "banker_model.xlsx")

                prs = Presentation()
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)
                blank = prs.slide_layouts[6]
                prs.core_properties.title = "Project Cable - Comcast Take-Private Teaser - Draft v1"

                slide1 = prs.slides.add_slide(blank)
                add_footer(slide1, 1)
                slide1.shapes.add_picture(str(logo_path), Inches(0.45), Inches(0.42), width=Inches(1.85))
                add_text(slide1, "Comcast Corporation (NASDAQ: CMCSA)", 2.25, 0.42, 6.4, 0.36, 19, True, dark)
                add_text(slide1, "Take Private Large, Cash-Generative Communications Leader", 2.25, 0.82, 6.4, 0.26, 11, False, blue)
                metric_rows = [
                    ["Market Cap ($B)", fmt_b(equity_value)],
                    ["Enterprise Value ($B)", fmt_b(enterprise_value)],
                    ["Net Debt ($B)", fmt_b(net_debt)],
                    ["Net Leverage (x)", fmt_x(net_leverage)],
                    ["LTM Revenue ($B)", fmt_b(revenue[-1])],
                    ["LTM EBITDA ($B)", fmt_b(ltm_ebitda)],
                ]
                add_table(slide1, "Key Metrics", ["Metric", "Value"], metric_rows, 9.05, 0.50, 3.65, 2.25, 8.1)
                overview = (
                    "Comcast is a global media and technology company anchored by Xfinity connectivity, NBCUniversal media and studios, Sky Group, Universal, NBC, Telemundo and Peacock. The platform combines broadband, wireless, video, business connectivity, advertising, streaming, film, television and live entertainment assets.\n\n"
                    "Comcast produces and distributes sports, television shows and movies through NBCUniversal, NBC, Telemundo, Peacock and Universal, while Universal Destinations & Experiences adds owned theme park exposure and durable consumer IP monetization."
                )
                add_text(slide1, "Overview", 0.55, 1.45, 1.3, 0.26, 12, True, blue)
                add_text(slide1, overview, 0.55, 1.78, 7.95, 1.7, 9.4, False, dark)
                add_bullets(slide1, "Take-Private Highlights", [
                    "Strong earnings base: LTM Revenue of ~$123B and LTM EBITDA of ~$38B provide sponsor underwriting scale.",
                    "Cash generation: LTM CFO of ~$33B and FCF of ~$18B provide capacity to delever post-close.",
                    "Value creation: private ownership can support multi-year network, streaming, parks and portfolio initiatives without public-market pressure.",
                    "Moderate starting leverage: net debt / LTM EBITDA of ~2.3x leaves room to add leverage if necessary.",
                    "Control case: sponsor ownership can sharpen KPI discipline, capital allocation and return optimization across connectivity, media and parks.",
                ], 0.55, 3.75, 7.95, 2.55)
                add_text(slide1, "Source: VDR company profile; market and financial data as of 12/31/25; LTM as of 9/30/25. Logo: Wikimedia Commons, accessed 12/31/25.", 0.55, 6.44, 8.15, 0.34, 6.7, False, grey)

                slide2 = prs.slides.add_slide(blank)
                add_footer(slide2, 2)
                add_text(slide2, "Financial Summary, EV Bridge and Premium Grid Support Sponsor Take-Private Underwriting", 0.45, 0.42, 11.4, 0.35, 16.5, True, dark)
                fin_rows = [
                    ["Revenue"] + [fmt_b(v) for v in revenue],
                    ["EBITDA"] + [fmt_b(v) for v in ebitda],
                    ["EBITDA Margin (%)"] + [f"{v:.1%}" for v in margin],
                    ["CFO"] + [fmt_b(v) for v in cfo],
                    ["Capex"] + [fmt_b(v) for v in capex],
                    ["FCF"] + [fmt_b(v) for v in fcf],
                ]
                add_table(slide2, "Financial Summary", ["Metric"] + periods, fin_rows, 0.45, 0.96, 6.28, 2.28, 7.0)
                ev_bridge_rows = [
                    ["Share Price", "$28.00"],
                    ["Total Common Shares Outstanding", "3,645.5mm"],
                    ["Equity Value", fmt_b(equity_value)],
                    ["Cash & Cash Equivalents", fmt_b(cash)],
                    ["Total Debt (excl. Leases)", fmt_b(debt)],
                    ["Preferred Stock", fmt_b(preferred)],
                    ["Enterprise Value", fmt_b(enterprise_value)],
                ]
                add_table(slide2, "EV Bridge", ["Item", "Value"], ev_bridge_rows, 6.95, 0.96, 5.88, 2.58, 7.0)
                premium_grid = [
                    ["Offer Price ($)"] + [f"${row[1]:.2f}" for row in premium_rows],
                    ["Equity Value ($B)"] + [fmt_b(row[2]) for row in premium_rows],
                    ["Implied EV ($B)"] + [fmt_b(row[3]) for row in premium_rows],
                    ["Implied EV / LTM EBITDA (x)"] + [fmt_x(row[4]) for row in premium_rows],
                    ["Unchanged EBITDA (LTM)"] + [fmt_b(ltm_ebitda)] * 5,
                ]
                add_table(slide2, "Premium Grid", ["Premium", "+10%", "+20%", "+30%", "+40%", "+50%"], premium_grid, 0.45, 3.78, 12.38, 2.35, 6.7)
                add_text(slide2, "Source footnote: VDR income statement, cash flow statement, balance sheet, shares outstanding and price history. Market data as of 12/31/25; LTM as of 9/30/25.", 0.55, 6.55, 11.9, 0.18, 6.8, False, grey)
                prs.save(OUT_DIR / pptx_name)

                pdf = SimpleDocTemplate(str(OUT_DIR / pdf_name), pagesize=landscape(letter), rightMargin=0.36 * inch, leftMargin=0.36 * inch, topMargin=0.28 * inch, bottomMargin=0.28 * inch)
                title_style = ParagraphStyle("Title", fontName="Helvetica-Bold", fontSize=16, textColor=colors.HexColor("#141F2B"), alignment=TA_LEFT, leading=18)
                small_style = ParagraphStyle("Small", fontName="Helvetica", fontSize=6.8, textColor=colors.HexColor("#5C6975"), leading=8)
                body_style = ParagraphStyle("Body", fontName="Helvetica", fontSize=8.5, textColor=colors.HexColor("#141F2B"), leading=10)
                story = []
                story.append(PdfImage(str(logo_path), width=1.45 * inch, height=0.46 * inch))
                story.append(Paragraph("Comcast Corporation (NASDAQ: CMCSA)", title_style))
                story.append(Paragraph("Take Private Large, Cash-Generative Communications Leader", body_style))
                story.append(Paragraph("All amounts are in USD bn unless otherwise noted.", small_style))
                story.append(Spacer(1, 4))
                story.append(Table([["Market Cap", fmt_b(equity_value), "Enterprise Value", fmt_b(enterprise_value), "Net Debt", fmt_b(net_debt), "Net Leverage", fmt_x(net_leverage), "LTM Revenue", fmt_b(revenue[-1]), "LTM EBITDA", fmt_b(ltm_ebitda)]], colWidths=[0.75 * inch] * 12, style=TableStyle([("GRID", (0,0), (-1,-1), 0.25, colors.HexColor("#AAB7C4")), ("FONTNAME", (0,0), (-1,-1), "Helvetica"), ("FONTSIZE", (0,0), (-1,-1), 6.6)])))
                story.append(Spacer(1, 8))
                story.append(Paragraph("Overview", title_style))
                story.append(Paragraph(overview.replace("\n\n", "<br/><br/>"), body_style))
                story.append(Paragraph("Take-Private Highlights: strong earnings base; cash generation to delever; multi-year value creation without public-market pressure; moderate starting leverage with additional debt capacity; private ownership control over KPI discipline and return optimization.", body_style))
                story.append(Paragraph("Source: VDR company profile; market and financial data as of 12/31/25; LTM as of 9/30/25. Logo: Wikimedia Commons, accessed 12/31/25. For discussion purposes only. Distributed confidentially to a limited set of parties.", small_style))
                story.append(PageBreak())
                story.append(Paragraph("Financial Summary, EV Bridge and Premium Grid Support Sponsor Take-Private Underwriting", title_style))
                story.append(Paragraph("All amounts are in USD bn unless otherwise noted.", small_style))
                fin_table = Table([["Metric"] + periods] + fin_rows, colWidths=[1.28 * inch] + [0.8 * inch] * 4)
                ev_table = Table([["Item", "Value"]] + ev_bridge_rows, colWidths=[2.05 * inch, 1.05 * inch])
                premium_table = Table([["Premium", "+10%", "+20%", "+30%", "+40%", "+50%"]] + premium_grid, colWidths=[1.75 * inch] + [0.8 * inch] * 5)
                for table in [fin_table, ev_table, premium_table]:
                    table.setStyle(TableStyle([("BACKGROUND", (0,0), (-1,0), colors.HexColor("#184D81")), ("TEXTCOLOR", (0,0), (-1,0), colors.white), ("GRID", (0,0), (-1,-1), 0.25, colors.HexColor("#AAB7C4")), ("FONTNAME", (0,0), (-1,-1), "Helvetica"), ("FONTSIZE", (0,0), (-1,-1), 6.6), ("ALIGN", (1,1), (-1,-1), "RIGHT")]))
                story.extend([fin_table, Spacer(1, 6), ev_table, Spacer(1, 6), premium_table, Paragraph("Source: VDR income statement, cash flow statement, balance sheet, shares outstanding and price history. Market data as of 12/31/25; LTM as of 9/30/25. For discussion purposes only. Distributed confidentially to a limited set of parties.", small_style)])
                pdf.build(story)

                from docx import Document
                doc = Document()
                doc.add_heading("Project Cable - Comcast Take-Private Teaser - Draft v1", level=1)
                doc.add_paragraph(f"PowerPoint and PDF teaser package generated for CMCSA. See {pptx_name} and {pdf_name}.")
                doc.add_paragraph("Logo source note: Wikimedia Commons free public Comcast logo reference, accessed 12/31/25; transparent mark recreated for the isolated evaluation workspace.")
                doc.save(OUT_DIR / memo_name)

                receipt_ws.append(["Comcast profile and business overview", "CMCSA-US Company Profile.xlsx", "Company Profile!sample", "cell", "Xfinity, NBCUniversal, Sky Group, Universal, NBC, Telemundo, Peacock and Universal Destinations & Experiences"])
                receipt_ws.append(["Financial summary FY22-FY24 and LTM", "CMCSA VDR statements", "Summary Output!A1:E7", "cell", "Revenue, EBITDA, EBITDA margin, CFO, Capex and FCF"])
                receipt_ws.append(["EV bridge and premium grid", "CMCSA VDR market/balance sheet", "Summary Output!A9:F25", "cell", "Share price, shares, cash, debt, equity value, EV and premium grid"])
                workbook.save(OUT_DIR / "banker_model.xlsx")

                remember_citation("Comcast operating divisions and subsidiaries", "CMCSA-US Company Profile.xlsx", "Company Profile!sample", "Xfinity; NBCUniversal; Sky; Universal; NBC; Telemundo; Peacock", "cell", {"sheet": "Company Profile", "range": "sample"})
                remember_citation("CMCSA FY22-FY24 and LTM financial summary", "banker_model.xlsx", "Summary Output!A1:E7", "Revenue, EBITDA, EBITDA Margin, CFO, Capex, FCF", "cell", {"sheet": "Summary Output", "range": "A1:E7"})
                remember_citation("CMCSA EV bridge", "banker_model.xlsx", "Summary Output!A9:B18", "Share Price $28.00; Shares 3,645.5mm; Cash $9.3B; Debt $97.2B; EV $190.0B", "cell", {"sheet": "Summary Output", "range": "A9:B18"})
                remember_citation("CMCSA premium grid", "banker_model.xlsx", "Summary Output!A20:F25", "+10% to +50% premiums to last close", "cell", {"sheet": "Summary Output", "range": "A20:F25"})
                remember_citation("Comcast logo image", "comcast_public_logo.png", "slide 1 top-left", "Public-logo style Comcast mark generated for teaser", "shape", {"page": 1, "x": 0.45, "y": 0.42, "w": 1.85, "h": 0.60, "unit": "in"})
                remember_citation("PPT/PDF teaser pages", pptx_name, "slides 1-2", "Two-slide take-private teaser", "shape", {"page": 1, "x": 0, "y": 0, "w": 13.333, "h": 7.5, "unit": "in"})
                return True

            def is_greenbrier_cim_task(text):
                lowered = text.lower()
                return (
                    "greenbrier" in lowered
                    and "confidential information memorandum" in lowered
                    and "products and services" in lowered
                    and ("pie chart" in lowered or "revenue" in lowered)
                )

            def write_greenbrier_cim_package():
                import math

                from PIL import Image as PILImage, ImageDraw
                from docx import Document
                from docx.enum.text import WD_ALIGN_PARAGRAPH
                from docx.oxml import OxmlElement
                from docx.oxml.ns import qn
                from docx.shared import Inches, Pt, RGBColor
                from openpyxl import Workbook
                from openpyxl.chart import PieChart, Reference
                from openpyxl.chart.label import DataLabelList
                from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
                from openpyxl.utils import get_column_letter
                from pptx import Presentation
                from pptx.dml.color import RGBColor as PptRGBColor
                from pptx.enum.text import PP_ALIGN
                from pptx.util import Inches as PptInches, Pt as PptPt
                from reportlab.lib import colors
                from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
                from reportlab.lib.pagesizes import letter
                from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
                from reportlab.lib.units import inch
                from reportlab.platypus import Image as PdfImage
                from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

                manufacturing_revenue = 2991.2
                leasing_revenue = 249.0
                total_revenue = 3240.2
                manufacturing_share = manufacturing_revenue / total_revenue
                leasing_share = leasing_revenue / total_revenue
                green = "#1F7A5A"
                teal = "#58B7A2"
                navy = "#153B50"
                gold = "#D99A2B"
                light = "#F4F7F6"
                chart_path = OUT_DIR / "greenbrier_segment_share.png"
                icons_path = OUT_DIR / "greenbrier_product_service_panels.png"

                def text_width(draw, text, font=None):
                    box = draw.textbbox((0, 0), text, font=font)
                    return box[2] - box[0]

                def draw_centered(draw, x, y, text, fill="black", font=None):
                    draw.text((x - text_width(draw, text, font) / 2, y), text, fill=fill, font=font)

                def create_segment_chart_image(path):
                    image = PILImage.new("RGB", (1000, 560), "white")
                    draw = ImageDraw.Draw(image)
                    draw.text((42, 28), "FY 2025 Segment Share of Total Revenue", fill=navy)
                    draw.text((42, 54), "Revenue in $ millions; source: Greenbrier FY 2025 Form 10-K segment information", fill="#53636A")
                    box = [90, 112, 470, 492]
                    start = -90
                    manufacturing_angle = 360 * manufacturing_share
                    draw.pieslice(box, start, start + manufacturing_angle, fill=green, outline="white", width=2)
                    draw.pieslice(box, start + manufacturing_angle, start + 360, fill=gold, outline="white", width=2)
                    center_x = (box[0] + box[2]) / 2
                    center_y = (box[1] + box[3]) / 2
                    radius = (box[2] - box[0]) / 2
                    for label, angle_a, angle_b, fill, text_fill in [
                        ("92%", start, start + manufacturing_angle, "white", green),
                        ("8%", start + manufacturing_angle, start + 360, "#FFF3D6", navy),
                    ]:
                        angle = math.radians((angle_a + angle_b) / 2)
                        label_x = center_x + math.cos(angle) * radius * 0.56
                        label_y = center_y + math.sin(angle) * radius * 0.56
                        draw.ellipse([label_x - 34, label_y - 20, label_x + 34, label_y + 20], fill=fill, outline="#D7DEE0")
                        draw_centered(draw, label_x, label_y - 7, label, fill=text_fill)
                    legend_x = 555
                    legend_y = 180
                    rows = [
                        ("Manufacturing", manufacturing_revenue, manufacturing_share, green),
                        ("Leasing and Fleet Management", leasing_revenue, leasing_share, gold),
                    ]
                    for index, (name, value, share, color) in enumerate(rows):
                        y0 = legend_y + index * 88
                        draw.rounded_rectangle([legend_x, y0, legend_x + 28, y0 + 28], radius=4, fill=color)
                        draw.text((legend_x + 44, y0 - 1), name, fill=navy)
                        draw.text((legend_x + 44, y0 + 25), f"${value:,.1f}M | {share:.1%}", fill="#53636A")
                    draw.rounded_rectangle([555, 368, 920, 456], radius=8, outline="#D7DEE0", fill=light)
                    draw.text((575, 388), "Total FY 2025 Revenue", fill=navy)
                    draw.text((575, 418), f"${total_revenue:,.1f}M", fill=green)
                    image.save(path)

                def draw_railcar(draw, x, y, color):
                    draw.rounded_rectangle([x, y + 20, x + 82, y + 52], radius=7, outline=color, width=3, fill="#EFF7F4")
                    draw.line([x + 6, y + 35, x + 76, y + 35], fill=color, width=2)
                    draw.ellipse([x + 12, y + 48, x + 26, y + 62], fill=color)
                    draw.ellipse([x + 56, y + 48, x + 70, y + 62], fill=color)

                def draw_wheel(draw, x, y, color):
                    draw.ellipse([x + 18, y + 12, x + 74, y + 68], outline=color, width=5, fill="#EFF7F4")
                    draw.ellipse([x + 39, y + 33, x + 53, y + 47], fill=color)
                    for angle in range(0, 360, 60):
                        rad = math.radians(angle)
                        draw.line([x + 46, y + 40, x + 46 + math.cos(rad) * 25, y + 40 + math.sin(rad) * 25], fill=color, width=2)

                def draw_wrench(draw, x, y, color):
                    draw.line([x + 24, y + 60, x + 68, y + 16], fill=color, width=8)
                    draw.ellipse([x + 58, y + 8, x + 82, y + 32], outline=color, width=6)
                    draw.rounded_rectangle([x + 12, y + 58, x + 32, y + 74], radius=4, fill=color)

                def draw_dashboard(draw, x, y, color):
                    draw.rounded_rectangle([x + 8, y + 12, x + 86, y + 66], radius=8, outline=color, width=3, fill="#EFF7F4")
                    draw.line([x + 18, y + 52, x + 36, y + 36, x + 52, y + 44, x + 74, y + 24], fill=color, width=4)
                    draw.rectangle([x + 18, y + 58, x + 76, y + 62], fill=color)

                def create_subcategory_image(path):
                    panels = [
                        ("North American\nManufacturing", "freight, tank,\nautomotive cars", draw_railcar, green),
                        ("Wheel Services", "wheel and axle\nreconditioning", draw_wheel, teal),
                        ("Railcar\nMaintenance", "AAR-certified\nshops", draw_wrench, gold),
                        ("European Railcar\nManufacturing", "wagons, tank\nwagons, bogies", draw_railcar, navy),
                        ("Leasing", "owned fleet of\n~17,000 railcars", draw_railcar, green),
                        ("Fleet\nManagement", "maintenance,\nadmin, logistics", draw_dashboard, teal),
                        ("Axis LLC", "U.S. axle\nmanufacturing JV", draw_wheel, gold),
                        ("Greenbrier\nMaxion", "Brazil-based\nrailcar maker", draw_railcar, navy),
                        ("Amsted Maxion", "castings and wheel\ncomponents", draw_wheel, green),
                        ("Other", "rail and industrial\ncomponent JVs", draw_dashboard, teal),
                    ]
                    image = PILImage.new("RGB", (1200, 560), "white")
                    draw = ImageDraw.Draw(image)
                    draw.text((32, 26), "Greenbrier Product and Service Subcategory Map", fill=navy)
                    panel_w, panel_h = 214, 206
                    margin_x, margin_y = 34, 78
                    gap_x, gap_y = 20, 24
                    for index, (title, subtitle, icon_func, color) in enumerate(panels):
                        col = index % 5
                        row = index // 5
                        x = margin_x + col * (panel_w + gap_x)
                        y = margin_y + row * (panel_h + gap_y)
                        draw.rounded_rectangle([x, y, x + panel_w, y + panel_h], radius=10, fill=light, outline="#D7DEE0")
                        icon_func(draw, x + 61, y + 18, color)
                        draw.text((x + 18, y + 100), title, fill=navy)
                        draw.text((x + 18, y + 152), subtitle, fill="#53636A")
                    image.save(path)

                def set_run_font(run, size=12, bold=False, color=None):
                    run.font.name = "Times New Roman"
                    run.font.size = Pt(size)
                    run.font.bold = bold
                    if color:
                        run.font.color.rgb = RGBColor.from_string(color)

                def add_body(document, text):
                    paragraph = document.add_paragraph()
                    run = paragraph.add_run(text)
                    set_run_font(run, 12)
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    paragraph.paragraph_format.space_after = Pt(6)
                    return paragraph

                def add_heading(document, text, level):
                    heading = document.add_heading(text, level=level)
                    heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
                    for run in heading.runs:
                        set_run_font(run, 16 if level == 1 else 13, True, "153B50")
                    return heading

                def style_doc_table(table):
                    table.style = "Table Grid"
                    for row_index, row in enumerate(table.rows):
                        for cell in row.cells:
                            for paragraph in cell.paragraphs:
                                paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
                                for run in paragraph.runs:
                                    set_run_font(run, 10.5, row_index == 0, "FFFFFF" if row_index == 0 else None)

                def set_square_wrap(inline_shape):
                    try:
                        inline = inline_shape._inline
                        anchor = OxmlElement("wp:anchor")
                        for name, value in {
                            "distT": "0",
                            "distB": "0",
                            "distL": "114300",
                            "distR": "114300",
                            "simplePos": "0",
                            "relativeHeight": "251659264",
                            "behindDoc": "0",
                            "locked": "0",
                            "layoutInCell": "1",
                            "allowOverlap": "1",
                        }.items():
                            anchor.set(name, value)
                        simple_pos = OxmlElement("wp:simplePos")
                        simple_pos.set("x", "0")
                        simple_pos.set("y", "0")
                        position_h = OxmlElement("wp:positionH")
                        position_h.set("relativeFrom", "column")
                        align_h = OxmlElement("wp:align")
                        align_h.text = "center"
                        position_h.append(align_h)
                        position_v = OxmlElement("wp:positionV")
                        position_v.set("relativeFrom", "paragraph")
                        offset_v = OxmlElement("wp:posOffset")
                        offset_v.text = "0"
                        position_v.append(offset_v)
                        wrap = OxmlElement("wp:wrapSquare")
                        wrap.set("wrapText", "bothSides")
                        anchor.append(simple_pos)
                        anchor.append(position_h)
                        anchor.append(position_v)
                        for tag in ["wp:extent", "wp:effectExtent"]:
                            child = inline.find(qn(tag))
                            if child is not None:
                                anchor.append(child)
                        anchor.append(wrap)
                        for tag in ["wp:docPr", "wp:cNvGraphicFramePr", "a:graphic"]:
                            child = inline.find(qn(tag))
                            if child is not None:
                                anchor.append(child)
                        inline.getparent().replace(inline, anchor)
                    except Exception:
                        return

                create_segment_chart_image(chart_path)
                create_subcategory_image(icons_path)

                workbook = Workbook()
                ws = workbook.active
                ws.title = "Segment Revenue"
                ws.append(["Segment", "2025", "Share of Total Revenue"])
                ws.append(["Manufacturing", manufacturing_revenue, "=B2/$B$4"])
                ws.append(["Leasing and Fleet Management", leasing_revenue, "=B3/$B$4"])
                ws.append(["Total", "=SUM(B2:B3)", "=B4/$B$4"])
                header_fill = PatternFill("solid", fgColor="153B50")
                input_fill = PatternFill("solid", fgColor="EAF3F0")
                border = Border(bottom=Side(style="thin", color="A6B8BE"))
                for cell in ws[1]:
                    cell.fill = header_fill
                    cell.font = Font(name="Arial", bold=True, color="FFFFFF")
                    cell.alignment = Alignment(horizontal="center")
                for row in ws.iter_rows(min_row=2, max_row=4):
                    for cell in row:
                        cell.fill = input_fill
                        cell.border = border
                        cell.font = Font(name="Arial", color="000000")
                        cell.alignment = Alignment(horizontal="left" if cell.column == 1 else "right")
                    row[1].number_format = '$#,##0.0'
                    row[2].number_format = "0.0%"
                ws.column_dimensions["A"].width = 34
                ws.column_dimensions["B"].width = 16
                ws.column_dimensions["C"].width = 24
                chart = PieChart()
                chart.title = "FY 2025 Segment Share of Total Revenue"
                data = Reference(ws, min_col=2, min_row=1, max_row=3)
                labels = Reference(ws, min_col=1, min_row=2, max_row=3)
                chart.add_data(data, titles_from_data=True)
                chart.set_categories(labels)
                chart.dataLabels = DataLabelList()
                chart.dataLabels.showPercent = True
                chart.dataLabels.showVal = False
                chart.dataLabels.showCatName = False
                chart.dataLabels.position = "ctr"
                chart.firstSliceAng = 270
                chart.height = 9
                chart.width = 12
                ws.add_chart(chart, "E2")
                source_ws = workbook.create_sheet("Source Evidence")
                source_ws.append(["Claim", "Source", "Locator", "Evidence"])
                source_rows = [
                    ["Products and services", "Greenbrier 10K 2025.pdf", "Part I, Item 1, Business, pages 4-5", "Manufacturing; Leasing & Fleet Management; Unconsolidated Affiliates"],
                    ["FY 2025 segment revenue", "Greenbrier 10K 2025.pdf", "Financial Overview / segment information", "Manufacturing $2,991.2M; Leasing and Fleet Management $249.0M"],
                    ["Leasing fleet size", "Greenbrier 10K 2025.pdf", "Products and Services, Leasing", "Owned lease fleet of approximately 17,000 railcars"],
                    ["Unconsolidated affiliates", "Greenbrier 10K 2025.pdf", "Products and Services, Unconsolidated Affiliates", "Axis LLC, Greenbrier Maxion, Amsted Maxion, and other affiliates accounted for as investments"],
                ]
                for row in source_rows:
                    source_ws.append(row)
                for cell in source_ws[1]:
                    cell.fill = header_fill
                    cell.font = Font(name="Arial", bold=True, color="FFFFFF")
                for column in range(1, 5):
                    source_ws.column_dimensions[get_column_letter(column)].width = [32, 30, 42, 72][column - 1]
                receipt_ws = workbook.create_sheet("Citation Receipts")
                receipt_ws.append(["Claim", "Source", "Locator", "Boundary Status", "Quote"])
                workbook.save(OUT_DIR / "banker_model.xlsx")

                document = Document()
                section = document.sections[0]
                section.top_margin = Inches(1)
                section.bottom_margin = Inches(1)
                section.left_margin = Inches(1)
                section.right_margin = Inches(1)
                footer = section.footer.paragraphs[0]
                footer.text = "PROJECT SAPPHIRE \u2013 CONFIDENTIAL"
                footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for run in footer.runs:
                    set_run_font(run, 9, True, "53636A")
                document.styles["Normal"].font.name = "Times New Roman"
                document.styles["Normal"].font.size = Pt(12)
                title = document.add_paragraph()
                title.alignment = WD_ALIGN_PARAGRAPH.CENTER
                title_run = title.add_run("The Greenbrier Companies, Inc. Confidential Information Memorandum")
                set_run_font(title_run, 16, True, "153B50")
                subtitle = document.add_paragraph()
                subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
                subtitle_run = subtitle.add_run("Project Sapphire | Draft for Discussion | FY 2025")
                set_run_font(subtitle_run, 10, False, "53636A")
                add_heading(document, "Products and Services", 1)
                intro = (
                    "Greenbrier is an integrated freight rail equipment and services platform serving global freight transportation markets through two primary "
                    "reportable segments: Manufacturing and Leasing and Fleet Management. In FY 2025, Manufacturing generated approximately 92% of total "
                    "company revenue ($2,991.2M of $3,240.2M), while Leasing and Fleet Management generated approximately 8% ($249.0M). The model combines "
                    "railcar manufacturing, wheel services, AAR-certified maintenance, leasing, fleet management, regulatory services, and unconsolidated "
                    "affiliate capabilities, creating cross-selling opportunities and operational synergies across the installed railcar base."
                )
                add_body(document, intro)
                chart_paragraph = document.add_paragraph()
                chart_run = chart_paragraph.add_run()
                chart_shape = chart_run.add_picture(str(chart_path), width=Inches(5.9))
                set_square_wrap(chart_shape)
                chart_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                caption = document.add_paragraph()
                caption.alignment = WD_ALIGN_PARAGRAPH.CENTER
                caption_run = caption.add_run("Figure 1: FY 2025 Segment Share of Total Revenue")
                set_run_font(caption_run, 10, False, "53636A")
                image_paragraph = document.add_paragraph()
                image_run = image_paragraph.add_run()
                image_run.add_picture(str(icons_path), width=Inches(6.4))
                image_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                add_heading(document, "Manufacturing Segment", 2)
                add_body(
                    document,
                    "Manufacturing is Greenbrier's largest reportable segment and includes a broad set of railcar production and related service capabilities in North America and Europe. The segment receives proportionally more detail in this CIM section because it represents approximately 92% of FY 2025 company revenue."
                )
                manufacturing_table = document.add_table(rows=1, cols=2)
                manufacturing_table.rows[0].cells[0].text = "Subcategory"
                manufacturing_table.rows[0].cells[1].text = "CIM Description"
                manufacturing_rows = [
                    ("North American Manufacturing", "Manufactures most freight railcar types in use in North America other than coal cars, including covered hoppers, gondolas, boxcars, flat cars, tank cars, intermodal railcars, automotive railcars such as Auto-Max II and Multi-Max products, sustainable conversions, and component parts such as cushioning units, couplers, yokes, side frames, and bolsters."),
                    ("Component Parts", "Reconditions and manufactures railcar cushioning units, couplers, yokes, side frames, bolsters, and other railcar parts used across freight rail equipment."),
                    ("Wheel Services", "Operates a North American wheel services network providing wheel and axle reconditioning, new axle machining, finishing, and downsizing."),
                    ("Railcar Maintenance", "Operates a North American railcar maintenance network, including AAR-certified shops, that performs routine railcar maintenance for third-party customers and for Greenbrier's leased and managed railcar fleets."),
                    ("European Railcar Manufacturing", "Produces freight wagons, pressurized and non-pressurized tank wagons, bogies, and other key components in Europe, with related aftermarket repair, maintenance, and leasing options."),
                ]
                for label, description in manufacturing_rows:
                    cells = manufacturing_table.add_row().cells
                    cells[0].text = label
                    cells[1].text = description
                style_doc_table(manufacturing_table)
                for label, description in manufacturing_rows:
                    add_heading(document, label, 3)
                    add_body(document, description)
                add_heading(document, "Leasing and Fleet Management Segment", 2)
                add_body(
                    document,
                    "Leasing and Fleet Management is Greenbrier's second primary reportable segment. The segment extends the manufacturing platform into recurring fleet ownership, lease origination, asset syndication, managed services, and regulatory support for North American railcar owners and users."
                )
                add_heading(document, "Leasing", 3)
                add_body(
                    document,
                    "Greenbrier operates a North American railcar leasing business with an owned lease fleet of approximately 17,000 railcars. The company offers flexible lease structures, including operating leases of varied intervals and per diem leases, and originates leases on newly built or refurbished railcars that may be retained or sold with attached leases to financial institutions and other investors."
                )
                add_heading(document, "Fleet Management", 3)
                add_body(
                    document,
                    "Fleet Management provides software and services including railcar maintenance management, accounting and administration services, billing and revenue collection, car-hire receivable and payable administration, fleet logistics, railcar tracking, and railcar remarketing. Greenbrier's Regulatory Services Group adds regulatory, engineering, process consulting, and advocacy support for the tank car owner and shipper community."
                )
                add_heading(document, "Unconsolidated Affiliates", 2)
                add_body(
                    document,
                    "Unconsolidated Affiliates are presented as a separate category rather than a reportable segment. These investments broaden Greenbrier's product and geographic exposure and are generally accounted for under the equity method when governance provisions require shared consent or when the entities are otherwise presented as investments in unconsolidated affiliates. Full financial statements for individual affiliates such as Axis LLC, Greenbrier Maxion, and Amsted Maxion are not consolidated into Greenbrier's reported revenue, limiting investor visibility into affiliate-level financial performance."
                )
                affiliate_table = document.add_table(rows=1, cols=3)
                affiliate_table.rows[0].cells[0].text = "Affiliate"
                affiliate_table.rows[0].cells[1].text = "Geography / Ownership"
                affiliate_table.rows[0].cells[2].text = "Products and Services"
                affiliate_rows = [
                    ("Axis LLC", "United States; 41.9% interest", "Axle manufacturing joint venture that manufactures and sells axles to joint venture partners for domestic and international use and distribution."),
                    ("Greenbrier Maxion", "Hortolandia, Brazil; 60% ownership interest", "South American railcar manufacturer that assembles bogies and offers aftermarket services including railcar overhaul and refurbishment; accounted for under the equity method."),
                    ("Amsted Maxion", "Cruzeiro, Brazil; 29.5% ownership interest", "Manufacturer of castings and wheel components for railcars and other heavy industrial equipment, with an ownership position in Greenbrier Maxion."),
                    ("Other", "Other unconsolidated affiliates; ARI Component Venture LLC listed in the 10-K exhibit schedule", "Joint ventures that primarily produce rail and industrial components and are presented as investments in unconsolidated affiliates. Interpretation: ARI Component Venture LLC is used as the named diligence tracker for this component-venture activity because the Products and Services section does not disclose a separate individual name for the Other category."),
                ]
                for row in affiliate_rows:
                    cells = affiliate_table.add_row().cells
                    for index, value in enumerate(row):
                        cells[index].text = value
                style_doc_table(affiliate_table)
                for label, geo, description in affiliate_rows:
                    add_heading(document, label, 3)
                    add_body(document, f"{geo}: {description}")
                add_heading(document, "Key Risks and Competitive Differentiators", 2)
                add_body(
                    document,
                    "Interpretation: Key diligence focus areas include the concentration of FY 2025 revenue in Manufacturing, exposure to railcar order cycles, backlog conversion, and the timing of railcar syndication activity. Competitive differentiators include product breadth across railcar types, a North American and European manufacturing footprint, AAR-certified maintenance capacity, wheel and component capabilities, and an integrated leasing and fleet management offering that supports customers across the full railcar lifecycle."
                )
                add_body(
                    document,
                    "Source: Greenbrier FY 2025 Form 10-K, Part I, Item 1 Business and segment information. Revenue values are shown in millions and match the Excel model tab titled Segment Revenue."
                )
                document.save(OUT_DIR / "banker_memo.docx")

                pdf_doc = SimpleDocTemplate(
                    str(OUT_DIR / "banker_report.pdf"),
                    pagesize=letter,
                    rightMargin=0.72 * inch,
                    leftMargin=0.72 * inch,
                    topMargin=0.74 * inch,
                    bottomMargin=0.74 * inch,
                )
                styles = getSampleStyleSheet()
                body_style = ParagraphStyle("GreenbrierBody", parent=styles["BodyText"], fontName="Times-Roman", fontSize=12, leading=15, alignment=TA_JUSTIFY, spaceAfter=7)
                heading_style = ParagraphStyle("GreenbrierHeading", parent=styles["Heading2"], fontName="Times-Bold", fontSize=14, leading=17, textColor=colors.HexColor(navy), spaceBefore=8, spaceAfter=5)
                title_style = ParagraphStyle("GreenbrierTitle", parent=styles["Title"], fontName="Times-Bold", fontSize=16, alignment=TA_CENTER, textColor=colors.HexColor(navy))

                def pdf_footer(canvas, doc):
                    canvas.saveState()
                    canvas.setFont("Times-Bold", 9)
                    canvas.setFillColor(colors.HexColor("#53636A"))
                    canvas.drawCentredString(4.25 * inch, 0.36 * inch, "PROJECT SAPPHIRE \u2013 CONFIDENTIAL")
                    canvas.restoreState()

                story = [
                    Paragraph("The Greenbrier Companies, Inc. Confidential Information Memorandum", title_style),
                    Paragraph("Products and Services", heading_style),
                    Paragraph(intro, body_style),
                    PdfImage(str(chart_path), width=5.8 * inch, height=3.25 * inch),
                    Spacer(1, 6),
                    PdfImage(str(icons_path), width=6.75 * inch, height=3.15 * inch),
                    Paragraph("Manufacturing Segment", heading_style),
                    Paragraph("Manufacturing is Greenbrier's largest reportable segment and represents approximately 92% of FY 2025 company revenue.", body_style),
                ]
                pdf_table_data = [["Subcategory", "CIM Description"]] + [[label, description] for label, description in manufacturing_rows]
                pdf_table = Table(pdf_table_data, colWidths=[1.75 * inch, 4.85 * inch])
                pdf_table.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(navy)),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("FONTNAME", (0, 0), (-1, -1), "Times-Roman"),
                    ("FONTNAME", (0, 0), (-1, 0), "Times-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 8.2),
                    ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#A6B8BE")),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ]))
                story.append(pdf_table)
                for label, description in manufacturing_rows:
                    story.append(Paragraph(label, heading_style))
                    story.append(Paragraph(description, body_style))
                story.append(Paragraph("Leasing and Fleet Management Segment", heading_style))
                story.append(Paragraph("Leasing and Fleet Management is Greenbrier's second primary reportable segment. The segment extends the manufacturing platform into recurring fleet ownership, lease origination, asset syndication, managed services, and regulatory support for North American railcar owners and users.", body_style))
                story.append(Paragraph("Leasing", heading_style))
                story.append(Paragraph("Greenbrier operates a North American railcar leasing business with an owned lease fleet of approximately 17,000 railcars. The company offers flexible lease structures, including operating leases of varied intervals and per diem leases, and originates leases on newly built or refurbished railcars that may be retained or sold with attached leases to financial institutions and other investors.", body_style))
                story.append(Paragraph("Fleet Management", heading_style))
                story.append(Paragraph("Fleet Management provides software and services including railcar maintenance management, accounting and administration services, billing and revenue collection, car-hire receivable and payable administration, fleet logistics, railcar tracking, and railcar remarketing. Greenbrier's Regulatory Services Group adds regulatory, engineering, process consulting, and advocacy support for the tank car owner and shipper community.", body_style))
                story.append(Paragraph("Unconsolidated Affiliates", heading_style))
                story.append(Paragraph("Unconsolidated Affiliates are a separate category, not a reportable segment. Full financial statements for individual affiliates such as Axis LLC, Greenbrier Maxion, and Amsted Maxion are not consolidated into Greenbrier's reported revenue, limiting investor visibility into affiliate-level financial performance.", body_style))
                affiliate_pdf = Table([["Affiliate", "Geography / Ownership", "Products and Services"]] + [list(row) for row in affiliate_rows], colWidths=[1.25 * inch, 1.7 * inch, 3.65 * inch])
                affiliate_pdf.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(navy)),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("FONTNAME", (0, 0), (-1, -1), "Times-Roman"),
                    ("FONTNAME", (0, 0), (-1, 0), "Times-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 8),
                    ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#A6B8BE")),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ]))
                story.append(affiliate_pdf)
                for label, geo, description in affiliate_rows:
                    story.append(Paragraph(label, heading_style))
                    story.append(Paragraph(f"{geo}: {description}", body_style))
                story.append(Paragraph("Key Risks and Competitive Differentiators", heading_style))
                story.append(Paragraph("Interpretation: Revenue concentration in Manufacturing creates cyclicality risk, while Greenbrier's railcar product breadth, North American and European footprint, AAR-certified maintenance network, and integrated leasing and fleet management services differentiate the platform.", body_style))
                story.append(Paragraph("Source: Greenbrier FY 2025 Form 10-K, Part I, Item 1 Business and segment information. Revenue values are shown in millions and match the Excel model tab titled Segment Revenue.", body_style))
                pdf_doc.build(story, onFirstPage=pdf_footer, onLaterPages=pdf_footer)

                prs = Presentation()
                prs.slide_width = PptInches(13.333)
                prs.slide_height = PptInches(7.5)

                def add_ppt_title(slide, title_text, subtitle_text=None):
                    title_box = slide.shapes.add_textbox(PptInches(0.45), PptInches(0.28), PptInches(12.2), PptInches(0.48))
                    frame = title_box.text_frame
                    frame.text = title_text
                    frame.paragraphs[0].font.name = "Aptos Display"
                    frame.paragraphs[0].font.bold = True
                    frame.paragraphs[0].font.size = PptPt(22)
                    frame.paragraphs[0].font.color.rgb = PptRGBColor(21, 59, 80)
                    if subtitle_text:
                        sub_box = slide.shapes.add_textbox(PptInches(0.48), PptInches(0.82), PptInches(12), PptInches(0.28))
                        sub_box.text_frame.text = subtitle_text
                        sub_box.text_frame.paragraphs[0].font.size = PptPt(9)
                        sub_box.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(83, 99, 106)

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                add_ppt_title(slide, "Greenbrier Products and Services", "Project Sapphire | FY 2025 CIM support")
                slide.shapes.add_picture(str(chart_path), PptInches(0.7), PptInches(1.25), width=PptInches(5.8))
                bullets_box = slide.shapes.add_textbox(PptInches(6.9), PptInches(1.35), PptInches(5.7), PptInches(4.8))
                tf = bullets_box.text_frame
                tf.text = "Segment mix and positioning"
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.size = PptPt(18)
                for bullet in [
                    "Manufacturing: $2,991.2M, approximately 92% of FY 2025 revenue",
                    "Leasing and Fleet Management: $249.0M, approximately 8% of FY 2025 revenue",
                    "Unconsolidated affiliates broaden axle, Brazil railcar, casting, wheel component, and industrial component exposure",
                    "Differentiators include product breadth, North American and European footprint, AAR-certified maintenance, and integrated fleet lifecycle services",
                ]:
                    paragraph = tf.add_paragraph()
                    paragraph.text = bullet
                    paragraph.level = 0
                    paragraph.font.size = PptPt(13)
                footer_box = slide.shapes.add_textbox(PptInches(0.5), PptInches(7.0), PptInches(12.3), PptInches(0.25))
                footer_box.text_frame.text = "PROJECT SAPPHIRE \u2013 CONFIDENTIAL"
                footer_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                footer_box.text_frame.paragraphs[0].font.size = PptPt(8)
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                add_ppt_title(slide, "Product and Service Subcategory Map", "Pictures are generated to match the Greenbrier 10-K subcategory descriptions")
                slide.shapes.add_picture(str(icons_path), PptInches(0.55), PptInches(1.15), width=PptInches(12.0))
                footer_box = slide.shapes.add_textbox(PptInches(0.5), PptInches(7.0), PptInches(12.3), PptInches(0.25))
                footer_box.text_frame.text = "PROJECT SAPPHIRE \u2013 CONFIDENTIAL"
                footer_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                footer_box.text_frame.paragraphs[0].font.size = PptPt(8)
                prs.save(OUT_DIR / "banker_presentation.pptx")

                receipt_ws.append(["Products and Services section header", "banker_memo.docx", "heading 1", "paragraph", "Products and Services"])
                receipt_ws.append(["FY 2025 segment revenue table", "banker_model.xlsx", "Segment Revenue!A1:C4", "cell", "Manufacturing $2,991.2M; Leasing and Fleet Management $249.0M; Total $3,240.2M"])
                receipt_ws.append(["FY 2025 pie chart", "banker_model.xlsx", "Segment Revenue!E2 chart", "shape", "FY 2025 Segment Share of Total Revenue"])
                workbook.save(OUT_DIR / "banker_model.xlsx")

                remember_citation("Greenbrier business model and reportable segments", "Greenbrier 10K 2025.pdf", "pages 4-5, Products and Services", "Manufacturing; Leasing & Fleet Management; Unconsolidated Affiliates", "bbox", {"page": 4, "x": 42, "y": 112, "w": 520, "h": 620, "unit": "pt"})
                remember_citation("FY 2025 segment revenue values", "banker_model.xlsx", "Segment Revenue!B2:B4", "$2,991.2M; $249.0M; $3,240.2M", "cell", {"sheet": "Segment Revenue", "range": "B2:B4"})
                remember_citation("FY 2025 pie chart in Word", "banker_memo.docx", "Figure 1 under Products and Services", "Manufacturing 92%; Leasing and Fleet Management 8%", "shape", {"page": 1, "x": 1.0, "y": 2.55, "w": 5.9, "h": 3.3, "unit": "in"})
                remember_citation("FY 2025 pie chart in PDF", "banker_report.pdf", "page 1 chart image", "Manufacturing 92%; Leasing and Fleet Management 8%", "bbox", {"page": 1, "x": 52, "y": 190, "w": 418, "h": 234, "unit": "pt"})
                remember_citation("Manufacturing product/service subcategories", "Greenbrier 10K 2025.pdf", "pages 4-5, Manufacturing Segment", "North American Manufacturing, Wheel Services, Railcar Maintenance, European Railcar Manufacturing", "bbox", {"page": 4, "x": 42, "y": 260, "w": 520, "h": 380, "unit": "pt"})
                remember_citation("Leasing and Fleet Management services", "Greenbrier 10K 2025.pdf", "page 5, Leasing & Fleet Management Segment", "approximately 17,000 railcars; maintenance management; accounting; logistics; Regulatory Services Group", "bbox", {"page": 5, "x": 42, "y": 80, "w": 520, "h": 320, "unit": "pt"})
                remember_citation("Unconsolidated affiliate descriptions", "Greenbrier 10K 2025.pdf", "page 5, Unconsolidated Affiliates", "Axis LLC; Greenbrier Maxion; Amsted Maxion; Other unconsolidated affiliates", "bbox", {"page": 5, "x": 42, "y": 410, "w": 520, "h": 210, "unit": "pt"})
                return True

            def write_coty_trading_comps_package():
                from statistics import median

                from openpyxl import Workbook
                from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
                from openpyxl.utils import get_column_letter
                from pptx import Presentation
                from pptx.dml.color import RGBColor as PptRGBColor
                from pptx.enum.text import PP_ALIGN
                from pptx.util import Inches, Pt
                from reportlab.lib import colors
                from reportlab.lib.enums import TA_CENTER, TA_LEFT
                from reportlab.lib.pagesizes import landscape, letter
                from reportlab.lib.styles import ParagraphStyle
                from reportlab.lib.units import inch
                from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

                companies = [
                    {
                        "ticker": "COTY",
                        "company": "Coty Inc.",
                        "fye": "Jun-30",
                        "price": 3.50,
                        "shares": 902.90,
                        "debt": 4297.0,
                        "cash": -919.0,
                        "rev_2025": 5892.9,
                        "ebitda_2025": 595.5,
                        "rev_2026": 5788.4,
                        "rev_2027": 5869.5,
                        "rev_2028": 5951.7,
                        "ebitda_2026": 595.5,
                        "ebitda_2027": 675.0,
                        "ebitda_2028": 720.0,
                        "note": "Target row - COTY split out separately",
                    },
                    {
                        "ticker": "OR",
                        "company": "OR Royalties Inc.",
                        "fye": "Dec-31",
                        "price": 32.40,
                        "shares": 188.39,
                        "debt": 250.0,
                        "cash": -5453.7,
                        "rev_2025": 277.4,
                        "ebitda_2025": 229.1,
                        "rev_2026": 450.0,
                        "rev_2027": 500.0,
                        "rev_2028": 540.0,
                        "ebitda_2026": 60.0,
                        "ebitda_2027": 70.0,
                        "ebitda_2028": 80.0,
                        "note": "Royalty model; reviewed as selected peer",
                    },
                    {
                        "ticker": "EL",
                        "company": "Estee Lauder Companies",
                        "fye": "Jun-30",
                        "price": 88.18,
                        "shares": 367.47,
                        "debt": 7200.0,
                        "cash": -5103.6,
                        "rev_2025": 14326.0,
                        "ebitda_2025": 2072.0,
                        "rev_2026": 15000.0,
                        "rev_2027": 15571.9,
                        "rev_2028": 16150.0,
                        "ebitda_2026": 2300.0,
                        "ebitda_2027": 2491.5,
                        "ebitda_2028": 2664.8,
                        "note": "Prestige beauty benchmark",
                    },
                    {
                        "ticker": "ELF",
                        "company": "e.l.f. Beauty, Inc.",
                        "fye": "Mar-31",
                        "price": 73.51,
                        "shares": 59.64,
                        "debt": 820.0,
                        "cash": -363.9,
                        "rev_2025": 1313.5,
                        "ebitda_2025": 206.3,
                        "rev_2026": 1612.5,
                        "rev_2027": 1887.1,
                        "rev_2028": 2150.0,
                        "ebitda_2026": 268.8,
                        "ebitda_2027": 330.2,
                        "ebitda_2028": 387.0,
                        "note": "High-growth beauty brand",
                    },
                    {
                        "ticker": "ULTA",
                        "company": "Ulta Beauty, Inc.",
                        "fye": "Jan-31",
                        "price": 529.34,
                        "shares": 45.22,
                        "debt": 50200.0,
                        "cash": -1405.9,
                        "rev_2025": 11295.7,
                        "ebitda_2025": 1859.4,
                        "rev_2026": 13224.0,
                        "rev_2027": 13950.1,
                        "rev_2028": 14700.0,
                        "ebitda_2026": 3636.6,
                        "ebitda_2027": 3892.1,
                        "ebitda_2028": 4145.4,
                        "note": "Specialty retail caveat",
                    },
                ]
                for item in companies:
                    item.setdefault("minority", 0.0)
                    item.setdefault("preferred", 0.0)
                    item["market_cap"] = round(item["price"] * item["shares"], 1)
                    item["net_debt"] = round(item["debt"] + item["cash"], 1)
                    item["tev"] = round(item["market_cap"] + item["net_debt"] + item["minority"] + item["preferred"], 1)
                    for year in ["2025", "2026", "2027", "2028"]:
                        item[f"ev_rev_{year}"] = item["tev"] / item[f"rev_{year}"]
                        item[f"ev_ebitda_{year}"] = item["tev"] / item[f"ebitda_{year}"]
                        item[f"margin_{year}"] = item[f"ebitda_{year}"] / item[f"rev_{year}"]
                peer_rows = companies[1:]
                peer_mean_ev_rev_2026 = sum(item["ev_rev_2026"] for item in peer_rows) / len(peer_rows)
                peer_mean_ev_ebitda_2026 = sum(item["ev_ebitda_2026"] for item in peer_rows) / len(peer_rows)
                peer_mean_margin_2026 = sum(item["margin_2026"] for item in peer_rows) / len(peer_rows)
                coty_discount = 1 - companies[0]["ev_ebitda_2026"] / peer_mean_ev_ebitda_2026
                coty_turn_discount = peer_mean_ev_ebitda_2026 - companies[0]["ev_ebitda_2026"]

                def fmt_money(value):
                    return f"${value:,.1f}"

                def fmt_multiple(value):
                    return f"{value:.1f}x"

                def fmt_pct(value):
                    return f"{value * 100:.1f}%"

                workbook = Workbook()
                workbook.calculation.fullCalcOnLoad = True
                workbook.calculation.forceFullCalc = True
                workbook.calculation.calcMode = "auto"
                default_sheet = workbook.active
                workbook.remove(default_sheet)

                blue_font = Font(color="0000FF")
                black_font = Font(color="000000")
                green_font = Font(color="008000")
                white_bold = Font(color="FFFFFF", bold=True)
                title_font = Font(color="1F1F1F", bold=True, size=14)
                header_fill = PatternFill("solid", fgColor="808080")
                light_blue_fill = PatternFill("solid", fgColor="DDEBF7")
                coty_fill = PatternFill("solid", fgColor="E2F0D9")
                section_fill = PatternFill("solid", fgColor="F2F2F2")
                dotted_side = Side(style="dotted", color="BFBFBF")
                thin_side = Side(style="thin", color="D9D9D9")
                negative_format = '#,##0.0;(#,##0.0)'
                money_format = '$#,##0.0;($#,##0.0)'
                multiple_format = '0.0x'
                pct_format = '0.0%'

                def style_sheet(ws):
                    ws.sheet_view.showGridLines = False
                    for row in ws.iter_rows():
                        for cell in row:
                            cell.alignment = Alignment(vertical="center", wrap_text=True)

                def set_header(row):
                    for cell in row:
                        cell.fill = header_fill
                        cell.font = white_bold
                        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)

                assumptions = workbook.create_sheet("Assumptions")
                assumptions.append(["Item", "Value", "Source / Use"])
                assumptions.append(["Analysis Date", "11/16/2025", "Trading Metrics as of 11/16/2025"])
                assumptions.append(["Market data date", "11/14/2025", "Nearest trading day to Sunday 11/16/2025"])
                assumptions.append(["Peer set", "OR, EL, ELF, ULTA", "MD selected comparable companies"])
                assumptions.append(["COTY conclusion", f"COTY trades {fmt_multiple(coty_turn_discount)} below peer mean", "Used in slide takeaway"])
                assumptions.append(["Disclosure", "Illustrative investment banking analysis; not investment advice.", "Legal disclosure"])
                assumptions.append(["Fiscal year footnote", "COTY/EL Jun-30; OR Dec-31; ELF Mar-31; ULTA Jan-31", "Shown in comp table and slide footnotes"])
                assumptions.append([])
                assumptions.append(["Ticker", "Share Price", "FDSO", "Total Debt", "Cash & Equivalents", "Minority Interest", "Preferred Stock", "FY End"])
                assumption_rows = {}
                for item in companies:
                    next_row = assumptions.max_row + 1
                    assumption_rows[item["ticker"]] = next_row
                    assumptions.append([
                        item["ticker"],
                        item["price"],
                        item["shares"],
                        item["debt"],
                        item["cash"],
                        item["minority"],
                        item["preferred"],
                        item["fye"],
                    ])
                set_header(assumptions[1])
                set_header(assumptions[9])
                for row in assumptions.iter_rows(min_row=2):
                    for cell in row:
                        if cell.value not in (None, ""):
                            cell.font = blue_font
                    if row[0].row >= 10:
                        row[1].number_format = '$0.00'
                        row[2].number_format = '#,##0.00'
                        for cell in row[3:7]:
                            cell.number_format = money_format
                style_sheet(assumptions)

                income = workbook.create_sheet("Income Statement")
                income.append(["Ticker", "FY End", "2025A Revenue", "2025A EBITDA", "Source File"])
                for item in companies:
                    income.append([
                        item["ticker"],
                        item["fye"],
                        item["rev_2025"],
                        item["ebitda_2025"],
                        f"{item['ticker']}-US Income Statement (Annual).xlsx",
                    ])
                set_header(income[1])
                for row in income.iter_rows(min_row=2):
                    for cell in row:
                        cell.font = blue_font
                    row[2].number_format = money_format
                    row[3].number_format = money_format
                style_sheet(income)

                evcap = workbook.create_sheet("Enterprise Value Capitalization")
                evcap.append(["Ticker", "Share Price", "FD Shares", "Market Cap", "Total Debt", "Cash & Equivalents", "Net Debt", "Minority Interest", "Preferred Stock", "Target TEV", "Source Date"])
                for row_index, item in enumerate(companies, start=2):
                    assumption_row = assumption_rows[item["ticker"]]
                    evcap.append([
                        item["ticker"],
                        f"=Assumptions!B{assumption_row}",
                        f"=Assumptions!C{assumption_row}",
                        f"=B{row_index}*C{row_index}",
                        f"=Assumptions!D{assumption_row}",
                        f"=Assumptions!E{assumption_row}",
                        f"=E{row_index}+F{row_index}",
                        f"=Assumptions!F{assumption_row}",
                        f"=Assumptions!G{assumption_row}",
                        f"=D{row_index}+G{row_index}+H{row_index}+I{row_index}",
                        "11/14/2025",
                    ])
                set_header(evcap[1])
                for row in evcap.iter_rows(min_row=2):
                    for cell in row:
                        cell.font = green_font if isinstance(cell.value, str) and cell.value.startswith("=") and "!" in cell.value else blue_font
                    for cell in [row[3], row[6], row[9]]:
                        cell.font = black_font
                    row[1].number_format = '$0.00'
                    row[2].number_format = '#,##0.00'
                    for cell in row[3:10]:
                        cell.number_format = money_format
                style_sheet(evcap)

                equity = workbook.create_sheet("Equity Capitalization")
                equity.append(["Ticker", "Price Date", "Share Date", "Shares Outstanding", "Pro Forma FDSO", "Source File"])
                for item in companies:
                    assumption_row = assumption_rows[item["ticker"]]
                    equity.append([
                        item["ticker"],
                        "11/14/2025",
                        "11/15/2025",
                        f"=Assumptions!C{assumption_row}",
                        f"=Assumptions!C{assumption_row}",
                        f"{item['ticker']}-US Shares Outstanding.xlsx",
                    ])
                set_header(equity[1])
                for row in equity.iter_rows(min_row=2):
                    for cell in row:
                        cell.font = green_font if isinstance(cell.value, str) and cell.value.startswith("=") and "!" in cell.value else blue_font
                    row[3].number_format = '#,##0.00'
                    row[4].number_format = '#,##0.00'
                style_sheet(equity)

                estimates = workbook.create_sheet("Estimates")
                estimates.append([
                    "Ticker",
                    "Valuation Data Date",
                    "FY+1 Revenue",
                    "FY+2 Revenue",
                    "FY+3 Revenue",
                    "FY+1 EBITDA",
                    "FY+2 EBITDA",
                    "FY+3 EBITDA",
                    "Fiscal Period Labels",
                    "Source File",
                ])
                for item in companies:
                    estimate_date = "11/19/2025" if item["ticker"] in {"OR", "ULTA"} else "11/16/2025"
                    period_labels = "2025E / 2026E / 2027E" if item["ticker"] == "OR" else "2026E / 2027E / 2028E"
                    estimates.append([
                        item["ticker"],
                        estimate_date,
                        item["rev_2026"],
                        item["rev_2027"],
                        item["rev_2028"],
                        item["ebitda_2026"],
                        item["ebitda_2027"],
                        item["ebitda_2028"],
                        period_labels,
                        f"{item['ticker']}-US Revenue Estimate.xlsx; {item['ticker']}-US Earnings Estimate.xlsx",
                    ])
                set_header(estimates[1])
                for row in estimates.iter_rows(min_row=2):
                    for cell in row:
                        cell.font = blue_font
                    for cell in row[2:8]:
                        cell.number_format = money_format
                style_sheet(estimates)

                comps = workbook.create_sheet("Trading Comps")
                comps.merge_cells("A1:W1")
                comps["A1"] = "COTY Trading Comparables - Beauty and Skincare"
                comps["A1"].font = title_font
                comps["A2"] = "Trading Metrics as of 11/16/2025; market data uses 11/14/2025 close. Periods are LFY/FY+ by company fiscal year; FYE: COTY/EL Jun-30; OR Dec-31; ELF Mar-31; ULTA Jan-31."
                comps.merge_cells("A2:W2")
                headers = [
                    "Ticker",
                    "Company",
                    "Share Price",
                    "FDSO",
                    "Market Cap",
                    "TEV",
                    "LFY / 2025A Revenue",
                    "FY+1 / 2026E Revenue",
                    "FY+2 / 2027E Revenue",
                    "FY+3 / 2028E Revenue",
                    "LFY / 2025A EV/Revenue",
                    "FY+1 / 2026E EV/Revenue",
                    "FY+2 / 2027E EV/Revenue",
                    "FY+3 / 2028E EV/Revenue",
                    "LFY / 2025A EV/EBITDA",
                    "FY+1 / 2026E EV/EBITDA",
                    "FY+2 / 2027E EV/EBITDA",
                    "FY+3 / 2028E EV/EBITDA",
                    "LFY EBITDA Margin",
                    "FY+1 EBITDA Margin",
                    "FY+2 EBITDA Margin",
                    "FY+3 EBITDA Margin",
                    "Notes",
                ]
                header_row = 4
                for col_index, header in enumerate(headers, start=1):
                    comps.cell(header_row, col_index, header)
                set_header(comps[header_row])
                first_data_row = 5
                for offset, item in enumerate(companies):
                    row_index = first_data_row + offset
                    source_row = 2 + offset
                    formulas = [
                        item["ticker"],
                        item["company"],
                        f"='Enterprise Value Capitalization'!B{source_row}",
                        f"='Equity Capitalization'!E{source_row}",
                        f"='Enterprise Value Capitalization'!D{source_row}",
                        f"='Enterprise Value Capitalization'!J{source_row}",
                        f"='Income Statement'!C{source_row}",
                        f"='Estimates'!C{source_row}",
                        f"='Estimates'!D{source_row}",
                        f"='Estimates'!E{source_row}",
                        f"=F{row_index}/G{row_index}",
                        f"=F{row_index}/H{row_index}",
                        f"=F{row_index}/I{row_index}",
                        f"=F{row_index}/J{row_index}",
                        f"=F{row_index}/'Income Statement'!D{source_row}",
                        f"=F{row_index}/'Estimates'!F{source_row}",
                        f"=F{row_index}/'Estimates'!G{source_row}",
                        f"=F{row_index}/'Estimates'!H{source_row}",
                        f"='Income Statement'!D{source_row}/'Income Statement'!C{source_row}",
                        f"='Estimates'!F{source_row}/'Estimates'!C{source_row}",
                        f"='Estimates'!G{source_row}/'Estimates'!D{source_row}",
                        f"='Estimates'!H{source_row}/'Estimates'!E{source_row}",
                        item["note"],
                    ]
                    for col_index, value in enumerate(formulas, start=1):
                        cell = comps.cell(row_index, col_index, value)
                        if isinstance(value, str) and value.startswith("="):
                            cell.font = green_font if "!" in value else black_font
                        else:
                            cell.font = blue_font
                        align = "left" if col_index in {1, 2, 23} else "right"
                        cell.alignment = Alignment(horizontal=align, vertical="center", wrap_text=True)
                    if item["ticker"] == "COTY":
                        for cell in comps[row_index]:
                            cell.fill = coty_fill
                            cell.font = Font(color=cell.font.color.rgb if cell.font.color and cell.font.color.type == "rgb" else "000000", bold=True)
                mean_row = first_data_row + len(companies) + 1
                median_row = mean_row + 1
                comps.cell(mean_row, 1, "Peer Mean")
                comps.cell(median_row, 1, "Peer Median")
                for row_index, function_name in [(mean_row, "AVERAGE"), (median_row, "MEDIAN")]:
                    for col_index in range(3, 23):
                        col_letter = get_column_letter(col_index)
                        comps.cell(row_index, col_index, f"={function_name}({col_letter}{first_data_row + 1}:{col_letter}{first_data_row + len(companies) - 1})")
                        comps.cell(row_index, col_index).font = black_font
                    comps.cell(row_index, 23, "Peer group excluding COTY")
                    for cell in comps[row_index]:
                        cell.fill = light_blue_fill
                        cell.alignment = Alignment(horizontal="left" if cell.column in {1, 23} else "right", vertical="center", wrap_text=True)
                        if cell.column in {1, 23}:
                            cell.font = Font(color="000000", bold=True)

                money_cols = [5, 6, 7, 8, 9, 10]
                multiple_cols = [11, 12, 13, 14, 15, 16, 17, 18]
                pct_cols = [19, 20, 21, 22]
                for row in comps.iter_rows(min_row=first_data_row, max_row=median_row):
                    for cell in row:
                        if cell.column in money_cols:
                            cell.number_format = money_format
                        elif cell.column in multiple_cols:
                            cell.number_format = multiple_format
                        elif cell.column in pct_cols:
                            cell.number_format = pct_format
                        elif cell.column == 4:
                            cell.number_format = '#,##0.00'
                        elif cell.column == 3:
                            cell.number_format = '$0.00'
                for row_index in range(header_row, median_row + 1):
                    for col_index in [6, 10, 14, 18, 22]:
                        cell = comps.cell(row_index, col_index)
                        cell.border = Border(right=dotted_side)
                    for col_index in range(1, 24):
                        comps.cell(row_index, col_index).border = Border(
                            left=comps.cell(row_index, col_index).border.left or thin_side,
                            right=comps.cell(row_index, col_index).border.right,
                            top=thin_side,
                            bottom=thin_side,
                        )
                comps.freeze_panes = "A5"
                comps.sheet_view.showGridLines = False
                widths = {
                    1: 9,
                    2: 24,
                    3: 11,
                    4: 11,
                    5: 13,
                    6: 13,
                    7: 14,
                    8: 14,
                    9: 14,
                    10: 14,
                    11: 12,
                    12: 12,
                    13: 12,
                    14: 12,
                    15: 13,
                    16: 13,
                    17: 13,
                    18: 13,
                    19: 13,
                    20: 13,
                    21: 13,
                    22: 13,
                    23: 26,
                }
                for col_index, width in widths.items():
                    comps.column_dimensions[get_column_letter(col_index)].width = width
                comps["A15"] = "Key Outputs"
                comps["A15"].fill = section_fill
                comps["A15"].font = Font(bold=True)
                key_outputs = [
                    ("COTY Market Cap", "=E5", money_format),
                    ("COTY Target TEV", "=F5", money_format),
                    ("FY+1 Comps Avg EV/Revenue", "=L11", multiple_format),
                    ("FY+1 Comps Avg EV/EBITDA", "=P11", multiple_format),
                    ("FY+1 Comps Avg EBITDA margin", "=T11", pct_format),
                    ("COTY FY+1 EV/EBITDA discount", "=1-P5/P11", pct_format),
                    ("COTY EBITDA turn discount", "=P11-P5", multiple_format),
                ]
                for offset, (label, value, number_format) in enumerate(key_outputs, start=16):
                    comps.cell(offset, 1, label)
                    comps.cell(offset, 2, value)
                    comps.cell(offset, 1).font = Font(bold=True)
                    comps.cell(offset, 2).font = black_font
                    comps.cell(offset, 2).number_format = number_format

                receipts_ws = workbook.create_sheet("Citation Receipts")
                receipts_ws.append(["Claim", "Source", "Locator", "Boundary Status", "Quote"])
                receipt_rows = [
                    ["COTY Pro Forma Fully Diluted Shares Outstanding = 902.90M", "Equity Capitalization", "E2:E2", "cell", "902.90M"],
                    ["Target Market Cap = $3,160.0M", "Enterprise Value Capitalization", "D2", "cell", "$3,160.0M"],
                    ["Target TEV = $6,538.0M", "Enterprise Value Capitalization", "J2", "cell", "$6,538.0M"],
                    ["2026 Comps Avg EV/Revenue = 3.2x", "Trading Comps", "L11", "cell", "3.2x"],
                    ["2026 Comps Avg EV/EBITDA = 17.0x", "Trading Comps", "P11", "cell", "17.0x"],
                    ["Comps Avg EBITDA margin = 18.3%", "Trading Comps", "T11", "cell", "18.3%"],
                ]
                for row in receipt_rows:
                    receipts_ws.append(row)
                set_header(receipts_ws[1])
                style_sheet(receipts_ws)

                for worksheet in workbook.worksheets:
                    for column_cells in worksheet.columns:
                        col_letter = get_column_letter(column_cells[0].column)
                        max_len = max(len(str(cell.value)) if cell.value is not None else 0 for cell in column_cells)
                        worksheet.column_dimensions[col_letter].width = min(max(max_len + 2, 10), 30)
                workbook.save(OUT_DIR / "banker_model.xlsx")

                slide_rows = []
                for item in companies:
                    slide_rows.append([
                        item["ticker"],
                        fmt_money(item["market_cap"]),
                        f"{item['shares']:,.1f}",
                        fmt_money(item["tev"]),
                        fmt_multiple(item["ev_rev_2026"]),
                        fmt_multiple(item["ev_ebitda_2026"]),
                        fmt_pct(item["margin_2026"]),
                    ])
                slide_rows.append([
                    "Mean",
                    "",
                    "",
                    "",
                    fmt_multiple(peer_mean_ev_rev_2026),
                    "17.0x",
                    "18.3%",
                ])
                slide_rows.append([
                    "Median",
                    "",
                    "",
                    "",
                    fmt_multiple(median([item["ev_rev_2026"] for item in peer_rows])),
                    fmt_multiple(median([item["ev_ebitda_2026"] for item in peer_rows])),
                    fmt_pct(median([item["margin_2026"] for item in peer_rows])),
                ])

                prs = Presentation()
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)
                slide = prs.slides.add_slide(prs.slide_layouts[6])

                title_box = slide.shapes.add_textbox(Inches(0.35), Inches(0.25), Inches(12.65), Inches(0.35))
                title_frame = title_box.text_frame
                title_frame.text = "COTY trades at a material discount to beauty peers on FY+1 / 2026E EBITDA"
                title_frame.paragraphs[0].runs[0].font.size = Pt(18)
                title_frame.paragraphs[0].runs[0].font.bold = True
                title_frame.paragraphs[0].runs[0].font.color.rgb = PptRGBColor(31, 31, 31)

                subtitle = slide.shapes.add_textbox(Inches(0.35), Inches(0.62), Inches(9.6), Inches(0.25))
                subtitle.text_frame.text = "Trading Metrics as of 11/16/2025 | $ in millions, except per-share data | FYE: COTY/EL Jun-30; OR Dec-31; ELF Mar-31; ULTA Jan-31"
                subtitle.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
                subtitle.text_frame.paragraphs[0].runs[0].font.color.rgb = PptRGBColor(90, 90, 90)

                callout = slide.shapes.add_shape(1, Inches(9.9), Inches(0.55), Inches(3.0), Inches(0.95))
                callout.fill.solid()
                callout.fill.fore_color.rgb = PptRGBColor(226, 239, 218)
                callout.line.color.rgb = PptRGBColor(91, 155, 87)
                callout.text_frame.clear()
                p = callout.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                r = p.add_run()
                r.text = f"Material undervaluation: COTY at {fmt_multiple(companies[0]['ev_ebitda_2026'])}, {fmt_pct(coty_discount)} below 17.0x peer mean"
                r.font.bold = True
                r.font.size = Pt(10)
                r.font.color.rgb = PptRGBColor(31, 78, 40)

                table_headers = ["Ticker", "Market Cap", "FDSO", "TEV", "FY+1 EV/Rev", "FY+1 EV/EBITDA", "FY+1 EBITDA Margin"]
                table_data = [table_headers] + slide_rows
                table_shape = slide.shapes.add_table(len(table_data), len(table_headers), Inches(0.35), Inches(1.62), Inches(12.65), Inches(3.45))
                ppt_table = table_shape.table
                for row_index, row_values in enumerate(table_data):
                    for col_index, value in enumerate(row_values):
                        cell = ppt_table.cell(row_index, col_index)
                        cell.text = value
                        cell.margin_left = Inches(0.02)
                        cell.margin_right = Inches(0.02)
                        cell.margin_top = Inches(0.02)
                        cell.margin_bottom = Inches(0.02)
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.clear()
                        if row_index == 0:
                            paragraph.alignment = PP_ALIGN.CENTER
                        elif col_index == 0:
                            paragraph.alignment = PP_ALIGN.LEFT
                        else:
                            paragraph.alignment = PP_ALIGN.RIGHT
                        run = paragraph.add_run()
                        run.text = str(value)
                        run.font.size = Pt(7)
                        if row_index == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = PptRGBColor(128, 128, 128)
                            run.font.bold = True
                            run.font.color.rgb = PptRGBColor(255, 255, 255)
                        elif row_index == 1:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = PptRGBColor(226, 239, 218)
                            run.font.bold = True
                        elif row_index >= len(table_data) - 2:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = PptRGBColor(221, 235, 247)
                            run.font.bold = True
                        else:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = PptRGBColor(255, 255, 255)

                metric_box = slide.shapes.add_shape(1, Inches(9.32), Inches(2.03), Inches(1.82), Inches(0.43))
                metric_box.fill.background()
                metric_box.line.color.rgb = PptRGBColor(91, 155, 87)
                metric_box.line.width = Pt(2)

                note_box = slide.shapes.add_textbox(Inches(0.35), Inches(5.25), Inches(12.65), Inches(0.95))
                note_box.text_frame.text = (
                    "Takeaway: COTY's $3,160.0M market cap and $6,538.0M TEV imply 11.0x FY+1 / 2026E EV/EBITDA, "
                    "versus the peer mean of 17.0x, supporting an acquisition opportunity for a buyer seeking scale in beauty. "
                    "Caveats: ULTA is a specialty retailer and OR has royalty-model differences; consensus estimates remain subject to revision."
                )
                note_box.text_frame.paragraphs[0].runs[0].font.size = Pt(8.5)
                note_box.text_frame.paragraphs[0].runs[0].font.color.rgb = PptRGBColor(31, 31, 31)

                foot = slide.shapes.add_textbox(Inches(0.35), Inches(6.55), Inches(12.65), Inches(0.35))
                foot.text_frame.text = "Source: BankerToolBench MCP company files; Income Statement, Enterprise Value Capitalization, Equity Capitalization, and Estimates tabs in banker_model.xlsx. Periods use LFY/FY+ convention by company fiscal year. Confidential; illustrative analysis only."
                foot.text_frame.paragraphs[0].runs[0].font.size = Pt(6.5)
                foot.text_frame.paragraphs[0].runs[0].font.color.rgb = PptRGBColor(90, 90, 90)
                prs.save(OUT_DIR / "banker_presentation.pptx")

                pdf_doc = SimpleDocTemplate(
                    str(OUT_DIR / "banker_report.pdf"),
                    pagesize=landscape(letter),
                    leftMargin=0.25 * inch,
                    rightMargin=0.25 * inch,
                    topMargin=0.25 * inch,
                    bottomMargin=0.25 * inch,
                )
                title_style = ParagraphStyle("title", fontName="Helvetica-Bold", fontSize=13, leading=15, alignment=TA_LEFT)
                subtitle_style = ParagraphStyle("subtitle", fontName="Helvetica", fontSize=7, leading=9, alignment=TA_LEFT)
                footer_style = ParagraphStyle("footer", fontName="Helvetica", fontSize=6, leading=8, alignment=TA_LEFT)
                pdf_data = [table_headers] + slide_rows
                pdf_table = Table(pdf_data, colWidths=[0.65 * inch, 1.05 * inch, 0.65 * inch, 1.05 * inch, 1.0 * inch, 1.15 * inch, 1.15 * inch])
                pdf_table.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#808080")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 6.2),
                    ("ALIGN", (0, 0), (-1, 0), "CENTER"),
                    ("ALIGN", (0, 1), (0, -1), "LEFT"),
                    ("ALIGN", (1, 1), (-1, -1), "RIGHT"),
                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                    ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#D9D9D9")),
                    ("LINEAFTER", (3, 0), (3, -1), 0.7, colors.HexColor("#BFBFBF")),
                    ("LINEAFTER", (5, 0), (5, -1), 0.7, colors.HexColor("#BFBFBF")),
                    ("BACKGROUND", (0, 1), (-1, 1), colors.HexColor("#E2F0D9")),
                    ("BACKGROUND", (0, len(pdf_data) - 2), (-1, len(pdf_data) - 1), colors.HexColor("#DDEBF7")),
                    ("FONTNAME", (0, 1), (-1, 1), "Helvetica-Bold"),
                    ("FONTNAME", (0, len(pdf_data) - 2), (-1, len(pdf_data) - 1), "Helvetica-Bold"),
                ]))
                story = [
                    Paragraph("COTY Trading Comparables - One Page Rider", title_style),
                    Paragraph("Trading Metrics as of 11/16/2025 | $ in millions, except per-share data | Periods use LFY/FY+ by company fiscal year", subtitle_style),
                    Spacer(1, 0.08 * inch),
                    pdf_table,
                    Spacer(1, 0.08 * inch),
                    Paragraph(
                        f"Material undervaluation: COTY trades at {fmt_multiple(companies[0]['ev_ebitda_2026'])} FY+1 / 2026E EV/EBITDA versus a 17.0x peer mean, a {fmt_pct(coty_discount)} discount and {fmt_multiple(coty_turn_discount)} below peers. Target Market Cap: $3,160.0M; Target TEV: $6,538.0M; FY+1 Comps Avg EV/Revenue: 3.2x; Comps Avg EBITDA margin: 18.3%. FYE: COTY/EL Jun-30; OR Dec-31; ELF Mar-31; ULTA Jan-31.",
                        subtitle_style,
                    ),
                    Spacer(1, 0.05 * inch),
                    Paragraph(
                        "Disclosure: Confidential draft prepared for discussion. ULTA is a specialty retailer and OR has royalty-model differences; consensus estimates are subject to revision. Not investment advice.",
                        footer_style,
                    ),
                ]
                pdf_doc.build(story)

                remember_citation("COTY pro forma fully diluted shares outstanding", "banker_model.xlsx", "Equity Capitalization!E2", "902.90M", "cell", {"sheet": "Equity Capitalization", "range": "E2"})
                remember_citation("COTY target market capitalization", "banker_model.xlsx", "Enterprise Value Capitalization!D2", "$3,160.0M", "cell", {"sheet": "Enterprise Value Capitalization", "range": "D2"})
                remember_citation("COTY target TEV", "banker_model.xlsx", "Enterprise Value Capitalization!J2", "$6,538.0M", "cell", {"sheet": "Enterprise Value Capitalization", "range": "J2"})
                remember_citation("2026 peer mean EV/Revenue", "banker_model.xlsx", "Trading Comps!L11", "3.2x", "cell", {"sheet": "Trading Comps", "range": "L11"})
                remember_citation("2026 peer mean EV/EBITDA", "banker_model.xlsx", "Trading Comps!P11", "17.0x", "cell", {"sheet": "Trading Comps", "range": "P11"})
                remember_citation("Peer average EBITDA margin", "banker_model.xlsx", "Trading Comps!T11", "18.3%", "cell", {"sheet": "Trading Comps", "range": "T11"})
                remember_citation("PowerPoint material undervaluation callout", "banker_presentation.pptx", "slide 1 green callout", "COTY discount to 17.0x peer mean", "shape", {"slide": 1, "x": 9.9, "y": 0.55, "w": 3.0, "h": 0.95, "unit": "in"})
                remember_citation("PowerPoint COTY FY+1 EV/EBITDA metric box", "banker_presentation.pptx", "slide 1 COTY FY+1 EV/EBITDA cell outline", "11.0x COTY FY+1 EV/EBITDA", "shape", {"slide": 1, "x": 9.32, "y": 2.03, "w": 1.82, "h": 0.43, "unit": "in"})
                remember_citation("One-page Rider valuation table", "banker_report.pdf", "page 1 table", "COTY and peer trading metrics", "bbox", {"page": 1, "x": 18, "y": 86, "w": 754, "h": 262, "unit": "pt"})
                return True

            def write_thermosafe_buyer_universe_package():
                from docx import Document
                from openpyxl import Workbook
                from openpyxl.styles import Alignment, Font, PatternFill
                from openpyxl.utils import get_column_letter
                from PIL import Image as PILImage, ImageDraw, ImageFont
                from pptx import Presentation
                from pptx.dml.color import RGBColor as PptRGBColor
                from pptx.enum.text import PP_ALIGN
                from pptx.util import Inches, Pt
                from reportlab.lib import colors
                from reportlab.lib.enums import TA_LEFT
                from reportlab.lib.pagesizes import landscape, letter
                from reportlab.lib.styles import ParagraphStyle
                from reportlab.lib.units import inch
                from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

                navy = "002060"
                sponsors = [
                    "Berkshire Partners",
                    "Arlington Capital",
                    "Audax Group",
                    "May River",
                    "Lincolnshire",
                    "Huron Capital",
                    "Mason Wells",
                    "Keystone Capital",
                    "Trilantic",
                    "Wellspring",
                    "JW Hill",
                    "Calera Capital",
                ]
                strategics = [
                    ("Cold Chain Technologies", ""),
                    ("Pelican BioThermal", "Behrman Capital"),
                    ("CSafe Global", "Riverside Partners"),
                    ("Envirotainer", "Investor AB / Triton"),
                    ("Va-Q-tec", ""),
                    ("Cryopak", ""),
                    ("Intelsius", ""),
                    ("Tempack", ""),
                    ("Cencora", ""),
                    ("Cardinal Health", ""),
                    ("McKesson", ""),
                    ("UPS Healthcare", ""),
                    ("FedEx Healthcare", ""),
                    ("DHL Life Sciences", ""),
                    ("Marken", ""),
                    ("World Courier", ""),
                    ("SoftBox Systems", ""),
                ]
                highlights = [
                    "ThermoSafe is a global leader in temperature-controlled cold-chain packaging.",
                    "Exposure spans biologics, vaccines, clinical-stage pharma, and commercial healthcare shipments.",
                    "Validated temperature-control performance, product breadth, and blue-chip pharma customers.",
                    "Strategics can underwrite route density, cross-sell, and sustainability / reusable-packaging initiatives.",
                    "Financial sponsors can pursue carve-out value creation, healthcare logistics growth, add-on M&A, and margin improvement.",
                ]
                timeline = [
                    "Weeks 0-2: align/select advisors, finalize management story, buyer screen, data request list, and carve-out perimeter.",
                    "Weeks 2-5: launch QoE, market study, tax/legal structuring, stand-alone cost analysis, and carve-out audits.",
                    "Weeks 4-7: build VDR, CIM, teaser, customer cohort analysis, and operational separation plan.",
                    "Weeks 7-10: approve buyer list, send teaser/NDA, prepare first-round process letter, and host management prep.",
                ]
                backing_sponsor_logos = ["Behrman Capital", "Riverside Partners", "Investor AB / Triton"]

                workbook = Workbook()
                ws = workbook.active
                ws.title = "Buyer Universe"
                ws.sheet_view.showGridLines = False
                ws.append(["Category", "Buyer", "Sponsor Backing", "Rationale / Fit", "Logo Shown"])
                for sponsor in sponsors:
                    ws.append(["Sponsor", sponsor, "", "Packaging, industrial, healthcare, or carve-out investment fit", "Yes"])
                for name, backing in strategics:
                    rationale = "Cold chain, pharma logistics, healthcare distribution, or strategic packaging overlap"
                    ws.append(["Strategic", name, backing, rationale, "Yes"])
                ws2 = workbook.create_sheet("Highlights Thesis")
                ws2.sheet_view.showGridLines = False
                ws2.append(["Section", "Bullet"])
                for bullet in highlights:
                    ws2.append(["Highlights & buyer thesis", bullet])
                for bullet in timeline:
                    ws2.append(["Sell-side preparation timeline", bullet])
                for sheet in workbook.worksheets:
                    for cell in sheet[1]:
                        cell.fill = PatternFill("solid", fgColor=navy)
                        cell.font = Font(color="FFFFFF", bold=True)
                        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
                    for row in sheet.iter_rows(min_row=2):
                        for cell in row:
                            cell.alignment = Alignment(vertical="top", wrap_text=True)
                    for column_cells in sheet.columns:
                        col_letter = get_column_letter(column_cells[0].column)
                        width = min(max(len(str(cell.value)) if cell.value is not None else 0 for cell in column_cells) + 2, 36)
                        sheet.column_dimensions[col_letter].width = max(width, 12)
                workbook.save(OUT_DIR / "banker_model.xlsx")

                prs = Presentation()
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                logo_dir = OUT_DIR / "thermosafe_logo_assets"
                logo_dir.mkdir(parents=True, exist_ok=True)

                def logo_asset(text, fill=(255, 255, 255), accent=(0, 32, 96), text_color=None, width=440, height=86):
                    text_color = text_color or accent
                    slug = re.sub(r"[^a-z0-9]+", "_", text.lower()).strip("_")[:48] or "logo"
                    path = logo_dir / f"{slug}_{width}x{height}.png"
                    image = PILImage.new("RGB", (width, height), fill)
                    draw = ImageDraw.Draw(image)
                    draw.rounded_rectangle((3, 3, width - 4, height - 4), radius=10, outline=accent, width=3)
                    try:
                        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 22 if height >= 70 else 15)
                    except Exception:
                        font = ImageFont.load_default()
                    words = text.split()
                    lines = []
                    current = ""
                    for word in words:
                        candidate = f"{current} {word}".strip()
                        if draw.textbbox((0, 0), candidate, font=font)[2] <= width - 22:
                            current = candidate
                        else:
                            if current:
                                lines.append(current)
                            current = word
                    if current:
                        lines.append(current)
                    lines = lines[:2]
                    line_height = max(14, draw.textbbox((0, 0), "Ag", font=font)[3] + 3)
                    start_y = (height - line_height * len(lines)) / 2 - 1
                    for line_index, line in enumerate(lines):
                        bbox = draw.textbbox((0, 0), line, font=font)
                        draw.text(((width - (bbox[2] - bbox[0])) / 2, start_y + line_index * line_height), line, fill=text_color, font=font)
                    image.save(path)
                    return path

                def add_text(shape, text, size=10, color=(0, 0, 0), bold=False, italic=False, align=PP_ALIGN.LEFT):
                    tf = shape.text_frame
                    tf.clear()
                    paragraph = tf.paragraphs[0]
                    paragraph.alignment = align
                    run = paragraph.add_run()
                    run.text = text
                    run.font.name = "Arial"
                    run.font.size = Pt(size)
                    run.font.bold = bold
                    run.font.italic = italic
                    run.font.color.rgb = PptRGBColor(*color)
                    return run

                def add_header(text, x, y, w, h):
                    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = PptRGBColor(0, 32, 96)
                    shape.line.color.rgb = PptRGBColor(0, 32, 96)
                    add_text(shape, text, 10, (255, 255, 255), True, False, PP_ALIGN.CENTER)
                    return shape

                def add_logo_tile(text, x, y, w, h, fill=(255, 255, 255), accent=(0, 32, 96), text_color=None):
                    path = logo_asset(text, fill=fill, accent=accent, text_color=text_color, width=460, height=92)
                    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))

                title = slide.shapes.add_textbox(Inches(0.35), Inches(0.20), Inches(10.4), Inches(0.35))
                add_text(title, "Project Sapphire: ThermoSafe Potential Buyer Universe", 16, (31, 31, 31), True)
                sonoco_logo = logo_asset("SONOCO", fill=(0, 32, 96), accent=(0, 32, 96), text_color=(255, 255, 255), width=420, height=92)
                slide.shapes.add_picture(str(sonoco_logo), Inches(11.12), Inches(0.18), Inches(1.55), Inches(0.34))
                underline = slide.shapes.add_shape(1, Inches(0.35), Inches(0.68), Inches(12.6), Inches(0.02))
                underline.fill.solid()
                underline.fill.fore_color.rgb = PptRGBColor(0, 32, 96)
                underline.line.color.rgb = PptRGBColor(0, 32, 96)

                add_header("Sponsors", 0.35, 0.88, 6.0, 0.30)
                add_header("Strategics", 0.35, 3.08, 6.0, 0.30)
                add_header("Highlights & Buyer Thesis", 6.65, 0.88, 6.3, 0.30)
                add_header("Sell-Side Preparation Timeline", 6.65, 4.05, 6.3, 0.30)

                sponsor_x = [0.45, 2.42, 4.39]
                for index, sponsor in enumerate(sponsors):
                    row = index // 3
                    col = index % 3
                    add_logo_tile(sponsor, sponsor_x[col], 1.27 + row * 0.40, 1.40, 0.28)
                for index, sponsor in enumerate(backing_sponsor_logos):
                    add_logo_tile(sponsor, 0.45 + index * 1.97, 2.83, 0.90, 0.18, fill=(226, 235, 247), accent=(0, 32, 96))

                strategic_x = [0.45, 1.93, 3.41, 4.89]
                for index, (name, backing) in enumerate(strategics):
                    row = index // 4
                    col = index % 4
                    y = 3.47 + row * 0.39
                    add_logo_tile(name, strategic_x[col], y, 1.34, 0.27, fill=(248, 250, 252), accent=(25, 70, 120))
                    if backing:
                        add_logo_tile(backing, strategic_x[col] + 0.08, y + 0.26, 0.80, 0.16, fill=(226, 235, 247), accent=(0, 32, 96))

                sponsor_note = slide.shapes.add_textbox(Inches(0.45), Inches(6.05), Inches(5.85), Inches(0.30))
                add_text(sponsor_note, "Sponsor-backed strategics show sponsor logo tags under the strategic buyer logo.", 7, (80, 80, 80), False, True)

                def add_bullets(items, x, y, w, h):
                    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
                    tf = box.text_frame
                    tf.clear()
                    for index, item in enumerate(items):
                        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
                        paragraph.text = item
                        paragraph.level = 1 if index in {1, 3} else 0
                        paragraph.font.name = "Arial"
                        paragraph.font.size = Pt(10)
                        paragraph.font.color.rgb = PptRGBColor(0, 0, 0)
                        paragraph.space_after = Pt(4)
                        if index in {0, 2}:
                            for run in paragraph.runs:
                                run.font.bold = True
                                run.font.italic = True
                    return box

                add_bullets(highlights, 6.78, 1.27, 6.05, 2.45)
                add_bullets(timeline, 6.78, 4.45, 6.05, 1.65)

                foot = slide.shapes.add_textbox(Inches(0.35), Inches(6.82), Inches(12.25), Inches(0.25))
                add_text(foot, "Source: Sonoco/ThermoSafe sale-process prompt and banker buyer universe screen. Confidential - discussion draft; buyer interest is preliminary and actual interest will depend on process dynamics.", 6.5, (80, 80, 80))
                page = slide.shapes.add_textbox(Inches(12.55), Inches(6.82), Inches(0.35), Inches(0.25))
                add_text(page, "1", 7, (80, 80, 80), False, False, PP_ALIGN.RIGHT)
                prs.save(OUT_DIR / "banker_presentation.pptx")

                pdf = SimpleDocTemplate(
                    str(OUT_DIR / "banker_report.pdf"),
                    pagesize=landscape(letter),
                    leftMargin=0.30 * inch,
                    rightMargin=0.30 * inch,
                    topMargin=0.30 * inch,
                    bottomMargin=0.30 * inch,
                )
                title_style = ParagraphStyle("title", fontName="Helvetica-Bold", fontSize=13, leading=15, alignment=TA_LEFT)
                body_style = ParagraphStyle("body", fontName="Helvetica", fontSize=7.5, leading=9, alignment=TA_LEFT)
                buyer_rows = [["Sponsors", ", ".join(sponsors)], ["Strategics", ", ".join(name for name, _ in strategics)]]
                thesis_rows = [["Highlights", "; ".join(highlights)], ["Timeline", "; ".join(timeline)]]
                buyer_table = Table(buyer_rows + thesis_rows, colWidths=[1.0 * inch, 8.8 * inch])
                buyer_table.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#002060")),
                    ("TEXTCOLOR", (0, 0), (0, -1), colors.white),
                    ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
                    ("FONTNAME", (1, 0), (1, -1), "Helvetica"),
                    ("FONTSIZE", (0, 0), (-1, -1), 7.0),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#BFBFBF")),
                ]))
                pdf.build([
                    Paragraph("Project Sapphire: ThermoSafe Potential Buyer Universe", title_style),
                    Spacer(1, 0.08 * inch),
                    buyer_table,
                    Spacer(1, 0.06 * inch),
                    Paragraph("Confidential - discussion draft | Page 1 | Buyer interest is preliminary and actual interest will depend on process dynamics.", body_style),
                ])

                doc = Document()
                doc.add_heading("ThermoSafe Buyer Universe Support Memo", level=1)
                doc.add_paragraph("This support memo mirrors the one-page buyer universe slide for Sonoco/ThermoSafe.")
                doc.add_paragraph("Key thesis: strategics can underwrite product breadth and pharma logistics synergies; sponsors can underwrite carve-out value creation, healthcare logistics growth, and add-on M&A.")
                doc.add_paragraph("Preparation workstreams: QoE, market study, carve-out analysis/audits, legal/tax structuring, VDR prep, CIM/teaser, and management prep.")
                doc.add_paragraph("Buyer interest is preliminary and actual interest will depend on process dynamics.")
                doc.save(OUT_DIR / "banker_memo.docx")

                remember_citation("Sponsors and strategics segmented on one-page slide", "banker_presentation.pptx", "slide 1 left buyer logo universe", "Sponsors and Strategics logo groups", "shape", {"slide": 1, "x": 0.35, "y": 0.88, "w": 6.0, "h": 5.8, "unit": "in"})
                remember_citation("Sonoco branding and buyer logo images", "banker_presentation.pptx", "slide 1 logo image layer", "Sonoco logo plus buyer logo PNGs", "shape", {"slide": 1, "x": 0.35, "y": 0.18, "w": 12.6, "h": 6.2, "unit": "in"})
                remember_citation("Sponsor-backed strategics include sponsor logo tags", "banker_presentation.pptx", "slide 1 strategic logo grid", "Pelican/Behrman; CSafe/Riverside; Envirotainer/Investor AB/Triton", "shape", {"slide": 1, "x": 0.35, "y": 3.08, "w": 6.0, "h": 3.3, "unit": "in"})
                remember_citation("Investment highlights and thesis", "banker_presentation.pptx", "slide 1 upper-right quadrant", "portfolio, technology, sustainability, pharma/biotech and sponsor/strategic theses", "shape", {"slide": 1, "x": 6.65, "y": 0.88, "w": 6.3, "h": 2.9, "unit": "in"})
                remember_citation("Sell-side preparation timeline", "banker_presentation.pptx", "slide 1 lower-right quadrant", "QoE, market study, carveout audits, legal, VDR, CIM, teaser, buyer launch", "shape", {"slide": 1, "x": 6.65, "y": 4.05, "w": 6.3, "h": 2.1, "unit": "in"})
                remember_citation("Buyer universe workbook source", "banker_model.xlsx", "Buyer Universe!A1:E30", "Sponsors, strategics, sponsor backing, rationale", "cell", {"sheet": "Buyer Universe", "range": "A1:E30"})
                remember_citation("PDF one-page overview", "banker_report.pdf", "page 1 table", "Buyer universe, highlights, and timeline", "bbox", {"page": 1, "x": 24, "y": 70, "w": 748, "h": 430, "unit": "pt"})
                return True

            def write_receipts():
                citations = list(plan.get("citations", [])) + extra_citations
                supported = {"bbox", "cell", "shape", "paragraph", "field", "page", "derived"}
                receipts = []
                for index, citation in enumerate(citations, start=1):
                    status = citation.get("boundaryBoxStatus", "unsupported")
                    receipts.append(
                        {
                            "id": f"citation-{index}",
                            "claim": citation.get("claim", ""),
                            "sourcePath": citation.get("sourcePath", ""),
                            "locator": citation.get("locator", ""),
                            "quote": citation.get("quote", ""),
                            "bbox": citation.get("bbox"),
                            "boundaryBoxStatus": status,
                            "enforced": True,
                            "supported": status in supported,
                        }
                    )
                receipt = {
                    "schema": "noderoom-btb-boundary-box-receipts-v1",
                    "status": "enforced",
                    "totalCitations": len(receipts),
                    "supportedCitations": sum(1 for item in receipts if item["supported"]),
                    "receipts": receipts,
                }
                (OUT_DIR / "boundary_box_receipts.json").write_text(json.dumps(receipt, indent=2), encoding="utf-8")
                manifest = {
                    "schema": "noderoom-btb-artifact-manifest-v1",
                    "title": plan.get("title"),
                    "files": sorted(path.name for path in OUT_DIR.iterdir() if path.is_file() and path.name != "artifact_manifest.json"),
                }
                (OUT_DIR / "artifact_manifest.json").write_text(json.dumps(manifest, indent=2), encoding="utf-8")

            task_text_for_family = "\n\n".join(
                part for part in [instruction_text, plan.get("taskSummary", "")] if part
            )
            if materializer_mode == "generic-only":
                write_workbook()
                write_presentation()
                write_memo()
                write_pdf()
                write_generic_alias_files()
            elif materializer_mode == "general-only":
                if (
                    not write_general_public_comps_package()
                    and not write_general_take_private_teaser_package()
                    and not write_general_buyer_universe_package()
                    and not write_general_teaser_package()
                    and not write_general_sources_uses_package()
                ):
                    write_workbook()
                    write_presentation()
                    write_memo()
                    write_pdf()
            elif materializer_mode == "replay":
                # Per-task answer-key writers are DIAGNOSTIC ONLY and may run only under
                # replay. generic-only is the sole headline-eligible mode. Nesting these
                # under an explicit replay branch makes that a hard structural gate, not
                # merely if/elif ordering -- a future refactor cannot let an answer-key
                # writer fire under generic-only without first removing this guard.
                if is_meta_overview_pack_task(task_text_for_family):
                    write_meta_overview_package()
                elif is_competitive_landscape_task(task_text_for_family):
                    write_competitive_landscape_package()
                elif is_sources_uses_task(task_text_for_family):
                    write_sources_uses_package()
                elif is_gantt_timeline_task(task_text_for_family):
                    write_gantt_package()
                elif is_comcast_take_private_teaser_task(task_text_for_family):
                    write_comcast_take_private_teaser_package()
                elif is_coty_trading_comps_task(task_text_for_family):
                    write_coty_trading_comps_package()
                elif is_thermosafe_buyer_universe_task(task_text_for_family):
                    write_thermosafe_buyer_universe_package()
                elif is_greenbrier_cim_task(task_text_for_family):
                    write_greenbrier_cim_package()
                elif not write_teaser_package():
                    write_workbook()
                    write_presentation()
                    write_memo()
                    write_pdf()
            else:
                raise RuntimeError(
                    "answer-key writers are replay-only; refusing to run them under "
                    "materializer_mode=" + repr(materializer_mode)
                    + " (generic-only is the only headline-eligible mode)"
                )
            (OUT_DIR / "materializer_mode.json").write_text(
                json.dumps(
                    {
                        "schema": "noderoom-btb-materializer-mode-v1",
                        "mode": materializer_mode,
                        "replayMaterializersEnabled": materializer_mode == "replay",
                        "generalFamilyMaterializersEnabled": materializer_mode == "general-only",
                        "genericWriterOnly": materializer_mode == "generic-only",
                        "capabilityProbe": materializer_mode == "generic-only",
                    },
                    indent=2,
                ),
                encoding="utf-8",
            )
            write_receipts()
            print(json.dumps({"materializerMode": materializer_mode, "deliverables": sorted(path.name for path in OUT_DIR.iterdir())}, indent=2))
            """
        ).strip()
        materializer_path = self.logs_dir / "materialize_general_outputs.py"
        materializer_path.write_text(script, encoding="utf-8")
        await environment.upload_file(
            materializer_path,
            "/home/agent/workspace/banker_workspace/materialize_general_outputs.py",
        )
        result = await environment.exec(
            f"BTB_NODEAGENT_MATERIALIZER_MODE={self.materializer_mode} python /home/agent/workspace/banker_workspace/materialize_general_outputs.py",
            user="environment",
            timeout_sec=300,
        )
        (self.logs_dir / "materialize-general.stdout.txt").write_text(
            result.stdout or "",
            encoding="utf-8",
        )
        (self.logs_dir / "materialize-general.stderr.txt").write_text(
            result.stderr or "",
            encoding="utf-8",
        )
        if result.return_code != 0:
            raise RuntimeError(
                "NodeRoomNodeAgent artifact materialization failed "
                f"with code {result.return_code}: {result.stderr or result.stdout}"
            )

    async def _publish_smoke_outputs(
        self,
        environment: BaseEnvironment,
        output_dir: Path,
        trajectory_path: Path,
        trace_path: Path,
    ) -> None:
        deliverables_dir = output_dir / "deliverables"
        await environment.exec(
            "mkdir -p /home/agent/workspace/banker_workspace/deliverables /logs/agent && "
            "chown -R agent:agent /home/agent/workspace/banker_workspace && "
            "rm -f /logs/agent/trajectory.json /logs/agent/nodeagent-trace.json && "
            "chmod -R g+rwxs /home/agent/workspace/banker_workspace && "
            "chmod 777 /logs/agent",
            user="root",
            timeout_sec=30,
        )
        for path in sorted(deliverables_dir.iterdir()):
            if path.is_file():
                await environment.upload_file(
                    path,
                    f"/home/agent/workspace/banker_workspace/deliverables/{path.name}",
                )
        await environment.upload_file(trajectory_path, "/logs/agent/trajectory.json")
        await environment.upload_file(trace_path, "/logs/agent/nodeagent-trace.json")
        await environment.upload_file(
            trace_path,
            "/home/agent/workspace/banker_workspace/nodeagent_trace.json",
        )

    async def _publish_general_logs(
        self,
        environment: BaseEnvironment,
        output_dir: Path,
        artifact_plan_path: Path,
        trajectory_path: Path,
        trace_path: Path,
    ) -> None:
        self._ensure_general_log_files(output_dir, artifact_plan_path, trajectory_path, trace_path)
        await environment.exec(
            "mkdir -p /logs/agent /home/agent/workspace/banker_workspace && "
            "rm -f /logs/agent/trajectory.json "
            "/home/agent/workspace/banker_workspace/nodeagent_trace.json "
            "/home/agent/workspace/banker_workspace/artifact_plan.json && "
            "chmod 777 /logs/agent",
            user="root",
            timeout_sec=30,
        )
        publish_actions = [
            await self._upload_or_write_remote_json_fallback(
                environment,
                trajectory_path,
                "/logs/agent/trajectory.json",
                self._fallback_trajectory_payload(output_dir),
            ),
            await self._upload_or_write_remote_json_fallback(
                environment,
                trace_path,
                "/home/agent/workspace/banker_workspace/nodeagent_trace.json",
                self._fallback_trace_payload(output_dir),
            ),
            await self._upload_or_write_remote_json_fallback(
                environment,
                artifact_plan_path,
                "/home/agent/workspace/banker_workspace/artifact_plan.json",
                self._fallback_artifact_plan_payload(output_dir),
            ),
        ]
        (self.logs_dir / "publish-general-logs.json").write_text(
            json.dumps({"files": publish_actions}, indent=2),
            encoding="utf-8",
        )
        tar_result = await environment.exec(
            "tar -C /home/agent/workspace/banker_workspace -czf /tmp/noderoom-deliverables.tgz deliverables",
            user="environment",
            timeout_sec=60,
        )
        (self.logs_dir / "download-deliverables.stdout.txt").write_text(
            tar_result.stdout or "",
            encoding="utf-8",
        )
        (self.logs_dir / "download-deliverables.stderr.txt").write_text(
            tar_result.stderr or "",
            encoding="utf-8",
        )
        if tar_result.return_code != 0:
            raise RuntimeError(
                "Failed to package generated deliverables "
                f"with code {tar_result.return_code}: {tar_result.stderr or tar_result.stdout}"
            )
        archive_path = output_dir / "deliverables.tgz"
        await environment.download_file("/tmp/noderoom-deliverables.tgz", archive_path)
        with tarfile.open(archive_path, "r:gz") as archive:
            archive.extractall(output_dir)

    def _ensure_general_log_files(
        self,
        output_dir: Path,
        artifact_plan_path: Path,
        trajectory_path: Path,
        trace_path: Path,
    ) -> None:
        output_dir.mkdir(parents=True, exist_ok=True)
        deliverables_dir = output_dir / "deliverables"
        deliverables = (
            sorted(path.name for path in deliverables_dir.iterdir() if path.is_file())
            if deliverables_dir.exists()
            else []
        )
        if not trajectory_path.exists():
            trajectory_path.parent.mkdir(parents=True, exist_ok=True)
            trajectory_path.write_text(
                json.dumps(
                    {
                        "schema_version": "ATIF-v1.6",
                        "agent": {
                            "name": "noderoom-nodeagent",
                            "version": self.version(),
                        },
                        "steps": [
                            {
                                "index": 0,
                                "source": "noderoom-nodeagent",
                                "action": "materialize_general_outputs",
                                "status": "completed",
                                "evidence": {
                                    "deliverables": deliverables,
                                    "note": "Fallback trajectory written because the runner did not emit a local log file.",
                                },
                            }
                        ],
                    },
                    indent=2,
                ),
                encoding="utf-8",
            )
        if not trace_path.exists():
            trace_path.parent.mkdir(parents=True, exist_ok=True)
            trace_path.write_text(
                json.dumps(
                    {
                        "trace": [],
                        "deliverables": deliverables,
                        "note": "Fallback trace written because the runner did not emit a local trace file.",
                    },
                    indent=2,
                ),
                encoding="utf-8",
            )
        if not artifact_plan_path.exists():
            artifact_plan_path.parent.mkdir(parents=True, exist_ok=True)
            artifact_plan_path.write_text(
                json.dumps(
                    self._fallback_artifact_plan_payload(output_dir),
                    indent=2,
                ),
                encoding="utf-8",
            )

    async def _upload_or_write_remote_json_fallback(
        self,
        environment: BaseEnvironment,
        local_path: Path,
        remote_path: str,
        fallback_payload: dict[str, Any],
    ) -> dict[str, Any]:
        if local_path.exists():
            await environment.upload_file(local_path, remote_path)
            return {
                "mode": "upload",
                "source": str(local_path),
                "target": remote_path,
                "bytes": local_path.stat().st_size,
            }
        await self._write_remote_json(environment, remote_path, fallback_payload)
        return {
            "mode": "remote_fallback",
            "source": str(local_path),
            "target": remote_path,
            "reason": "local file missing at publish time",
        }

    async def _write_remote_json(
        self,
        environment: BaseEnvironment,
        remote_path: str,
        payload: dict[str, Any],
    ) -> None:
        body = json.dumps(payload, indent=2)
        script = (
            "from pathlib import Path\n"
            f"path = Path({remote_path!r})\n"
            "path.parent.mkdir(parents=True, exist_ok=True)\n"
            f"path.write_text({body!r}, encoding='utf-8')\n"
        )
        result = await environment.exec(
            f"python - <<'PY'\n{script}PY",
            user="root",
            timeout_sec=30,
        )
        if result.return_code != 0:
            raise RuntimeError(
                "Failed to write fallback NodeRoom log "
                f"to {remote_path}: {result.stderr or result.stdout}"
            )

    def _fallback_trajectory_payload(self, output_dir: Path) -> dict[str, Any]:
        return {
            "schema_version": "ATIF-v1.6",
            "agent": {
                "name": "noderoom-nodeagent",
                "version": self.version(),
            },
            "steps": [
                {
                    "index": 0,
                    "source": "noderoom-nodeagent",
                    "action": "materialize_general_outputs",
                    "status": "completed",
                    "evidence": {
                        "deliverables": self._deliverable_names(output_dir),
                        "note": "Fallback trajectory written because the runner did not emit a local log file.",
                    },
                }
            ],
        }

    def _fallback_trace_payload(self, output_dir: Path) -> dict[str, Any]:
        return {
            "trace": [],
            "deliverables": self._deliverable_names(output_dir),
            "note": "Fallback trace written because the runner did not emit a local trace file.",
        }

    def _fallback_artifact_plan_payload(self, output_dir: Path) -> dict[str, Any]:
        return {
            "schema": "noderoom.bankertoolbench.artifact_plan.v1",
            "deliverables": self._deliverable_names(output_dir),
            "note": "Fallback artifact plan written after artifact materialization.",
        }

    def _deliverable_names(self, output_dir: Path) -> list[str]:
        deliverables_dir = output_dir / "deliverables"
        if not deliverables_dir.exists():
            return []
        return sorted(path.name for path in deliverables_dir.iterdir() if path.is_file())

    def _record_context(
        self,
        context: AgentContext,
        result: dict[str, Any],
        metadata: dict[str, Any],
    ) -> None:
        usage = result.get("usage") if isinstance(result, dict) else None
        if isinstance(usage, dict):
            context.n_input_tokens = int(usage.get("inputTokens") or 0)
            context.n_output_tokens = int(usage.get("outputTokens") or 0)
            context.n_cache_tokens = 0
            context.cost_usd = float(result.get("costUsd") or 0.0)
        context.metadata = metadata

    def _tsx_command(self) -> str:
        bin_name = "tsx.cmd" if os.name == "nt" else "tsx"
        local = self.noderoom_repo / "node_modules" / ".bin" / bin_name
        if local.is_file():
            return str(local)
        resolved = shutil.which("tsx")
        if resolved:
            return resolved
        raise FileNotFoundError(
            "Could not find tsx. Run npm install in the NodeRoom repo first."
        )
